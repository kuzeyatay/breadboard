"""
shells_seed/stat_trio_row.py — Three-stat hero row: each stat has a value, label, delta, and caption.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'metric_dashboard'
DESCRIPTION = 'Three-stat hero row: each stat has a value, label, delta, and caption.'
ARCHETYPE = 'data'
MOOD = ['boardroom', 'bold', 'technical']
DENSITY = 'balanced'
STYLE_TAGS = ['three_up', 'stat_row', 'trio', 'kpi']

SLOTS = [{'name': 'section_label', 'kind': 'text', 'max_chars': 30, 'style': 'caption', 'required': False}, {'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title', 'required': True}, {'name': 'stats', 'kind': 'bullet_list', 'bullet_capacity': 3, 'required': True}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    section_label = get_slot(slots, "section_label")
    if section_label:
        sl = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(8), Inches(0.4))
        set_textbox_text(sl, section_label, theme, "caption", "accent", align="left")
        sl.name = "section_label"

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(1.0), Inches(13.333 - 2*margin), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"
    add_theme_entrance(slide, hl, theme, delay_ms=100, index=0)

    add_hairline(slide, margin, 2.1, 13.333 - 2*margin, theme, "border")

    stats = get_slot(slots, "stats", default=[]) or []
    if isinstance(stats, str):
        stats = [x.strip() for x in stats.split("\n") if x.strip()]
    stats = (stats + [""]*3)[:3]
    col_w = (13.333 - 2*margin) / 3
    col_y = 2.6
    for i, entry in enumerate(stats):
        x = margin + i * col_w
        # Parse "VALUE | LABEL | DELTA | CAPTION"
        if isinstance(entry, dict):
            val = entry.get("value", ""); lbl = entry.get("label", "")
            dlt = entry.get("delta", ""); cap = entry.get("caption", "")
        else:
            parts = [p.strip() for p in str(entry).split("|")]
            parts = (parts + [""]*4)[:4]
            val, lbl, dlt, cap = parts

        vb = slide.shapes.add_textbox(Inches(x+0.2), Inches(col_y),
                                       Inches(col_w-0.4), Inches(1.8))
        set_textbox_text(vb, val, theme, "metric_xl", "accent", align="left")
        vb.name = f"stat_value_{i}"
        lb = slide.shapes.add_textbox(Inches(x+0.2), Inches(col_y+1.7),
                                       Inches(col_w-0.4), Inches(0.6))
        set_textbox_text(lb, lbl, theme, "body_bold", "text", align="left")
        lb.name = f"stat_label_{i}"
        if dlt:
            db = slide.shapes.add_textbox(Inches(x+0.2), Inches(col_y+2.3),
                                           Inches(col_w-0.4), Inches(0.4))
            set_textbox_text(db, dlt, theme, "body_bold", "success", align="left")
            db.name = f"stat_delta_{i}"
        if cap:
            cb = slide.shapes.add_textbox(Inches(x+0.2), Inches(col_y+2.8),
                                           Inches(col_w-0.4), Inches(1.2))
            set_textbox_text(cb, truncate_to(cap, 100), theme, "body", "muted", align="left")
            cb.name = f"stat_caption_{i}"
        add_theme_entrance(slide, vb, theme, delay_ms=300+i*150, index=i+1)
