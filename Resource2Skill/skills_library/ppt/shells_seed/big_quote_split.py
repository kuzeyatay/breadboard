"""
shells_seed/big_quote_split.py — Large quote on left 65%, attribution card with role + company on right 35%.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'hero_quote'
DESCRIPTION = 'Large quote on left 65%, attribution card with role + company on right 35%.'
ARCHETYPE = 'brand'
MOOD = ['editorial', 'cinematic', 'bold']
DENSITY = 'sparse'
STYLE_TAGS = ['quote', 'split', 'big_quote', 'attribution_card']

SLOTS = [{'name': 'quote', 'kind': 'text', 'max_chars': 220, 'style': 'accent_serif', 'required': True}, {'name': 'attribution', 'kind': 'text', 'max_chars': 50, 'style': 'body_bold', 'required': False}, {'name': 'role', 'kind': 'text', 'max_chars': 60, 'style': 'caption', 'required': False}, {'name': 'company', 'kind': 'text', 'max_chars': 40, 'style': 'caption', 'required': False}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    # Big serif quote glyph top-left
    g = slide.shapes.add_textbox(Inches(margin), Inches(0.5), Inches(1.8), Inches(1.6))
    g.text_frame.text = "\u201C"
    r = g.text_frame.paragraphs[0].runs[0]
    r.font.size = Pt(140); r.font.italic = True
    r.font.name = theme.get("typography", {}).get("accent_serif", {}).get("font", "Playfair Display")
    r.font.color.rgb = palette_color(theme, "accent")
    g.name = "quote_glyph"

    quote = truncate_to(get_slot(slots, "quote", required=True), 220)
    qb = slide.shapes.add_textbox(Inches(margin), Inches(2.0),
                                   Inches(8.2), Inches(4.5))
    set_textbox_text(qb, quote, theme, "accent_serif", "text", align="left")
    qb.name = "quote"
    add_theme_entrance(slide, qb, theme, delay_ms=200, index=0)

    # Right attribution card
    card_x = 9.3; card_y = 2.0; card_w = 13.333 - card_x - margin; card_h = 4.5
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(card_x), Inches(card_y),
                                   Inches(card_w), Inches(card_h))
    card.fill.solid(); card.fill.fore_color.rgb = palette_color(theme, "panel")
    card.line.color.rgb = palette_color(theme, "accent")
    card.line.width = Pt(1.0); card.shadow.inherit = False
    card.name = "attribution_card"
    # Hairline at top
    hairline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(card_x + 0.3), Inches(card_y + 0.6),
                                       Inches(0.8), Inches(0.04))
    hairline.fill.solid(); hairline.fill.fore_color.rgb = palette_color(theme, "accent")
    hairline.line.fill.background(); hairline.shadow.inherit = False

    attrib = get_slot(slots, "attribution") or ""
    role = get_slot(slots, "role") or ""
    company = get_slot(slots, "company") or ""
    ab = slide.shapes.add_textbox(Inches(card_x + 0.3), Inches(card_y + 1.0),
                                   Inches(card_w - 0.6), Inches(0.7))
    set_textbox_text(ab, attrib, theme, "body_bold", "text", align="left")
    ab.name = "attribution"
    rb = slide.shapes.add_textbox(Inches(card_x + 0.3), Inches(card_y + 1.8),
                                   Inches(card_w - 0.6), Inches(0.5))
    set_textbox_text(rb, role, theme, "caption", "muted", align="left")
    rb.name = "role"
    cb = slide.shapes.add_textbox(Inches(card_x + 0.3), Inches(card_y + 2.3),
                                   Inches(card_w - 0.6), Inches(0.5))
    set_textbox_text(cb, company, theme, "caption", "accent", align="left")
    cb.name = "company"

    add_theme_entrance(slide, card, theme, delay_ms=400, index=1)
