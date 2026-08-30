"""
shells_seed/metric_dashboard.py — 3-4 KPI metrics with big numbers.

Slots:
  section_label (text, optional, caption style)
  headline      (text, required, ≤60 chars, title style)
  metrics       (bullet_list, required, 3-4 items, each {value, label, delta?})
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, apply_type_style, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'metric_dashboard'
DESCRIPTION = 'Multi-metric 2x2 KPI tile grid with section_label + headline stack.'
ARCHETYPE = 'data'
MOOD = ['boardroom', 'technical']
DENSITY = 'dense'
STYLE_TAGS = ['grid', 'data_tiles', 'two_by_two']

SLOTS = [
    {"name": "section_label", "kind": "text",        "max_chars": 30, "style": "caption", "required": False},
    {"name": "headline",      "kind": "text",        "max_chars": 60, "style": "title",   "required": True},
    {"name": "metrics",       "kind": "bullet_list", "bullet_capacity": 4, "required": True},
]


def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    gutter = theme.get("spacing", {}).get("gutter", 0.3)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    y = margin + 0.2

    # Eyebrow
    label = get_slot(slots, "section_label")
    if label:
        lb = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(6), Inches(0.35))
        set_textbox_text(lb, truncate_to(label, 30), theme, "caption", color_key="accent")
        reveal.append(lb); y += 0.55
        if theme.get("motif", {}).get("type") == "thin_rule":
            add_hairline(slide, margin, y, 1.8, theme); y += 0.15

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 60)
    head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(1.1))
    set_textbox_text(head, headline, theme, "title", color_key="text")
    reveal.append(head); y += 1.5

    # Metrics row
    metrics = get_slot(slots, "metrics", required=True)
    if not isinstance(metrics, list) or len(metrics) < 2:
        raise ValueError(f"'metrics' must be a list of 2-4 items, got {metrics!r}")
    metrics = metrics[:4]

    total_w = 13.333 - 2 * margin
    card_w = (total_w - (len(metrics) - 1) * gutter) / len(metrics)
    card_h = 3.6

    for i, m in enumerate(metrics):
        if not isinstance(m, dict):
            m = {"value": str(m), "label": ""}
        x = margin + i * (card_w + gutter)

        # Card
        card = add_solid_rect(
            slide, x, y, card_w, card_h, theme,
            color_key="panel", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        card.adjustments[0] = 0.04

        # Top accent bar
        add_solid_rect(slide, x + 0.3, y + 0.35, 0.7, 0.07, theme, color_key="accent")

        # Big metric value — use full card width so 5-char values don't wrap
        value_text = truncate_to(str(m.get("value", "")), 10)
        value_box = slide.shapes.add_textbox(
            Inches(x + 0.05), Inches(y + 0.65), Inches(card_w - 0.1), Inches(1.6),
        )
        value_tf = value_box.text_frame
        value_tf.word_wrap = False  # prevent wrapping on long metrics
        value_tf.margin_left = value_tf.margin_right = Inches(0.02)
        p = value_tf.paragraphs[0]
        p.text = value_text
        if p.runs:
            style_token = "metric_xl" if "metric_xl" in theme.get("typography", {}) else "title_xl"
            apply_type_style(p.runs[0], theme, style_token, color_key="accent")
            from pptx.util import Pt
            from pptx.enum.text import PP_ALIGN
            base = theme.get("typography", {}).get(style_token, {}).get("size", 56)
            # Card-width-aware sizing. Card_w in inches; assume 0.55*pt per char.
            # Effective char budget ≈ (card_w - 0.2) inches × 72 pt/in / 0.55 / size
            # Solve: size ≤ (card_w - 0.2) × 72 / (0.55 × len)
            max_size_for_width = (card_w - 0.2) * 72 / (0.55 * max(len(value_text), 1))
            chosen = min(base, max_size_for_width, 64)
            p.runs[0].font.size = Pt(max(chosen, 22))
            p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(y + 2.25), Inches(card_w - 0.5), Inches(0.5),
        )
        set_textbox_text(lbl_box, truncate_to(m.get("label", ""), 40), theme, "body", color_key="text")

        # Delta (optional)
        delta = m.get("delta")
        if delta:
            delta_color = "success" if str(delta).startswith(("+", "↑")) else \
                          "warning" if str(delta).startswith(("-", "↓")) else "muted"
            d_box = slide.shapes.add_textbox(
                Inches(x + 0.25), Inches(y + 2.8), Inches(card_w - 0.5), Inches(0.4),
            )
            set_textbox_text(d_box, str(delta), theme, "caption", color_key=delta_color)
            reveal.append(d_box)

        reveal.extend([card, value_box, lbl_box])

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
