"""
shells_seed/dense_roadmap_quarters.py — 4-quarter roadmap, dense.

Slots:
  section_label (text, optional, caption)
  headline      (text, required, title)
  quarters      (bullet_list, required, list of 4 dicts OR 4 strings where each
                  is a title + "|" + up to 3 milestone lines)

Layout: 4 vertical column cards, each with quarter label + milestone bullets
inside. Information-dense roadmap/pipeline view.
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = "timeline_horizontal"
DESCRIPTION = "Dense 4-column roadmap: each column shows a quarter/phase + 3 milestone bullets."
ARCHETYPE = "boardroom"
MOOD = ["boardroom", "restrained", "technical", "editorial"]
DENSITY = "dense"
STYLE_TAGS = ["four_column", "roadmap", "milestones"]

SLOTS = [
    {"name": "section_label", "kind": "text", "max_chars": 30, "style": "caption", "required": False},
    {"name": "headline",      "kind": "text", "max_chars": 100, "style": "title", "required": True},
    {"name": "quarters",      "kind": "bullet_list", "bullet_capacity": 4, "required": True},
]


def _parse_quarter(entry):
    """Accept either a dict {label, items:[...]} or a 'label | item1 | item2 | item3' string."""
    if isinstance(entry, dict):
        return entry.get("label") or entry.get("title") or "", \
               list(entry.get("items") or entry.get("milestones") or [])[:3]
    if isinstance(entry, str):
        parts = [x.strip() for x in entry.split("|")]
        if len(parts) >= 2:
            return parts[0], parts[1:4]
        return parts[0], []
    return "", []


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    section_label = get_slot(slots, "section_label")
    if section_label:
        sl = slide.shapes.add_textbox(Inches(margin), Inches(margin),
                                       Inches(7.0), Inches(0.4))
        set_textbox_text(sl, section_label, theme, "caption", "accent", align="left")
        sl.name = "section_label"

    headline = truncate_to(get_slot(slots, "headline", required=True), 100)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(1.1),
                                   Inches(13.333 - 2 * margin), Inches(1.0))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"
    add_theme_entrance(slide, hl, theme, delay_ms=120, index=0)

    # Horizontal rule
    add_hairline(slide, margin, 2.25, 13.333 - 2 * margin, theme, "border")

    entries = get_slot(slots, "quarters", default=[]) or []
    if isinstance(entries, str):
        entries = [entries]
    entries = entries[:4]
    n = len(entries)
    if n == 0:
        return

    col_y = 2.6
    col_h = 4.5
    gap = 0.3
    usable_w = 13.333 - 2 * margin - (n - 1) * gap
    col_w = max(1.0, usable_w / n)

    reveals = []
    for i, entry in enumerate(entries):
        label, items = _parse_quarter(entry)
        x = margin + i * (col_w + gap)
        # Column card background
        card = add_solid_rect(slide, x, col_y, col_w, col_h, theme,
                              color_key="panel", line=True)
        # Accent top bar
        add_solid_rect(slide, x, col_y, col_w, 0.08, theme, color_key="accent")
        # Label
        lb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(col_y + 0.25),
                                       Inches(col_w - 0.3), Inches(0.5))
        set_textbox_text(lb, label, theme, "body_bold", "accent", align="left")
        lb.name = f"quarter_label_{i}"
        # Items
        items = list(items)[:3]
        for j, it in enumerate(items):
            y = col_y + 0.9 + j * 1.05
            tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y),
                                           Inches(col_w - 0.35),
                                           Inches(0.95))
            set_textbox_text(tb, f"• {it}", theme, "body", "text", align="left")
            tb.name = f"q{i}_item_{j}"
        reveals.append(card)

    add_sequential_reveal(slide, reveals, theme, start_delay_ms=300, step_ms=150)
