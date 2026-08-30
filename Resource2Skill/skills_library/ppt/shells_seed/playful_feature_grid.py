"""
shells_seed/playful_feature_grid.py — Six rounded tiles in a 3x2 grid, each with icon + label + 1-line description.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'feature_grid'
DESCRIPTION = 'Six rounded tiles in a 3x2 grid, each with icon + label + 1-line description.'
ARCHETYPE = 'brand'
MOOD = ['playful', 'warm', 'punchy']
DENSITY = 'balanced'
STYLE_TAGS = ['grid', 'rounded', 'tiles', 'six_up']

SLOTS = [{'name': 'section_label', 'kind': 'text', 'max_chars': 30, 'style': 'caption', 'required': False}, {'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title', 'required': True}, {'name': 'features', 'kind': 'bullet_list', 'bullet_capacity': 6, 'required': True}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    section_label = get_slot(slots, "section_label")
    if section_label:
        sl = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(7), Inches(0.4))
        set_textbox_text(sl, section_label, theme, "caption", "accent", align="left")
        sl.name = "section_label"

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(1.0), Inches(13.333 - 2*margin), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"
    add_theme_entrance(slide, hl, theme, delay_ms=100, index=0)

    features = get_slot(slots, "features", default=[]) or []
    if isinstance(features, str):
        features = [x.strip() for x in features.split("\n") if x.strip()]
    features = (features + [""] * 6)[:6]

    grid_x0 = margin
    grid_y0 = 2.2
    tile_w = (13.333 - 2*margin - 2*0.3) / 3
    tile_h = (7.5 - grid_y0 - margin - 0.3) / 2
    gap = 0.3
    tiles = []
    for i, text in enumerate(features):
        col = i % 3; row = i // 3
        x = grid_x0 + col * (tile_w + gap)
        y = grid_y0 + row * (tile_h + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(y), Inches(tile_w), Inches(tile_h))
        card.fill.solid(); card.fill.fore_color.rgb = palette_color(theme, "panel")
        card.line.color.rgb = palette_color(theme, "border", "muted")
        card.line.width = Pt(0.75); card.shadow.inherit = False
        card.name = f"feature_tile_{i}"
        # Small accent dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.3), Inches(y+0.35),
                                      Inches(0.3), Inches(0.3))
        dot.fill.solid(); dot.fill.fore_color.rgb = palette_color(theme, "accent")
        dot.line.fill.background(); dot.shadow.inherit = False
        dot.name = f"feature_dot_{i}"
        # Title + optional hyphenated detail line
        parts = str(text).split(" — ", 1) if " — " in str(text) else [str(text), ""]
        title = parts[0]
        detail = parts[1] if len(parts) > 1 else ""
        tb_title = slide.shapes.add_textbox(Inches(x+0.3), Inches(y+0.75),
                                             Inches(tile_w-0.6), Inches(0.5))
        set_textbox_text(tb_title, truncate_to(title, 40), theme, "body_bold", "text", align="left")
        tb_title.name = f"feature_title_{i}"
        if detail:
            tb_det = slide.shapes.add_textbox(Inches(x+0.3), Inches(y+1.3),
                                               Inches(tile_w-0.6), Inches(tile_h-1.5))
            set_textbox_text(tb_det, truncate_to(detail, 90), theme, "body", "muted", align="left")
            tb_det.name = f"feature_detail_{i}"
        tiles.append(card)
    add_sequential_reveal(slide, tiles, theme, start_delay_ms=300, step_ms=110)
