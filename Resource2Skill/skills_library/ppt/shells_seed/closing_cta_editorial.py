"""
shells_seed/closing_cta_editorial.py — editorial closing with serif accent.

Slots:
  headline (text, required, ≤100 chars, title_xl)
  cta      (text, required, ≤60 chars, subtitle)
  url      (text, optional, ≤60 chars, caption)

Layout: centered stack, decorative serif numeral "3." or "end.", calm editorial feel.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_theme_entrance, get_slot, palette_color,
    set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = "closing_cta"
DESCRIPTION = "Editorial centered closing: serif end-glyph, calm headline stack, subtle URL caption."
ARCHETYPE = "narrative"
MOOD = ["editorial", "calm", "cinematic"]
DENSITY = "sparse"
STYLE_TAGS = ["centered", "serif_accent", "calm"]

SLOTS = [
    {"name": "headline", "kind": "text", "max_chars": 100, "style": "title_xl", "required": True},
    {"name": "cta",      "kind": "text", "max_chars":  60, "style": "subtitle", "required": True},
    {"name": "url",      "kind": "text", "max_chars":  60, "style": "caption",  "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    # Decorative end glyph (large period)
    glyph = slide.shapes.add_textbox(Inches(6.0), Inches(0.6),
                                      Inches(1.333), Inches(1.4))
    glyph.text_frame.text = "."
    run = glyph.text_frame.paragraphs[0].runs[0]
    acc = theme.get("typography", {}).get("accent_serif", {})
    run.font.name = acc.get("font", "Playfair Display")
    run.font.size = Pt(140)
    run.font.color.rgb = palette_color(theme, "accent")
    glyph.name = "end_glyph"

    # Hairline
    add_hairline(slide, 3.333, 2.2, 6.667, theme, "border")

    headline = truncate_to(get_slot(slots, "headline", required=True), 100)
    hl = slide.shapes.add_textbox(Inches(1.5), Inches(2.7),
                                   Inches(10.333), Inches(2.0))
    set_textbox_text(hl, headline, theme, "title_xl", "text", align="center")
    hl.name = "headline"
    add_theme_entrance(slide, hl, theme, delay_ms=200, index=0)

    cta = truncate_to(get_slot(slots, "cta", required=True), 60)
    ct = slide.shapes.add_textbox(Inches(2.5), Inches(5.1),
                                   Inches(8.333), Inches(0.8))
    set_textbox_text(ct, cta, theme, "subtitle", "accent", align="center")
    ct.name = "cta"
    add_theme_entrance(slide, ct, theme, delay_ms=420, index=1)

    url = get_slot(slots, "url")
    if url:
        u = slide.shapes.add_textbox(Inches(2.5), Inches(6.0),
                                      Inches(8.333), Inches(0.4))
        set_textbox_text(u, url, theme, "caption", "muted", align="center")
        u.name = "url"
        add_theme_entrance(slide, u, theme, delay_ms=620, index=2)

    add_hairline(slide, 3.333, 6.7, 6.667, theme, "border")
