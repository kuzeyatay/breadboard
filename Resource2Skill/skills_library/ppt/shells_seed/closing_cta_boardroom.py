"""
shells_seed/closing_cta_boardroom.py — restrained boardroom closing.

Slots:
  section_label (text, optional, ≤30 chars, caption style)
  headline      (text, required, ≤100 chars, title_xl)
  cta           (text, required, ≤80 chars, subtitle)
  footer        (text, optional, ≤60 chars, caption)

Layout: left-aligned text stack, hairline rules top + bottom, muted footer.
Feels like a board resolution page.
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = "closing_cta"
DESCRIPTION = "Restrained boardroom closing: left-aligned headline + action, hairline rules, muted footer."
ARCHETYPE = "boardroom"
MOOD = ["boardroom", "restrained"]
DENSITY = "sparse"
STYLE_TAGS = ["left_aligned", "hairline", "resolution"]

SLOTS = [
    {"name": "section_label", "kind": "text", "max_chars":  30, "style": "caption",  "required": False},
    {"name": "headline",      "kind": "text", "max_chars": 100, "style": "title_xl", "required": True},
    {"name": "cta",           "kind": "text", "max_chars":  80, "style": "subtitle", "required": True},
    {"name": "footer",        "kind": "text", "max_chars":  60, "style": "caption",  "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)
    # bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    # Top hairline
    add_hairline(slide, margin, 1.1, 13.333 - 2 * margin, theme, "border")

    # Section label
    section_label = get_slot(slots, "section_label") or "DECISION"
    sl = slide.shapes.add_textbox(Inches(margin), Inches(0.6),
                                   Inches(6.0), Inches(0.4))
    set_textbox_text(sl, section_label, theme, "caption",
                     "accent", align="left")
    sl.name = "section_label"
    add_theme_entrance(slide, sl, theme, delay_ms=120, index=0)

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 100)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(2.2),
                                   Inches(13.333 - 2 * margin), Inches(2.2))
    set_textbox_text(hl, headline, theme, "title_xl", "text", align="left")
    hl.name = "headline"
    add_theme_entrance(slide, hl, theme, delay_ms=250, index=1)

    # Accent rule above CTA
    add_solid_rect(slide, margin, 4.6, 0.8, 0.05, theme, color_key="accent")

    # CTA
    cta = truncate_to(get_slot(slots, "cta", required=True), 80)
    ct = slide.shapes.add_textbox(Inches(margin), Inches(4.85),
                                   Inches(13.333 - 2 * margin), Inches(0.9))
    set_textbox_text(ct, cta, theme, "subtitle", "accent", align="left")
    ct.name = "cta"
    add_theme_entrance(slide, ct, theme, delay_ms=420, index=2)

    # Bottom hairline
    add_hairline(slide, margin, 7.5 - margin - 0.2,
                 13.333 - 2 * margin, theme, "border")

    # Footer
    footer = get_slot(slots, "footer")
    if footer:
        fo = slide.shapes.add_textbox(Inches(margin),
                                       Inches(7.5 - margin),
                                       Inches(13.333 - 2 * margin),
                                       Inches(0.3))
        set_textbox_text(fo, footer, theme, "caption", "muted", align="left")
        fo.name = "footer"
        add_theme_entrance(slide, fo, theme, delay_ms=560, index=3)
