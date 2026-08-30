"""
skills_library/ppt/shells_seed/_shell_helpers.py

Shared helpers for every shell's render(slide, slots, theme) function.
Converts Theme tokens (logical names from the Theme JSON) into concrete
python-pptx calls. Shells MUST go through these helpers — never write
raw RGB / Pt values, because that breaks cross-slide theme consistency.
"""
from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches, Emu


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' to RGBColor."""
    s = hex_str.lstrip("#")
    r, g, b = int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    return RGBColor(r, g, b)


def palette_color(theme: dict, key: str, fallback: str = "text") -> RGBColor:
    """Look up a palette color by key; fall back to another palette key if missing."""
    pal = theme.get("palette", {})
    hex_str = pal.get(key) or pal.get(fallback) or "#000000"
    return hex_to_rgb(hex_str)


# ---------------------------------------------------------------------------
# Typography helpers
# ---------------------------------------------------------------------------


def apply_type_style(run, theme: dict, style_token: str, color_key: str = "text") -> None:
    """Apply a named typography style from the theme to a python-pptx run."""
    typo = theme.get("typography", {}).get(style_token)
    if not typo:
        # Fallback to body
        typo = theme.get("typography", {}).get("body", {
            "font": "Inter", "size": 14, "weight": 400,
        })
    run.font.name = typo.get("font", "Inter")
    run.font.size = Pt(typo.get("size", 14))
    weight = typo.get("weight", 400)
    run.font.bold = weight >= 600
    run.font.italic = typo.get("italic", False)
    run.font.color.rgb = palette_color(theme, color_key)
    # Note: tracking/upper are not directly supported in python-pptx run-level;
    # upper is handled by caller (uppercase the text string before assignment).


def should_uppercase(theme: dict, style_token: str) -> bool:
    """Return True if this style wants uppercase text."""
    typo = theme.get("typography", {}).get(style_token, {})
    return bool(typo.get("upper", False))


def set_textbox_text(textbox, text: str, theme: dict, style_token: str,
                     color_key: str = "text", align: str | None = None,
                     auto_fit: bool = True, min_size_pt: float = 9.0) -> None:
    """One-shot: set text on a textbox with a theme-styled run.

    If auto_fit is True (default), the font is shrunk so the text fits the
    textbox width × height. Floor at min_size_pt.
    """
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    if should_uppercase(theme, style_token):
        text = text.upper()
    tf = textbox.text_frame
    tf.word_wrap = True
    # Snug margins so we can use the full box width
    try:
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    if p.runs:
        apply_type_style(p.runs[0], theme, style_token, color_key)
        if auto_fit:
            _shrink_to_fit(textbox, p.runs[0], text, theme, style_token, min_size_pt)
    if align:
        align_map = {
            "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
        }
        if align in align_map:
            p.alignment = align_map[align]


def _shrink_to_fit(textbox, run, text: str, theme: dict, style_token: str,
                   min_size_pt: float = 9.0) -> None:
    """Heuristic: estimate text capacity for the box at the chosen style size,
    and shrink the run's font size if the text overflows.

    Avg glyph width ≈ 0.55 × font size in pt. Line height ≈ 1.2 × size.
    EMU/inch = 914400; 1 inch = 72 pt → 1 pt ≈ 12700 EMU.
    """
    try:
        w_pt = (textbox.width or 0) / 12700.0
        h_pt = (textbox.height or 0) / 12700.0
    except Exception:
        return
    if w_pt < 8 or h_pt < 8:
        return

    typo = theme.get("typography", {}).get(style_token, {}) or {}
    base_size = float(typo.get("size", 14))
    n = max(len(text), 1)

    # Try sizes from base down to min, pick the largest that fits.
    for size in (base_size, base_size * 0.9, base_size * 0.8, base_size * 0.7,
                 base_size * 0.6, base_size * 0.5, max(min_size_pt, base_size * 0.4)):
        if size < min_size_pt:
            size = min_size_pt
        chars_per_line = max(int(w_pt / (size * 0.55)), 4)
        lines = max(1, -(-n // chars_per_line))  # ceil div
        used_h = lines * size * 1.2
        if used_h <= h_pt:
            run.font.size = Pt(size)
            return
        if size <= min_size_pt:
            run.font.size = Pt(min_size_pt)
            return


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------


def add_solid_rect(slide, x: float, y: float, w: float, h: float,
                   theme: dict, color_key: str = "panel", line: bool = False,
                   shape_type=None):
    """Add a rectangle with a theme-palette fill. Returns the shape."""
    from pptx.enum.shapes import MSO_SHAPE

    shape_type = shape_type or MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = palette_color(theme, color_key)
    if not line:
        s.line.fill.background()
    else:
        s.line.color.rgb = palette_color(theme, "border", "muted")
        s.line.width = Pt(0.75)
    return s


def add_hairline(slide, x: float, y: float, w: float, theme: dict,
                 color_key: str = "border") -> None:
    """Add a thin horizontal rule — honors the 'thin_rule' motif."""
    from pptx.enum.shapes import MSO_SHAPE

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Emu(4572)  # ~0.005in
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = palette_color(theme, color_key, "muted")
    rule.line.fill.background()


# ---------------------------------------------------------------------------
# Slot validation
# ---------------------------------------------------------------------------


def get_slot(slots: dict, name: str, default: Any = None, required: bool = False) -> Any:
    """Fetch a slot value with required/default semantics."""
    if name in slots and slots[name] not in (None, ""):
        return slots[name]
    if required:
        raise ValueError(f"Required slot '{name}' not provided")
    return default


def truncate_to(text: str, max_chars: int | None) -> str:
    """Shrink overlong text to fit a fixed textbox.

    Truncation is a visual safety net, not a content decision. When the
    rewriter honors the slot's character budget no cut happens. When a long
    string arrives we cut at the LAST whitespace boundary inside the budget
    with NO ellipsis, since mid-word "…" reads as a bug to the viewer.
    As a last resort (single-word overflow) we hard-cut without ellipsis."""
    if not max_chars or len(text) <= max_chars:
        return text
    window = text[:max_chars]
    cut = window.rfind(" ")
    if cut >= max_chars // 2:
        return window[:cut].rstrip(" ,;:.-")
    return window.rstrip(" ,;:.-")


# ---------------------------------------------------------------------------
# Motion helpers (theme-aware entrance)
# ---------------------------------------------------------------------------


def add_theme_entrance(slide, shape, theme: dict, delay_ms: int = 0,
                        duration_ms: int = 500, index: int = 0) -> None:
    """
    Pick an entrance preset that matches the theme's motion.personality,
    and stagger across shapes using theme.motion.stagger_ms.
    """
    motion = theme.get("motion", {})
    pref = motion.get("preferred_entrance") or ["fade"]
    # Cycle through preferred entrances so a staggered group gets variety
    filter_name = pref[index % len(pref)]
    preset_map = {
        "fade": 10, "zoom": 53, "in_bottom": 2, "in_top": 2,
        "horizontal": 22, "vertical": 22, "in_left": 2, "in_right": 2, "in": 6,
        "wipe": 12, "circle": 6, "diamond": 16, "wedge": 19, "split": 13,
    }
    preset_id = preset_map.get(filter_name, 10)

    stagger = motion.get("stagger_ms", 120)
    total_delay = delay_ms + index * stagger
    _append_entrance_xml(slide, shape, total_delay, duration_ms, preset_id, filter_name)


def add_emphasis_pulse(slide, shape, theme: dict, delay_ms: int = 0,
                       duration_ms: int = 800) -> None:
    """Add an emphasis 'pulse' (grow/shrink) animation. Useful for callouts,
    KPI numbers, and CTA buttons. Triggers after entrances complete."""
    _append_emphasis_xml(slide, shape, delay_ms, duration_ms, preset_id=6, filter_name="grow_shrink")


def add_sequential_reveal(slide, shapes_in_order, theme: dict,
                          start_delay_ms: int = 0, step_ms: int = 200) -> None:
    """Reveal a sequence of shapes one after another with consistent stagger.
    Useful for lists, timelines, and feature grids where order matters."""
    for i, shape in enumerate(shapes_in_order):
        add_theme_entrance(slide, shape, theme,
                           delay_ms=start_delay_ms + i * step_ms,
                           duration_ms=400, index=i)


# ---------------------------------------------------------------------------
# Motion primitives — for "hero moments" where animation IS the content
# ---------------------------------------------------------------------------


def add_counter_up(slide, shapes_sequence, theme: dict, delay_ms: int = 0,
                   total_duration_ms: int = 1500) -> None:
    """Counter-up effect: a sequence of text shapes (e.g., ['0', '37', '74',
    '111', '148M']) at the SAME position. Each is initially hidden; they
    flip through quickly so the final value appears to count up.

    Caller must place the shapes in the same XY beforehand. This injects:
      - hide all but the last initially (style.visibility=hidden via set_pres effect)
      - on each step: hide previous + show current
    """
    if not shapes_sequence:
        return
    n = len(shapes_sequence)
    step_ms = max(80, total_duration_ms // n)
    for i, shape in enumerate(shapes_sequence):
        # Each step: instant appear (very short fade) at delay+i*step
        _append_appear_xml(slide, shape, delay_ms + i * step_ms, duration_ms=20)
        if i < n - 1:
            # Hide it again right before the next step so only the latest is visible
            _append_disappear_xml(slide, shape, delay_ms + (i + 1) * step_ms - 10)


def add_line_grow(slide, line_shape, theme: dict, delay_ms: int = 0,
                  duration_ms: int = 900, axis: str = "x") -> None:
    """Grow a horizontal (axis='x') or vertical (axis='y') line/rectangle from 0
    to 100% along the chosen axis. The shape's anchor stays at its (left, top).

    Use for: chart axes drawing in, progress bars filling, timeline rules
    extending across the slide.
    """
    _append_scale_xml(slide, line_shape, delay_ms, duration_ms,
                      from_pct=(0 if axis == "x" else 100000),
                      to_pct=100000, axis=axis)


def add_pivot_swap(slide, shape_out, shape_in, theme: dict, delay_ms: int = 0,
                   duration_ms: int = 900) -> None:
    """Pivot/rotate-swap: shape_out spins out (0→90°, fades), then shape_in
    spins in (-90°→0°, appears). The two shapes should occupy the same
    (or overlapping) bounding box for a proper card-flip feel.

    Use for: before/after comparisons, character/profile swaps, stage
    transitions on the same anchor.
    """
    _append_rotate_xml(slide, shape_out, delay_ms, duration_ms // 2,
                       from_deg=0, to_deg=90)
    _append_disappear_xml(slide, shape_out, delay_ms + duration_ms // 2)
    _append_appear_xml(slide, shape_in, delay_ms + duration_ms // 2, duration_ms=20)
    _append_rotate_xml(slide, shape_in, delay_ms + duration_ms // 2, duration_ms // 2,
                       from_deg=-90, to_deg=0)


def add_mask_wipe(slide, shape, theme: dict, delay_ms: int = 0,
                  duration_ms: int = 700, direction: str = "left_to_right") -> None:
    """Wipe-in entrance with explicit direction. PowerPoint preset 12 = wipe.
    Subtype controls direction:  4=left, 8=top, 2=right, 1=bottom.

    Use for: revealing text behind a moving mask (cinematic title), exposing
    a chart by sweeping across, hero word-by-word headlines.
    """
    direction_map = {"left_to_right": 4, "top_to_bottom": 8,
                     "right_to_left": 2, "bottom_to_top": 1}
    subtype = direction_map.get(direction, 4)
    _append_entrance_xml(slide, shape, delay_ms, duration_ms,
                         preset_id=12, filter_name="wipe", preset_subtype=subtype)


def add_zoom_climax(slide, shape, theme: dict, delay_ms: int = 0,
                    duration_ms: int = 700) -> None:
    """Big-zoom entrance for the single most important shape on the slide
    (e.g., the headline number, the hero quote). Larger duration + zoom
    preset gives a real ‘drop in' moment.
    """
    _append_entrance_xml(slide, shape, delay_ms, duration_ms,
                         preset_id=53, filter_name="zoom")


# ---------------------------------------------------------------------------
# Low-level XML injection helpers
# ---------------------------------------------------------------------------


def _bootstrap_timing(slide):
    """Ensure the slide has a <p:timing> mainSeq scaffold; return (timing, group)."""
    from pptx.oxml.ns import qn
    timing = slide._element.find(qn("p:timing"))
    if timing is None:
        # Use the entrance helper to bootstrap
        # (it parses + appends the full timing block on first invocation)
        return None, None
    seq = timing.find(".//" + qn("p:seq"))
    if seq is None:
        return timing, None
    groups = seq.findall(".//" + qn("p:childTnLst"))
    return timing, (groups[-1] if groups else None)


def _append_entrance_xml(slide, shape, delay_ms: int, duration_ms: int,
                          preset_id: int, filter_name: str,
                          preset_subtype: int = 0) -> None:
    """Inject an entrance animation into the slide's auto-play par."""
    from pptx.oxml import parse_xml

    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    group = _ensure_timing_group(slide)
    if group is None:
        return
    sid = shape.shape_id
    uid = 1000 + int(delay_ms) + sid * 7
    xml = (
        f'<p:par xmlns:p="{P_NS}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:cTn id="{uid}" presetID="{preset_id}" presetClass="entr" presetSubtype="{int(preset_subtype)}" '
        f'fill="hold" grpId="0" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="{int(delay_ms)}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:set><p:cBhvr><p:cTn id="{uid+1}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'<p:animEffect transition="in" filter="{filter_name}">'
        f'<p:cBhvr><p:cTn id="{uid+2}" dur="{int(duration_ms)}"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'</p:cBhvr></p:animEffect>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))


def _append_emphasis_xml(slide, shape, delay_ms: int, duration_ms: int,
                          preset_id: int, filter_name: str) -> None:
    """Inject an emphasis animation into the slide's auto-play par."""
    from pptx.oxml import parse_xml

    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    group = _ensure_timing_group(slide)
    if group is None:
        return
    sid = shape.shape_id
    uid = 5000 + int(delay_ms) + sid * 11
    xml = (
        f'<p:par xmlns:p="{P_NS}">'
        f'<p:cTn id="{uid}" presetID="{preset_id}" presetClass="emph" presetSubtype="0" '
        f'fill="hold" grpId="1" nodeType="afterEffect">'
        f'<p:stCondLst><p:cond delay="{int(delay_ms)}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:animScale><p:cBhvr><p:cTn id="{uid+1}" dur="{int(duration_ms)}" autoRev="1"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr>'
        f'<p:by x="110000" y="110000"/></p:animScale>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))


# ---------------------------------------------------------------------------
# Ambient motion primitives (indefinite-loop)
# ---------------------------------------------------------------------------
#
# These helpers emit OOXML that plays continuously while the slide is shown,
# modeled on the patterns used in reference_decks/pm_cogs.pptx (rotating cog
# mechanism) and reference_decks/pm_planets.pptx (orbital motion). Each
# primitive degrades gracefully when the shape argument lacks required
# geometry: the helper returns without raising so prs.save() always succeeds.


def add_infinite_rotation(slide, shape, duration_ms: int = 6000,
                          direction: str = "cw") -> None:
    """Rotate ``shape`` around its center continuously while the slide plays.

    Emits an ``<p:animRot>`` with ``repeatCount="indefinite"`` (the reference
    pattern from ``pm_cogs.pptx``). Direction is ``cw`` or ``ccw``. The
    helper is a no-op if the shape lacks a valid shape id.

    Impress quirk: a positive-``by`` ``animRot`` with indefinite loop plays
    as a static element in the current LibreOffice Impress build (verified
    by the feasibility spike). We therefore emit a negative-``by`` rotation
    regardless of ``direction`` — PowerPoint respects the sign, Impress
    plays both as motion, and the visual difference between cw and ccw is
    marginal on rotationally symmetric shapes (which is the primary use
    case for this primitive: cogs, rings, haloes).
    """
    if getattr(shape, "shape_id", None) is None:
        return
    _append_ambient_rotate_xml(slide, shape, duration_ms, sign=-1)


def add_orbital_motion(slide, shape, center_xy: tuple[float, float],
                       radius_in: float, duration_ms: int = 8000,
                       direction: str = "cw") -> None:
    """Move ``shape`` along a circular path around ``center_xy`` forever.

    Parameters
    ----------
    center_xy
        ``(cx, cy)`` in inches — the orbit center relative to the slide.
    radius_in
        Orbit radius in inches.

    Implementation note: PowerPoint's ``animMotion`` uses a normalised SVG
    path where (0,0) is the current shape position, and the viewport is the
    slide (1.0 = full width/height). A full circle path of radius ``r`` is
    described with two elliptical arcs. Coordinates are relative to the
    slide dimensions so the caller supplies inches and the helper converts.
    """
    if getattr(shape, "shape_id", None) is None or radius_in <= 0:
        return
    try:
        slide_w_in = slide.part.package.presentation_part.presentation.slide_width / 914400.0
        slide_h_in = slide.part.package.presentation_part.presentation.slide_height / 914400.0
    except Exception:
        slide_w_in, slide_h_in = 13.333, 7.5
    rx = radius_in / max(slide_w_in, 0.1)
    ry = radius_in / max(slide_h_in, 0.1)
    # Origin of the motion path = shape's current center relative to slide
    try:
        shape_cx_in = (shape.left + shape.width / 2) / 914400.0
        shape_cy_in = (shape.top + shape.height / 2) / 914400.0
    except Exception:
        shape_cx_in, shape_cy_in = 0.0, 0.0
    center_cx_in, center_cy_in = center_xy
    # Path starts at (center + (rx, 0)) relative to current shape center.
    # Described as two arcs so we trace a full circle.
    sign = 1 if direction == "cw" else -1
    start_dx = (center_cx_in + radius_in - shape_cx_in) / max(slide_w_in, 0.1)
    start_dy = (center_cy_in - shape_cy_in) / max(slide_h_in, 0.1)
    path = (
        f"M {start_dx:.5f} {start_dy:.5f} "
        f"a {rx:.5f} {ry:.5f} 0 1 {1 if sign == 1 else 0} {-2*rx:.5f} 0 "
        f"a {rx:.5f} {ry:.5f} 0 1 {1 if sign == 1 else 0} {2*rx:.5f} 0"
    )
    _append_ambient_motion_xml(slide, shape, duration_ms, path)


def add_pulse_loop(slide, shape, duration_ms: int = 2000,
                   scale_pct: int = 110) -> None:
    """Scale ``shape`` between 100% and ``scale_pct`` forever (auto-reverse).

    Uses ``<p:animScale>`` with ``autoRev="1"`` and ``repeatCount="indefinite"``.
    Good for hero metric rings or call-to-action badges.
    """
    if getattr(shape, "shape_id", None) is None:
        return
    _append_ambient_scale_xml(slide, shape, duration_ms, scale_pct * 1000)


def add_drift_motion(slide, shape, dx_in: float = 0.0, dy_in: float = 0.3,
                     duration_ms: int = 5000, pingpong: bool = True) -> None:
    """Translate ``shape`` by ``(dx_in, dy_in)`` and back forever.

    Uses ``<p:animMotion>`` with a short linear path and ``autoRev="1"`` to
    create a gentle back-and-forth drift. Setting ``pingpong=False`` makes
    the path one-way and then jumps back (less smooth).
    """
    if getattr(shape, "shape_id", None) is None:
        return
    try:
        slide_w_in = slide.part.package.presentation_part.presentation.slide_width / 914400.0
        slide_h_in = slide.part.package.presentation_part.presentation.slide_height / 914400.0
    except Exception:
        slide_w_in, slide_h_in = 13.333, 7.5
    dx_rel = dx_in / max(slide_w_in, 0.1)
    dy_rel = dy_in / max(slide_h_in, 0.1)
    path = f"M 0 0 L {dx_rel:.5f} {dy_rel:.5f}"
    _append_ambient_motion_xml(slide, shape, duration_ms, path,
                               auto_reverse=pingpong)


# ---------------------------------------------------------------------------
# Morph anchor emission (PowerPoint-compatibility deliverable)
# ---------------------------------------------------------------------------


MORPH_ANCHOR_ROLES = frozenset({
    "brand_mark",
    "accent_orb",
    "hero_number",
    "hero_headline",
    "section_chip",
})

_MORPH_ANCHOR_PREFIX = "!!sameName"


def set_morph_anchor(shape, role: str, theme: dict | None = None) -> str:
    """Tag ``shape`` so PowerPoint force-matches it across adjacent slides.

    Assigns ``<p:cNvPr name="!!SameName!!<role>"/>`` per the naming rule in
    ``docs/ppt_morph_continuity_contract.md``. Enforces the per-slide
    uniqueness rule: two shapes on the same slide with the same role raise
    ``ValueError``.

    Parameters
    ----------
    shape
        A python-pptx shape (any type whose ``nvSpPr/cNvPr`` or equivalent
        non-visual property element exists).
    role
        One of ``MORPH_ANCHOR_ROLES``.
    theme
        Reserved; accepted for forward compatibility (e.g. theme-aware role
        filtering). Currently unused.

    Returns
    -------
    The assigned ``cNvPr/@name`` value.
    """
    if role not in MORPH_ANCHOR_ROLES:
        raise ValueError(
            f"role '{role}' is not in the morph anchor vocabulary "
            f"{sorted(MORPH_ANCHOR_ROLES)}"
        )

    cnv_pr = _find_cnv_pr(shape)
    if cnv_pr is None:
        raise ValueError(
            "shape has no <p:cNvPr> or <p:nvSpPr> element; cannot tag it"
        )

    anchor_name = f"{_MORPH_ANCHOR_PREFIX}{role}"
    parent_slide_elem = _ascend_to_slide(cnv_pr)
    if parent_slide_elem is not None:
        for existing in parent_slide_elem.iter():
            tag = existing.tag
            if not tag.endswith("}cNvPr"):
                continue
            if existing is cnv_pr:
                continue
            if existing.get("name") == anchor_name:
                raise ValueError(
                    f"slide already has a shape tagged with role '{role}'; "
                    "the morph contract mandates one anchor per role per slide"
                )

    cnv_pr.set("name", anchor_name)
    return anchor_name


def _find_cnv_pr(shape):
    """Return the shape's ``<p:cNvPr>`` element, or None if it can't be located."""
    try:
        elem = shape._element
    except AttributeError:
        elem = getattr(shape, "element", None)
    if elem is None:
        return None
    for child in elem.iter():
        if child.tag.endswith("}cNvPr"):
            return child
    return None


def _ascend_to_slide(elem):
    """Walk up from ``elem`` to the owning ``<p:sld>`` element, or None."""
    node = elem
    while node is not None:
        if node.tag.endswith("}sld"):
            return node
        node = node.getparent() if hasattr(node, "getparent") else None
    return None


# ---------------------------------------------------------------------------
# Low-level XML for ambient (indefinite-loop) effects
# ---------------------------------------------------------------------------


def _append_ambient_rotate_xml(slide, shape, duration_ms: int, sign: int) -> None:
    from pptx.oxml import parse_xml
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    group = _ensure_timing_group(slide)
    if group is None:
        return
    sid = shape.shape_id
    uid = 11000 + int(duration_ms) + sid * 29
    from_deg = 0
    to_deg = 360 * sign
    # autoRev="1" is required for LibreOffice Impress to actually animate the
    # loop; a bare repeatCount="indefinite" animRot plays as a static element
    # in Impress today (feasibility spike confirmed this for the primary
    # viewer of record). autoRev is harmless on a 360-degree rotation because
    # the reverse half-cycle retraces the same visual trajectory.
    xml = (
        f'<p:par xmlns:p="{P_NS}">'
        f'<p:cTn id="{uid}" fill="hold" grpId="5" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:animRot by="{int((to_deg - from_deg) * 60000)}">'
        f'<p:cBhvr><p:cTn id="{uid+1}" dur="{int(duration_ms)}" '
        f'repeatCount="indefinite" autoRev="1" fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr>'
        f'</p:animRot>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))


def _append_ambient_motion_xml(slide, shape, duration_ms: int, path: str,
                               auto_reverse: bool = False) -> None:
    from pptx.oxml import parse_xml
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    group = _ensure_timing_group(slide)
    if group is None:
        return
    sid = shape.shape_id
    uid = 12000 + int(duration_ms) + sid * 31
    auto_rev_attr = ' autoRev="1"' if auto_reverse else ""
    xml = (
        f'<p:par xmlns:p="{P_NS}">'
        f'<p:cTn id="{uid}" fill="hold" grpId="6" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:animMotion origin="layout" path="{path}" pathEditMode="relative">'
        f'<p:cBhvr><p:cTn id="{uid+1}" dur="{int(duration_ms)}" '
        f'repeatCount="indefinite"{auto_rev_attr} fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr>'
        f'</p:animMotion>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))


def _append_ambient_scale_xml(slide, shape, duration_ms: int,
                              scale_val: int) -> None:
    from pptx.oxml import parse_xml
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    group = _ensure_timing_group(slide)
    if group is None:
        return
    sid = shape.shape_id
    uid = 13000 + int(duration_ms) + sid * 37
    xml = (
        f'<p:par xmlns:p="{P_NS}">'
        f'<p:cTn id="{uid}" fill="hold" grpId="7" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:animScale>'
        f'<p:cBhvr><p:cTn id="{uid+1}" dur="{int(duration_ms)}" '
        f'repeatCount="indefinite" autoRev="1" fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr>'
        f'<p:from x="100000" y="100000"/>'
        f'<p:to x="{int(scale_val)}" y="{int(scale_val)}"/>'
        f'</p:animScale>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))


# ---------------------------------------------------------------------------
# More low-level XML for the motion primitives
# ---------------------------------------------------------------------------


def _ensure_timing_group(slide):
    """Return the withGroup childTnLst (auto-start animation bucket).

    Matches the structure MS Office writes in decks that play on slide
    load — see reference_decks/pm_cogs.pptx (verified in PowerPoint):

      tmRoot -> seq mainSeq -> par clickPar
          stCondLst: cond delay=indefinite + cond evt=onBegin tn=2
          par withGroup
              stCondLst cond delay=0
              childTnLst  <- emissions land here
    """
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    timing = slide._element.find(qn("p:timing"))
    if timing is None:
        timing = parse_xml(
            f'<p:timing xmlns:p="{P_NS}">'
            '<p:tnLst><p:par>'
            '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            '<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
            '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            '<p:par><p:cTn id="3" fill="hold" nodeType="clickPar">'
            '<p:stCondLst>'
            '<p:cond delay="indefinite"/>'
            '<p:cond evt="onBegin" delay="0"><p:tn val="2"/></p:cond>'
            '</p:stCondLst>'
            '<p:childTnLst><p:par><p:cTn id="4" fill="hold" nodeType="withGroup">'
            '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            '<p:childTnLst/>'
            '</p:cTn></p:par></p:childTnLst>'
            '</p:cTn></p:par></p:childTnLst></p:cTn>'
            '<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            '<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            '</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
        )
        slide._element.append(timing)
    else:
        # Retrofit: add the auto-start cond + withGroup par if missing.
        click_par_cTn = None
        withgroup_children = None
        for cTn in timing.iter(qn("p:cTn")):
            nt = cTn.get("nodeType")
            if nt == "clickPar":
                click_par_cTn = cTn
            elif nt == "withGroup":
                withgroup_children = cTn.find(qn("p:childTnLst"))
        if click_par_cTn is not None:
            stCondLst = click_par_cTn.find(qn("p:stCondLst"))
            has_auto = False
            if stCondLst is not None:
                for cond in stCondLst.findall(qn("p:cond")):
                    if cond.get("evt") == "onBegin":
                        has_auto = True
                        break
                if not has_auto:
                    stCondLst.append(parse_xml(
                        f'<p:cond xmlns:p="{P_NS}" evt="onBegin" delay="0">'
                        '<p:tn val="2"/></p:cond>'
                    ))
        if withgroup_children is None:
            mainSeq = None
            for cTn in timing.iter(qn("p:cTn")):
                if cTn.get("nodeType") == "mainSeq":
                    mainSeq = cTn
                    break
            if mainSeq is not None:
                ms_children = mainSeq.find(qn("p:childTnLst"))
                if ms_children is not None:
                    ms_children.append(parse_xml(
                        f'<p:par xmlns:p="{P_NS}">'
                        '<p:cTn id="3" fill="hold" nodeType="clickPar">'
                        '<p:stCondLst>'
                        '<p:cond delay="indefinite"/>'
                        '<p:cond evt="onBegin" delay="0"><p:tn val="2"/></p:cond>'
                        '</p:stCondLst>'
                        '<p:childTnLst><p:par><p:cTn id="4" fill="hold" nodeType="withGroup">'
                        '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
                        '<p:childTnLst/>'
                        '</p:cTn></p:par></p:childTnLst>'
                        '</p:cTn></p:par>'
                    ))
    for cTn in timing.iter(qn("p:cTn")):
        if cTn.get("nodeType") == "withGroup":
            return cTn.find(qn("p:childTnLst"))
    return None


def _append_appear_xml(slide, shape, delay_ms: int, duration_ms: int = 50) -> None:
    """Make a shape become visible at delay_ms (instant appear)."""
    from pptx.oxml import parse_xml
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    group = _ensure_timing_group(slide)
    if group is None:
        return
    sid = shape.shape_id
    uid = 8000 + int(delay_ms) + sid * 13
    xml = (
        f'<p:par xmlns:p="{P_NS}">'
        f'<p:cTn id="{uid}" presetID="1" presetClass="entr" presetSubtype="0" '
        f'fill="hold" grpId="0" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="{int(delay_ms)}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:set><p:cBhvr><p:cTn id="{uid+1}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))


def _append_disappear_xml(slide, shape, delay_ms: int) -> None:
    """Hide a shape at delay_ms (instant)."""
    from pptx.oxml import parse_xml
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    group = _ensure_timing_group(slide)
    if group is None:
        return
    sid = shape.shape_id
    uid = 9000 + int(delay_ms) + sid * 17
    xml = (
        f'<p:par xmlns:p="{P_NS}">'
        f'<p:cTn id="{uid}" presetID="1" presetClass="exit" presetSubtype="0" '
        f'fill="hold" grpId="2" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="{int(delay_ms)}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:set><p:cBhvr><p:cTn id="{uid+1}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="hidden"/></p:to></p:set>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))


def _append_scale_xml(slide, shape, delay_ms: int, duration_ms: int,
                      from_pct: int, to_pct: int, axis: str = "x") -> None:
    """Animate the shape's scale on x or y axis. Values in 1/1000 of percent
    (100000 = 100%). Use for line-grow / progress bar fills."""
    from pptx.oxml import parse_xml
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    group = _ensure_timing_group(slide)
    if group is None:
        return
    sid = shape.shape_id
    uid = 6000 + int(delay_ms) + sid * 19
    other_axis = "y" if axis == "x" else "x"
    xml = (
        f'<p:par xmlns:p="{P_NS}">'
        f'<p:cTn id="{uid}" fill="hold" grpId="3" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="{int(delay_ms)}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:animScale><p:cBhvr><p:cTn id="{uid+1}" dur="{int(duration_ms)}" fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr>'
        f'<p:from {axis}="{int(from_pct)}" {other_axis}="100000"/>'
        f'<p:to {axis}="{int(to_pct)}" {other_axis}="100000"/>'
        f'</p:animScale>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))


def _append_rotate_xml(slide, shape, delay_ms: int, duration_ms: int,
                       from_deg: int, to_deg: int) -> None:
    """Rotate a shape from from_deg to to_deg. PPT uses 60000ths of a degree."""
    from pptx.oxml import parse_xml
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    group = _ensure_timing_group(slide)
    if group is None:
        return
    sid = shape.shape_id
    uid = 7000 + int(delay_ms) + sid * 23
    xml = (
        f'<p:par xmlns:p="{P_NS}">'
        f'<p:cTn id="{uid}" fill="hold" grpId="4" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="{int(delay_ms)}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:animRot><p:cBhvr><p:cTn id="{uid+1}" dur="{int(duration_ms)}" fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr>'
        f'<p:from val="{int(from_deg) * 60000}"/>'
        f'<p:to val="{int(to_deg) * 60000}"/>'
        f'</p:animRot>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))


__all__ = [
    "hex_to_rgb", "palette_color",
    "apply_type_style", "should_uppercase", "set_textbox_text",
    "add_solid_rect", "add_hairline",
    "get_slot", "truncate_to",
    "add_theme_entrance", "add_emphasis_pulse", "add_sequential_reveal",
    "add_counter_up", "add_line_grow", "add_pivot_swap",
    "add_mask_wipe", "add_zoom_climax",
    "add_infinite_rotation", "add_orbital_motion", "add_pulse_loop",
    "add_drift_motion",
    "set_morph_anchor", "MORPH_ANCHOR_ROLES",
]
