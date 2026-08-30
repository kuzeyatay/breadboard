"""
shells_seed/bullet_card_list.py — 3-column bullet card list.

Slots:
  section_label (text, optional, caption style)            — "Problem", "Key points", etc.
  headline      (text, required, ≤80 chars, title style)   — slide title
  bullets       (bullet_list, required, 3 items exactly)   — each item: {title, body}
    where each item is a dict {"title": str, "body": str}

Layout: header band on top 1.5in, 3 equal-width cards below sharing 1 row.
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'bullet_card_list'
DESCRIPTION = 'Four-card grid of short bullets with a section_label + headline stack.'
ARCHETYPE = 'narrative'
MOOD = ['editorial', 'boardroom']
DENSITY = 'balanced'
STYLE_TAGS = ['grid', 'card', 'two_by_two']

SLOTS = [
    {"name": "section_label", "kind": "text",        "max_chars": 30, "style": "caption",  "required": False},
    {"name": "headline",      "kind": "text",        "max_chars": 80, "style": "title",    "required": True},
    {"name": "bullets",       "kind": "bullet_list", "bullet_capacity": 3,                   "required": True},
]


def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    gutter = theme.get("spacing", {}).get("gutter", 0.3)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    y = margin
    reveal = []

    # Section label (kicker)
    label = get_slot(slots, "section_label")
    if label:
        lb = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(6), Inches(0.35))
        set_textbox_text(lb, truncate_to(label, 30), theme, "caption", color_key="accent")
        reveal.append(lb)
        y += 0.55
        if theme.get("motif", {}).get("type") == "thin_rule":
            add_hairline(slide, margin, y, 1.8, theme)
            y += 0.2

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(1.2))
    set_textbox_text(head, headline, theme, "title", color_key="text")
    reveal.append(head)
    y += 1.5

    # 3 cards
    bullets = get_slot(slots, "bullets", required=True)
    if not isinstance(bullets, list) or len(bullets) != 3:
        raise ValueError(f"'bullets' slot must be a list of 3 items, got {bullets!r}")

    card_h = 3.4
    total_w = 13.333 - 2 * margin
    card_w = (total_w - 2 * gutter) / 3
    for i, item in enumerate(bullets):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
        x = margin + i * (card_w + gutter)

        # Card panel
        card = add_solid_rect(
            slide, x, y, card_w, card_h, theme,
            color_key="panel", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        card.adjustments[0] = 0.05

        # Accent bar at top of card
        bar = add_solid_rect(
            slide, x + 0.25, y + 0.25, 0.6, 0.08, theme,
            color_key="accent", line=False,
        )

        # Card title
        t_box = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(y + 0.55), Inches(card_w - 0.5), Inches(0.8),
        )
        set_textbox_text(
            t_box, truncate_to(item.get("title", ""), 40), theme, "body_bold", color_key="text",
        )

        # Card body
        b_box = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(y + 1.35), Inches(card_w - 0.5), Inches(card_h - 1.6),
        )
        set_textbox_text(
            b_box, truncate_to(item.get("body", ""), 160), theme, "body", color_key="muted",
        )

        reveal.extend([card, bar, t_box, b_box])

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=80, duration_ms=500, index=i)
