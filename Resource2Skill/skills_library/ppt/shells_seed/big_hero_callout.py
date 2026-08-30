"""
shells_seed/big_hero_callout.py — Full-bleed accent background with giant headline callout + URL.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'closing_cta'
DESCRIPTION = 'Full-bleed accent background with giant headline callout + URL.'
ARCHETYPE = 'brand'
MOOD = ['bold', 'cinematic', 'punchy']
DENSITY = 'sparse'
STYLE_TAGS = ['full_bleed', 'hero_callout', 'big_text', 'cta']

SLOTS = [{'name': 'eyebrow', 'kind': 'text', 'max_chars': 30, 'style': 'caption', 'required': False}, {'name': 'headline', 'kind': 'text', 'max_chars': 90, 'style': 'title_xl', 'required': True}, {'name': 'cta', 'kind': 'text', 'max_chars': 60, 'style': 'subtitle', 'required': True}, {'name': 'url', 'kind': 'text', 'max_chars': 60, 'style': 'caption', 'required': False}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    # Full-bleed accent
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "accent")
    bg.line.fill.background(); bg.shadow.inherit = False
    bg.name = "hero_bg"

    # Decorative oval
    halo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1.5),
                                   Inches(4.5), Inches(4.5))
    halo.fill.solid(); halo.fill.fore_color.rgb = palette_color(theme, "accent2")
    halo.line.fill.background(); halo.shadow.inherit = False
    halo.name = "accent_orb"

    eyebrow = get_slot(slots, "eyebrow") or ""
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(margin), Inches(margin + 0.2),
                                       Inches(13.333 - 2*margin), Inches(0.5))
        set_textbox_text(eb, eyebrow, theme, "caption", "bg", align="left")
        eb.name = "eyebrow"

    headline = truncate_to(get_slot(slots, "headline", required=True), 90)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(2.1),
                                   Inches(13.333 - 2*margin), Inches(2.8))
    set_textbox_text(hl, headline, theme, "title_xl", "bg", align="left")
    hl.name = "hero_headline"
    add_theme_entrance(slide, hl, theme, delay_ms=180, index=0)

    # CTA pill
    cta = truncate_to(get_slot(slots, "cta", required=True), 60)
    ctb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(margin), Inches(5.2),
                                  Inches(4.0), Inches(0.7))
    ctb.fill.solid(); ctb.fill.fore_color.rgb = palette_color(theme, "bg")
    ctb.line.fill.background(); ctb.shadow.inherit = False
    ctb.name = "cta"
    ctt = slide.shapes.add_textbox(Inches(margin), Inches(5.25),
                                    Inches(4.0), Inches(0.6))
    set_textbox_text(ctt, cta, theme, "body_bold", "accent", align="center")
    ctt.name = "cta_text"
    add_theme_entrance(slide, ctb, theme, delay_ms=420, index=1)

    # URL
    url = get_slot(slots, "url") or ""
    if url:
        ub = slide.shapes.add_textbox(Inches(margin), Inches(6.2),
                                       Inches(13.333 - 2*margin), Inches(0.5))
        set_textbox_text(ub, url, theme, "caption", "bg", align="left")
        ub.name = "url"
