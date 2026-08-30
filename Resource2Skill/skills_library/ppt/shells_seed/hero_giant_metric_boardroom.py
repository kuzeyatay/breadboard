"""
shells_seed/hero_giant_metric_boardroom.py — restrained boardroom KPI hero.

Slots:
  eyebrow (text, optional, ≤30 chars, caption)
  value   (text, required, ≤12 chars, metric_xl)
  delta   (text, optional, ≤20 chars, body_bold — e.g. "+38% YoY")
  label   (text, required, ≤40 chars, subtitle)
  context (text, optional, ≤160 chars, body)

Layout: left 60% — value + delta stacked with hairline underneath, right
40% — context paragraph. Boardroom/financial feel.
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_theme_entrance, get_slot, palette_color,
    set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = "hero_giant_metric"
DESCRIPTION = "Boardroom KPI hero: left value+delta with hairline, right-side context paragraph."
ARCHETYPE = "boardroom"
MOOD = ["boardroom", "restrained", "technical"]
DENSITY = "balanced"
STYLE_TAGS = ["split", "giant_numeral", "kpi_delta"]

SLOTS = [
    {"name": "eyebrow", "kind": "text", "max_chars":  30, "style": "caption",   "required": False},
    {"name": "value",   "kind": "text", "max_chars":  12, "style": "metric_xl", "required": True},
    {"name": "delta",   "kind": "text", "max_chars":  20, "style": "body_bold", "required": False},
    {"name": "label",   "kind": "text", "max_chars":  40, "style": "subtitle",  "required": True},
    {"name": "context", "kind": "text", "max_chars": 160, "style": "body",      "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    # Eyebrow (top left)
    eyebrow = get_slot(slots, "eyebrow")
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(margin), Inches(margin),
                                       Inches(6), Inches(0.4))
        set_textbox_text(eb, eyebrow, theme, "caption", "accent", align="left")
        eb.name = "eyebrow"
        add_theme_entrance(slide, eb, theme, delay_ms=100, index=0)

    # Value (huge numeral)
    value = truncate_to(get_slot(slots, "value", required=True), 12)
    val = slide.shapes.add_textbox(Inches(margin), Inches(1.6),
                                    Inches(7.5), Inches(3.2))
    set_textbox_text(val, value, theme, "metric_xl", "text", align="left")
    val.name = "value"
    add_theme_entrance(slide, val, theme, delay_ms=250, index=1)

    # Delta chip
    delta = get_slot(slots, "delta")
    if delta:
        dl = slide.shapes.add_textbox(Inches(margin), Inches(4.85),
                                       Inches(3.5), Inches(0.5))
        set_textbox_text(dl, delta, theme, "body_bold", "success",
                         align="left")
        dl.name = "delta"
        add_theme_entrance(slide, dl, theme, delay_ms=420, index=2)

    # Hairline under value
    add_hairline(slide, margin, 5.5, 7.5, theme, "border")

    # Label
    label = truncate_to(get_slot(slots, "label", required=True), 40)
    lb = slide.shapes.add_textbox(Inches(margin), Inches(5.7),
                                   Inches(7.5), Inches(0.6))
    set_textbox_text(lb, label, theme, "subtitle", "muted", align="left")
    lb.name = "label"
    add_theme_entrance(slide, lb, theme, delay_ms=540, index=3)

    # Right-side context
    context = get_slot(slots, "context")
    if context:
        ctx = slide.shapes.add_textbox(Inches(8.4), Inches(2.0),
                                        Inches(4.3), Inches(4.0))
        set_textbox_text(ctx, context, theme, "body", "text", align="left")
        ctx.name = "context"
        add_theme_entrance(slide, ctx, theme, delay_ms=380, index=4)
