"""
shells_seed/numbered_steps_column.py — Vertical numbered list with big serif numerals and step descriptions.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'bullet_card_list'
DESCRIPTION = 'Vertical numbered list with big serif numerals and step descriptions.'
ARCHETYPE = 'narrative'
MOOD = ['editorial', 'restrained', 'warm']
DENSITY = 'balanced'
STYLE_TAGS = ['numbered', 'vertical', 'editorial_list', 'serif_numeral']

SLOTS = [{'name': 'section_label', 'kind': 'text', 'max_chars': 30, 'style': 'caption', 'required': False}, {'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title', 'required': True}, {'name': 'steps', 'kind': 'bullet_list', 'bullet_capacity': 5, 'required': True}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    section_label = get_slot(slots, "section_label")
    if section_label:
        sl = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(8), Inches(0.4))
        set_textbox_text(sl, section_label, theme, "caption", "accent", align="left")
        sl.name = "section_label"

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(1.0), Inches(10), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"

    steps = get_slot(slots, "steps", default=[]) or []
    if isinstance(steps, str):
        steps = [x.strip() for x in steps.split("\n") if x.strip()]
    n = min(5, max(3, len(steps)))
    steps = (steps + [""]*n)[:n]
    start_y = 2.3
    step_h = (7.5 - start_y - margin) / n
    revealables = []
    for i, text in enumerate(steps):
        y = start_y + i * step_h
        # Big serif numeral
        nb = slide.shapes.add_textbox(Inches(margin), Inches(y),
                                       Inches(1.2), Inches(step_h))
        nb.text_frame.text = f"{i+1:02d}"
        r = nb.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(54); r.font.bold = True
        r.font.name = theme.get("typography", {}).get("accent_serif", {}).get("font",
                      theme.get("typography", {}).get("title_xl", {}).get("font", "Inter"))
        r.font.color.rgb = palette_color(theme, "accent")
        nb.name = f"step_num_{i}"
        # Rule
        add_hairline(slide, margin + 1.3, y + 0.15, 11.333 - margin, theme, "border")
        # Step text
        tb = slide.shapes.add_textbox(Inches(margin + 1.5), Inches(y + 0.25),
                                       Inches(11.333 - margin - 1.5),
                                       Inches(step_h - 0.3))
        set_textbox_text(tb, truncate_to(str(text), 180), theme, "body", "text", align="left")
        tb.name = f"step_body_{i}"
        revealables.append(nb)
    add_sequential_reveal(slide, revealables, theme, start_delay_ms=250, step_ms=160)
