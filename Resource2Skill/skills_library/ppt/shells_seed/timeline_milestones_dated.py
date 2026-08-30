"""
shells_seed/timeline_milestones_dated.py — Horizontal dated timeline with 3-5 milestones, each with date/label/detail.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'timeline_horizontal'
DESCRIPTION = 'Horizontal dated timeline with 3-5 milestones, each with date/label/detail.'
ARCHETYPE = 'narrative'
MOOD = ['editorial', 'restrained', 'technical']
DENSITY = 'balanced'
STYLE_TAGS = ['timeline', 'milestones', 'dated', 'connector']

SLOTS = [{'name': 'section_label', 'kind': 'text', 'max_chars': 30, 'style': 'caption', 'required': False}, {'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title', 'required': True}, {'name': 'milestones', 'kind': 'bullet_list', 'bullet_capacity': 5, 'required': True}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    section_label = get_slot(slots, "section_label")
    if section_label:
        sl = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(8), Inches(0.4))
        set_textbox_text(sl, section_label, theme, "caption", "accent", align="left")
        sl.name = "section_label"

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(1.0), Inches(13.333 - 2*margin), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"
    add_theme_entrance(slide, hl, theme, delay_ms=100, index=0)

    milestones = get_slot(slots, "milestones", default=[]) or []
    if isinstance(milestones, str):
        milestones = [x.strip() for x in milestones.split("\n") if x.strip()]
    n = max(1, min(5, len(milestones)))
    milestones = milestones[:n]

    # Horizontal rule across
    rule_y = 3.8
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(margin), Inches(rule_y - 0.02),
                                   Inches(13.333 - 2*margin), Inches(0.04))
    rule.fill.solid(); rule.fill.fore_color.rgb = palette_color(theme, "border", "muted")
    rule.line.fill.background(); rule.name = "timeline_rule"

    segment_w = (13.333 - 2*margin) / max(n, 1)
    dots = []
    for i, entry in enumerate(milestones):
        if isinstance(entry, dict):
            date = entry.get("date", ""); label = entry.get("label", "")
            detail = entry.get("detail", "")
        else:
            parts = [p.strip() for p in str(entry).split("|")]
            parts = (parts + [""]*3)[:3]
            date, label, detail = parts
        cx = margin + segment_w * (i + 0.5)
        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.15),
                                      Inches(rule_y - 0.15), Inches(0.3), Inches(0.3))
        dot.fill.solid(); dot.fill.fore_color.rgb = palette_color(theme, "accent")
        dot.line.fill.background(); dot.shadow.inherit = False
        dot.name = f"milestone_dot_{i}"
        # Date above
        db = slide.shapes.add_textbox(Inches(cx - segment_w/2 + 0.1),
                                       Inches(rule_y - 1.1),
                                       Inches(segment_w - 0.2), Inches(0.4))
        set_textbox_text(db, date, theme, "caption", "accent", align="center")
        db.name = f"milestone_date_{i}"
        # Label
        lb = slide.shapes.add_textbox(Inches(cx - segment_w/2 + 0.1),
                                       Inches(rule_y + 0.3),
                                       Inches(segment_w - 0.2), Inches(0.5))
        set_textbox_text(lb, label, theme, "body_bold", "text", align="center")
        lb.name = f"milestone_label_{i}"
        # Detail
        if detail:
            dt = slide.shapes.add_textbox(Inches(cx - segment_w/2 + 0.1),
                                           Inches(rule_y + 0.95),
                                           Inches(segment_w - 0.2), Inches(1.6))
            set_textbox_text(dt, truncate_to(detail, 120), theme, "body", "muted", align="center")
            dt.name = f"milestone_detail_{i}"
        dots.append(dot)
    add_sequential_reveal(slide, dots, theme, start_delay_ms=360, step_ms=200)
