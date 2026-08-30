"""
shells_seed/dense_bullet_list.py — 8-10 bullet list, information-heavy.

Slots:
  section_label (text, optional, caption)
  headline      (text, required, title)
  items         (bullet_list, required, 6-10 short strings)
  footer        (text, optional, caption)

Layout: left-aligned headline block, right ~60% holds two columns of bullets.
Density=dense: designed to look like a proper briefing slide, not a hero.
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = "bullet_card_list"
DESCRIPTION = "Two-column dense bullet list with 6-10 items + headline + caption footer."
ARCHETYPE = "boardroom"
MOOD = ["boardroom", "restrained", "editorial"]
DENSITY = "dense"
STYLE_TAGS = ["two_column", "bullets", "dense_list"]

SLOTS = [
    {"name": "section_label", "kind": "text",        "max_chars": 30,  "style": "caption", "required": False},
    {"name": "headline",      "kind": "text",        "max_chars": 100, "style": "title",   "required": True},
    {"name": "items",         "kind": "bullet_list", "bullet_capacity": 10, "required": True},
    {"name": "footer",        "kind": "text",        "max_chars": 80,  "style": "caption", "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    # Section label
    section_label = get_slot(slots, "section_label")
    if section_label:
        sl = slide.shapes.add_textbox(Inches(margin), Inches(margin),
                                       Inches(7.0), Inches(0.4))
        set_textbox_text(sl, section_label, theme, "caption", "accent", align="left")
        sl.name = "section_label"

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 100)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(1.1),
                                   Inches(13.333 - 2 * margin), Inches(1.0))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"
    add_theme_entrance(slide, hl, theme, delay_ms=120, index=0)

    # Hairline under headline
    add_hairline(slide, margin, 2.25, 13.333 - 2 * margin, theme, "border")

    # Items — two-column layout
    items = get_slot(slots, "items", default=[]) or []
    if isinstance(items, str):
        items = [x.strip() for x in items.split("\n") if x.strip()]
    items = items[:10]
    n = len(items)
    if n == 0:
        return
    # Split into two columns
    left_col = items[: (n + 1) // 2]
    right_col = items[(n + 1) // 2:]
    col_w = (13.333 - 2 * margin - 0.4) / 2
    col_y = 2.5
    col_h = 4.4
    item_h = col_h / max(len(left_col), 1)

    reveals = []
    for i, text in enumerate(left_col):
        y = col_y + i * item_h
        # bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(margin),
                                      Inches(y + 0.1), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = palette_color(theme, "accent")
        dot.line.fill.background(); dot.shadow.inherit = False
        # text
        tb = slide.shapes.add_textbox(Inches(margin + 0.25), Inches(y),
                                       Inches(col_w - 0.25),
                                       Inches(item_h - 0.05))
        set_textbox_text(tb, str(text), theme, "body", "text", align="left")
        tb.name = f"item_left_{i}"
        reveals.append(tb)

    for i, text in enumerate(right_col):
        y = col_y + i * item_h
        x0 = margin + col_w + 0.4
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x0),
                                      Inches(y + 0.1), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = palette_color(theme, "accent2")
        dot.line.fill.background(); dot.shadow.inherit = False
        tb = slide.shapes.add_textbox(Inches(x0 + 0.25), Inches(y),
                                       Inches(col_w - 0.25),
                                       Inches(item_h - 0.05))
        set_textbox_text(tb, str(text), theme, "body", "text", align="left")
        tb.name = f"item_right_{i}"
        reveals.append(tb)

    add_sequential_reveal(slide, reveals, theme, start_delay_ms=300, step_ms=90)

    # Footer
    footer = get_slot(slots, "footer")
    if footer:
        fo = slide.shapes.add_textbox(Inches(margin), Inches(7.5 - margin - 0.2),
                                       Inches(13.333 - 2 * margin), Inches(0.3))
        set_textbox_text(fo, footer, theme, "caption", "muted", align="left")
        fo.name = "footer"
