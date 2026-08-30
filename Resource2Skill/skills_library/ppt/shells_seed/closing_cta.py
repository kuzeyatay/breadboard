"""
shells_seed/closing_cta.py — Closing slide with CTA.

Slots:
  eyebrow       (text, optional)                 — "Thank you" / "Let's talk"
  headline      (text, required, ≤80 chars)      — main closing message
  cta_text      (text, optional, ≤40 chars)      — button / CTA label
  contact_lines (bullet_list, optional, max 3)   — email / url / phone lines

Layout: centered stack, restrained. Single focal point.
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'closing_cta'
DESCRIPTION = 'Centered closing slide with big headline and single CTA chip / URL.'
ARCHETYPE = 'product'
MOOD = ['bold', 'punchy']
DENSITY = 'sparse'
STYLE_TAGS = ['centered', 'cta', 'big_text']

SLOTS = [
    {"name": "eyebrow",       "kind": "text",        "max_chars": 30, "style": "caption", "required": False},
    {"name": "headline",      "kind": "text",        "max_chars": 80, "style": "title_xl","required": True},
    {"name": "cta_text",      "kind": "text",        "max_chars": 40, "style": "body_bold","required": False},
    {"name": "contact_lines", "kind": "bullet_list", "bullet_capacity": 3, "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    center_x = 13.333 / 2
    content_w = 9.0
    content_x = center_x - content_w / 2

    y = 2.3

    # Eyebrow
    eyebrow = get_slot(slots, "eyebrow")
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(content_x), Inches(y), Inches(content_w), Inches(0.4))
        set_textbox_text(eb, truncate_to(eyebrow, 30), theme, "caption", color_key="accent", align="center")
        reveal.append(eb)
        y += 0.6
        if theme.get("motif", {}).get("type") == "thin_rule":
            add_hairline(slide, center_x - 0.6, y, 1.2, theme)
            y += 0.2

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hb = slide.shapes.add_textbox(Inches(content_x), Inches(y), Inches(content_w), Inches(1.6))
    set_textbox_text(hb, headline, theme, "title_xl", color_key="text", align="center")
    reveal.append(hb)
    y += 1.8

    # CTA button
    cta_text = get_slot(slots, "cta_text")
    if cta_text:
        cta_w, cta_h = 3.0, 0.7
        cta_x = center_x - cta_w / 2
        btn = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cta_x), Inches(y), Inches(cta_w), Inches(cta_h),
        )
        btn.fill.solid(); btn.fill.fore_color.rgb = palette_color(theme, "accent")
        btn.line.fill.background()
        btn.adjustments[0] = 0.4

        cta_tb = slide.shapes.add_textbox(
            Inches(cta_x), Inches(y + 0.15), Inches(cta_w), Inches(cta_h - 0.3),
        )
        set_textbox_text(cta_tb, truncate_to(cta_text, 40), theme, "body_bold", color_key="bg", align="center")
        reveal.extend([btn, cta_tb])
        y += cta_h + 0.5

    # Contact lines
    contact_lines = get_slot(slots, "contact_lines")
    if contact_lines:
        if not isinstance(contact_lines, list):
            contact_lines = [str(contact_lines)]
        contact_lines = contact_lines[:3]
        for i, line in enumerate(contact_lines):
            cb = slide.shapes.add_textbox(
                Inches(content_x), Inches(y + i * 0.35), Inches(content_w), Inches(0.3),
            )
            set_textbox_text(cb, truncate_to(str(line), 80), theme, "body", color_key="muted", align="center")
            reveal.append(cb)

    # Entrance — closing uses theme's preferred entrance, light stagger
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=150, duration_ms=700, index=i)
