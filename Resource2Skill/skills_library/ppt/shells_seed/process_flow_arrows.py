"""
shells_seed/process_flow_arrows.py — 4-step process flow with chevron arrows between steps.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'timeline_horizontal'
DESCRIPTION = '4-step process flow with chevron arrows between steps.'
ARCHETYPE = 'product'
MOOD = ['bold', 'technical', 'punchy']
DENSITY = 'balanced'
STYLE_TAGS = ['process', 'flow', 'arrows', 'chevron']

SLOTS = [{'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title', 'required': True}, {'name': 'steps', 'kind': 'bullet_list', 'bullet_capacity': 4, 'required': True}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2*margin), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"

    steps = get_slot(slots, "steps", default=[]) or []
    if isinstance(steps, str):
        steps = [x.strip() for x in steps.split("\n") if x.strip()]
    steps = (steps + [""]*4)[:4]
    n = 4
    row_y = 3.0; row_h = 2.5
    arrow_w = 0.5; gap = 0.15
    chev_w = (13.333 - 2*margin - (n-1)*(arrow_w+2*gap)) / n
    xs = []
    for i in range(n):
        x = margin + i * (chev_w + arrow_w + 2*gap)
        xs.append(x)
    revealables = []
    for i, text in enumerate(steps):
        x = xs[i]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(row_y),
                                       Inches(chev_w), Inches(row_h))
        card.fill.solid(); card.fill.fore_color.rgb = palette_color(theme, "panel")
        card.line.color.rgb = palette_color(theme, "accent")
        card.line.width = Pt(1.25); card.shadow.inherit = False
        card.name = f"step_card_{i}"
        # Step number
        sb = slide.shapes.add_textbox(Inches(x+0.25), Inches(row_y+0.2),
                                       Inches(chev_w-0.5), Inches(0.5))
        set_textbox_text(sb, f"0{i+1}", theme, "caption", "accent", align="left")
        sb.name = f"step_num_{i}"
        # Label
        lb = slide.shapes.add_textbox(Inches(x+0.25), Inches(row_y+0.9),
                                       Inches(chev_w-0.5), Inches(row_h-1.0))
        set_textbox_text(lb, truncate_to(str(text), 120), theme, "body_bold", "text", align="left")
        lb.name = f"step_label_{i}"
        revealables.append(card)
        # Arrow between
        if i < n-1:
            ax = x + chev_w + gap
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            Inches(ax), Inches(row_y + row_h/2 - 0.25),
                                            Inches(arrow_w), Inches(0.5))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = palette_color(theme, "accent2", "accent")
            arrow.line.fill.background(); arrow.shadow.inherit = False
            arrow.name = f"step_arrow_{i}"
    add_sequential_reveal(slide, revealables, theme, start_delay_ms=300, step_ms=180)
