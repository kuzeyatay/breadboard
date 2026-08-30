"""
shells_seed/hero_ambient_ring.py — Hero metric with an ambient ring.

Seed shell demonstrating ambient rotation paired with a hero numeric figure.
The decorative ring rotates continuously while the number drops in. The
`hero_number` shape and the ring's `accent_orb` both carry morph anchors so
adjacent slides can morph the same number in-place or slide the orb across.

Exists as a distiller-calibration reference. Not hard-coded into prompts.

Slots:
  eyebrow  (text, optional, ≤24 chars, caption style)
  number   (text, required, ≤12 chars, title_xl style)  — hero metric
  caption  (text, optional, ≤80 chars, body style)      — one-liner below
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_infinite_rotation, add_zoom_climax, get_slot, palette_color,
    set_morph_anchor, set_textbox_text, truncate_to,
)

AMBIENT_CAPABLE = True

# ---- LLM-readable shell metadata ----
ROLE = 'hero_giant_metric'
DESCRIPTION = 'Centered headline with ambient concentric ring glyph behind.'
ARCHETYPE = 'brand'
MOOD = ['editorial', 'cinematic']
DENSITY = 'sparse'
STYLE_TAGS = ['ring', 'centered', 'ambient_accent']

SLOTS = [
    {"name": "eyebrow", "kind": "text", "max_chars": 24, "style": "caption",  "required": False},
    {"name": "number",  "kind": "text", "max_chars": 12, "style": "title_xl", "required": True},
    {"name": "caption", "kind": "text", "max_chars": 80, "style": "body",     "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    bg.shadow.inherit = False

    cx, cy = 6.666, 3.6
    ring_r = 2.4

    # Ring (a thin-bordered oval) that rotates continuously.
    ring = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - ring_r), Inches(cy - ring_r),
        Inches(ring_r * 2), Inches(ring_r * 2),
    )
    ring.fill.background()
    ring.line.color.rgb = palette_color(theme, "accent")
    ring.line.width = Inches(0.035)
    set_morph_anchor(ring, "accent_orb")
    add_infinite_rotation(slide, ring, duration_ms=9000, direction="cw")

    # Eyebrow above the number.
    eyebrow = get_slot(slots, "eyebrow")
    if eyebrow:
        eyebrow = truncate_to(eyebrow, 24)
        eb = slide.shapes.add_textbox(
            Inches(cx - 2.5), Inches(cy - 1.9),
            Inches(5.0), Inches(0.35),
        )
        set_textbox_text(eb, eyebrow, theme, "caption", color_key="accent",
                         align="center")

    # Hero number (the focus of the slide).
    number = truncate_to(get_slot(slots, "number", required=True), 12)
    num_box = slide.shapes.add_textbox(
        Inches(cx - 3.0), Inches(cy - 1.2),
        Inches(6.0), Inches(2.2),
    )
    set_textbox_text(num_box, number, theme, "title_xl", color_key="text",
                     align="center")
    set_morph_anchor(num_box, "hero_number")
    add_zoom_climax(slide, num_box, theme, delay_ms=200, duration_ms=700)

    # Caption below.
    caption = get_slot(slots, "caption")
    if caption:
        caption = truncate_to(caption, 80)
        cap_box = slide.shapes.add_textbox(
            Inches(cx - 4.0), Inches(cy + 1.4),
            Inches(8.0), Inches(0.7),
        )
        set_textbox_text(cap_box, caption, theme, "body", color_key="muted",
                         align="center")
