"""
shells_seed/agenda_numbered.py — Agenda slide with numbered section items.

Slots:
  eyebrow  (text, optional, caption style)               — e.g. "Agenda" / "What's inside"
  headline (text, required, ≤60 chars, title style)      — title
  items    (bullet_list, required, 3-6 items)            — each: {title, body?}
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, apply_type_style,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'agenda'
DESCRIPTION = 'Numbered agenda with 3-6 items stacked vertically and an eyebrow/headline block.'
ARCHETYPE = 'narrative'
MOOD = ['editorial', 'restrained']
DENSITY = 'balanced'
STYLE_TAGS = ['numbered', 'serif_accent', 'list_left']

SLOTS = [
    {"name": "eyebrow",  "kind": "text",        "max_chars": 30, "style": "caption", "required": False},
    {"name": "headline", "kind": "text",        "max_chars": 60, "style": "title",   "required": True},
    {"name": "items",    "kind": "bullet_list", "bullet_capacity": 5, "required": True,
     "item_schema": {"title": "agenda item label (≤50 chars)",
                      "body":  "optional one-line description (≤100 chars, may be '')"}},
]


def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    gutter = theme.get("spacing", {}).get("gutter", 0.3)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    y = margin + 0.3

    # Eyebrow (or default "AGENDA")
    eyebrow = get_slot(slots, "eyebrow")
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(5), Inches(0.35))
        set_textbox_text(eb, truncate_to(eyebrow, 30), theme, "caption", color_key="accent")
        reveal.append(eb)
        y += 0.55
        if theme.get("motif", {}).get("type") == "thin_rule":
            add_hairline(slide, margin, y, 1.5, theme)
            y += 0.2

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 60)
    head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(1.2))
    set_textbox_text(head, headline, theme, "title", color_key="text")
    reveal.append(head)
    y += 1.6

    # Items
    items = get_slot(slots, "items", required=True)
    if not isinstance(items, list):
        raise ValueError(f"'items' must be a list, got {type(items).__name__}")
    items = items[:6]

    item_h = min(0.85, (6.9 - y) / max(len(items), 1))
    for i, item in enumerate(items):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
        item_y = y + i * (item_h + 0.12)

        # Number badge — use a bounded font size so the glyph fits the box
        # height cleanly. accent_serif is usually 72pt italic which overflows
        # a 0.75in box and visually stacks digits across adjacent items.
        num = f"{i+1:02d}"
        num_box = slide.shapes.add_textbox(
            Inches(margin), Inches(item_y + 0.05), Inches(1.0), Inches(item_h),
        )
        num_tf = num_box.text_frame
        num_tf.word_wrap = False
        num_tf.margin_top = 0
        num_tf.margin_bottom = 0
        p = num_tf.paragraphs[0]
        p.text = num
        if p.runs:
            style_token = "accent_serif" if "accent_serif" in theme.get("typography", {}) else "title"
            apply_type_style(p.runs[0], theme, style_token, color_key="accent")
            # Clamp the font size to fit the row so the glyph never crosses
            # into the next row's band. 36pt ≈ 0.5in cap-height.
            max_pt = int(min(48, max(24, (item_h - 0.12) * 72)))
            p.runs[0].font.size = Pt(max_pt)
            p.runs[0].font.italic = True
        reveal.append(num_box)

        # Title
        t_box = slide.shapes.add_textbox(
            Inches(margin + 1.1), Inches(item_y + 0.12), Inches(6), Inches(item_h - 0.1),
        )
        set_textbox_text(t_box, truncate_to(item.get("title", ""), 60), theme, "body_bold", color_key="text")
        reveal.append(t_box)

        # Body (right side)
        body = item.get("body", "")
        if body:
            b_box = slide.shapes.add_textbox(
                Inches(margin + 7.3), Inches(item_y + 0.18), Inches(13.333 - margin - 7.5),
                Inches(item_h - 0.1),
            )
            set_textbox_text(b_box, truncate_to(body, 120), theme, "body", color_key="muted")
            reveal.append(b_box)

        # Divider line between items
        if i < len(items) - 1:
            add_hairline(slide, margin, item_y + item_h + 0.04, 13.333 - 2 * margin, theme, color_key="border")

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=80, duration_ms=450, index=i)
