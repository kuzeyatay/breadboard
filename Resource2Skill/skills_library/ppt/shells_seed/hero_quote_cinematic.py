"""
shells_seed/hero_quote_cinematic.py — cinematic full-bleed quote slide.

Slots:
  context      (text, optional, ≤40 chars, caption)
  quote        (text, required, ≤180 chars, title_xl)
  attribution  (text, optional, ≤60 chars, caption)

Layout: left 40% accent panel with eyebrow + attribution stack, right 60%
large bold quote. Cinematic product/brand feel.
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot, palette_color,
    set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = "hero_quote"
DESCRIPTION = "Cinematic split: left accent panel with context + attribution, right side large bold quote."
ARCHETYPE = "brand"
MOOD = ["cinematic", "bold", "editorial"]
DENSITY = "sparse"
STYLE_TAGS = ["split", "bold_quote", "asymmetric"]

SLOTS = [
    {"name": "context",     "kind": "text", "max_chars":  40, "style": "caption",  "required": False},
    {"name": "quote",       "kind": "text", "max_chars": 180, "style": "title_xl", "required": True},
    {"name": "attribution", "kind": "text", "max_chars":  60, "style": "caption",  "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)
    # bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    # Left accent panel (40% width)
    panel_w = 5.1
    add_solid_rect(slide, 0, 0, panel_w, 7.5, theme, color_key="panel")

    # Left: context eyebrow + attribution stacked bottom
    context = get_slot(slots, "context")
    if context:
        eb = slide.shapes.add_textbox(Inches(margin), Inches(margin),
                                       Inches(panel_w - 2 * margin), Inches(0.5))
        set_textbox_text(eb, context, theme, "caption", "accent", align="left")
        eb.name = "context"
        add_theme_entrance(slide, eb, theme, delay_ms=100, index=0)

    attrib = get_slot(slots, "attribution")
    if attrib:
        at = slide.shapes.add_textbox(Inches(margin),
                                       Inches(7.5 - margin - 0.6),
                                       Inches(panel_w - 2 * margin),
                                       Inches(0.5))
        set_textbox_text(at, attrib, theme, "caption", "muted", align="left")
        at.name = "attribution"
        add_theme_entrance(slide, at, theme, delay_ms=500, index=2)

    # Right: big quote
    quote_text = truncate_to(get_slot(slots, "quote", required=True), 180)
    q = slide.shapes.add_textbox(Inches(panel_w + margin), Inches(1.4),
                                  Inches(13.333 - panel_w - 2 * margin),
                                  Inches(4.8))
    set_textbox_text(q, quote_text, theme, "title_xl",
                     color_key="text", align="left")
    q.name = "quote"
    add_theme_entrance(slide, q, theme, delay_ms=250, index=1)
