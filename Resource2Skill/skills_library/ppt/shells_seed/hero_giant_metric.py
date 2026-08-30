"""
shells_seed/hero_giant_metric.py — One single huge number, dramatic drop-in.

Use case: the deck's headline KPI ($148M ARR, 3.2x faster, 99.99% uptime).
Renders the final number as a single hero shape with a zoom-in entrance and
a follow-up emphasis pulse — the visual climax of the deck.

Slots:
  eyebrow      (text, optional, caption — small section label above)
  value        (text, required — the final value, e.g. "$148M")
  label        (text, required — what the number means, below)
  context      (text, optional, body — one short sentence under the label)
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from _shell_helpers import (
    add_emphasis_pulse, add_theme_entrance, add_zoom_climax, apply_type_style,
    get_slot, palette_color, set_morph_anchor, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'hero_giant_metric'
DESCRIPTION = 'Single huge numeric value + label + eyebrow; one-metric hero slide.'
ARCHETYPE = 'data'
MOOD = ['boardroom', 'bold']
DENSITY = 'sparse'
STYLE_TAGS = ['giant_numeral', 'centered', 'single_metric']

SLOTS = [
    {"name": "eyebrow", "kind": "text", "max_chars": 40,  "style": "caption",   "required": False},
    {"name": "value",   "kind": "text", "max_chars": 12,  "style": "metric_xl", "required": True},
    {"name": "label",   "kind": "text", "max_chars": 60,  "style": "subtitle",  "required": True},
    {"name": "context", "kind": "text", "max_chars": 140, "style": "body",      "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    eyebrow = get_slot(slots, "eyebrow")
    value = truncate_to(get_slot(slots, "value", required=True), 12)
    label = truncate_to(get_slot(slots, "label", required=True), 60)
    context = get_slot(slots, "context")

    # Eyebrow at top-left
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(0.4))
        set_textbox_text(eb, truncate_to(eyebrow, 40), theme, "caption", color_key="accent")
        add_theme_entrance(slide, eb, theme, delay_ms=200, index=0)

    # Hero number — single giant text, centered
    counter_box_w, counter_box_h = 12.5, 3.0
    cx = (13.333 - counter_box_w) / 2
    cy = (7.5 - counter_box_h) / 2 - 0.3
    hero = slide.shapes.add_textbox(
        Inches(cx), Inches(cy), Inches(counter_box_w), Inches(counter_box_h),
    )
    tf = hero.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = value
    if p.runs:
        style_token = "metric_xl" if "metric_xl" in theme.get("typography", {}) else "title_xl"
        apply_type_style(p.runs[0], theme, style_token, color_key="accent")
        # Force a giant size — this is THE hero of the deck
        n = len(value)
        p.runs[0].font.size = Pt(180 if n <= 4 else 150 if n <= 6 else 120 if n <= 8 else 96)
        p.alignment = PP_ALIGN.CENTER

    # Climactic zoom entrance + emphasis pulse — the aha moment
    add_zoom_climax(slide, hero, theme, delay_ms=600, duration_ms=700)
    add_emphasis_pulse(slide, hero, theme, delay_ms=1500, duration_ms=600)
    set_morph_anchor(hero, "hero_number")

    # Label below the hero number
    lab = slide.shapes.add_textbox(
        Inches(1.0), Inches(cy + counter_box_h + 0.2),
        Inches(11.333), Inches(0.7),
    )
    set_textbox_text(lab, label, theme, "subtitle", color_key="text", align="center")
    add_theme_entrance(slide, lab, theme, delay_ms=1700, index=0)

    # Optional context line under that
    if context:
        ctx = slide.shapes.add_textbox(
            Inches(2.0), Inches(cy + counter_box_h + 1.05),
            Inches(9.333), Inches(0.7),
        )
        set_textbox_text(ctx, truncate_to(context, 140), theme, "body",
                         color_key="muted", align="center")
        add_theme_entrance(slide, ctx, theme, delay_ms=2000, index=1)
