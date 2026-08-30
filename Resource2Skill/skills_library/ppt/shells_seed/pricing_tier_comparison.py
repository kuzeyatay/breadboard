"""
shells_seed/pricing_tier_comparison.py — Three pricing tier cards with tier name, price, features, and CTA.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'comparison_split'
DESCRIPTION = 'Three pricing tier cards with tier name, price, features, and CTA.'
ARCHETYPE = 'product'
MOOD = ['bold', 'punchy', 'boardroom']
DENSITY = 'dense'
STYLE_TAGS = ['pricing', 'tiers', 'three_up', 'comparison_cards']

SLOTS = [{'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title', 'required': True}, {'name': 'tiers', 'kind': 'bullet_list', 'bullet_capacity': 3, 'required': True}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2*margin), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="center")
    hl.name = "headline"

    tiers = get_slot(slots, "tiers", default=[]) or []
    if isinstance(tiers, str):
        tiers = [x.strip() for x in tiers.split("\n") if x.strip()]
    tiers = (tiers + [""]*3)[:3]

    card_y = 1.85; card_h = 5.3
    gap = 0.3
    card_w = (13.333 - 2*margin - 2*gap) / 3
    cards = []
    for i, entry in enumerate(tiers):
        if isinstance(entry, dict):
            tname = entry.get("name", "")
            price = entry.get("price", "")
            features = entry.get("features", [])
            cta = entry.get("cta", "")
        else:
            parts = [p.strip() for p in str(entry).split("|")]
            parts = (parts + [""]*4)[:4]
            tname, price, feats_str, cta = parts
            features = [f.strip() for f in (feats_str or "").split(";") if f.strip()]
        x = margin + i * (card_w + gap)
        highlight = (i == 1)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(card_y),
                                       Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = palette_color(theme, "accent" if highlight else "panel")
        card.line.color.rgb = palette_color(theme, "border", "muted")
        card.line.width = Pt(1.25); card.shadow.inherit = False
        card.name = f"tier_card_{i}"
        text_color = "bg" if highlight else "text"
        # Tier name
        nb = slide.shapes.add_textbox(Inches(x+0.3), Inches(card_y+0.3),
                                       Inches(card_w-0.6), Inches(0.6))
        set_textbox_text(nb, tname, theme, "body_bold", text_color, align="center")
        nb.name = f"tier_name_{i}"
        # Price
        pb = slide.shapes.add_textbox(Inches(x+0.3), Inches(card_y+1.0),
                                       Inches(card_w-0.6), Inches(1.2))
        set_textbox_text(pb, price, theme, "metric_xl",
                         text_color if highlight else "accent", align="center")
        pb.name = f"tier_price_{i}"
        # Features
        features = (features + [""]*5)[:5]
        for j, f in enumerate(features):
            if not f:
                continue
            fb = slide.shapes.add_textbox(Inches(x+0.4),
                                           Inches(card_y+2.3+j*0.4),
                                           Inches(card_w-0.8), Inches(0.4))
            set_textbox_text(fb, f"- {f}", theme, "body",
                             "bg" if highlight else "muted", align="left")
            fb.name = f"tier_feat_{i}_{j}"
        # CTA pill
        if cta:
            cta_w = card_w - 1.0
            ctb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(x+0.5), Inches(card_y+card_h-0.8),
                                          Inches(cta_w), Inches(0.5))
            ctb.fill.solid()
            ctb.fill.fore_color.rgb = palette_color(theme, "bg" if highlight else "accent")
            ctb.line.fill.background(); ctb.shadow.inherit = False
            ctb.name = f"tier_cta_{i}"
            cbt = slide.shapes.add_textbox(Inches(x+0.5), Inches(card_y+card_h-0.75),
                                            Inches(cta_w), Inches(0.45))
            set_textbox_text(cbt, cta, theme, "body_bold",
                             "accent" if highlight else "bg", align="center")
            cbt.name = f"tier_cta_text_{i}"
        cards.append(card)
    add_sequential_reveal(slide, cards, theme, start_delay_ms=280, step_ms=160)
