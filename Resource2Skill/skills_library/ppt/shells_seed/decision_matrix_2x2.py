"""
shells_seed/decision_matrix_2x2.py — 2x2 decision matrix with x/y axis labels and 4 quadrant items.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'comparison_split'
DESCRIPTION = '2x2 decision matrix with x/y axis labels and 4 quadrant items.'
ARCHETYPE = 'boardroom'
MOOD = ['boardroom', 'restrained', 'editorial']
DENSITY = 'balanced'
STYLE_TAGS = ['matrix', '2x2', 'decision', 'axis']

SLOTS = [{'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title', 'required': True}, {'name': 'x_axis', 'kind': 'text', 'max_chars': 40, 'style': 'caption', 'required': True}, {'name': 'y_axis', 'kind': 'text', 'max_chars': 40, 'style': 'caption', 'required': True}, {'name': 'quadrants', 'kind': 'bullet_list', 'bullet_capacity': 4, 'required': True}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2*margin), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"

    matrix_x = 2.0; matrix_y = 2.0
    matrix_w = 9.333; matrix_h = 4.8
    # Draw grid
    for row in range(2):
        for col in range(2):
            x = matrix_x + col * (matrix_w/2)
            y = matrix_y + row * (matrix_h/2)
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(x), Inches(y),
                                           Inches(matrix_w/2), Inches(matrix_h/2))
            cell.fill.solid()
            cell.fill.fore_color.rgb = palette_color(theme, "panel" if (row+col)%2 else "bg")
            cell.line.color.rgb = palette_color(theme, "border", "muted")
            cell.line.width = Pt(0.75); cell.shadow.inherit = False
            cell.name = f"matrix_cell_{row}_{col}"

    quads = get_slot(slots, "quadrants", default=[]) or []
    if isinstance(quads, str):
        quads = [x.strip() for x in quads.split("\n") if x.strip()]
    quads = (quads + [""]*4)[:4]
    positions = [(0,0),(0,1),(1,0),(1,1)]
    for i, (row, col) in enumerate(positions):
        x = matrix_x + col * (matrix_w/2)
        y = matrix_y + row * (matrix_h/2)
        tb = slide.shapes.add_textbox(Inches(x+0.25), Inches(y+0.25),
                                       Inches(matrix_w/2 - 0.5),
                                       Inches(matrix_h/2 - 0.5))
        set_textbox_text(tb, truncate_to(str(quads[i]), 120),
                         theme, "body", "text", align="left")
        tb.name = f"matrix_text_{i}"

    x_axis = get_slot(slots, "x_axis", required=True)
    xb = slide.shapes.add_textbox(Inches(matrix_x),
                                   Inches(matrix_y + matrix_h + 0.15),
                                   Inches(matrix_w), Inches(0.4))
    set_textbox_text(xb, f"-> {x_axis} ->", theme, "caption", "accent", align="center")
    xb.name = "x_axis"

    y_axis = get_slot(slots, "y_axis", required=True)
    yb = slide.shapes.add_textbox(Inches(0.4), Inches(matrix_y),
                                   Inches(1.4), Inches(matrix_h))
    set_textbox_text(yb, y_axis, theme, "caption", "accent", align="left")
    yb.name = "y_axis"
