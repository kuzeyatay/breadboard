"""
shells_seed/hero_giant_metric_brand.py — brand-bold KPI hero with halo.

Slots:
  eyebrow (text, optional, ≤30 chars, caption)
  value   (text, required, ≤12 chars, metric_xl)
  label   (text, required, ≤40 chars, subtitle)
  context (text, optional, ≤140 chars, body)

Layout: centered value with an accent-filled circular halo behind,
eyebrow above, label + context below. Product/brand launch feel.
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, add_pulse_loop, get_slot, palette_color,
    set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = "hero_giant_metric"
DESCRIPTION = "Bold brand KPI hero: accent halo behind centered numeral, calm caption above, label + context below."
ARCHETYPE = "brand"
MOOD = ["bold", "cinematic", "punchy"]
DENSITY = "sparse"
STYLE_TAGS = ["centered", "halo", "ambient_accent"]

SLOTS = [
    {"name": "eyebrow", "kind": "text", "max_chars":  30, "style": "caption",   "required": False},
    {"name": "value",   "kind": "text", "max_chars":  12, "style": "metric_xl", "required": True},
    {"name": "label",   "kind": "text", "max_chars":  40, "style": "subtitle",  "required": True},
    {"name": "context", "kind": "text", "max_chars": 140, "style": "body",      "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    # Halo (large accent-filled oval behind center)
    halo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.4), Inches(1.4),
                                   Inches(4.5), Inches(4.5))
    halo.fill.solid(); halo.fill.fore_color.rgb = palette_color(theme, "accent2")
    halo.line.fill.background(); halo.shadow.inherit = False
    halo.name = "accent_halo"
    try:
        add_pulse_loop(slide, halo, duration_ms=3200, scale_pct=106)
    except Exception:
        pass

    # Eyebrow
    eyebrow = get_slot(slots, "eyebrow")
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(2.5), Inches(0.9),
                                       Inches(8.333), Inches(0.5))
        set_textbox_text(eb, eyebrow, theme, "caption", "accent",
                         align="center")
        eb.name = "eyebrow"
        add_theme_entrance(slide, eb, theme, delay_ms=120, index=0)

    # Value centered over halo
    value = truncate_to(get_slot(slots, "value", required=True), 12)
    val = slide.shapes.add_textbox(Inches(2.5), Inches(2.4),
                                    Inches(8.333), Inches(2.8))
    set_textbox_text(val, value, theme, "metric_xl", "text",
                     align="center")
    val.name = "value"
    add_theme_entrance(slide, val, theme, delay_ms=280, index=1)

    # Label
    label = truncate_to(get_slot(slots, "label", required=True), 40)
    lb = slide.shapes.add_textbox(Inches(2.5), Inches(5.5),
                                   Inches(8.333), Inches(0.6))
    set_textbox_text(lb, label, theme, "subtitle", "muted",
                     align="center")
    lb.name = "label"
    add_theme_entrance(slide, lb, theme, delay_ms=460, index=2)

    # Context
    context = get_slot(slots, "context")
    if context:
        ctx = slide.shapes.add_textbox(Inches(3.2), Inches(6.3),
                                        Inches(6.933), Inches(0.8))
        set_textbox_text(ctx, context, theme, "body", "text",
                         align="center")
        ctx.name = "context"
        add_theme_entrance(slide, ctx, theme, delay_ms=620, index=3)
