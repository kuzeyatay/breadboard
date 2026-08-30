"""
shells_seed/hero_before_after.py — Before/after pivot-swap reveal.

Use case: dramatic comparison (old way → our way, slow query → fast query,
manual workflow → automated). Two cards occupy the SAME center anchor.
The "before" card spins out, the "after" card spins in — a true card-flip moment.

Slots:
  context       (text, optional, caption — "BEFORE / AFTER")
  before_label  (text, required, ≤30, body_bold — short label e.g. "Old way")
  before_value  (text, required, ≤80, title — what it was)
  after_label   (text, required, ≤30, body_bold)
  after_value   (text, required, ≤80, title)
  takeaway      (text, optional, ≤120, body — summary line at the bottom)
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from _shell_helpers import (
    add_emphasis_pulse, add_mask_wipe, add_pivot_swap, add_solid_rect,
    add_theme_entrance, add_zoom_climax, apply_type_style, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'comparison_split'
DESCRIPTION = 'Left/right before-vs-after panels with takeaway strip below.'
ARCHETYPE = 'narrative'
MOOD = ['editorial', 'restrained']
DENSITY = 'balanced'
STYLE_TAGS = ['split', 'comparison', 'symmetric']

SLOTS = [
    {"name": "context",      "kind": "text", "max_chars": 40,  "style": "caption",   "required": False},
    {"name": "before_label", "kind": "text", "max_chars": 30,  "style": "body_bold", "required": True},
    {"name": "before_value", "kind": "text", "max_chars": 80,  "style": "title",     "required": True},
    {"name": "after_label",  "kind": "text", "max_chars": 30,  "style": "body_bold", "required": True},
    {"name": "after_value",  "kind": "text", "max_chars": 80,  "style": "title",     "required": True},
    {"name": "takeaway",     "kind": "text", "max_chars": 120, "style": "body",      "required": False},
]


def _build_card(slide, x, y, w, h, theme, label, value, label_color, value_color, fill_key):
    """Add a card group: rounded rect + label at top + value below. Returns the card shape."""
    from pptx.enum.text import PP_ALIGN
    card = add_solid_rect(slide, x, y, w, h, theme,
                          color_key=fill_key, line=True,
                          shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.04
    # Label
    lab = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 0.5), Inches(w - 0.8), Inches(0.5))
    set_textbox_text(lab, label, theme, "body_bold", color_key=label_color, align="center")
    # Value
    val = slide.shapes.add_textbox(
        Inches(x + 0.4), Inches(y + h * 0.42), Inches(w - 0.8), Inches(h * 0.5),
    )
    tf = val.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = value
    if p.runs:
        apply_type_style(p.runs[0], theme, "title", color_key=value_color)
        # Aggressive shrink for long values
        n = len(value)
        size = 36 if n <= 30 else 28 if n <= 60 else 22
        p.runs[0].font.size = Pt(size)
        p.alignment = PP_ALIGN.CENTER
    return card


def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    context = get_slot(slots, "context")
    before_label = truncate_to(get_slot(slots, "before_label", required=True), 30)
    before_value = truncate_to(get_slot(slots, "before_value", required=True), 80)
    after_label  = truncate_to(get_slot(slots, "after_label",  required=True), 30)
    after_value  = truncate_to(get_slot(slots, "after_value",  required=True), 80)
    takeaway     = get_slot(slots, "takeaway")

    # Context label centered top
    if context:
        cb = slide.shapes.add_textbox(Inches(1.0), Inches(0.7), Inches(11.333), Inches(0.45))
        set_textbox_text(cb, truncate_to(context, 40), theme, "caption",
                         color_key="accent", align="center")
        add_theme_entrance(slide, cb, theme, delay_ms=200, index=0)

    # Side-by-side cards with arrow between — static-friendly layout
    card_w, card_h = 5.4, 4.2
    gap = 0.6
    total = card_w * 2 + gap
    left_x = (13.333 - total) / 2
    right_x = left_x + card_w + gap
    cy = (7.5 - card_h) / 2 - 0.1

    before_card = _build_card(
        slide, left_x, cy, card_w, card_h, theme,
        before_label, before_value,
        label_color="muted", value_color="muted", fill_key="panel",
    )
    after_card = _build_card(
        slide, right_x, cy, card_w, card_h, theme,
        after_label, after_value,
        label_color="accent", value_color="text", fill_key="panel",
    )

    # Arrow between
    arrow_w = gap - 0.1
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left_x + card_w + 0.05), Inches(cy + card_h / 2 - 0.25),
        Inches(arrow_w), Inches(0.5),
    )
    arrow.fill.solid(); arrow.fill.fore_color.rgb = palette_color(theme, "accent")
    arrow.line.fill.background()

    # Choreography: before fades in, arrow grows, after lands with emphasis
    add_theme_entrance(slide, before_card, theme, delay_ms=600,  index=0)
    add_mask_wipe(slide, arrow, theme, delay_ms=1100, duration_ms=600,
                  direction="left_to_right")
    add_zoom_climax(slide, after_card, theme, delay_ms=1500, duration_ms=600)
    add_emphasis_pulse(slide, after_card, theme, delay_ms=2200, duration_ms=600)

    # Takeaway under the cards
    if takeaway:
        tb = slide.shapes.add_textbox(
            Inches(2.0), Inches(cy + card_h + 0.4), Inches(9.333), Inches(0.7),
        )
        set_textbox_text(tb, truncate_to(takeaway, 120), theme, "body",
                         color_key="muted", align="center")
        add_theme_entrance(slide, tb, theme, delay_ms=2500, index=0)
