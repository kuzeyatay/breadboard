"""
shells_seed/cover_ambient_orbit.py — Cover with ambient orbital motion.

Seed shell demonstrating how an ambient primitive plus a morph anchor come
together on a cover slide. The small accent orb orbits slowly while the
headline holds still; the orb is tagged with the `accent_orb` morph role so
the next slide can morph the orb in-place.

Exists as a distiller-calibration reference (the distiller will learn from
this shape of code when fed `animation/` frames). It is NOT hard-coded into
the agent prompt — selection flows through `select_shell`.

Slots:
  eyebrow   (text, optional, ≤30 chars, caption style)
  headline  (text, required, ≤60 chars, title_xl style)
  subhead   (text, optional, ≤120 chars, subtitle style)
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_orbital_motion, add_pulse_loop, add_theme_entrance,
    get_slot, palette_color, set_morph_anchor, set_textbox_text, truncate_to,
)

AMBIENT_CAPABLE = True

# ---- LLM-readable shell metadata ----
ROLE = 'cover'
DESCRIPTION = 'Cover with orbital/halo accent on right, text stack on left; ambient feel.'
ARCHETYPE = 'brand'
MOOD = ['editorial', 'cinematic', 'calm']
DENSITY = 'sparse'
STYLE_TAGS = ['orbital', 'accent_ring', 'asymmetric']

SLOTS = [
    {"name": "eyebrow",  "kind": "text", "max_chars":  30, "style": "caption",  "required": False},
    {"name": "headline", "kind": "text", "max_chars":  60, "style": "title_xl", "required": True},
    {"name": "subhead",  "kind": "text", "max_chars": 120, "style": "subtitle", "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    bg.shadow.inherit = False

    # Accent orb that orbits around a center point. Small shape, high contrast.
    orb_radius_in = 0.35
    center_x, center_y = 10.2, 3.6
    orbit_r = 1.8
    orb = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x + orbit_r - orb_radius_in),
        Inches(center_y - orb_radius_in),
        Inches(orb_radius_in * 2),
        Inches(orb_radius_in * 2),
    )
    orb.fill.solid()
    orb.fill.fore_color.rgb = palette_color(theme, "accent")
    orb.line.fill.background()
    set_morph_anchor(orb, "accent_orb")
    add_orbital_motion(slide, orb, center_xy=(center_x, center_y),
                       radius_in=orbit_r, duration_ms=7000, direction="cw")

    # Secondary breathing node for depth (pulses in place).
    node = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x - 0.18),
        Inches(center_y - 0.18),
        Inches(0.36), Inches(0.36),
    )
    node.fill.solid()
    node.fill.fore_color.rgb = palette_color(theme, "accent2", fallback="accent")
    node.line.fill.background()
    add_pulse_loop(slide, node, duration_ms=2200, scale_pct=120)

    # Left text stack
    text_x = margin + 0.2
    text_w = 6.5
    y_cursor = 1.4
    reveal_shapes = []

    eyebrow = get_slot(slots, "eyebrow")
    if eyebrow:
        eyebrow = truncate_to(eyebrow, 30)
        eb = slide.shapes.add_textbox(Inches(text_x), Inches(y_cursor),
                                      Inches(text_w), Inches(0.35))
        set_textbox_text(eb, eyebrow, theme, "caption", color_key="accent")
        reveal_shapes.append(eb)
        y_cursor += 0.55
        if theme.get("motif", {}).get("type") == "thin_rule":
            add_hairline(slide, text_x, y_cursor, min(text_w, 1.8), theme)
            y_cursor += 0.25

    headline = truncate_to(get_slot(slots, "headline", required=True), 60)
    head_box = slide.shapes.add_textbox(Inches(text_x), Inches(y_cursor),
                                        Inches(text_w), Inches(2.2))
    set_textbox_text(head_box, headline, theme, "title_xl", color_key="text")
    set_morph_anchor(head_box, "hero_headline")
    reveal_shapes.append(head_box)
    y_cursor += 2.3

    subhead = get_slot(slots, "subhead")
    if subhead:
        subhead = truncate_to(subhead, 120)
        sub_box = slide.shapes.add_textbox(Inches(text_x), Inches(y_cursor),
                                           Inches(text_w), Inches(1.4))
        set_textbox_text(sub_box, subhead, theme, "subtitle", color_key="muted")
        reveal_shapes.append(sub_box)

    for i, shape in enumerate(reveal_shapes):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
