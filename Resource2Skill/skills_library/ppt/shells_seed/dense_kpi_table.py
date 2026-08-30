"""
shells_seed/dense_kpi_table.py — 6-row KPI table, data-dense.

Slots:
  section_label (text, optional, caption)
  headline      (text, required, title)
  kpis          (bullet_list, required, list of "label | value | delta | note"
                  strings OR dicts with those keys; 4-6 rows)

Layout: single-column table with columns [label | value | delta | note].
Reads like a board report. Density=dense.
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = "metric_dashboard"
DESCRIPTION = "Dense 4-column KPI table (label / value / delta / note) with 4-6 rows."
ARCHETYPE = "boardroom"
MOOD = ["boardroom", "technical", "restrained"]
DENSITY = "dense"
STYLE_TAGS = ["table", "kpi_rows", "boardroom_style"]

SLOTS = [
    {"name": "section_label", "kind": "text", "max_chars": 30,  "style": "caption", "required": False},
    {"name": "headline",      "kind": "text", "max_chars": 100, "style": "title", "required": True},
    {"name": "kpis",          "kind": "bullet_list", "bullet_capacity": 6, "required": True,
     "item_schema": {"label": "metric name (≤30 chars)",
                      "value": "headline number (≤10 chars, e.g. '94%' or '$1.2M')",
                      "delta": "delta or trend (≤10 chars, e.g. '+3.1%' or '↑')",
                      "note":  "one-clause context (≤60 chars)"}},
]


def _parse_row(entry):
    if isinstance(entry, dict):
        return (
            str(entry.get("label") or ""),
            str(entry.get("value") or ""),
            str(entry.get("delta") or ""),
            str(entry.get("note") or ""),
        )
    if isinstance(entry, str):
        parts = [x.strip() for x in entry.split("|")]
        return tuple((parts + ["", "", "", ""])[:4])
    return "", "", "", ""


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    section_label = get_slot(slots, "section_label")
    if section_label:
        sl = slide.shapes.add_textbox(Inches(margin), Inches(margin),
                                       Inches(7.0), Inches(0.4))
        set_textbox_text(sl, section_label, theme, "caption", "accent", align="left")
        sl.name = "section_label"

    headline = truncate_to(get_slot(slots, "headline", required=True), 100)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(1.0),
                                   Inches(13.333 - 2 * margin), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"
    add_theme_entrance(slide, hl, theme, delay_ms=120, index=0)

    rows = get_slot(slots, "kpis", default=[]) or []
    if isinstance(rows, str):
        rows = [rows]
    rows = rows[:6]
    n = len(rows)
    if n == 0:
        return

    # Table header row
    header_y = 2.2
    col_xs = [margin, margin + 4.5, margin + 7.5, margin + 9.5]
    col_ws = [4.3, 2.9, 1.9, 13.333 - margin - (margin + 9.5) - 0.2]
    headers = ["METRIC", "VALUE", "DELTA", "NOTE"]
    for x, w, h in zip(col_xs, col_ws, headers):
        tb = slide.shapes.add_textbox(Inches(x), Inches(header_y),
                                       Inches(w), Inches(0.4))
        set_textbox_text(tb, h, theme, "caption", "muted", align="left")
    add_hairline(slide, margin, header_y + 0.5,
                 13.333 - 2 * margin, theme, "border")

    # Rows
    row_start = 2.85
    row_h = (7.5 - row_start - margin) / max(n, 1)
    for i, entry in enumerate(rows):
        label, value, delta, note = _parse_row(entry)
        y = row_start + i * row_h
        # Subtle alternating row background
        if i % 2 == 0:
            add_solid_rect(slide, margin, y, 13.333 - 2 * margin,
                           row_h - 0.05, theme, color_key="panel")
        for x, w, style, color, txt, name in [
            (col_xs[0], col_ws[0], "body_bold", "text", label, f"kpi_label_{i}"),
            (col_xs[1], col_ws[1], "title", "accent", value, f"kpi_value_{i}"),
            (col_xs[2], col_ws[2], "body_bold", "success", delta, f"kpi_delta_{i}"),
            (col_xs[3], col_ws[3], "body", "muted", note, f"kpi_note_{i}"),
        ]:
            if not txt:
                continue
            tb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.15),
                                           Inches(w - 0.2),
                                           Inches(row_h - 0.2))
            set_textbox_text(tb, truncate_to(txt, 80), theme, style, color, align="left")
            tb.name = name
