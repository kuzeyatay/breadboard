"""
shells_seed/timeline_horizontal.py — 4-step horizontal timeline.

Slots:
  section_label (text, optional, caption style)
  headline      (text, required, ≤60 chars, title style)
  steps         (bullet_list, required, 3-5 items, each {label, body, marker?})
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Inches, Pt, Emu

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'timeline_horizontal'
DESCRIPTION = 'Horizontal row of sequential steps with connecting rule; process/pipeline view.'
ARCHETYPE = 'narrative'
MOOD = ['editorial', 'technical']
DENSITY = 'balanced'
STYLE_TAGS = ['sequential', 'timeline', 'horizontal']

SLOTS = [
    {"name": "section_label", "kind": "text",        "max_chars": 30, "style": "caption", "required": False},
    {"name": "headline",      "kind": "text",        "max_chars": 60, "style": "title",   "required": True},
    {"name": "steps",         "kind": "bullet_list", "bullet_capacity": 5, "required": True},
]


def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    y = margin + 0.2

    # Eyebrow
    label = get_slot(slots, "section_label")
    if label:
        lb = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(6), Inches(0.35))
        set_textbox_text(lb, truncate_to(label, 30), theme, "caption", color_key="accent")
        reveal.append(lb); y += 0.55
        if theme.get("motif", {}).get("type") == "thin_rule":
            add_hairline(slide, margin, y, 1.8, theme); y += 0.15

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 60)
    head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(1.1))
    set_textbox_text(head, headline, theme, "title", color_key="text")
    reveal.append(head); y += 1.6

    # Timeline axis
    steps = get_slot(slots, "steps", required=True)
    if not isinstance(steps, list) or len(steps) < 2:
        raise ValueError(f"'steps' must be a list of 2-5 items, got {steps!r}")
    steps = steps[:5]
    n = len(steps)

    axis_y = y + 1.6  # vertical center of the timeline axis
    axis_left = margin + 0.3
    axis_right = 13.333 - margin - 0.3
    axis_w = axis_right - axis_left

    # Horizontal axis line (thin rule in accent)
    axis_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(axis_left), Inches(axis_y - 0.02), Inches(axis_w), Inches(0.04),
    )
    axis_rect.fill.solid(); axis_rect.fill.fore_color.rgb = palette_color(theme, "accent")
    axis_rect.line.fill.background()
    reveal.append(axis_rect)

    # Step dots + labels
    for i, step in enumerate(steps):
        if not isinstance(step, dict):
            step = {"label": str(step), "body": ""}
        x_center = axis_left + axis_w * i / max(n - 1, 1)

        # Dot
        dot_r = 0.22
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x_center - dot_r), Inches(axis_y - dot_r), Inches(dot_r * 2), Inches(dot_r * 2),
        )
        dot.fill.solid(); dot.fill.fore_color.rgb = palette_color(theme, "accent")
        dot.line.color.rgb = palette_color(theme, "bg")
        dot.line.width = Pt(3)
        reveal.append(dot)

        # Marker text inside dot (optional)
        marker = step.get("marker") or f"{i+1}"
        m_box = slide.shapes.add_textbox(
            Inches(x_center - 0.3), Inches(axis_y - 0.18), Inches(0.6), Inches(0.36),
        )
        set_textbox_text(m_box, truncate_to(str(marker), 3), theme, "caption", color_key="bg", align="center")
        reveal.append(m_box)

        # Label above (or below, alternating to save space)
        label_above = i % 2 == 0
        label_y = axis_y - 1.4 if label_above else axis_y + 0.5
        body_y = label_y + 0.38

        col_w = min(axis_w / n - 0.1, 2.6)
        col_x = x_center - col_w / 2
        # Clamp column text box to canvas bounds (edge-step columns can overflow)
        col_x = max(margin, min(col_x, 13.333 - margin - col_w))

        lbl_box = slide.shapes.add_textbox(
            Inches(col_x), Inches(label_y), Inches(col_w), Inches(0.4),
        )
        set_textbox_text(lbl_box, truncate_to(step.get("label", ""), 24), theme, "body_bold", color_key="text", align="center")
        reveal.append(lbl_box)

        body = step.get("body", "")
        if body:
            b_box = slide.shapes.add_textbox(
                Inches(col_x), Inches(body_y), Inches(col_w), Inches(0.7),
            )
            set_textbox_text(b_box, truncate_to(body, 100), theme, "body", color_key="muted", align="center")
            reveal.append(b_box)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=80, duration_ms=420, index=i)
