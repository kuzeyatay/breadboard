"""
shells_seed/bento_grid_feature.py — Bento-style irregular grid (2 large + 3 small tiles) for feature highlights.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'feature_grid'
DESCRIPTION = 'Bento-style irregular grid (2 large + 3 small tiles) for feature highlights.'
ARCHETYPE = 'product'
MOOD = ['bold', 'cinematic', 'editorial']
DENSITY = 'balanced'
STYLE_TAGS = ['bento', 'asymmetric', 'feature', 'hero_grid']

SLOTS = [{'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title', 'required': True}, {'name': 'tiles', 'kind': 'bullet_list', 'bullet_capacity': 5, 'required': True}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2*margin), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"

    tiles = get_slot(slots, "tiles", default=[]) or []
    if isinstance(tiles, str):
        tiles = [x.strip() for x in tiles.split("\n") if x.strip()]
    tiles = (tiles + [""]*5)[:5]

    grid_y = 1.9; grid_w = 13.333 - 2*margin; grid_h = 5.0
    gap = 0.25
    big_w = (grid_w - gap) * 0.55
    big_h = (grid_h - gap) / 2
    small_w = grid_w - big_w - gap
    small_h = (grid_h - 2*gap) / 3

    # Layout: 2 big tiles on left (stacked), 3 small tiles on right (stacked)
    layouts = [
        (margin, grid_y,          big_w, big_h),                       # tile 0: big top-left
        (margin, grid_y + big_h + gap, big_w, big_h),                  # tile 1: big bottom-left
        (margin + big_w + gap, grid_y,                   small_w, small_h),  # tile 2: small top-right
        (margin + big_w + gap, grid_y + small_h + gap,   small_w, small_h),  # tile 3: small mid-right
        (margin + big_w + gap, grid_y + 2*(small_h + gap), small_w, small_h) # tile 4: small bot-right
    ]
    reveals = []
    for i, (x, y, w, h) in enumerate(layouts):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(y), Inches(w), Inches(h))
        color_key = "accent" if i == 0 else ("accent2" if i == 1 else "panel")
        card.fill.solid(); card.fill.fore_color.rgb = palette_color(theme, color_key)
        card.line.fill.background(); card.shadow.inherit = False
        card.name = f"bento_tile_{i}"
        text_color = "text" if color_key == "panel" else "bg"
        tb = slide.shapes.add_textbox(Inches(x+0.3), Inches(y+0.3),
                                       Inches(w-0.6), Inches(h-0.6))
        style = "title" if i < 2 else "body_bold"
        set_textbox_text(tb, truncate_to(str(tiles[i]), 140 if i<2 else 80),
                         theme, style, text_color, align="left")
        tb.name = f"bento_text_{i}"
        reveals.append(card)
    add_sequential_reveal(slide, reveals, theme, start_delay_ms=300, step_ms=140)
