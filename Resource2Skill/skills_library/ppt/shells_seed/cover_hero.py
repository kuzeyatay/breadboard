"""
shells_seed/cover_hero.py — Cover slide shell.

Slots:
  eyebrow     (text, optional, ≤30 chars, caption style)      — kicker line above title
  headline    (text, required, ≤60 chars, title_xl style)     — main title
  subhead     (text, optional, ≤120 chars, subtitle style)    — supporting sentence
  hero_image  (image, optional, fills right ~55% of canvas)   — optional image path

Layout: title stack on left 45%, hero image or accent panel on right 55%.
Works on both dark and light themes (color via theme palette only).
"""
from __future__ import annotations

from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_morph_anchor, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'cover'
DESCRIPTION = 'Left-aligned title stack + right panel accent or hero image; bold cover.'
ARCHETYPE = 'product'
MOOD = ['bold', 'editorial']
DENSITY = 'balanced'
STYLE_TAGS = ['split', 'asymmetric', 'hero_image']

SLOTS = [
    {"name": "eyebrow",    "kind": "text",  "max_chars":  30, "style": "caption",   "required": False},
    {"name": "headline",   "kind": "text",  "max_chars":  60, "style": "title_xl",  "required": True},
    {"name": "subhead",    "kind": "text",  "max_chars": 120, "style": "subtitle",  "required": False},
    {"name": "hero_image", "kind": "image", "aspect": "4:5",  "required": False, "fallback": "solid_fill:accent"},
]


def render(slide, slots: dict, theme: dict) -> None:
    """Render a cover slide."""
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background fill (full canvas bg color — assumes blank layout 6)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    bg.shadow.inherit = False

    # --- Right panel: hero image or accent block (56% of width) ---
    right_x = 7.3
    right_w = 13.333 - right_x - margin
    hero_path = get_slot(slots, "hero_image")
    if hero_path:
        try:
            slide.shapes.add_picture(
                hero_path, Inches(right_x), Inches(margin),
                Inches(right_w), Inches(7.5 - 2 * margin),
            )
        except Exception:
            hero_path = None  # fall through to solid panel
    if not hero_path:
        # Fallback: accent-filled panel
        add_solid_rect(
            slide, right_x, margin, right_w, 7.5 - 2 * margin,
            theme, color_key="accent", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )

    # --- Left text stack ---
    text_x = margin + 0.2
    text_w = right_x - text_x - 0.3

    y_cursor = 1.4
    reveal_shapes = []

    # Eyebrow (kicker)
    eyebrow = get_slot(slots, "eyebrow")
    if eyebrow:
        eyebrow = truncate_to(eyebrow, 30)
        eb = slide.shapes.add_textbox(Inches(text_x), Inches(y_cursor), Inches(text_w), Inches(0.35))
        set_textbox_text(eb, eyebrow, theme, "caption", color_key="accent")
        reveal_shapes.append(eb)
        y_cursor += 0.55

        # Motif: editorial themes get a hairline rule below the eyebrow
        if theme.get("motif", {}).get("type") == "thin_rule":
            add_hairline(slide, text_x, y_cursor, min(text_w, 1.8), theme)
            y_cursor += 0.25

    # Headline (required)
    headline = truncate_to(get_slot(slots, "headline", required=True), 60)
    head_box = slide.shapes.add_textbox(Inches(text_x), Inches(y_cursor), Inches(text_w), Inches(2.2))
    set_textbox_text(head_box, headline, theme, "title_xl", color_key="text")
    set_morph_anchor(head_box, "hero_headline")
    reveal_shapes.append(head_box)
    y_cursor += 2.3

    # Subhead
    subhead = get_slot(slots, "subhead")
    if subhead:
        subhead = truncate_to(subhead, 120)
        sub_box = slide.shapes.add_textbox(Inches(text_x), Inches(y_cursor), Inches(text_w), Inches(1.4))
        set_textbox_text(sub_box, subhead, theme, "subtitle", color_key="muted")
        reveal_shapes.append(sub_box)

    # --- Motion: theme-driven entrance, staggered ---
    for i, shape in enumerate(reveal_shapes):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
