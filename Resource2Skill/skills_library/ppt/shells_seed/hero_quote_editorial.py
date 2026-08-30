"""
shells_seed/hero_quote_editorial.py — editorial pull-quote slide.

Slots:
  context      (text, optional, ≤40 chars, caption style)
  quote        (text, required, ≤200 chars, accent_serif/title style)
  attribution  (text, optional, ≤60 chars, body style)

Layout: huge serif quotation mark top-left, quote body centered, attribution
below-right. Editorial/research feel.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_theme_entrance, get_slot, palette_color,
    set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = "hero_quote"
DESCRIPTION = "Editorial pull quote with large serif glyph, centered body, right-side attribution."
ARCHETYPE = "research"
MOOD = ["editorial", "calm", "restrained"]
DENSITY = "sparse"
STYLE_TAGS = ["centered", "serif_accent", "pull_quote"]

SLOTS = [
    {"name": "context",     "kind": "text", "max_chars":  40, "style": "caption",      "required": False},
    {"name": "quote",       "kind": "text", "max_chars": 200, "style": "accent_serif", "required": True},
    {"name": "attribution", "kind": "text", "max_chars":  60, "style": "body",         "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    margin = theme.get("spacing", {}).get("margin", 0.6)
    # background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    # Decorative large quotation mark
    glyph = slide.shapes.add_textbox(Inches(margin), Inches(0.9),
                                      Inches(2.5), Inches(2.5))
    glyph.text_frame.text = "\u201C"
    run = glyph.text_frame.paragraphs[0].runs[0]
    fallback = theme.get("typography", {}).get("accent_serif", {})
    run.font.name = fallback.get("font", "Playfair Display")
    run.font.size = Pt(160)
    run.font.italic = True
    run.font.color.rgb = palette_color(theme, "accent")
    glyph.name = "quote_glyph"

    # Context eyebrow (top)
    context = get_slot(slots, "context")
    if context:
        eb = slide.shapes.add_textbox(Inches(margin), Inches(0.7),
                                       Inches(6.0), Inches(0.4))
        set_textbox_text(eb, context, theme, "caption", "muted", align="left")
        eb.name = "context_eyebrow"
        add_theme_entrance(slide, eb, theme, delay_ms=120, index=0)

    # Hairline separator
    add_hairline(slide, margin, 1.1, 5.5, theme, "border")

    # Quote body — centered big block
    quote_text = truncate_to(get_slot(slots, "quote", required=True), 200)
    body = slide.shapes.add_textbox(Inches(1.5), Inches(2.6),
                                     Inches(10.333), Inches(3.0))
    set_textbox_text(body, quote_text, theme, "accent_serif",
                     color_key="text", align="center")
    body.name = "quote"
    add_theme_entrance(slide, body, theme, delay_ms=300, index=1)

    # Attribution
    attrib = get_slot(slots, "attribution")
    if attrib:
        at = slide.shapes.add_textbox(Inches(6.0), Inches(5.9),
                                       Inches(6.5), Inches(0.5))
        set_textbox_text(at, f"— {attrib}", theme, "body",
                         color_key="muted", align="right")
        at.name = "attribution"
        add_theme_entrance(slide, at, theme, delay_ms=600, index=2)

    # Bottom hairline
    add_hairline(slide, margin, 7.5 - margin - 0.1,
                 13.333 - 2 * margin, theme, "border")
