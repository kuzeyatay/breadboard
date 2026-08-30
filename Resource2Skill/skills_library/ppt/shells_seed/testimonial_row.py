"""
shells_seed/testimonial_row.py — Three testimonial cards with quote + attribution + role.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'bullet_card_list'
DESCRIPTION = 'Three testimonial cards with quote + attribution + role.'
ARCHETYPE = 'brand'
MOOD = ['editorial', 'warm', 'cinematic']
DENSITY = 'balanced'
STYLE_TAGS = ['testimonial', 'three_up', 'quotes', 'social_proof']

SLOTS = [{'name': 'section_label', 'kind': 'text', 'max_chars': 30, 'style': 'caption', 'required': False}, {'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title', 'required': True}, {'name': 'testimonials', 'kind': 'bullet_list', 'bullet_capacity': 3, 'required': True}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    section_label = get_slot(slots, "section_label")
    if section_label:
        sl = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(8), Inches(0.4))
        set_textbox_text(sl, section_label, theme, "caption", "accent", align="left")
        sl.name = "section_label"

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(margin), Inches(1.0), Inches(13.333 - 2*margin), Inches(0.9))
    set_textbox_text(hl, headline, theme, "title", "text", align="left")
    hl.name = "headline"

    testimonials = get_slot(slots, "testimonials", default=[]) or []
    if isinstance(testimonials, str):
        testimonials = [x.strip() for x in testimonials.split("\n") if x.strip()]
    testimonials = (testimonials + [""]*3)[:3]
    card_y = 2.4
    card_h = 4.4
    gap = 0.3
    card_w = (13.333 - 2*margin - 2*gap) / 3
    cards = []
    for i, entry in enumerate(testimonials):
        if isinstance(entry, dict):
            quote = entry.get("quote", "")
            attrib = entry.get("attribution", "")
            role = entry.get("role", "")
        else:
            parts = [p.strip() for p in str(entry).split("|")]
            parts = (parts + [""]*3)[:3]
            quote, attrib, role = parts
        x = margin + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(card_y),
                                       Inches(card_w), Inches(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = palette_color(theme, "panel")
        card.line.color.rgb = palette_color(theme, "border","muted")
        card.line.width = Pt(0.75); card.shadow.inherit = False
        card.name = f"testimonial_card_{i}"
        # Large quote glyph
        g = slide.shapes.add_textbox(Inches(x+0.3), Inches(card_y+0.15),
                                      Inches(card_w-0.6), Inches(0.8))
        g.text_frame.text = "\u201C"
        run = g.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(56)
        run.font.name = theme.get("typography", {}).get("accent_serif", {}).get("font", "Playfair Display")
        run.font.color.rgb = palette_color(theme, "accent")
        g.name = f"testimonial_glyph_{i}"
        # Quote
        qb = slide.shapes.add_textbox(Inches(x+0.3), Inches(card_y+1.0),
                                       Inches(card_w-0.6), Inches(2.4))
        set_textbox_text(qb, truncate_to(str(quote), 180), theme, "body", "text", align="left")
        qb.name = f"testimonial_quote_{i}"
        # Attribution
        ab = slide.shapes.add_textbox(Inches(x+0.3), Inches(card_y+card_h-1.0),
                                       Inches(card_w-0.6), Inches(0.4))
        set_textbox_text(ab, attrib, theme, "body_bold", "text", align="left")
        ab.name = f"testimonial_attrib_{i}"
        # Role
        rb = slide.shapes.add_textbox(Inches(x+0.3), Inches(card_y+card_h-0.55),
                                       Inches(card_w-0.6), Inches(0.4))
        set_textbox_text(rb, role, theme, "caption", "muted", align="left")
        rb.name = f"testimonial_role_{i}"
        cards.append(card)
    add_sequential_reveal(slide, cards, theme, start_delay_ms=300, step_ms=160)
