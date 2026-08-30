"""
shells_seed/photo_caption_story.py — Left 55% photo with overlaid caption block on right 45% — editorial feature page.

Auto-authored by scripts/distill_more_shells.py on 2026-04-19. Hand-
crafted (not LLM-generated) for robustness; no empty placeholder boxes.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from _shell_helpers import (
    add_hairline, add_sequential_reveal, add_solid_rect, add_theme_entrance,
    get_slot, palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'cover'
DESCRIPTION = 'Left 55% photo with overlaid caption block on right 45% — editorial feature page.'
ARCHETYPE = 'narrative'
MOOD = ['editorial', 'cinematic', 'calm']
DENSITY = 'sparse'
STYLE_TAGS = ['photo', 'editorial', 'split', 'caption']

SLOTS = [{'name': 'eyebrow', 'kind': 'text', 'max_chars': 30, 'style': 'caption', 'required': False}, {'name': 'headline', 'kind': 'text', 'max_chars': 80, 'style': 'title_xl', 'required': True}, {'name': 'story', 'kind': 'text', 'max_chars': 200, 'style': 'body', 'required': False}, {'name': 'byline', 'kind': 'text', 'max_chars': 40, 'style': 'caption', 'required': False}, {'name': 'hero_image', 'kind': 'image', 'aspect': '4:5', 'required': False, 'fallback': 'solid_fill:accent'}]


def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background(); bg.shadow.inherit = False

    # Left panel (photo or accent)
    left_w = 7.1
    hero_path = get_slot(slots, "hero_image")
    if hero_path:
        try:
            slide.shapes.add_picture(hero_path, Inches(0), Inches(0),
                                      Inches(left_w), Inches(7.5))
        except Exception:
            hero_path = None
    if not hero_path:
        pn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     Inches(left_w), Inches(7.5))
        pn.fill.solid(); pn.fill.fore_color.rgb = palette_color(theme, "accent")
        pn.line.fill.background(); pn.shadow.inherit = False
        pn.name = "hero_image"

    right_x = left_w + margin
    right_w = 13.333 - right_x - margin
    eyebrow = get_slot(slots, "eyebrow")
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(right_x), Inches(margin + 0.2),
                                       Inches(right_w), Inches(0.4))
        set_textbox_text(eb, eyebrow, theme, "caption", "accent", align="left")
        eb.name = "eyebrow"
        add_theme_entrance(slide, eb, theme, delay_ms=100, index=0)

    add_hairline(slide, right_x, 1.4, right_w, theme, "border")

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    hl = slide.shapes.add_textbox(Inches(right_x), Inches(1.7),
                                   Inches(right_w), Inches(2.4))
    set_textbox_text(hl, headline, theme, "title_xl", "text", align="left")
    hl.name = "headline"
    add_theme_entrance(slide, hl, theme, delay_ms=260, index=1)

    story = get_slot(slots, "story")
    if story:
        st = slide.shapes.add_textbox(Inches(right_x), Inches(4.3),
                                       Inches(right_w), Inches(2.0))
        set_textbox_text(st, story, theme, "body", "muted", align="left")
        st.name = "story"
        add_theme_entrance(slide, st, theme, delay_ms=460, index=2)

    byline = get_slot(slots, "byline")
    if byline:
        by = slide.shapes.add_textbox(Inches(right_x), Inches(7.5 - margin - 0.3),
                                       Inches(right_w), Inches(0.3))
        set_textbox_text(by, byline, theme, "caption", "muted", align="left")
        by.name = "byline"
