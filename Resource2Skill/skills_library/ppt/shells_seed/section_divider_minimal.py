"""
shells_seed/section_divider_minimal.py — Minimalist section divider.

Slots:
  section_number (text, optional, e.g. "01")   — big decorative numeral
  section_label  (text, required, ≤40 chars)   — "PROBLEM" / "PRODUCT" etc.
  tagline        (text, optional, ≤160 chars)  — supporting sentence below
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_theme_entrance, apply_type_style, get_slot,
    palette_color, set_morph_anchor, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'section_divider'
DESCRIPTION = 'Full-bleed section break with decorative number and label.'
ARCHETYPE = 'narrative'
MOOD = ['editorial', 'restrained']
DENSITY = 'sparse'
STYLE_TAGS = ['minimal', 'serif_numeral', 'centered']

SLOTS = [
    {"name": "section_number", "kind": "text", "max_chars":  4,  "style": "accent_serif", "required": False},
    {"name": "section_label",  "kind": "text", "max_chars": 40,  "style": "title_xl",     "required": True},
    {"name": "tagline",        "kind": "text", "max_chars": 160, "style": "subtitle",     "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Decorative section number (huge, serif, accent color) — left-aligned big
    section_number = get_slot(slots, "section_number")
    if section_number:
        num_box = slide.shapes.add_textbox(Inches(margin), Inches(1.1), Inches(5.0), Inches(3.2))
        num_tf = num_box.text_frame
        num_tf.word_wrap = True
        p = num_tf.paragraphs[0]
        p.text = truncate_to(section_number, 4)
        if p.runs:
            style_token = "accent_serif" if "accent_serif" in theme.get("typography", {}) else "metric_xl"
            apply_type_style(p.runs[0], theme, style_token, color_key="accent")
        reveal.append(num_box)

    # Hairline rule above label (editorial themes get this)
    label_y = 4.6
    if theme.get("motif", {}).get("type") == "thin_rule":
        add_hairline(slide, margin + 0.1, label_y - 0.25, 1.6, theme, color_key="accent")

    # Section label (big)
    label = truncate_to(get_slot(slots, "section_label", required=True), 40)
    label_box = slide.shapes.add_textbox(
        Inches(margin), Inches(label_y), Inches(13.333 - 2 * margin), Inches(1.4),
    )
    set_textbox_text(label_box, label, theme, "title_xl", color_key="text")
    set_morph_anchor(label_box, "section_chip")
    reveal.append(label_box)

    # Tagline
    tagline = get_slot(slots, "tagline")
    if tagline:
        tag_box = slide.shapes.add_textbox(
            Inches(margin), Inches(label_y + 1.3), Inches(10), Inches(1.5),
        )
        set_textbox_text(tag_box, truncate_to(tagline, 160), theme, "subtitle", color_key="muted")
        reveal.append(tag_box)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=700, index=i)
