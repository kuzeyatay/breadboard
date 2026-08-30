"""
shells_seed/hero_quote_reveal.py — Full-bleed quote with cinematic mask wipe.

Use case: dramatic single-page testimonial or thesis statement.
Built around add_mask_wipe so the quote text is revealed by a sweeping
mask, then the attribution fades in.

Slots:
  quote        (text, required, max 220 chars — the hero quote)
  attribution  (text, required, max 80 chars — name, title)
  context      (text, optional, max 60 chars, caption — e.g. "WHY WE BELIEVE")
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from _shell_helpers import (
    add_mask_wipe, add_theme_entrance, apply_type_style, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

# ---- LLM-readable shell metadata ----
ROLE = 'hero_quote'
DESCRIPTION = 'Large editorial quote with attribution; minimalist hero-quote layout.'
ARCHETYPE = 'research'
MOOD = ['editorial', 'calm']
DENSITY = 'sparse'
STYLE_TAGS = ['centered', 'editorial', 'serif_accent']

SLOTS = [
    {"name": "context",     "kind": "text", "max_chars": 60,  "style": "caption", "required": False},
    {"name": "quote",       "kind": "text", "max_chars": 220, "style": "title",   "required": True},
    {"name": "attribution", "kind": "text", "max_chars": 80,  "style": "subtitle","required": True},
]


def render(slide, slots: dict, theme: dict) -> None:
    # Background — use the theme's panel color for moody contrast against bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "panel", "bg")
    bg.line.fill.background()

    # Optional context label, top-center
    context = get_slot(slots, "context")
    if context:
        cb = slide.shapes.add_textbox(Inches(1.0), Inches(0.7), Inches(11.333), Inches(0.45))
        set_textbox_text(cb, truncate_to(context, 60), theme, "caption",
                         color_key="accent", align="center")
        add_theme_entrance(slide, cb, theme, delay_ms=200, index=0)

    # The quote — large, centered, full-width
    quote = truncate_to(get_slot(slots, "quote", required=True), 220)
    qb = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.333), Inches(3.5))
    tf = qb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = f"\u201c{quote}\u201d"  # smart quotes
    if p.runs:
        apply_type_style(p.runs[0], theme, "title_xl" if "title_xl" in theme.get("typography", {}) else "title",
                         color_key="text")
        # Bigger than title for hero feel
        from pptx.util import Pt
        n = len(quote)
        size = 44 if n <= 100 else 36 if n <= 160 else 30
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.italic = True
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER

    # Mask wipe so the quote sweeps in left-to-right
    add_mask_wipe(slide, qb, theme, delay_ms=600, duration_ms=900, direction="left_to_right")

    # Attribution, bottom-right
    attribution = truncate_to(get_slot(slots, "attribution", required=True), 80)
    ab = slide.shapes.add_textbox(Inches(7.0), Inches(6.2), Inches(5.333), Inches(0.6))
    set_textbox_text(ab, f"\u2014 {attribution}", theme, "subtitle",
                     color_key="muted", align="right")
    add_theme_entrance(slide, ab, theme, delay_ms=1700, index=1)

    # Hairline above the attribution for editorial polish
    hr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(11.333), Inches(6.05), Inches(1.0), Inches(0.02),
    )
    hr.fill.solid(); hr.fill.fore_color.rgb = palette_color(theme, "accent")
    hr.line.fill.background()
    add_theme_entrance(slide, hr, theme, delay_ms=1500, index=0)
