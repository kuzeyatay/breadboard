from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin_x = theme.get("spacing", {}).get("margin", 1.0)
    margin_y = theme.get("spacing", {}).get("margin", 1.0)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline_text = get_slot(slots, "headline", default="Click to add title")
    head = slide.shapes.add_textbox(Inches(margin_x), Inches(margin_y), Inches(13.333 - 2 * margin_x), Inches(1.2))
    set_textbox_text(head, truncate_to(headline_text, 80), theme, "title", color_key="text")
    reveal.append(head)

    # Bullets
    bullets = get_slot(slots, "bullets", required=True)
    if not isinstance(bullets, list):
        bullets = [{"title": "Bullet 1"}, {"title": "Bullet 2"}, {"title": "Bullet 3"}]

    y_offset = margin_y + 1.5
    bullet_spacing = 0.6

    # Render each bullet as a separate textbox to allow staggered animation as seen in the source image
    for i, item in enumerate(bullets[:7]):
        if isinstance(item, dict):
            text = item.get("title", "")
        else:
            text = str(item)

        bullet_text = f"•  {truncate_to(text, 120)}"
        b_box = slide.shapes.add_textbox(Inches(margin_x + 0.2), Inches(y_offset), Inches(13.333 - 2 * margin_x - 0.2), Inches(0.5))
        set_textbox_text(b_box, bullet_text, theme, "body", color_key="text")
        reveal.append(b_box)
        y_offset += bullet_spacing

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=400, index=i)
