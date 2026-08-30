"""
Curved Split Cover.

Slots:
  hero_image (image, required)
  headline   (text, required, title_xl)
  subtitle   (text, optional, subtitle)
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []

    # Hero Image (Left side)
    img_slot = get_slot(slots, "hero_image")
    if img_slot:
        pic = slide.shapes.add_picture(img_slot, Inches(0), Inches(0), width=Inches(8.5), height=Inches(7.5))
        reveal.append(pic)
    else:
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(8.5), Inches(7.5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = palette_color(theme, "muted")
        placeholder.line.fill.background()
        reveal.append(placeholder)

    # Curved overlay (Right side)
    # We use a large oval to create the curved split effect
    overlay = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.5), Inches(-2.5), Inches(12), Inches(12.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = palette_color(theme, "bg")
    overlay.line.fill.background()
    reveal.append(overlay)

    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(6.5), Inches(2.0), Inches(6.0), Inches(2.0))
    set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title_xl", color_key="text")
    reveal.append(head_box)

    # Subtitle
    subtitle_text = get_slot(slots, "subtitle")
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(6.5), Inches(4.5), Inches(6.0), Inches(1.5))
        set_textbox_text(sub_box, truncate_to(subtitle_text, 100), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
