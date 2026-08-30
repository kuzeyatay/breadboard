from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    y = margin

    # Headline with panel background
    headline = get_slot(slots, "headline", required=True)
    if headline:
        hl_bg = add_solid_rect(slide, margin, y, 8.5, 0.7, theme, color_key="panel")
        reveal.append(hl_bg)
        
        head_box = slide.shapes.add_textbox(Inches(margin + 0.15), Inches(y + 0.05), Inches(8.2), Inches(0.6))
        set_textbox_text(head_box, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head_box)
        y += 0.7

    # Subtitle with panel background
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_bg = add_solid_rect(slide, margin, y, 8.5, 0.5, theme, color_key="panel")
        reveal.append(sub_bg)

        sub_box = slide.shapes.add_textbox(Inches(margin + 0.15), Inches(y + 0.05), Inches(8.2), Inches(0.4))
        set_textbox_text(sub_box, truncate_to(subtitle, 120), theme, "subtitle", color_key="text")
        reveal.append(sub_box)
        y += 0.7
    else:
        y += 0.2

    # Chart Title
    chart_title = get_slot(slots, "chart_title")
    if chart_title:
        ct_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(6.0), Inches(0.6))
        set_textbox_text(ct_box, truncate_to(chart_title, 80), theme, "body_bold", color_key="text")
        reveal.append(ct_box)
        y += 0.7

    # Chart Placeholder
    chart_h = 7.0 - y - 0.2
    chart_w = 13.333 - (2 * margin)
    chart_ph = add_solid_rect(slide, margin, y, chart_w, chart_h, theme, color_key="panel")
    reveal.append(chart_ph)

    # Source
    source = get_slot(slots, "source")
    if source:
        src_box = slide.shapes.add_textbox(Inches(margin), Inches(7.0), Inches(10.0), Inches(0.3))
        set_textbox_text(src_box, truncate_to(source, 100), theme, "caption", color_key="muted")
        reveal.append(src_box)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
