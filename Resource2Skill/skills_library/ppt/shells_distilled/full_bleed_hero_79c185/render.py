from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_theme_entrance, get_slot, set_textbox_text, truncate_to, palette_color

def render(slide, slots: dict, theme: dict) -> None:
    # 1. Full Bleed Image
    img_path = get_slot(slots, "hero_image")
    if img_path:
        pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        add_theme_entrance(slide, pic, theme, index=0)
    else:
        # Fallback background if image is missing
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = palette_color(theme, "panel")
        bg.line.fill.background()
        add_theme_entrance(slide, bg, theme, index=0)

    # 2. Optional Text Overlay
    headline = get_slot(slots, "headline")
    subtitle = get_slot(slots, "subtitle")

    if headline or subtitle:
        margin = theme.get("spacing", {}).get("margin", 0.8)
        y_pos = Inches(4.8) # Lower third positioning
        reveal_shapes = []

        if headline:
            head_box = slide.shapes.add_textbox(Inches(margin), y_pos, Inches(13.333 - 2 * margin), Inches(1.5))
            # Use 'bg' color key for light text on dark image assumption
            set_textbox_text(head_box, truncate_to(headline, 60), theme, "title_xl", color_key="bg")
            reveal_shapes.append(head_box)
            y_pos += Inches(1.5)

        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(margin), y_pos, Inches(13.333 - 2 * margin), Inches(1.0))
            set_textbox_text(sub_box, truncate_to(subtitle, 100), theme, "subtitle", color_key="bg")
            reveal_shapes.append(sub_box)

        for i, shape in enumerate(reveal_shapes):
            add_theme_entrance(slide, shape, theme, delay_ms=400 + (i * 150), index=i+1)