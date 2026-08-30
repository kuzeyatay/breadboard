from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Perspective Grid
    vx, vy = 6.666, 3.75
    lines = [
        (0, vy, 13.333, vy),          # Horizon
        (0, 6.5, 13.333, 6.5),        # Bottom horizontal
        (vx, vy, vx, 7.5),            # Vertical center
        (vx, vy, 4.5, 7.5),           # Left radiating 1
        (vx, vy, 2.0, 7.5),           # Left radiating 2
        (vx, vy, 0, 6.0),             # Left radiating 3
        (vx, vy, 8.8, 7.5),           # Right radiating 1
        (vx, vy, 11.3, 7.5),          # Right radiating 2
        (vx, vy, 13.333, 6.0),        # Right radiating 3
    ]
    
    for x1, y1, x2, y2 in lines:
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.line.color.rgb = palette_color(theme, "muted")
        reveal.append(connector)
        
    # Content Box (Label)
    box_x, box_y = 1.0, 5.5
    box_w, box_h = 4.0, 0.8
    
    panel = add_solid_rect(slide, box_x, box_y, box_w, box_h, theme, color_key="panel", line=False)
    reveal.append(panel)
    
    headline = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(box_x + 0.2), Inches(box_y + 0.1), Inches(box_w - 0.4), Inches(0.6))
    set_textbox_text(head_box, truncate_to(headline, 40), theme, "title", color_key="text")
    reveal.append(head_box)
    
    # Staggered entrance for grid lines and label
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
