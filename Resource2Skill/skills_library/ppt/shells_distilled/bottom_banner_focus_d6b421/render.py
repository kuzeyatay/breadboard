from __future__ import annotations
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.5)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Center area (body)
    body_text = get_slot(slots, "body")
    if body_text:
        panel_w = 8.0
        panel_h = 3.0
        panel_x = (13.333 - panel_w) / 2
        panel_y = (7.5 - panel_h) / 2 - 0.5
        
        panel = add_solid_rect(slide, panel_x, panel_y, panel_w, panel_h, theme, color_key="panel", line=True)
        reveal.append(panel)
        
        body_box = slide.shapes.add_textbox(Inches(panel_x + 0.5), Inches(panel_y + 0.5), Inches(panel_w - 1.0), Inches(panel_h - 1.0))
        set_textbox_text(body_box, truncate_to(body_text, 200), theme, "body", color_key="text")
        reveal.append(body_box)

    # Bottom Banner
    headline = get_slot(slots, "headline", required=True)
    banner_w = 9.0
    banner_h = 1.2
    banner_x = margin
    banner_y = 7.5 - margin - banner_h
    
    banner = add_solid_rect(slide, banner_x, banner_y, banner_w, banner_h, theme, color_key="accent")
    reveal.append(banner)
    
    head_box = slide.shapes.add_textbox(Inches(banner_x + 0.3), Inches(banner_y + 0.1), Inches(banner_w - 0.6), Inches(banner_h - 0.2))
    set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="bg")
    reveal.append(head_box)
    
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
