from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    gutter = theme.get("spacing", {}).get("gutter", 0.4)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(1.0))
    set_textbox_text(head, headline, theme, "title", color_key="text")
    reveal.append(head)
    
    # Content Area
    content_y = margin + 1.2
    content_h = 7.5 - content_y - margin
    col_w = (13.333 - 2 * margin - gutter) / 2
    img_h = min(col_w * 0.75, content_h) # 4:3 aspect ratio constraint
    
    # Left Image
    img_path = get_slot(slots, "hero_image")
    if img_path:
        try:
            pic = slide.shapes.add_picture(img_path, Inches(margin), Inches(content_y), width=Inches(col_w), height=Inches(img_h))
            reveal.append(pic)
        except Exception:
            pic = add_solid_rect(slide, margin, content_y, col_w, img_h, theme, color_key="muted")
            reveal.append(pic)
    else:
        pic = add_solid_rect(slide, margin, content_y, col_w, img_h, theme, color_key="muted")
        reveal.append(pic)
        
    # Right Body
    body = get_slot(slots, "body")
    if body:
        body_box = slide.shapes.add_textbox(Inches(margin + col_w + gutter), Inches(content_y), Inches(col_w), Inches(content_h))
        set_textbox_text(body_box, truncate_to(body, 500), theme, "body", color_key="text")
        reveal.append(body_box)
        
    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
