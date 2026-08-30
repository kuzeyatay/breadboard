from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import palette_color, set_textbox_text, add_theme_entrance, get_slot, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    y = margin
    
    # Optional Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.8))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)
        y += 1.0
        
    # Hero Image (Screenshot)
    img_path = get_slot(slots, "hero_image", required=True)
    if img_path:
        remaining_h = 7.5 - y - margin - 0.5
        pic = slide.shapes.add_picture(img_path, Inches(margin), Inches(y), width=Inches(13.333 - 2 * margin), height=Inches(remaining_h))
        reveal.append(pic)
        y += remaining_h + 0.1
        
    # Optional Caption
    caption = get_slot(slots, "caption")
    if caption:
        cap = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.4))
        set_textbox_text(cap, truncate_to(caption, 100), theme, "caption", color_key="muted")
        reveal.append(cap)
        
    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
