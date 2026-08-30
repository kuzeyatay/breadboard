from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "headline", "kind": "text", "max_chars": 100, "style": "title_xl", "required": True},
    {"name": "subtitle", "kind": "text", "max_chars": 150, "style": "subtitle", "required": False},
]


def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 100)
    head_box = slide.shapes.add_textbox(Inches(margin), Inches(2.8), Inches(13.333 - 2 * margin), Inches(1.5))
    set_textbox_text(head_box, headline, theme, "title_xl", color_key="text")
    # Center align text
    for paragraph in head_box.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    reveal.append(head_box)

    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(margin), Inches(4.3), Inches(13.333 - 2 * margin), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subtitle, 150), theme, "subtitle", color_key="muted")
        # Center align text
        for paragraph in sub_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        reveal.append(sub_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
