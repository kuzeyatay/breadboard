from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    slide_w, slide_h = 13.333, 7.5
    w1, w2, w3, w4 = 4.333, 3.0, 3.0, 3.0
    margin = 0.6
    reveal = []

    bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(w1), Inches(slide_h))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = palette_color(theme, "bg")
    bg1.line.fill.background()
    reveal.append(bg1)

    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(margin), Inches(2.0), Inches(w1 - 2*margin), Inches(1.5))
        set_textbox_text(head_box, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head_box)

    intro_body = get_slot(slots, "intro_body")
    if intro_body:
        body_box = slide.shapes.add_textbox(Inches(margin), Inches(3.5), Inches(w1 - 2*margin), Inches(2.5))
        set_textbox_text(body_box, truncate_to(intro_body, 200), theme, "body", color_key="muted")
        reveal.append(body_box)

    img_slot = get_slot(slots, "hero_image")
    if img_slot and isinstance(img_slot, str):
        try:
            pic = slide.shapes.add_picture(img_slot, Inches(w1), Inches(0), Inches(w2), Inches(slide_h))
            reveal.append(pic)
        except Exception:
            pic = add_solid_rect(slide, w1, 0, w2, slide_h, theme, "muted")
            reveal.append(pic)
    else:
        pic = add_solid_rect(slide, w1, 0, w2, slide_h, theme, "muted")
        reveal.append(pic)

    num1 = slide.shapes.add_textbox(Inches(w1 + w2 - 1.0), Inches(slide_h - 1.2), Inches(0.8), Inches(0.8))
    set_textbox_text(num1, "01", theme, "body", color_key="bg")
    reveal.append(num1)

    features = get_slot(slots, "features", default=[])
    def get_feature_text(idx):
        if idx < len(features):
            item = features[idx]
            if isinstance(item, dict):
                return item.get("body", item.get("title", ""))
            return str(item)
        return ""

    bg3 = add_solid_rect(slide, w1 + w2, 0, w3, slide_h, theme, "panel")
    reveal.append(bg3)

    acc3 = add_solid_rect(slide, w1 + w2 + 0.5, 1.0, 0.4, 0.05, theme, "accent")
    reveal.append(acc3)

    body3_text = get_feature_text(0)
    if body3_text:
        body3 = slide.shapes.add_textbox(Inches(w1 + w2 + 0.5), Inches(3.0), Inches(w3 - 1.0), Inches(3.0))
        set_textbox_text(body3, truncate_to(body3_text, 150), theme, "body", color_key="text")
        reveal.append(body3)

    num2 = slide.shapes.add_textbox(Inches(w1 + w2 + w3 - 1.0), Inches(slide_h - 1.2), Inches(0.8), Inches(0.8))
    set_textbox_text(num2, "02", theme, "body", color_key="text")
    reveal.append(num2)

    bg4 = add_solid_rect(slide, w1 + w2 + w3, 0, w4, slide_h, theme, "text")
    reveal.append(bg4)

    acc4 = add_solid_rect(slide, w1 + w2 + w3 + 0.5, 1.0, 0.4, 0.05, theme, "accent")
    reveal.append(acc4)

    body4_text = get_feature_text(1)
    if body4_text:
        body4 = slide.shapes.add_textbox(Inches(w1 + w2 + w3 + 0.5), Inches(3.0), Inches(w4 - 1.0), Inches(3.0))
        set_textbox_text(body4, truncate_to(body4_text, 150), theme, "body", color_key="bg")
        reveal.append(body4)

    num3 = slide.shapes.add_textbox(Inches(w1 + w2 + w3 + w4 - 1.0), Inches(slide_h - 1.2), Inches(0.8), Inches(0.8))
    set_textbox_text(num3, "03", theme, "body", color_key="bg")
    reveal.append(num3)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
