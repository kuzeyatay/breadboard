from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Headline
    headline = get_slot(slots, "headline", required=True)
    head = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(11.333), Inches(1.2))
    set_textbox_text(head, truncate_to(headline, 50), theme, "title", color_key="text")
    
    items = get_slot(slots, "items", required=True)
    if not isinstance(items, list) or len(items) != 3:
        items = [{"title": f"Item {i+1}", "body": "Description goes here."} for i in range(3)]

    D = 4.2
    overlap = 0.9
    start_x = (13.333 - (3 * D - 2 * overlap)) / 2
    y = 2.0

    # Draw order: Left (0), Right (2), Middle (1) so middle is on top visually
    positions = [
        (0, start_x),
        (2, start_x + 2 * (D - overlap)),
        (1, start_x + D - overlap)
    ]

    item_shapes = {0: [], 1: [], 2: []}

    for idx, x in positions:
        item = items[idx]
        
        # Circle
        circle = add_solid_rect(
            slide, x, y, D, D, theme,
            color_key="panel", line=True, shape_type=MSO_SHAPE.OVAL
        )
        item_shapes[idx].append(circle)

        # Inscribed text box area
        box_size = D * 0.65
        offset = (D - box_size) / 2
        tx = x + offset
        ty = y + offset + 0.2

        t_box = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(box_size), Inches(0.6))
        set_textbox_text(t_box, truncate_to(item.get("title", ""), 30), theme, "body_bold", color_key="text")
        item_shapes[idx].append(t_box)

        b_box = slide.shapes.add_textbox(Inches(tx), Inches(ty + 0.6), Inches(box_size), Inches(box_size - 0.6))
        set_textbox_text(b_box, truncate_to(item.get("body", ""), 80), theme, "body", color_key="muted")
        item_shapes[idx].append(b_box)

    # Animation order: Headline, then Left, Middle, Right
    reveal = [head]
    reveal.extend(item_shapes[0])
    reveal.extend(item_shapes[1])
    reveal.extend(item_shapes[2])

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
