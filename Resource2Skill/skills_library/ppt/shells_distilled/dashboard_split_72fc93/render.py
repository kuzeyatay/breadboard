from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.5)
    reveal = []

    # Background
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")
    reveal.append(bg)

    y = margin

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(12.333), Inches(0.8))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)
        y += 1.0

    # Hero Image (Dashboard Diagram)
    img_path = get_slot(slots, "hero_image")
    img_h = 3.5
    img_w = 13.333 - 2 * margin
    if img_path:
        try:
            pic = slide.shapes.add_picture(img_path, Inches(margin), Inches(y), width=Inches(img_w), height=Inches(img_h))
            reveal.append(pic)
        except Exception:
            placeholder = add_solid_rect(slide, margin, y, img_w, img_h, theme, color_key="panel")
            reveal.append(placeholder)
    else:
        placeholder = add_solid_rect(slide, margin, y, img_w, img_h, theme, color_key="panel")
        reveal.append(placeholder)

    y += img_h + 0.4

    # Data Table (Bottom)
    data_list = get_slot(slots, "data_list")
    if data_list and isinstance(data_list, list):
        # Table Header
        header_h = 0.4
        header_bg = add_solid_rect(slide, margin, y, img_w, header_h, theme, color_key="accent")
        reveal.append(header_bg)

        col1_w = img_w * 0.3
        col2_w = img_w * 0.7

        h_title = slide.shapes.add_textbox(Inches(margin + 0.1), Inches(y), Inches(col1_w - 0.2), Inches(header_h))
        set_textbox_text(h_title, "Parameter", theme, "body_bold", color_key="bg")
        reveal.append(h_title)

        h_val = slide.shapes.add_textbox(Inches(margin + col1_w + 0.1), Inches(y), Inches(col2_w - 0.2), Inches(header_h))
        set_textbox_text(h_val, "Value / Status", theme, "body_bold", color_key="bg")
        reveal.append(h_val)

        y += header_h

        # Table Rows
        row_h = 0.35
        max_rows = 5
        for i, item in enumerate(data_list[:max_rows]):
            if not isinstance(item, dict):
                item = {"title": str(item), "body": ""}

            # Alternating row background
            row_color = "panel" if i % 2 == 0 else "bg"
            row_bg = add_solid_rect(slide, margin, y, img_w, row_h, theme, color_key=row_color)
            
            # Add a subtle border to the row
            row_bg.line.fill.solid()
            row_bg.line.fill.fore_color.rgb = palette_color(theme, "muted")
            reveal.append(row_bg)

            r_title = slide.shapes.add_textbox(Inches(margin + 0.1), Inches(y), Inches(col1_w - 0.2), Inches(row_h))
            set_textbox_text(r_title, truncate_to(item.get("title", ""), 50), theme, "body", color_key="text")
            reveal.append(r_title)

            r_val = slide.shapes.add_textbox(Inches(margin + col1_w + 0.1), Inches(y), Inches(col2_w - 0.2), Inches(row_h))
            set_textbox_text(r_val, truncate_to(item.get("body", ""), 100), theme, "body", color_key="text")
            reveal.append(r_val)

            y += row_h

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
