from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []

    # Large Accent Rectangle
    outer_w = 9.0
    outer_h = 5.0
    outer_x = (13.333 - outer_w) / 2
    outer_y = (7.5 - outer_h) / 2

    outer_rect = add_solid_rect(
        slide, outer_x, outer_y, outer_w, outer_h, theme,
        color_key="accent", line=False, shape_type=MSO_SHAPE.RECTANGLE
    )
    reveal.append(outer_rect)

    # Inner Panel (Rounded Rectangle)
    inner_w = 4.5
    inner_h = 1.8
    inner_x = (13.333 - inner_w) / 2
    inner_y = (7.5 - inner_h) / 2

    inner_rect = add_solid_rect(
        slide, inner_x, inner_y, inner_w, inner_h, theme,
        color_key="panel", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    inner_rect.adjustments[0] = 0.15
    reveal.append(inner_rect)

    # Text
    headline = get_slot(slots, "headline", required=True)
    subhead = get_slot(slots, "subhead")

    text_y = inner_y + 0.3
    if not subhead:
        text_y = inner_y + 0.5

    head_box = slide.shapes.add_textbox(Inches(inner_x + 0.2), Inches(text_y), Inches(inner_w - 0.4), Inches(0.6))
    set_textbox_text(head_box, truncate_to(headline, 40), theme, "title", color_key="text")
    for p in head_box.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    reveal.append(head_box)

    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(inner_x + 0.2), Inches(text_y + 0.6), Inches(inner_w - 0.4), Inches(0.5))
        set_textbox_text(sub_box, truncate_to(subhead, 60), theme, "subtitle", color_key="muted")
        for p in sub_box.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        reveal.append(sub_box)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
