from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import palette_color, set_textbox_text, add_theme_entrance, get_slot, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)

    # Features
    features = get_slot(slots, "features", required=True)
    if not isinstance(features, list):
        features = []
    features = features[:3]

    for i, item in enumerate(features):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}

        # Staggered positioning
        cx = 3.166 + i * 3.5
        if i == 1:
            cy = 2.5
            y_text = 4.5
        else:
            cy = 3.5
            y_text = 5.5

        # Magnifying Glass Handle (drawn first to sit behind the circle border)
        handle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx + 0.5), Inches(cy + 0.5), Inches(0.4), Inches(1.5))
        handle.fill.solid()
        handle.fill.fore_color.rgb = palette_color(theme, "text")
        handle.line.fill.background()
        handle.rotation = 315
        reveal.append(handle)

        # Magnifying Glass Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 1.0), Inches(cy - 1.0), Inches(2.0), Inches(2.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = palette_color(theme, "accent")
        circle.line.color.rgb = palette_color(theme, "text")
        circle.line.width = Inches(0.1)
        reveal.append(circle)

        # Inner reflection highlight
        refl = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.7), Inches(cy - 0.7), Inches(0.4), Inches(0.4))
        refl.fill.solid()
        refl.fill.fore_color.rgb = palette_color(theme, "bg")
        refl.line.fill.background()
        reveal.append(refl)

        # Number
        num_box = slide.shapes.add_textbox(Inches(cx - 0.4), Inches(cy - 0.4), Inches(0.8), Inches(0.8))
        set_textbox_text(num_box, f"{i+1:02d}", theme, "title", color_key="text")
        reveal.append(num_box)

        # Text Title
        t_box = slide.shapes.add_textbox(Inches(cx - 1.5), Inches(y_text), Inches(3.0), Inches(0.5))
        set_textbox_text(t_box, truncate_to(item.get("title", ""), 40), theme, "body_bold", color_key="text")
        reveal.append(t_box)

        # Text Body
        b_box = slide.shapes.add_textbox(Inches(cx - 1.5), Inches(y_text + 0.5), Inches(3.0), Inches(1.0))
        set_textbox_text(b_box, truncate_to(item.get("body", ""), 100), theme, "body", color_key="muted")
        reveal.append(b_box)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50*i, index=i)
