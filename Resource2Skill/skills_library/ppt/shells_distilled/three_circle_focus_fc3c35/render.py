from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

SLOTS = [
    {"name": "headline", "kind": "text", "style": "title", "required": False},
    {"name": "labels", "kind": "bullet_list", "bullet_capacity": 3, "required": True},
    {"name": "description", "kind": "text", "style": "body", "required": False}
]

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        head.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(head)

    # Labels (3 circles)
    labels = get_slot(slots, "labels", required=True)
    if not isinstance(labels, list):
        labels = [{"title": "A"}, {"title": "B"}, {"title": "C"}]
    while len(labels) < 3:
        labels.append({"title": ""})
    labels = labels[:3]

    cx_list = [3.333, 6.666, 10.0]
    cy = 3.5
    outer_r = 1.6
    inner_r = 1.2

    for i, item in enumerate(labels):
        title_text = item.get("title", str(item)) if isinstance(item, dict) else str(item)
        cx = cx_list[i]

        # Outer ring
        outer = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - outer_r), Inches(cy - outer_r),
            Inches(outer_r * 2), Inches(outer_r * 2)
        )
        outer.fill.solid()
        outer.fill.fore_color.rgb = palette_color(theme, "panel")
        outer.line.color.rgb = palette_color(theme, "accent")
        outer.line.width = Pt(4)

        # Inner circle
        inner = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - inner_r), Inches(cy - inner_r),
            Inches(inner_r * 2), Inches(inner_r * 2)
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = palette_color(theme, "accent")
        inner.line.fill.background()

        # Small number badge (top left of outer circle)
        badge_size = 0.35
        badge_x = cx - outer_r + 0.1
        badge_y = cy - outer_r + 0.1
        
        badge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(badge_x), Inches(badge_y),
            Inches(badge_size), Inches(badge_size)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = palette_color(theme, "bg")
        badge.line.color.rgb = palette_color(theme, "muted")

        badge_text = slide.shapes.add_textbox(
            Inches(badge_x), Inches(badge_y),
            Inches(badge_size), Inches(badge_size)
        )
        set_textbox_text(badge_text, str(i + 1), theme, "caption", color_key="text")
        badge_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_text.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Main Text inside inner circle
        tb = slide.shapes.add_textbox(
            Inches(cx - inner_r), Inches(cy - inner_r),
            Inches(inner_r * 2), Inches(inner_r * 2)
        )
        set_textbox_text(tb, truncate_to(title_text, 3), theme, "metric_xl", color_key="bg")
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER

        reveal.extend([outer, inner, badge, badge_text, tb])

    # Description at bottom
    desc = get_slot(slots, "description")
    if desc:
        desc_box = slide.shapes.add_textbox(Inches(2.666), Inches(5.8), Inches(8.0), Inches(1.0))
        set_textbox_text(desc_box, truncate_to(desc, 150), theme, "body", color_key="text")
        desc_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(desc_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
