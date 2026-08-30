from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "muted")
    bg.line.fill.background()
    
    reveal = []
    
    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)

    steps = get_slot(slots, "steps", required=True)
    if not isinstance(steps, list):
        steps = []
    while len(steps) < 2:
        steps.append({"title": "Step Title", "body": "Description goes here."})
    steps = steps[:2]

    card_w = 4.8
    card_h = 1.6
    spacing = 0.8
    total_w = 2 * card_w + spacing
    start_x = (13.333 - total_w) / 2
    y = 3.0

    for i, item in enumerate(steps):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
            
        x = start_x + i * (card_w + spacing)
        
        # Ribbon tail (behind)
        ribbon = add_solid_rect(
            slide, x - 0.3, y + 0.3, 0.6, 1.0, theme,
            color_key="accent", line=False, shape_type=MSO_SHAPE.RECTANGLE
        )
        
        # Main card
        card = add_solid_rect(
            slide, x, y, card_w, card_h, theme,
            color_key="bg", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
        )
        card.adjustments[0] = 0.1
        
        # Circle
        circle = add_solid_rect(
            slide, x - 0.5, y + 0.3, 1.0, 1.0, theme,
            color_key="accent", line=False, shape_type=MSO_SHAPE.OVAL
        )
        
        # Step text inside circle
        step_box = slide.shapes.add_textbox(Inches(x - 0.5), Inches(y + 0.45), Inches(1.0), Inches(0.7))
        set_textbox_text(step_box, f"STEP\n0{i+1}", theme, "caption", color_key="bg")
        for paragraph in step_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            
        # Title
        t_box = slide.shapes.add_textbox(Inches(x + 0.7), Inches(y + 0.2), Inches(3.0), Inches(0.5))
        set_textbox_text(t_box, truncate_to(item.get("title", "TITLE"), 30), theme, "body_bold", color_key="accent")
        
        # Body
        b_box = slide.shapes.add_textbox(Inches(x + 0.7), Inches(y + 0.7), Inches(3.0), Inches(0.8))
        set_textbox_text(b_box, truncate_to(item.get("body", "Lorem ipsum dolor sit amet."), 100), theme, "body", color_key="text")
        
        # Decorative icon placeholder (right side)
        icon_ph = add_solid_rect(
            slide, x + 3.9, y + 0.5, 0.6, 0.6, theme,
            color_key="bg", line=True, shape_type=MSO_SHAPE.OVAL
        )
        icon_ph.line.color.rgb = palette_color(theme, "accent")
        icon_ph.line.width = Inches(0.02)
        
        icon_inner = add_solid_rect(
            slide, x + 4.0, y + 0.6, 0.4, 0.4, theme,
            color_key="bg", line=True, shape_type=MSO_SHAPE.OVAL
        )
        icon_inner.line.color.rgb = palette_color(theme, "accent")
        
        reveal.extend([ribbon, card, circle, step_box, t_box, b_box, icon_ph, icon_inner])

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50*i, index=i)
