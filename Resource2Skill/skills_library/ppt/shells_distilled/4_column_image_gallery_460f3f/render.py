from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, add_emphasis_pulse, add_sequential_reveal,
    get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.8)
    gutter = theme.get("spacing", {}).get("gutter", 0.3)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    y_offset = 0
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        add_theme_entrance(slide, head, theme, delay_ms=0, index=0)
        y_offset = 0.6

    img_slots = ["image_1", "image_2", "image_3", "image_4"]
    images = [get_slot(slots, name) for name in img_slots]
    
    n = 4
    total_w = 13.333 - 2 * margin
    img_w = (total_w - (n - 1) * gutter) / n
    img_h = img_w * 1.45  # Portrait aspect ratio
    
    y = (7.5 - img_h) / 2 + y_offset / 2
    
    shapes_to_reveal = []
    for i, img_path in enumerate(images):
        x = margin + i * (img_w + gutter)
        shape_added = False
        if img_path:
            try:
                pic = slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(img_w), Inches(img_h))
                shapes_to_reveal.append(pic)
                shape_added = True
            except Exception:
                pass
        
        if not shape_added:
            rect = add_solid_rect(slide, x, y, img_w, img_h, theme, color_key="panel")
            shapes_to_reveal.append(rect)
            
    add_sequential_reveal(slide, shapes_to_reveal, theme)
    
    for i, shp in enumerate(shapes_to_reveal):
        add_emphasis_pulse(slide, shp, theme, delay_ms=2000 + 100 * i)
