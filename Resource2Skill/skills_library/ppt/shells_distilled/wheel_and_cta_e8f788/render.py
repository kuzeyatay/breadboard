def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from _shell_helpers import palette_color, set_textbox_text, add_theme_entrance, get_slot, truncate_to

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Central Circle
    cx, cy = 6.666, 3.2
    radius = 2.6
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - radius), Inches(cy - radius), 
        Inches(radius * 2), Inches(radius * 2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = palette_color(theme, "panel")
    circle.line.fill.background()
    reveal.append(circle)
    
    line_color = palette_color(theme, "text")
    
    def add_spoke(x1, y1, x2, y2):
        spoke = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        spoke.line.color.rgb = line_color
        spoke.line.width = Inches(0.08)
        reveal.append(spoke)

    add_spoke(cx, cy - radius, cx, cy + radius)
    add_spoke(cx - radius, cy, cx + radius, cy)
    
    offset = radius * 0.7071
    add_spoke(cx - offset, cy - offset, cx + offset, cy + offset)
    add_spoke(cx - offset, cy + offset, cx + offset, cy - offset)
    
    # Outer Button (Pill)
    btn_w, btn_h = 5.5, 1.0
    btn_x = cx - btn_w / 2
    btn_y = 6.1
    
    outer_btn = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(btn_x), Inches(btn_y), Inches(btn_w), Inches(btn_h)
    )
    outer_btn.fill.solid()
    outer_btn.fill.fore_color.rgb = palette_color(theme, "bg")
    outer_btn.line.color.rgb = palette_color(theme, "accent")
    outer_btn.line.width = Inches(0.06)
    outer_btn.adjustments[0] = 0.5
    reveal.append(outer_btn)
    
    # Inner Button
    inner_w, inner_h = 2.8, 0.6
    inner_x = cx - inner_w / 2
    inner_y = btn_y + (btn_h - inner_h) / 2
    
    inner_btn = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(inner_x), Inches(inner_y), Inches(inner_w), Inches(inner_h)
    )
    inner_btn.fill.solid()
    inner_btn.fill.fore_color.rgb = palette_color(theme, "text")
    inner_btn.line.fill.background()
    inner_btn.adjustments[0] = 0.2
    reveal.append(inner_btn)
    
    # Button Text
    headline = get_slot(slots, "headline", default="SUBSCRIBE")
    text_box = slide.shapes.add_textbox(
        Inches(inner_x), Inches(inner_y + 0.05), Inches(inner_w), Inches(inner_h)
    )
    set_textbox_text(text_box, truncate_to(headline, 20), theme, "body_bold", color_key="bg")
    reveal.append(text_box)
    
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
