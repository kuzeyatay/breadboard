from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []
    
    # 1. Hero Image (Background)
    hero_img_path = get_slot(slots, "hero_image")
    if hero_img_path:
        try:
            pic = slide.shapes.add_picture(hero_img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            reveal.append(pic)
        except Exception:
            pass
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = palette_color(theme, "bg")
        bg.line.fill.background()
        reveal.append(bg)

    # 2. Logo
    logo_path = get_slot(slots, "logo")
    if logo_path:
        logo_w = 3.0
        logo_x = (13.333 - logo_w) / 2
        logo_y = 2.0
        try:
            logo_pic = slide.shapes.add_picture(logo_path, Inches(logo_x), Inches(logo_y), width=Inches(logo_w))
            reveal.append(logo_pic)
        except Exception:
            pass

    # 3. Headline
    headline_text = get_slot(slots, "headline")
    if headline_text:
        head_w = 10.0
        head_h = 1.2
        head_x = (13.333 - head_w) / 2
        head_y = 4.8
        head_box = slide.shapes.add_textbox(Inches(head_x), Inches(head_y), Inches(head_w), Inches(head_h))
        # Using 'bg' color key assuming a light theme where 'bg' is white, providing contrast against a dark photo
        set_textbox_text(head_box, truncate_to(headline_text, 40), theme, "title_xl", color_key="bg")
        for paragraph in head_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        reveal.append(head_box)

    # 4. Subhead
    subhead_text = get_slot(slots, "subhead")
    if subhead_text:
        sub_w = 8.0
        sub_h = 0.8
        sub_x = (13.333 - sub_w) / 2
        sub_y = 6.0
        sub_box = slide.shapes.add_textbox(Inches(sub_x), Inches(sub_y), Inches(sub_w), Inches(sub_h))
        set_textbox_text(sub_box, truncate_to(subhead_text, 60), theme, "subtitle", color_key="bg")
        for paragraph in sub_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        reveal.append(sub_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
