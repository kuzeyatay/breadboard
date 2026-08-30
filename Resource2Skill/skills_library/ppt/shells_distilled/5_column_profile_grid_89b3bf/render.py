from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin_x = 0.6
    margin_y = 1.5
    col_w = 2.1
    gutter = 0.4
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    # Top accent bar
    top_bar = add_solid_rect(slide, 0, 0, 13.333, 0.15, theme, color_key="accent", line=False)
    
    reveal = [top_bar]
    
    profiles = get_slot(slots, "profiles", required=True)
    if not isinstance(profiles, list):
        profiles = []
        
    for i in range(5):
        x_col = margin_x + i * (col_w + gutter)
        
        # Image
        img_slot_name = f"image_{i+1}"
        img_path = get_slot(slots, img_slot_name)
        img_w = 1.6
        x_img = x_col + (col_w - img_w) / 2
        y_img = margin_y
        
        if img_path:
            try:
                pic = slide.shapes.add_picture(img_path, Inches(x_img), Inches(y_img), width=Inches(img_w), height=Inches(img_w))
                reveal.append(pic)
            except Exception:
                placeholder = add_solid_rect(slide, x_img, y_img, img_w, img_w, theme, color_key="muted", line=False)
                reveal.append(placeholder)
        else:
            placeholder = add_solid_rect(slide, x_img, y_img, img_w, img_w, theme, color_key="muted", line=False)
            reveal.append(placeholder)
            
        # Text
        if i < len(profiles):
            item = profiles[i]
            if not isinstance(item, dict):
                item = {"title": str(item), "body": ""}
            
            # Title
            t_box = slide.shapes.add_textbox(Inches(x_col), Inches(y_img + img_w + 0.2), Inches(col_w), Inches(0.5))
            set_textbox_text(t_box, truncate_to(item.get("title", ""), 30), theme, "title", color_key="text")
            reveal.append(t_box)
            
            # Body
            b_box = slide.shapes.add_textbox(Inches(x_col), Inches(y_img + img_w + 0.8), Inches(col_w), Inches(2.5))
            set_textbox_text(b_box, truncate_to(item.get("body", ""), 150), theme, "body", color_key="text")
            reveal.append(b_box)
            
    # Footer
    footer_text = get_slot(slots, "footer")
    if footer_text:
        f_box = slide.shapes.add_textbox(Inches(margin_x), Inches(6.5), Inches(13.333 - 2 * margin_x), Inches(0.8))
        set_textbox_text(f_box, truncate_to(footer_text, 100), theme, "title", color_key="text")
        reveal.append(f_box)
        
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
