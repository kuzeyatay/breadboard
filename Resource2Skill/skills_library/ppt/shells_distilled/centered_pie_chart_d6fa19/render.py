from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(0.8))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)

    # Pie Chart Placeholder (Circle)
    pie_radius = 2.2
    pie_x = (13.333 / 2) - pie_radius
    pie_y = 3.2 - pie_radius
    pie = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(pie_x), Inches(pie_y), Inches(pie_radius * 2), Inches(pie_radius * 2))
    pie.fill.solid(); pie.fill.fore_color.rgb = palette_color(theme, "accent")
    pie.line.color.rgb = palette_color(theme, "bg")
    pie.line.width = Inches(0.05)
    reveal.append(pie)

    # Legend
    data_points = get_slot(slots, "data_points")
    if not data_points:
        data_points = [{"title": f"Item {i+1}"} for i in range(5)]

    num_items = len(data_points)
    if num_items > 0:
        legend_y = 6.0
        max_items_per_row = 7
        item_w = (13.333 - 2 * margin) / min(num_items, max_items_per_row)

        for i, item in enumerate(data_points):
            if not isinstance(item, dict):
                item = {"title": str(item)}

            row = i // max_items_per_row
            col = i % max_items_per_row

            items_in_this_row = min(num_items - row * max_items_per_row, max_items_per_row)
            row_width = items_in_this_row * item_w
            start_x = (13.333 - row_width) / 2

            x = start_x + col * item_w
            y = legend_y + row * 0.4

            # Legend color box
            box_size = 0.15
            box = add_solid_rect(slide, x, y + 0.1, box_size, box_size, theme, color_key="panel" if i % 2 == 0 else "muted")
            reveal.append(box)

            # Legend text
            tb = slide.shapes.add_textbox(Inches(x + box_size + 0.05), Inches(y), Inches(item_w - box_size - 0.05), Inches(0.35))
            set_textbox_text(tb, truncate_to(item.get("title", ""), 15), theme, "caption", color_key="text")
            reveal.append(tb)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
