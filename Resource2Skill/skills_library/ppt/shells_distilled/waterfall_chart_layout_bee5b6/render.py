from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "headline", "kind": "text", "max_chars": 80, "style": "title", "required": True},
    {"name": "chart_data", "kind": "chart", "required": True},
]

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = get_slot(slots, "headline", default="Financial Overview")
    head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(1.0))
    set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
    reveal.append(head)

    # Mock Waterfall Chart (Placeholder for chart_data)
    chart_y = margin + 1.2
    chart_h = 7.5 - chart_y - margin
    chart_w = 13.333 - 2 * margin

    # X-axis
    axis_y = chart_y + chart_h - 0.5
    axis = add_solid_rect(
        slide, margin, axis_y, chart_w, 0.02, theme,
        color_key="muted", shape_type=MSO_SHAPE.RECTANGLE, line=False
    )
    reveal.append(axis)

    # Bars
    num_bars = 6
    bar_w = 0.8
    spacing = (chart_w - (num_bars * bar_w)) / (num_bars + 1)

    scale = 4.0 / 100
    steps = [100, -20, -10, -15, -5, 50]
    labels = ["Start", "Factor A", "Factor B", "Factor C", "Factor D", "End"]
    current_val = 100

    for i, val in enumerate(steps):
        bx = margin + spacing + i * (bar_w + spacing)

        if i == 0:
            bh = val * scale
            by = axis_y - bh
            color = "accent"
            lbl_text = str(val)
        elif i == len(steps) - 1:
            bh = val * scale
            by = axis_y - bh
            color = "accent"
            lbl_text = str(val)
        else:
            bh = abs(val) * scale
            by = axis_y - (current_val * scale)
            color = "muted"
            current_val += val
            lbl_text = str(val)

        bar = add_solid_rect(
            slide, bx, by, bar_w, bh, theme,
            color_key=color, shape_type=MSO_SHAPE.RECTANGLE, line=False
        )
        reveal.append(bar)

        # Value label
        lbl_val = slide.shapes.add_textbox(Inches(bx - 0.2), Inches(by - 0.4), Inches(bar_w + 0.4), Inches(0.3))
        set_textbox_text(lbl_val, lbl_text, theme, "caption", color_key="text")
        reveal.append(lbl_val)

        # Step label
        lbl_step = slide.shapes.add_textbox(Inches(bx - 0.4), Inches(axis_y + 0.1), Inches(bar_w + 0.8), Inches(0.6))
        set_textbox_text(lbl_step, labels[i], theme, "caption", color_key="text")
        reveal.append(lbl_step)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
