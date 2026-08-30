from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Card dimensions
    card_w = 8.0
    card_h = 4.0
    x = (13.333 - card_w) / 2
    y = (7.5 - card_h) / 2

    # Card panel
    card = add_solid_rect(
        slide, x, y, card_w, card_h, theme,
        color_key="panel", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    card.adjustments[0] = 0.08
    reveal.append(card)

    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    subhead = get_slot(slots, "subhead")

    start_y = y + 1.2
    if subhead:
        start_y = y + 0.8

    head_box = slide.shapes.add_textbox(Inches(x + 0.8), Inches(start_y), Inches(card_w - 1.6), Inches(1.0))
    set_textbox_text(head_box, headline, theme, "title", color_key="text")
    reveal.append(head_box)

    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(x + 0.8), Inches(start_y + 1.2), Inches(card_w - 1.6), Inches(1.5))
        set_textbox_text(sub_box, truncate_to(subhead, 150), theme, "body", color_key="muted")
        reveal.append(sub_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
