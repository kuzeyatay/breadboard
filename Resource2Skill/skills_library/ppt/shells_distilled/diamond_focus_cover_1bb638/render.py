from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    cx, cy = 13.333 / 2, 7.5 / 2
    reveal = []

    # Background
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, "bg")
    reveal.append(bg)

    img_slot = get_slot(slots, "hero_image")
    has_img = False
    if img_slot:
        try:
            pic = slide.shapes.add_picture(img_slot, Inches(0), Inches(0), width=Inches(13.333))
            reveal.append(pic)
            has_img = True
        except Exception:
            pass

    S = 4.0
    D = 4.2
    offset = D * 0.7071

    if not has_img:
        # Adjacent diamonds (Grid pattern)
        adj_positions = [
            (cx + offset, cy - offset),
            (cx + offset, cy + offset),
            (cx - offset, cy - offset),
            (cx - offset, cy + offset),
            (cx + 2*offset, cy),
            (cx - 2*offset, cy),
        ]

        for ax, ay in adj_positions:
            adj = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ax - S/2), Inches(ay - S/2), Inches(S), Inches(S))
            adj.rotation = 45
            adj.adjustments[0] = 0.15
            adj.fill.solid()
            adj.fill.fore_color.rgb = palette_color(theme, "panel")
            adj.line.color.rgb = palette_color(theme, "bg")
            adj.line.width = Inches(0.05)
            reveal.append(adj)

    # Central Diamond
    center_diamond = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - S/2), Inches(cy - S/2), Inches(S), Inches(S))
    center_diamond.rotation = 45
    center_diamond.adjustments[0] = 0.15
    center_diamond.fill.solid()
    center_diamond.fill.fore_color.rgb = palette_color(theme, "bg" if has_img else "accent")
    center_diamond.line.color.rgb = palette_color(theme, "accent" if has_img else "bg")
    center_diamond.line.width = Inches(0.08)
    reveal.append(center_diamond)

    # Text
    text_w = 2.8
    text_color = "text" if has_img else "bg"
    muted_color = "muted" if has_img else "bg"
    
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(cx - text_w/2), Inches(cy - 0.8), Inches(text_w), Inches(1.0))
        set_textbox_text(head_box, truncate_to(headline, 40), theme, "title", color_key=text_color)
        for p in head_box.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        reveal.append(head_box)

    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(cx - text_w/2), Inches(cy + 0.2), Inches(text_w), Inches(0.8))
        set_textbox_text(sub_box, truncate_to(subtitle, 80), theme, "body", color_key=muted_color)
        for p in sub_box.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        reveal.append(sub_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
