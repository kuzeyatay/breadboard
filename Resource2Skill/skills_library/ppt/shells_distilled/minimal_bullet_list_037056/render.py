from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    y = margin
    
    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)
        y += 1.2
        
    # Bullets
    bullets = get_slot(slots, "bullets", required=True)
    if isinstance(bullets, list):
        for i, item in enumerate(bullets):
            if not isinstance(item, dict):
                item = {"title": str(item), "body": ""}
            
            title_text = item.get("title", "")
            if title_text:
                t_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.5))
                set_textbox_text(t_box, f"• {truncate_to(title_text, 100)}", theme, "body_bold", color_key="text")
                reveal.append(t_box)
                y += 0.4
                
            body_text = item.get("body", "")
            if body_text:
                b_box = slide.shapes.add_textbox(Inches(margin + 0.3), Inches(y), Inches(13.333 - 2 * margin - 0.3), Inches(0.5))
                set_textbox_text(b_box, truncate_to(body_text, 200), theme, "body", color_key="muted")
                reveal.append(b_box)
                y += 0.5
                
            y += 0.2  # Spacing between bullet items

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
