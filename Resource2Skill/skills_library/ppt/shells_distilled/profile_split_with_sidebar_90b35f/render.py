def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from _shell_helpers import palette_color, set_textbox_text, add_solid_rect, add_theme_entrance, get_slot, truncate_to

    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Photo (Right side)
    photo_x, photo_y, photo_w, photo_h = 9.0, 0.6, 3.7, 6.3
    photo_path = get_slot(slots, "photo")
    photo_shape = None
    if photo_path and isinstance(photo_path, str):
        try:
            photo_shape = slide.shapes.add_picture(photo_path, Inches(photo_x), Inches(photo_y), width=Inches(photo_w), height=Inches(photo_h))
        except Exception:
            pass
    
    if not photo_shape:
        photo_shape = add_solid_rect(slide, photo_x, photo_y, photo_w, photo_h, theme, color_key="muted")
    
    reveal.append(photo_shape)

    # Left Column: Roles
    primary_role = get_slot(slots, "primary_role", default="Primary Role")
    secondary_roles = get_slot(slots, "secondary_roles", default=[])
    
    pill_y = 3.2
    pill = add_solid_rect(slide, 0.6, pill_y, 2.8, 0.8, theme, color_key="accent", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    pill.adjustments[0] = 0.5
    reveal.append(pill)
    
    role_box = slide.shapes.add_textbox(Inches(0.9), Inches(pill_y + 0.15), Inches(2.2), Inches(0.5))
    set_textbox_text(role_box, truncate_to(primary_role.upper(), 30), theme, "body_bold", color_key="bg")
    reveal.append(role_box)
    
    if isinstance(secondary_roles, list):
        for i, role in enumerate(secondary_roles[:5]):
            role_text = role.get("title", "") if isinstance(role, dict) else str(role)
            if i == 0:
                y_pos = pill_y - 0.6
            else:
                y_pos = pill_y + 1.0 + (i - 1) * 0.6
                
            box = slide.shapes.add_textbox(Inches(0.9), Inches(y_pos), Inches(2.5), Inches(0.5))
            set_textbox_text(box, truncate_to(role_text.upper(), 30), theme, "caption", color_key="text")
            reveal.append(box)

    # Middle Column: Name & Bio
    first_name = get_slot(slots, "first_name", default="First Name")
    last_name = get_slot(slots, "last_name", default="Last Name")
    bio = get_slot(slots, "bio")
    
    fn_box = slide.shapes.add_textbox(Inches(4.0), Inches(2.8), Inches(4.5), Inches(1.0))
    set_textbox_text(fn_box, truncate_to(first_name.upper(), 20), theme, "title_xl", color_key="text")
    reveal.append(fn_box)
    
    ln_box = slide.shapes.add_textbox(Inches(5.0), Inches(3.6), Inches(4.5), Inches(1.0))
    set_textbox_text(ln_box, truncate_to(last_name.upper(), 20), theme, "title_xl", color_key="text")
    reveal.append(ln_box)
    
    if bio:
        bio_box = slide.shapes.add_textbox(Inches(4.0), Inches(4.8), Inches(4.5), Inches(2.0))
        set_textbox_text(bio_box, truncate_to(bio, 300), theme, "body", color_key="muted")
        reveal.append(bio_box)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
