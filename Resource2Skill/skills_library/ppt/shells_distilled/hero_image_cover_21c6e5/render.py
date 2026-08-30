from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "headline", "kind": "text", "max_chars": 80, "style": "title_xl", "required": True},
    {"name": "hero_image", "kind": "image", "required": True},
]

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = get_slot(slots, "headline", required=True)
    head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(1.5))
    set_textbox_text(head, truncate_to(headline, 80), theme, "title_xl", color_key="text")
    reveal.append(head)

    # Hero Image
    hero_image = get_slot(slots, "hero_image")
    if hero_image:
        img_w = 8.0
        img_h = 4.5
        img_x = (13.333 - img_w) / 2
        img_y = 2.5
        try:
            pic = slide.shapes.add_picture(hero_image, Inches(img_x), Inches(img_y), width=Inches(img_w), height=Inches(img_h))
            reveal.append(pic)
        except Exception:
            # Fallback if image path is invalid or missing in test environment
            rect = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="panel")
            reveal.append(rect)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
