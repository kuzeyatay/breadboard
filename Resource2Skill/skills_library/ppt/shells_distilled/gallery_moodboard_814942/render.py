from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    gutter = theme.get("spacing", {}).get("gutter", 0.2)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(0.8))
        set_textbox_text(head, truncate_to(headline, 60), theme, "title", color_key="text")
        reveal.append(head)

    y_start = 1.6
    content_h = 7.5 - y_start - margin
    hero_w = 7.5
    grid_w = 13.333 - 2 * margin - hero_w - gutter

    def place_image_slot(slot_name, x, y, w, h):
        img_path = get_slot(slots, slot_name)
        shape_added = False
        if img_path:
            try:
                pic = slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
                reveal.append(pic)
                shape_added = True
            except Exception:
                pass
        if not shape_added:
            # Fallback placeholder if image is missing or fails to load
            rect = add_solid_rect(slide, x, y, w, h, theme, color_key="panel")
            reveal.append(rect)

    # Hero Image
    hero_h = content_h - 0.5 # Leave room for caption
    place_image_slot("hero_image", margin, y_start, hero_w, hero_h)

    # Hero Caption
    caption = get_slot(slots, "hero_caption")
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(margin), Inches(y_start + hero_h + 0.1), Inches(hero_w), Inches(0.4))
        set_textbox_text(cap_box, truncate_to(caption, 80), theme, "caption", color_key="muted")
        reveal.append(cap_box)

    # Grid Images (2x2)
    grid_x_start = margin + hero_w + gutter
    img_w = (grid_w - gutter) / 2
    img_h = (content_h - gutter) / 2

    place_image_slot("grid_image_1", grid_x_start, y_start, img_w, img_h)
    place_image_slot("grid_image_2", grid_x_start + img_w + gutter, y_start, img_w, img_h)
    place_image_slot("grid_image_3", grid_x_start, y_start + img_h + gutter, img_w, img_h)
    place_image_slot("grid_image_4", grid_x_start + img_w + gutter, y_start + img_h + gutter, img_w, img_h)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50 * i, duration_ms=500, index=i)
