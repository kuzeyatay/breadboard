from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, add_emphasis_pulse,
    get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Left angled shapes (using rotated rectangles to create the '\' angle)
    # Dark blue shape (leftmost)
    shape1 = add_solid_rect(slide, -3.0, -2.0, 7.0, 11.5, theme, color_key="accent2")
    shape1.rotation = 345
    reveal.append(shape1)

    # Bright blue shape (middle)
    shape2 = add_solid_rect(slide, 3.5, -2.0, 3.5, 11.5, theme, color_key="accent")
    shape2.rotation = 345
    reveal.append(shape2)

    # Text Content
    headline_text = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(7.2), Inches(2.8), Inches(5.5), Inches(1.5))
    set_textbox_text(head_box, truncate_to(headline_text, 40), theme, "title_xl", color_key="text")
    reveal.append(head_box)

    subtitle_text = get_slot(slots, "subtitle")
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(7.2), Inches(4.4), Inches(5.5), Inches(0.8))
        set_textbox_text(sub_box, truncate_to(subtitle_text, 100), theme, "body", color_key="muted")
        reveal.append(sub_box)

    # Decorative block (acts as a visual anchor or faux CTA)
    dec_block = add_solid_rect(slide, 7.2, 5.5, 2.0, 0.5, theme, color_key="accent")
    reveal.append(dec_block)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)

    add_emphasis_pulse(slide, dec_block, theme, delay_ms=1500)
