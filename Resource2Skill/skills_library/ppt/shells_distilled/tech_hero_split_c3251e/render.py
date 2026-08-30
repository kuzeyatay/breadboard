from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_hairline, add_theme_entrance, add_emphasis_pulse,
    get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Logo
    logo_text = get_slot(slots, "logo")
    if logo_text:
        logo = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(3), Inches(0.4))
        set_textbox_text(logo, truncate_to(logo_text, 20), theme, "body_bold", color_key="text")
        reveal.append(logo)

    # Top Right Nav
    nav_text = get_slot(slots, "top_right_nav")
    if nav_text:
        nav = slide.shapes.add_textbox(Inches(13.333 - margin - 4), Inches(margin), Inches(4), Inches(0.4))
        set_textbox_text(nav, truncate_to(nav_text, 40), theme, "caption", color_key="muted")
        reveal.append(nav)

    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head = slide.shapes.add_textbox(Inches(margin), Inches(1.8), Inches(8.0), Inches(3.5))
    set_textbox_text(head, truncate_to(headline_text, 60), theme, "title_xl", color_key="text")
    reveal.append(head)

    # Subhead
    subhead_text = get_slot(slots, "subhead")
    if subhead_text:
        sub = slide.shapes.add_textbox(Inches(9.0), Inches(3.5), Inches(3.5), Inches(1.0))
        set_textbox_text(sub, truncate_to(subhead_text, 100), theme, "body", color_key="muted")
        reveal.append(sub)

    # CTA
    cta_text = get_slot(slots, "cta")
    cta_box = None
    if cta_text:
        cta_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), Inches(4.8), Inches(2.5), Inches(0.6))
        cta_box.fill.background()
        cta_box.line.color.rgb = palette_color(theme, "accent")
        cta_box.line.width = Inches(0.02)
        
        cta_tb = slide.shapes.add_textbox(Inches(9.0), Inches(4.8), Inches(2.5), Inches(0.6))
        set_textbox_text(cta_tb, truncate_to(cta_text, 20), theme, "body_bold", color_key="text")
        reveal.extend([cta_box, cta_tb])

    # Decorative elements
    add_hairline(slide, margin, 6.8, 13.333 - 2*margin, theme)
    
    dec_text = slide.shapes.add_textbox(Inches(margin), Inches(6.9), Inches(4), Inches(0.4))
    set_textbox_text(dec_text, "/// SYSTEM READY", theme, "caption", color_key="accent")
    reveal.append(dec_text)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
        
    if cta_box:
        add_emphasis_pulse(slide, cta_box, theme, delay_ms=1500)
