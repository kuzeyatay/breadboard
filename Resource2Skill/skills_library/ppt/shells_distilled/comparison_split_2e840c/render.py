from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")
    
    reveal = []

    # Headline
    headline_text = get_slot(slots, "headline", default="Product Comparison")
    head = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.0))
    set_textbox_text(head, headline_text, theme, "title", color_key="text")
    head.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.append(head)

    # Center Spine Background (adds depth for the comparison)
    spine = add_solid_rect(slide, 5.833, 1.8, 1.666, 5.2, theme, color_key="panel")
    reveal.append(spine)

    # --- Product A Header (Left) ---
    prod_a_name = get_slot(slots, "product_a_name", default="Product A")
    prod_a_desc = get_slot(slots, "product_a_desc", default="")

    badge_a = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.033), Inches(1.8), Inches(0.8), Inches(0.8))
    badge_a.fill.solid()
    badge_a.fill.fore_color.rgb = palette_color(theme, "accent")
    badge_a.line.fill.background()
    set_textbox_text(badge_a, "A", theme, "title", color_key="bg")
    badge_a.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.append(badge_a)

    box_a_name = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(3.3), Inches(0.4))
    set_textbox_text(box_a_name, prod_a_name, theme, "subtitle", color_key="accent")
    box_a_name.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    reveal.append(box_a_name)

    if prod_a_desc:
        box_a_desc = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(3.3), Inches(0.8))
        set_textbox_text(box_a_desc, truncate_to(prod_a_desc, 120), theme, "body", color_key="muted")
        box_a_desc.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        reveal.append(box_a_desc)

    # --- Product B Header (Right) ---
    prod_b_name = get_slot(slots, "product_b_name", default="Product B")
    prod_b_desc = get_slot(slots, "product_b_desc", default="")

    badge_b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(1.8), Inches(0.8), Inches(0.8))
    badge_b.fill.solid()
    badge_b.fill.fore_color.rgb = palette_color(theme, "accent2")
    badge_b.line.fill.background()
    set_textbox_text(badge_b, "B", theme, "title", color_key="bg")
    badge_b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.append(badge_b)

    box_b_name = slide.shapes.add_textbox(Inches(8.533), Inches(1.8), Inches(3.3), Inches(0.4))
    set_textbox_text(box_b_name, prod_b_name, theme, "subtitle", color_key="accent2")
    box_b_name.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    reveal.append(box_b_name)

    if prod_b_desc:
        box_b_desc = slide.shapes.add_textbox(Inches(8.533), Inches(2.2), Inches(3.3), Inches(0.8))
        set_textbox_text(box_b_desc, truncate_to(prod_b_desc, 120), theme, "body", color_key="muted")
        box_b_desc.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        reveal.append(box_b_desc)

    # --- Features / Comparison Rows ---
    features = get_slot(slots, "features", default=[])
    if not features:
        features = [{"label": "", "val_a": "", "val_b": ""} for _ in range(5)]

    y_start = 3.2
    row_h = 0.5
    gap = 0.3

    for i, feat in enumerate(features[:5]):
        if not isinstance(feat, dict):
            feat = {"label": str(feat), "val_a": "", "val_b": ""}

        y = y_start + i * (row_h + gap)

        # Left Bar (Product A)
        bar_a = add_solid_rect(slide, 2.0, y, 3.833, row_h, theme, color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        bar_a.adjustments[0] = 0.5
        reveal.append(bar_a)

        val_a = feat.get("val_a", "")
        if val_a:
            tb_a = slide.shapes.add_textbox(Inches(2.2), Inches(y), Inches(3.433), row_h)
            set_textbox_text(tb_a, truncate_to(val_a, 40), theme, "body", color_key="text")
            tb_a.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            reveal.append(tb_a)

        # Right Bar (Product B)
        bar_b = add_solid_rect(slide, 7.5, y, 3.833, row_h, theme, color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        bar_b.adjustments[0] = 0.5
        reveal.append(bar_b)

        val_b = feat.get("val_b", "")
        if val_b:
            tb_b = slide.shapes.add_textbox(Inches(7.7), Inches(y), Inches(3.433), row_h)
            set_textbox_text(tb_b, truncate_to(val_b, 40), theme, "body", color_key="text")
            tb_b.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            reveal.append(tb_b)

        # Center Label (Feature Name)
        label = feat.get("label", feat.get("title", ""))
        if label:
            tb_lbl = slide.shapes.add_textbox(Inches(5.833), Inches(y), Inches(1.666), row_h)
            set_textbox_text(tb_lbl, truncate_to(label, 25), theme, "body_bold", color_key="text")
            tb_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            reveal.append(tb_lbl)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
