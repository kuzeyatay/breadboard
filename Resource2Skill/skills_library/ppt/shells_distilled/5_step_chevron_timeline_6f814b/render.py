from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    y = margin
    reveal = []

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.8))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)
        y += 1.0
    else:
        y += 0.5

    # Steps
    steps = get_slot(slots, "steps", required=True)
    if not isinstance(steps, list):
        steps = []
    
    num_steps = min(len(steps), 5)
    if num_steps == 0:
        return

    total_w = 13.333 - 2 * margin
    overlap = 0.5
    step_w = (total_w + (num_steps - 1) * overlap) / num_steps
    step_h = 7.5 - y - margin
    
    start_x = margin

    for i in range(num_steps):
        item = steps[i]
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
            
        x = start_x + i * (step_w - overlap)
        
        # First shape is flat on the left, subsequent shapes are notched
        shape_type = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        
        shape = slide.shapes.add_shape(
            shape_type, Inches(x), Inches(y), Inches(step_w), Inches(step_h)
        )
        shape.fill.solid()
        
        # Cycle through theme colors to distinguish steps
        color_keys = ["primary", "secondary", "accent"]
        shape.fill.fore_color.rgb = palette_color(theme, color_keys[i % len(color_keys)])
        
        # Add a background-colored border to create a gap effect
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = palette_color(theme, "bg")
        
        reveal.append(shape)

        # Text positioning adjusts based on whether the shape has a left notch
        text_x = x + 0.2 if i == 0 else x + overlap + 0.1
        text_w = step_w - overlap - 0.4 if i == 0 else step_w - overlap * 1.5 - 0.2
        
        # Title
        t_box = slide.shapes.add_textbox(
            Inches(text_x), Inches(y + 0.5), Inches(text_w), Inches(0.8)
        )
        set_textbox_text(
            t_box, truncate_to(item.get("title", ""), 40), theme, "body_bold", color_key="bg"
        )
        reveal.append(t_box)

        # Body
        b_box = slide.shapes.add_textbox(
            Inches(text_x), Inches(y + 1.3), Inches(text_w), Inches(step_h - 1.8)
        )
        set_textbox_text(
            b_box, truncate_to(item.get("body", ""), 120), theme, "caption", color_key="bg"
        )
        reveal.append(b_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
