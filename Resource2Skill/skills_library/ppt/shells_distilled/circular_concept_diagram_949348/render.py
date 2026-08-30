import math
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    cx = 4.5
    cy = 3.75
    donut_radius = 2.5

    # Donut Shape
    donut = slide.shapes.add_shape(
        MSO_SHAPE.DONUT, 
        Inches(cx - donut_radius), 
        Inches(cy - donut_radius), 
        Inches(donut_radius * 2), 
        Inches(donut_radius * 2)
    )
    donut.fill.solid()
    donut.fill.fore_color.rgb = palette_color(theme, "primary")
    donut.line.color.rgb = palette_color(theme, "primary")
    try:
        donut.adjustments[0] = 0.5  # 50% hole to make a thick ring
    except:
        pass
    reveal.append(donut)

    # Arrow Shape
    arrow_w = 6.0
    arrow_h = 1.2
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(cx + 1.0),  # Overlap with donut to merge shapes visually
        Inches(cy - arrow_h / 2),
        Inches(arrow_w),
        Inches(arrow_h)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = palette_color(theme, "primary")
    arrow.line.color.rgb = palette_color(theme, "primary")
    reveal.append(arrow)

    # Main Label (inside Arrow)
    label_text = get_slot(slots, "main_label", default="CORE CONCEPT")
    t_box = slide.shapes.add_textbox(
        Inches(cx + 1.5), 
        Inches(cy - arrow_h / 2), 
        Inches(arrow_w - 2.0), 
        Inches(arrow_h)
    )
    set_textbox_text(t_box, truncate_to(label_text, 30), theme, "title", color_key="bg")
    for paragraph in t_box.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    reveal.append(t_box)

    # Items distributed around the Donut
    items = get_slot(slots, "items", required=True)
    if not isinstance(items, list):
        items = [{"title": "Item 1"}, {"title": "Item 2"}, {"title": "Item 3"}, {"title": "Item 4"}]
    
    num_items = len(items)
    if num_items > 0:
        text_radius = 1.875  # Position text in the middle of the thick donut ring
        angle_step = 2 * math.pi / num_items
        
        for i, item in enumerate(items):
            if not isinstance(item, dict):
                item = {"title": str(item)}
            
            text_val = truncate_to(item.get("title", ""), 15).upper()
            
            # Calculate position (start at top, go clockwise)
            angle = i * angle_step - math.pi / 2
            
            tw = 1.4
            th = 0.4
            tx = cx + text_radius * math.cos(angle) - tw / 2
            ty = cy + text_radius * math.sin(angle) - th / 2
            
            el_box = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
            set_textbox_text(el_box, text_val, theme, "caption", color_key="bg")
            
            for paragraph in el_box.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                
            # Rotate text to follow the curve, keeping it readable
            rot_deg = math.degrees(angle) + 90
            rot_deg = rot_deg % 360
            if 90 < rot_deg < 270:
                rot_deg -= 180
            el_box.rotation = rot_deg
                
            reveal.append(el_box)

    # Staggered entrance animation
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
