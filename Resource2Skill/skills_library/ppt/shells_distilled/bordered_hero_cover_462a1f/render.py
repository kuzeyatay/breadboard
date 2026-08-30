def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from _shell_helpers import (
        add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
    )

    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []

    # Hero Image
    hero_image = get_slot(slots, "hero_image")
    if hero_image:
        try:
            pic = slide.shapes.add_picture(
                hero_image, 
                Inches(margin), 
                Inches(margin), 
                Inches(13.333 - 2 * margin), 
                Inches(7.5 - 2 * margin)
            )
            reveal.append(pic)
        except Exception:
            fallback = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(margin), 
                Inches(margin), 
                Inches(13.333 - 2 * margin), 
                Inches(7.5 - 2 * margin)
            )
            fallback.fill.solid()
            fallback.fill.fore_color.rgb = palette_color(theme, "panel")
            fallback.line.fill.background()
            reveal.append(fallback)

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(
            Inches(margin + 0.5), 
            Inches(margin + 0.5), 
            Inches(13.333 - 2 * margin - 1.0), 
            Inches(2)
        )
        set_textbox_text(head, truncate_to(headline, 100), theme, "title_xl", color_key="text")
        reveal.append(head)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
