from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "headline", "kind": "text", "max_chars": 50, "style": "title", "required": True},
    {"name": "metrics", "kind": "bullet_list", "bullet_capacity": 3, "required": True},
]

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline_text = get_slot(slots, "headline", default="PERFORMANCE METRICS")
    head = slide.shapes.add_textbox(Inches(margin), Inches(0.8), Inches(13.333 - 2 * margin), Inches(1.0))
    set_textbox_text(head, truncate_to(headline_text, 50), theme, "title", color_key="text")
    head.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.append(head)

    # Metrics
    metrics = get_slot(slots, "metrics", required=True)
    if not isinstance(metrics, list) or len(metrics) == 0:
        metrics = [{"title": "0%", "body": "Metric"}] * 3
    while len(metrics) < 3:
        metrics.append(metrics[-1])
    metrics = metrics[:3]

    total_w = 13.333
    col_w = total_w / 3
    gauge_size = 2.8
    y_gauge = 2.8
    y_center = y_gauge + gauge_size / 2

    for i, item in enumerate(metrics):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}

        x_center = (i + 0.5) * col_w
        x_gauge = x_center - gauge_size / 2

        # Gauge Track (Donut)
        track = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(x_gauge), Inches(y_gauge), Inches(gauge_size), Inches(gauge_size))
        track.fill.solid()
        track.fill.fore_color.rgb = palette_color(theme, "accent")
        track.line.fill.background()
        try:
            track.adjustments[0] = 0.65  # Adjust hole size for a thicker arc
        except:
            pass
        reveal.append(track)

        # Mask bottom half to create a semi-circle gauge effect
        mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_gauge - 0.1), Inches(y_center), Inches(gauge_size + 0.2), Inches(gauge_size / 2 + 0.1))
        mask.fill.solid()
        mask.fill.fore_color.rgb = palette_color(theme, "bg")
        mask.line.fill.background()
        reveal.append(mask)

        # Metric Value (Title) - positioned inside the semi-circle
        val_box = slide.shapes.add_textbox(Inches(x_center - 1.0), Inches(y_center - 0.8), Inches(2.0), Inches(0.8))
        set_textbox_text(val_box, truncate_to(item.get("title", ""), 10), theme, "title", color_key="text")
        val_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(val_box)

        # Metric Label (Body) - positioned below the gauge
        lbl_box = slide.shapes.add_textbox(Inches(x_center - 1.5), Inches(y_center + 0.3), Inches(3.0), Inches(0.8))
        set_textbox_text(lbl_box, truncate_to(item.get("body", ""), 40), theme, "body", color_key="muted")
        lbl_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(lbl_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
