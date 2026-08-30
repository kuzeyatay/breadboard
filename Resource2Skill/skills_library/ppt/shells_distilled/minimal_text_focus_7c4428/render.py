from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Decorative text (left)
    dec_text = get_slot(slots, "decorative_text")
    if dec_text:
        dec_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.5), Inches(1.0), Inches(5.0))
        set_textbox_text(dec_box, truncate_to(dec_text, 100), theme, "body", color_key="muted")
        reveal.append(dec_box)

    # Main text (center)
    main_text = get_slot(slots, "main_text", required=True)
    main_box = slide.shapes.add_textbox(Inches(5.0), Inches(2.5), Inches(7.0), Inches(2.5))
    set_textbox_text(main_box, truncate_to(main_text, 200), theme, "title_xl", color_key="text")
    reveal.append(main_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
