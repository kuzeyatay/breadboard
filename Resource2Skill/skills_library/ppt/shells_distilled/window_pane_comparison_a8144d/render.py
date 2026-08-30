from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    y_start = margin

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(margin), Inches(y_start), Inches(13.333 - 2*margin), Inches(1.0))
        set_textbox_text(head_box, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head_box)
        y_start += 1.2

    # Items (Doors)
    items = get_slot(slots, "items", required=True)
    if not isinstance(items, list):
        items = []
    # Pad to exactly 2 items if needed
    while len(items) < 2:
        items.append({"title": "Feature Title", "body": "Description goes here."})
    items = items[:2]

    door_w = 4.5
    door_h = 5.0
    gap = 0.3
    total_w = door_w * 2 + gap
    start_x = (13.333 - total_w) / 2
    
    # Center vertically in remaining space
    available_h = 7.5 - y_start - margin
    start_y = y_start + (available_h - door_h) / 2
    if start_y < y_start:
        start_y = y_start

    for i, item in enumerate(items):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}

        x = start_x + i * (door_w + gap)

        # Outer frame (Accent color)
        frame = add_solid_rect(slide, x, start_y, door_w, door_h, theme, color_key="accent")
        reveal.append(frame)

        pane_margin = 0.3
        pane_w = door_w - 2 * pane_margin
        pane_h = (door_h - 3 * pane_margin) / 2

        # Top pane (Panel color)
        top_pane = add_solid_rect(slide, x + pane_margin, start_y + pane_margin, pane_w, pane_h, theme, color_key="panel")
        reveal.append(top_pane)

        title_text = item.get("title", "")
        if title_text:
            t_box = slide.shapes.add_textbox(
                Inches(x + pane_margin + 0.3), 
                Inches(start_y + pane_margin + 0.3), 
                Inches(pane_w - 0.6), 
                Inches(pane_h - 0.6)
            )
            set_textbox_text(t_box, truncate_to(title_text, 60), theme, "subtitle", color_key="text")
            reveal.append(t_box)

        # Bottom pane (Panel color)
        bottom_pane = add_solid_rect(
            slide, 
            x + pane_margin, 
            start_y + 2*pane_margin + pane_h, 
            pane_w, 
            pane_h, 
            theme, 
            color_key="panel"
        )
        reveal.append(bottom_pane)

        body_text = item.get("body", "")
        if body_text:
            b_box = slide.shapes.add_textbox(
                Inches(x + pane_margin + 0.3), 
                Inches(start_y + 2*pane_margin + pane_h + 0.3), 
                Inches(pane_w - 0.6), 
                Inches(pane_h - 0.6)
            )
            set_textbox_text(b_box, truncate_to(body_text, 150), theme, "body", color_key="text")
            reveal.append(b_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
