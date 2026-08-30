from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    palette_color, set_textbox_text, add_theme_entrance, get_slot, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Vertical Kicker (Left)
    kicker_text = get_slot(slots, "vertical_kicker")
    if kicker_text:
        # Force vertical stacking by joining characters with newlines
        stacked_text = "\n".join(list(truncate_to(kicker_text, 15)))
        kicker_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.5), Inches(0.8), Inches(4.5))
        set_textbox_text(kicker_box, stacked_text, theme, "body", color_key="accent")
        reveal.append(kicker_box)

    # Main Metric (Center)
    metric_text = get_slot(slots, "main_metric", required=True)
    metric_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.5), Inches(7.0), Inches(2.0))
    set_textbox_text(metric_box, truncate_to(metric_text, 20), theme, "title_xl", color_key="text")
    reveal.append(metric_box)

    # Caption (Bottom)
    caption_text = get_slot(slots, "caption")
    if caption_text:
        caption_box = slide.shapes.add_textbox(Inches(4.5), Inches(6.0), Inches(7.0), Inches(0.8))
        set_textbox_text(caption_box, truncate_to(caption_text, 60), theme, "caption", color_key="muted")
        reveal.append(caption_box)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
