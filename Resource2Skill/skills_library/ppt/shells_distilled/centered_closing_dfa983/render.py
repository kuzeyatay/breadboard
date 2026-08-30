from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    headline_text = truncate_to(headline_text, 50)

    tb_width = 11.0
    tb_height = 3.0
    x = (13.333 - tb_width) / 2
    y = (7.5 - tb_height) / 2

    head = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(tb_width), Inches(tb_height))
    set_textbox_text(head, headline_text, theme, "title_xl", color_key="text")

    # Force center alignment
    for paragraph in head.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

    add_theme_entrance(slide, head, theme, delay_ms=0, duration_ms=800, index=0)
