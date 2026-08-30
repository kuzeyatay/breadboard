from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.8)
    reveal = []

    # 1. Base Background (Left side)
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")
    reveal.append(bg)

    # Helper to create vertical waves by rotating standard horizontal waves 90 degrees.
    # By making them 10 inches long, they bleed off the top and bottom of the 7.5 inch slide,
    # hiding the flat edges and leaving only the continuous wavy intersection.
    def add_vertical_wave(x, y, w, h, color_key):
        wave = slide.shapes.add_shape(MSO_SHAPE.WAVE, Inches(x), Inches(y), Inches(w), Inches(h))
        wave.rotation = 90
        wave.fill.solid()
        color = palette_color(theme, color_key)
        wave.fill.fore_color.rgb = color
        # Match outline to fill to prevent visible seams
        wave.line.fill.solid()
        wave.line.fill.fore_color.rgb = color
        return wave

    # 2. Ribbon Shadow (Offset slightly left to create 3D depth)
    ribbon_shadow = add_vertical_wave(0.55, 3.0, 10, 1.5, "muted")
    reveal.append(ribbon_shadow)

    # 3. Dark Ribbon
    ribbon = add_vertical_wave(0.75, 3.0, 10, 1.5, "accent")
    reveal.append(ribbon)

    # 4. Right Block Shadow
    right_shadow = add_vertical_wave(5.6, -1.25, 10, 10, "muted")
    reveal.append(right_shadow)

    # 5. Right Block (Covers the entire right half of the slide)
    right_block = add_vertical_wave(5.8, -1.25, 10, 10, "panel")
    reveal.append(right_block)

    # Text Content (Left-aligned in the safe zone before the wave)
    y_offset = 2.5
    
    if theme.get("motif", {}).get("type") == "thin_rule":
        add_hairline(slide, margin, y_offset - 0.3, 1.5, theme)

    headline_text = get_slot(slots, "headline")
    if headline_text:
        head_box = slide.shapes.add_textbox(Inches(margin), Inches(y_offset), Inches(3.8), Inches(1.5))
        set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title_xl", color_key="text")
        reveal.append(head_box)
        y_offset += 1.6

    subhead_text = get_slot(slots, "subhead")
    if subhead_text:
        sub_box = slide.shapes.add_textbox(Inches(margin), Inches(y_offset), Inches(3.8), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subhead_text, 120), theme, "body", color_key="text")
        reveal.append(sub_box)

    # Staggered entrance animation
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
