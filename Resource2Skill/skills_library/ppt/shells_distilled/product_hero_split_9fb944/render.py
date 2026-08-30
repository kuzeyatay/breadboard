from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.8)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Left Column
    left_x = margin
    left_w = 5.0
    y = 2.0
    
    headline_text = get_slot(slots, "headline")
    if headline_text:
        head_box = slide.shapes.add_textbox(Inches(left_x), Inches(y), Inches(left_w), Inches(1.2))
        set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title_xl", color_key="text")
        reveal.append(head_box)
        y += 1.4
        
    subhead_text = get_slot(slots, "subhead")
    if subhead_text:
        sub_box = slide.shapes.add_textbox(Inches(left_x), Inches(y), Inches(left_w), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subhead_text, 120), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)
        y += 1.2
        
    cta_text = get_slot(slots, "cta_text")
    if cta_text:
        cta_w = 2.5
        cta_h = 0.6
        # Button background
        btn = add_solid_rect(
            slide, left_x, y, cta_w, cta_h, theme, 
            color_key="accent", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
        )
        btn.adjustments[0] = 0.5 # fully rounded
        
        # Button text
        btn_text = slide.shapes.add_textbox(Inches(left_x + 0.2), Inches(y + 0.1), Inches(cta_w - 0.4), Inches(cta_h - 0.2))
        set_textbox_text(btn_text, truncate_to(cta_text, 30), theme, "body_bold", color_key="bg")
        
        reveal.extend([btn, btn_text])
        
    # Right Column (Hero Image)
    hero_image = get_slot(slots, "hero_image")
    if hero_image:
        right_w = 6.5
        right_x = 13.333 - margin - right_w
        img_y = 1.5
        img_h = 4.5
        
        try:
            pic = slide.shapes.add_picture(hero_image, Inches(right_x), Inches(img_y), width=Inches(right_w))
            reveal.append(pic)
        except Exception:
            placeholder = add_solid_rect(slide, right_x, img_y, right_w, img_h, theme, color_key="panel")
            reveal.append(placeholder)
            
    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
