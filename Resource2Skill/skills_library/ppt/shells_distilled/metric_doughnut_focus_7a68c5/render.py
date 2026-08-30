from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")
    
    reveal = []
    
    # Left Panel
    panel = add_solid_rect(slide, 0, 0, 0.8, 7.5, theme, color_key="panel")
    reveal.append(panel)
    
    # Sidebar Text
    sidebar_text = get_slot(slots, "sidebar_text")
    if sidebar_text:
        sb_w, sb_h = 5.0, 0.5
        sb_x, sb_y = -2.1, 3.5
        sb_box = slide.shapes.add_textbox(Inches(sb_x), Inches(sb_y), Inches(sb_w), Inches(sb_h))
        sb_box.rotation = -90
        set_textbox_text(sb_box, truncate_to(sidebar_text, 30), theme, "caption", color_key="muted")
        reveal.append(sb_box)

    center_x = 13.333 / 2
    
    # Headline
    headline = get_slot(slots, "headline", required=True)
    head_w = 10.0
    head_x = center_x - head_w / 2
    head_y = 1.2
    head_box = slide.shapes.add_textbox(Inches(head_x), Inches(head_y), Inches(head_w), Inches(0.8))
    set_textbox_text(head_box, truncate_to(headline, 40), theme, "title", color_key="text")
    reveal.append(head_box)
    
    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_y = head_y + 0.8
        sub_box = slide.shapes.add_textbox(Inches(head_x), Inches(sub_y), Inches(head_w), Inches(0.5))
        set_textbox_text(sub_box, truncate_to(subtitle, 60), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)
        
    # Doughnut Chart Track
    donut_size = 3.5
    donut_x = center_x - donut_size / 2
    donut_y = 3.0
    
    track = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(donut_x), Inches(donut_y), Inches(donut_size), Inches(donut_size))
    track.fill.solid()
    track.fill.fore_color.rgb = palette_color(theme, "panel")
    track.line.fill.background()
    reveal.append(track)
    
    # Doughnut Chart Progress Arc
    arc = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, Inches(donut_x), Inches(donut_y), Inches(donut_size), Inches(donut_size))
    arc.fill.solid()
    arc.fill.fore_color.rgb = palette_color(theme, "accent")
    arc.line.fill.background()
    reveal.append(arc)
    
    # Metric Text
    metric = get_slot(slots, "metric", required=True)
    met_w = donut_size
    met_h = 1.0
    met_x = donut_x
    met_y = donut_y + (donut_size - met_h) / 2
    met_box = slide.shapes.add_textbox(Inches(met_x), Inches(met_y), Inches(met_w), Inches(met_h))
    set_textbox_text(met_box, truncate_to(metric, 10), theme, "metric_xl", color_key="text")
    reveal.append(met_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
