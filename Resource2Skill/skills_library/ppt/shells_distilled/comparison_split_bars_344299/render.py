from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")
    reveal = [bg]

    # Headline
    headline_text = get_slot(slots, "headline", default="Product Comparison")
    head = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.8))
    set_textbox_text(head, headline_text, theme, "title", color_key="text")
    head.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.append(head)

    # Product A Header
    pa_name = get_slot(slots, "product_a_name", default="Product A")
    pa_desc = get_slot(slots, "product_a_desc", default="")
    
    circ_a = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.8), Inches(1.6), Inches(0.5), Inches(0.5))
    circ_a.fill.solid()
    circ_a.fill.fore_color.rgb = palette_color(theme, "accent")
    circ_a.line.fill.background()
    set_textbox_text(circ_a, "A", theme, "title", color_key="bg")
    circ_a.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.append(circ_a)

    tb_a_name = slide.shapes.add_textbox(Inches(2.0), Inches(1.5), Inches(2.6), Inches(0.4))
    set_textbox_text(tb_a_name, pa_name, theme, "body_bold", color_key="accent")
    tb_a_name.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    reveal.append(tb_a_name)

    if pa_desc:
        tb_a_desc = slide.shapes.add_textbox(Inches(2.0), Inches(1.9), Inches(2.6), Inches(0.8))
        set_textbox_text(tb_a_desc, truncate_to(pa_desc, 100), theme, "caption", color_key="muted")
        tb_a_desc.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        reveal.append(tb_a_desc)

    # Product B Header
    pb_name = get_slot(slots, "product_b_name", default="Product B")
    pb_desc = get_slot(slots, "product_b_desc", default="")
    
    circ_b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.933), Inches(1.6), Inches(0.5), Inches(0.5))
    circ_b.fill.solid()
    circ_b.fill.fore_color.rgb = palette_color(theme, "accent2")
    circ_b.line.fill.background()
    set_textbox_text(circ_b, "B", theme, "title", color_key="bg")
    circ_b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.append(circ_b)

    tb_b_name = slide.shapes.add_textbox(Inches(8.733), Inches(1.5), Inches(2.6), Inches(0.4))
    set_textbox_text(tb_b_name, pb_name, theme, "body_bold", color_key="accent2")
    reveal.append(tb_b_name)

    if pb_desc:
        tb_b_desc = slide.shapes.add_textbox(Inches(8.733), Inches(1.9), Inches(2.6), Inches(0.8))
        set_textbox_text(tb_b_desc, truncate_to(pb_desc, 100), theme, "caption", color_key="muted")
        reveal.append(tb_b_desc)

    # Central Panel & Metrics
    panel_w = 2.4
    panel_x = 5.466
    panel_y = 2.8
    row_h = 0.6
    spacing = 0.15
    max_bar_w = 3.5
    
    metrics = get_slot(slots, "metrics")
    if not metrics:
        metrics = [
            {"title": "Target Market Penetration", "pct_a": 0.5, "pct_b": 0.95},
            {"title": "Market Share", "pct_a": 0.75, "pct_b": 0.65},
            {"title": "Customer Acquisition Cost", "pct_a": 1.0, "pct_b": 0.5},
            {"title": "Average Revenue Per User", "pct_a": 0.4, "pct_b": 0.85},
            {"title": "Customer Lifetime Value", "pct_a": 0.85, "pct_b": 1.0}
        ]
    
    num_metrics = min(len(metrics), 5)
    total_panel_h = num_metrics * (row_h + spacing)
    c_panel = add_solid_rect(slide, panel_x, panel_y - 0.1, panel_w, total_panel_h + 0.1, theme, color_key="panel")
    reveal.append(c_panel)

    for i, item in enumerate(metrics[:5]):
        if not isinstance(item, dict):
            item = {"title": str(item)}
            
        y = panel_y + i * (row_h + spacing)
        
        # Metric Title
        title = item.get("title", f"Metric {i+1}")
        tb_m = slide.shapes.add_textbox(Inches(panel_x + 0.1), Inches(y + 0.05), Inches(panel_w - 0.2), Inches(row_h - 0.1))
        set_textbox_text(tb_m, truncate_to(title, 40), theme, "caption", color_key="text")
        tb_m.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(tb_m)
        
        # Left Bar (Product A)
        pct_a = item.get("pct_a", 0.5)
        try:
            pct_a = float(pct_a)
        except:
            pct_a = 0.5
        val_a = item.get("value_a", f"{int(pct_a*100)}%")
        bar_w_a = max(0.8, max_bar_w * pct_a)
        bar_x_a = panel_x - bar_w_a
        
        track_a = add_solid_rect(slide, panel_x - max_bar_w, y + 0.1, max_bar_w, row_h - 0.2, theme, color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        reveal.append(track_a)
        
        bar_a = add_solid_rect(slide, bar_x_a, y + 0.1, bar_w_a, row_h - 0.2, theme, color_key="accent", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        reveal.append(bar_a)
        
        tb_val_a = slide.shapes.add_textbox(Inches(bar_x_a + 0.1), Inches(y + 0.1), Inches(0.8), Inches(row_h - 0.2))
        set_textbox_text(tb_val_a, truncate_to(val_a, 10), theme, "caption", color_key="bg")
        reveal.append(tb_val_a)
        
        # Right Bar (Product B)
        pct_b = item.get("pct_b", 0.5)
        try:
            pct_b = float(pct_b)
        except:
            pct_b = 0.5
        val_b = item.get("value_b", f"{int(pct_b*100)}%")
        bar_w_b = max(0.8, max_bar_w * pct_b)
        bar_x_b = panel_x + panel_w
        
        track_b = add_solid_rect(slide, bar_x_b, y + 0.1, max_bar_w, row_h - 0.2, theme, color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        reveal.append(track_b)
        
        bar_b = add_solid_rect(slide, bar_x_b, y + 0.1, bar_w_b, row_h - 0.2, theme, color_key="accent2", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        reveal.append(bar_b)
        
        tb_val_b = slide.shapes.add_textbox(Inches(bar_x_b + bar_w_b - 0.9), Inches(y + 0.1), Inches(0.8), Inches(row_h - 0.2))
        set_textbox_text(tb_val_b, truncate_to(val_b, 10), theme, "caption", color_key="bg")
        tb_val_b.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        reveal.append(tb_val_b)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
