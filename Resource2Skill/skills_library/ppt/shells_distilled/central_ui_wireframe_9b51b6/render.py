from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Optional Headline
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.333), Inches(0.8))
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="text")
        reveal.append(head_box)
        
    # Center coordinates
    center_x = 13.333 / 2
    
    # Large top rectangle (Media/Art placeholder)
    rect_w, rect_h = 3.5, 3.5
    rect_x = center_x - rect_w / 2
    rect_y = 1.0
    
    main_rect = add_solid_rect(
        slide, rect_x, rect_y, rect_w, rect_h, theme,
        color_key="panel", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    main_rect.adjustments[0] = 0.1
    reveal.append(main_rect)
    
    # Slider track
    slider_y = rect_y + rect_h + 0.3
    track = add_solid_rect(
        slide, rect_x, slider_y, rect_w, 0.1, theme,
        color_key="muted", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    reveal.append(track)
    
    # Slider fill
    fill_w = rect_w * 0.6
    fill = add_solid_rect(
        slide, rect_x, slider_y, fill_w, 0.1, theme,
        color_key="accent", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    reveal.append(fill)
    
    # Slider thumb
    thumb_size = 0.25
    thumb = add_solid_rect(
        slide, rect_x + fill_w - thumb_size/2, slider_y - (thumb_size - 0.1)/2, thumb_size, thumb_size, theme,
        color_key="text", line=False, shape_type=MSO_SHAPE.OVAL
    )
    reveal.append(thumb)
    
    # 3 Buttons
    btn_size = 0.8
    btn_spacing = (rect_w - 3 * btn_size) / 2
    btn_y = slider_y + 0.4
    
    for i in range(3):
        bx = rect_x + i * (btn_size + btn_spacing)
        btn = add_solid_rect(
            slide, bx, btn_y, btn_size, btn_size, theme,
            color_key="panel", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
        )
        btn.adjustments[0] = 0.2
        reveal.append(btn)
        
    # Bottom Action Button
    action_y = btn_y + btn_size + 0.3
    action_h = 0.8
    action_btn = add_solid_rect(
        slide, rect_x, action_y, rect_w, action_h, theme,
        color_key="panel", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    action_btn.adjustments[0] = 0.5
    reveal.append(action_btn)
    
    # Action Label
    action_label = get_slot(slots, "action_label")
    if action_label:
        lbl_box = slide.shapes.add_textbox(Inches(rect_x), Inches(action_y + 0.15), Inches(rect_w), Inches(0.5))
        set_textbox_text(lbl_box, truncate_to(action_label, 30), theme, "body_bold", color_key="text")
        reveal.append(lbl_box)
        
    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
