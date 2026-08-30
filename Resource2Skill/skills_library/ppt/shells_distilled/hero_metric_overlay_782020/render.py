from __future__ import annotations
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from _shell_helpers import add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

SLOTS = [
    {"name": "background_image", "kind": "image", "required": True, "aspect": "16:9"},
    {"name": "metric", "kind": "text", "required": True, "style": "metric_xl", "max_chars": 15},
    {"name": "label", "kind": "text", "required": False, "style": "title", "max_chars": 40}
]

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # Background Image
    bg_img_path = get_slot(slots, "background_image")
    if bg_img_path:
        try:
            pic = slide.shapes.add_picture(bg_img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            reveal.append(pic)
        except Exception:
            pass
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = palette_color(theme, "bg")
        bg.line.fill.background()
        reveal.append(bg)

    metric_text = truncate_to(get_slot(slots, "metric", required=True), 15)
    label_text = get_slot(slots, "label")

    # Layout math
    if label_text:
        metric_y = 2.2
        label_y = 4.5
    else:
        metric_y = 2.8

    # Metric
    metric_box = slide.shapes.add_textbox(Inches(1.66), Inches(metric_y), Inches(10.0), Inches(2.0))
    set_textbox_text(metric_box, metric_text, theme, "metric_xl", color_key="text")
    for p in metric_box.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    reveal.append(metric_box)

    # Label
    if label_text:
        label_box = slide.shapes.add_textbox(Inches(1.66), Inches(label_y), Inches(10.0), Inches(1.0))
        set_textbox_text(label_box, truncate_to(label_text, 40), theme, "title", color_key="text")
        for p in label_box.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        reveal.append(label_box)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=150, index=i)
