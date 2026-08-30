from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    y = margin

    # Headline
    headline = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(1.0))
    set_textbox_text(head_box, truncate_to(headline, 100), theme, "title", color_key="text")
    reveal.append(head_box)
    y += 1.2

    # Chart Title
    chart_title = get_slot(slots, "chart_title")
    if chart_title:
        ct_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.35))
        set_textbox_text(ct_box, truncate_to(chart_title, 60), theme, "body_bold", color_key="text")
        reveal.append(ct_box)
        y += 0.35

    # Chart Subtitle
    chart_subtitle = get_slot(slots, "chart_subtitle")
    if chart_subtitle:
        cs_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.35))
        set_textbox_text(cs_box, truncate_to(chart_subtitle, 100), theme, "body", color_key="muted")
        reveal.append(cs_box)
        y += 0.4

    # Hairline separator
    if chart_title or chart_subtitle:
        add_hairline(slide, margin, y, 13.333 - 2 * margin, theme)
        y += 0.2

    # Chart Placeholder
    chart_h = 6.8 - y
    chart_rect = add_solid_rect(
        slide, margin, y, 13.333 - 2 * margin, chart_h, theme,
        color_key="panel", line=True
    )
    reveal.append(chart_rect)

    # Footer
    footer = get_slot(slots, "footer")
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(margin), Inches(6.9), Inches(13.333 - 2 * margin), Inches(0.4))
        set_textbox_text(foot_box, truncate_to(footer, 80), theme, "caption", color_key="muted")
        reveal.append(foot_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=80, duration_ms=500, index=i)
