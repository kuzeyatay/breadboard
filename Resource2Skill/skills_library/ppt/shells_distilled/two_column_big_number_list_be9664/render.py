from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import palette_color, set_textbox_text, add_theme_entrance, get_slot, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    y_offset = margin

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(y_offset), Inches(13.333 - 2 * margin), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)
        y_offset += 1.2

    # Items
    items = get_slot(slots, "items", [])
    if not items:
        return

    num_items = len(items)
    # Split into two columns (left column fills first)
    items_per_col = (num_items + 1) // 2
    col_w = (13.333 - 2 * margin) / 2
    available_h = 7.5 - y_offset - margin
    row_h = available_h / max(items_per_col, 1)

    for i, item in enumerate(items):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}

        col = 0 if i < items_per_col else 1
        row = i if col == 0 else i - items_per_col

        x = margin + col * col_w
        y = y_offset + row * row_h

        # Big Number / Title
        title_text = truncate_to(item.get("title", ""), 10)
        num_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.5), Inches(row_h))
        set_textbox_text(num_box, title_text, theme, "metric_xl", color_key="accent")
        reveal.append(num_box)

        # Body Text
        body_text = item.get("body", "")
        if body_text:
            body_box = slide.shapes.add_textbox(Inches(x + 1.6), Inches(y + 0.1), Inches(col_w - 1.8), Inches(row_h - 0.2))
            set_textbox_text(body_box, truncate_to(body_text, 120), theme, "body", color_key="text")
            reveal.append(body_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
