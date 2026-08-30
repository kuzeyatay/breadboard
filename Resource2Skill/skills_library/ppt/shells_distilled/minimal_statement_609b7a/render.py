from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

SLOTS = [
    {"name": "statement", "kind": "text", "max_chars": 150, "style": "title", "required": True}
]

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "accent")
    bg.line.fill.background()

    statement_text = get_slot(slots, "statement", required=True)
    statement_text = truncate_to(statement_text, 150)

    # Text box
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.5))
    set_textbox_text(tb, statement_text, theme, "title", color_key="bg")

    add_theme_entrance(slide, tb, theme, delay_ms=100)
