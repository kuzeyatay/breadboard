from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    box_width = 6.0
    x = (13.333 - box_width) / 2
    y = 1.8

    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_width), Inches(1.2))
        set_textbox_text(head_box, truncate_to(headline, 100), theme, "title", color_key="text")
        reveal.append(head_box)
        y += 1.3

    subhead = get_slot(slots, "subhead")
    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_width), Inches(0.6))
        set_textbox_text(sub_box, truncate_to(subhead, 60), theme, "subtitle", color_key="accent")
        reveal.append(sub_box)
        y += 0.8

    body = get_slot(slots, "body")
    if body:
        body_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_width), Inches(2.5))
        set_textbox_text(body_box, truncate_to(body, 400), theme, "body", color_key="muted")
        reveal.append(body_box)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
