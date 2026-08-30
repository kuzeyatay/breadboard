from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "kicker", "kind": "text", "max_chars": 30, "style": "caption", "required": False},
    {"name": "headline", "kind": "text", "max_chars": 80, "style": "title", "required": False},
    {"name": "timeline_items", "kind": "bullet_list", "bullet_capacity": 3, "required": True},
]

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    y = margin
    
    # Kicker
    kicker = get_slot(slots, "kicker")
    if kicker:
        k_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(6), Inches(0.35))
        set_textbox_text(k_box, truncate_to(kicker, 30), theme, "caption", color_key="accent")
        reveal.append(k_box)
        y += 0.5
        
    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        h_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(1.0))
        set_textbox_text(h_box, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(h_box)
        
    # Timeline layout
    y_timeline = 3.75
    total_w = 13.333 - 2 * margin
    
    # Main horizontal line
    line = add_solid_rect(slide, margin, y_timeline, total_w, 0.02, theme, color_key="muted")
    reveal.append(line)
    
    items = get_slot(slots, "timeline_items", required=True)
    if not isinstance(items, list):
        items = []
        
    col_w = total_w / 3
    
    for i in range(3):
        if i < len(items):
            item = items[i]
            if not isinstance(item, dict):
                item = {"title": str(item), "body": ""}
        else:
            # Fallback if fewer than 3 items provided
            item = {"title": "", "body": ""}
            
        cx = margin + i * col_w + col_w / 2
        
        # Node (circle)
        node = add_solid_rect(
            slide, cx - 0.15, y_timeline - 0.15, 0.3, 0.3, theme, 
            color_key="accent", shape_type=MSO_SHAPE.OVAL
        )
        reveal.append(node)
        
        # Vertical drop line
        drop_line = add_solid_rect(
            slide, cx - 0.01, y_timeline + 0.15, 0.02, 0.4, theme, color_key="muted"
        )
        reveal.append(drop_line)
        
        # Title
        t_box = slide.shapes.add_textbox(
            Inches(cx - col_w/2 + 0.2), Inches(y_timeline + 0.6), Inches(col_w - 0.4), Inches(0.5)
        )
        set_textbox_text(t_box, truncate_to(item.get("title", ""), 40), theme, "body_bold", color_key="text")
        reveal.append(t_box)
        
        # Body
        b_box = slide.shapes.add_textbox(
            Inches(cx - col_w/2 + 0.2), Inches(y_timeline + 1.1), Inches(col_w - 0.4), Inches(1.5)
        )
        set_textbox_text(b_box, truncate_to(item.get("body", ""), 120), theme, "body", color_key="muted")
        reveal.append(b_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
