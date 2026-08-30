import math
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    reveal = [bg]
    
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2*margin), Inches(1.0))
        set_textbox_text(head, headline, theme, "title", color_key="text")
        reveal.append(head)
        
    # Draw 8-segment wheel
    center_x = 13.333 / 2
    center_y = 7.5 / 2
    radius = 2.2
    
    items = get_slot(slots, "wheel_items") or []
    
    for i in range(8):
        start_angle = i * 45
        end_angle = (i + 1) * 45
        
        pie = slide.shapes.add_shape(
            MSO_SHAPE.PIE, 
            Inches(center_x - radius), 
            Inches(center_y - radius), 
            Inches(radius * 2), 
            Inches(radius * 2)
        )
        pie.fill.solid()
        pie.fill.fore_color.rgb = palette_color(theme, "panel")
        
        # Thick background-colored border creates the visual gaps between slices
        pie.line.color.rgb = palette_color(theme, "bg")
        pie.line.width = Inches(0.08)
        
        try:
            pie.adjustments[0] = start_angle
            pie.adjustments[1] = end_angle
        except Exception:
            pass
            
        reveal.append(pie)
        
        if i < len(items):
            angle_rad = math.radians(start_angle + 22.5)
            tx = center_x + (radius + 0.9) * math.cos(angle_rad)
            ty = center_y + (radius + 0.9) * math.sin(angle_rad)
            
            tbox = slide.shapes.add_textbox(Inches(tx - 1.0), Inches(ty - 0.3), Inches(2.0), Inches(0.6))
            set_textbox_text(tbox, items[i].get("title", ""), theme, "body", color_key="text")
            reveal.append(tbox)

    # Call to action button
    cta_text = get_slot(slots, "call_to_action")
    if cta_text:
        btn_w = 3.5
        btn_h = 0.8
        btn_x = center_x - btn_w / 2
        btn_y = 7.5 - margin - btn_h
        
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(btn_x), Inches(btn_y), Inches(btn_w), Inches(btn_h))
        btn.fill.solid()
        btn.fill.fore_color.rgb = palette_color(theme, "accent")
        btn.line.fill.background()
        try:
            btn.adjustments[0] = 0.5
        except Exception:
            pass
            
        btn_tbox = slide.shapes.add_textbox(Inches(btn_x), Inches(btn_y + 0.15), Inches(btn_w), Inches(0.5))
        set_textbox_text(btn_tbox, cta_text, theme, "body_bold", color_key="bg")
        
        reveal.extend([btn, btn_tbox])
    
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
