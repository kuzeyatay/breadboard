from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

SLOTS = [
    {"name": "headline", "kind": "text", "required": True, "style": "title", "max_chars": 60},
    {"name": "subhead", "kind": "text", "required": False, "style": "body", "max_chars": 120},
    {"name": "hero_image", "kind": "image", "required": True}
]

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []

    # Accent Block (Middle band)
    block_y = 3.0
    block_h = 3.5
    block_w = 11.0
    block = add_solid_rect(
        slide, 0, block_y, block_w, block_h, theme,
        color_key="primary", line=False, shape_type=MSO_SHAPE.RECTANGLE
    )
    reveal.append(block)

    # Text inside the block
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(6.0), Inches(block_y + 0.6), Inches(4.5), Inches(1.0))
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="bg")
        reveal.append(head_box)

    subhead = get_slot(slots, "subhead")
    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(6.0), Inches(block_y + 1.6), Inches(4.5), Inches(1.5))
        set_textbox_text(sub_box, truncate_to(subhead, 120), theme, "body", color_key="bg")
        reveal.append(sub_box)

    # Hero Image (Left, overlapping block)
    img_slot = get_slot(slots, "hero_image")
    if img_slot:
        try:
            pic = slide.shapes.add_picture(img_slot, Inches(1.0), Inches(0.5), width=Inches(4.5))
            reveal.append(pic)
        except Exception:
            # Fallback placeholder
            placeholder = add_solid_rect(
                slide, 1.0, 0.5, 4.5, 6.5, theme,
                color_key="muted", line=False, shape_type=MSO_SHAPE.RECTANGLE
            )
            reveal.append(placeholder)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
