import math
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Left Graphic (Donut)
    donut_size = 5.5
    donut_x = 1.0
    donut_y = (7.5 - donut_size) / 2
    donut = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(donut_x), Inches(donut_y), Inches(donut_size), Inches(donut_size))
    donut.fill.solid()
    donut.fill.fore_color.rgb = palette_color(theme, "accent")
    donut.line.color.rgb = palette_color(theme, "bg")
    donut.line.width = Inches(0.05)
    
    if len(donut.adjustments) > 0:
        donut.adjustments[0] = 0.25  # Make ring thick enough for text
    reveal.append(donut)

    # Labels on Donut
    labels = get_slot(slots, "labels")
    if labels and isinstance(labels, list):
        cx = donut_x + donut_size / 2
        cy = donut_y + donut_size / 2
        # Radius to middle of the donut ring
        r = donut_size / 2 * 0.87 
        n = min(len(labels), 8)
        for i in range(n):
            item = labels[i]
            text = item.get("title", str(item)) if isinstance(item, dict) else str(item)
            
            angle = i * (2 * math.pi / n) - math.pi / 2
            lx = cx + r * math.cos(angle)
            ly = cy + r * math.sin(angle)
            
            lw, lh = 1.5, 0.4
            lbox = slide.shapes.add_textbox(Inches(lx - lw/2), Inches(ly - lh/2), Inches(lw), Inches(lh))
            set_textbox_text(lbox, truncate_to(text, 15), theme, "caption", color_key="bg")
            
            # Rotate text to follow the curve
            deg = angle * 180 / math.pi
            rot = (deg + 90) % 360
            # Flip text if it's on the bottom half to keep it readable
            if 90 < rot < 270:
                rot = (rot + 180) % 360
            lbox.rotation = rot
            
            reveal.append(lbox)

    # Right Text Area
    text_x = 7.0
    text_w = 5.5

    headline_text = get_slot(slots, "headline", default="Circular Process")
    head = slide.shapes.add_textbox(Inches(text_x), Inches(2.5), Inches(text_w), Inches(1.2))
    set_textbox_text(head, truncate_to(headline_text, 80), theme, "title", color_key="text")
    reveal.append(head)

    body_text = get_slot(slots, "body")
    if body_text:
        body = slide.shapes.add_textbox(Inches(text_x), Inches(3.8), Inches(text_w), Inches(3.0))
        set_textbox_text(body, truncate_to(body_text, 300), theme, "body", color_key="muted")
        reveal.append(body)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
