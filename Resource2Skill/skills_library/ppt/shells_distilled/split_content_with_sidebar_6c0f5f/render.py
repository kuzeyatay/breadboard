from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    reveal.append(bg)

    # Main Panel (Left 2/3)
    panel_w = 9.0
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(margin), Inches(margin), Inches(panel_w - margin), Inches(7.5 - 2*margin))
    panel.fill.solid()
    panel.fill.fore_color.rgb = palette_color(theme, "panel")
    panel.line.fill.background()
    reveal.append(panel)

    # Left Column: Headline & Lead Text
    left_x = margin + 0.4
    left_w = 3.5
    y_left = margin + 0.4

    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(left_x), Inches(y_left), Inches(left_w), Inches(1.5))
        set_textbox_text(head_box, truncate_to(headline, 80), theme, "title", color_key="accent")
        reveal.append(head_box)
        y_left += 1.8

    lead_text = get_slot(slots, "lead_text")
    if lead_text:
        lead_box = slide.shapes.add_textbox(Inches(left_x), Inches(y_left), Inches(left_w), Inches(4.0))
        set_textbox_text(lead_box, truncate_to(lead_text, 300), theme, "body_bold", color_key="text")
        reveal.append(lead_box)

    # Middle Column: Bullets
    mid_x = left_x + left_w + 0.4
    mid_w = panel_w - mid_x - 0.4
    y_mid = margin + 0.4

    bullets = get_slot(slots, "main_bullets")
    if bullets and isinstance(bullets, list):
        for i, item in enumerate(bullets):
            if not isinstance(item, dict):
                item = {"title": str(item), "body": ""}
            
            b_title = item.get("title", "")
            if b_title:
                bt_box = slide.shapes.add_textbox(Inches(mid_x), Inches(y_mid), Inches(mid_w), Inches(0.5))
                set_textbox_text(bt_box, truncate_to(b_title, 60), theme, "body_bold", color_key="accent")
                reveal.append(bt_box)
                y_mid += 0.5
            
            b_body = item.get("body", "")
            if b_body:
                bb_box = slide.shapes.add_textbox(Inches(mid_x), Inches(y_mid), Inches(mid_w), Inches(1.0))
                set_textbox_text(bb_box, truncate_to(b_body, 150), theme, "body", color_key="text")
                reveal.append(bb_box)
                y_mid += 1.2

    # Right Column: Sidebar (Outside panel)
    right_x = panel_w + 0.4
    right_w = 13.333 - right_x - margin
    y_right = margin + 0.4

    sidebar_title = get_slot(slots, "sidebar_title")
    if sidebar_title:
        st_box = slide.shapes.add_textbox(Inches(right_x), Inches(y_right), Inches(right_w), Inches(0.6))
        set_textbox_text(st_box, truncate_to(sidebar_title, 40), theme, "subtitle", color_key="accent")
        reveal.append(st_box)
        y_right += 0.8

    sidebar_text = get_slot(slots, "sidebar_text")
    if sidebar_text:
        sb_box = slide.shapes.add_textbox(Inches(right_x), Inches(y_right), Inches(right_w), Inches(4.0))
        set_textbox_text(sb_box, truncate_to(sidebar_text, 300), theme, "body", color_key="text")
        reveal.append(sb_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
