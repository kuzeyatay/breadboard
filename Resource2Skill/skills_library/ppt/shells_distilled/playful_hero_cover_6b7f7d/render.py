from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Top hanging accent block (inspired by the red rectangle)
    accent_w = 1.2
    accent_h = 3.2
    accent_x = 9.0
    accent_y = 0.0
    accent = add_solid_rect(
        slide, accent_x, accent_y, accent_w, accent_h, theme,
        color_key="accent", line=False
    )
    reveal.append(accent)

    # Text content (Left side)
    y_text = 2.5
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(margin), Inches(y_text), Inches(6.0), Inches(1.5))
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title_xl", color_key="text")
        reveal.append(head_box)
        y_text += 1.6

    subhead = get_slot(slots, "subhead")
    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(margin), Inches(y_text), Inches(6.0), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subhead, 100), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)

    # Hero Image (Bottom Right, inspired by the soccer ball)
    hero_img = get_slot(slots, "hero_image")
    if hero_img:
        img_w = 3.0
        img_h = 3.0
        img_x = 7.0
        img_y = 4.0
        try:
            pic = slide.shapes.add_picture(hero_img, Inches(img_x), Inches(img_y), Inches(img_w), Inches(img_h))
            reveal.append(pic)
        except Exception:
            placeholder = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="panel")
            reveal.append(placeholder)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
