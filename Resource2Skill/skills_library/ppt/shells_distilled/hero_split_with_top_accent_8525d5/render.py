from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Decorative accent block (hanging from top right)
    accent = add_solid_rect(
        slide, 10.0, 0.0, 1.5, 3.0, theme,
        color_key="accent", line=False
    )
    reveal.append(accent)

    # Headline (bottom left)
    headline_text = get_slot(slots, "headline", required=True)
    if headline_text:
        head_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(5.5), Inches(2.5))
        set_textbox_text(head_box, truncate_to(headline_text, 80), theme, "title_xl", color_key="text")
        reveal.append(head_box)

    # Hero Image (right side)
    image_path = get_slot(slots, "hero_image", required=True)
    if image_path:
        pic = slide.shapes.add_picture(image_path, Inches(7.0), Inches(2.0), width=Inches(5.5))
        reveal.append(pic)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
