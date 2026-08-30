from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

SLOTS = [
    {"name": "bullets", "kind": "bullet_list", "bullet_capacity": 5, "required": True},
]

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    bullets = get_slot(slots, "bullets", required=True)
    if not isinstance(bullets, list):
        bullets = []
        
    # Centered placeholder dimensions based on the image
    start_x = 2.5
    start_y = 2.0
    w = 8.333
    
    y = start_y
    
    for i, item in enumerate(bullets):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
            
        title_text = truncate_to(item.get("title", ""), 120)
        body_text = truncate_to(item.get("body", ""), 300)
        
        if title_text:
            t_box = slide.shapes.add_textbox(Inches(start_x), Inches(y), Inches(w), Inches(0.5))
            set_textbox_text(t_box, f"• {title_text}", theme, "title", color_key="text")
            reveal.append(t_box)
            y += 0.8
            
        if body_text:
            b_box = slide.shapes.add_textbox(Inches(start_x + 0.5), Inches(y), Inches(w - 0.5), Inches(0.5))
            set_textbox_text(b_box, body_text, theme, "body", color_key="muted")
            reveal.append(b_box)
            y += 0.6
            
        y += 0.3

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
