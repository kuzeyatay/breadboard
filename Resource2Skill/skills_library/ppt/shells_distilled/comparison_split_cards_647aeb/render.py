from __future__ import annotations
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    gutter = theme.get("spacing", {}).get("gutter", 0.4)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Overall Headline (Optional)
    headline_text = get_slot(slots, "headline")
    y_offset = margin
    if headline_text:
        head = slide.shapes.add_textbox(Inches(margin), Inches(y_offset), Inches(13.333 - 2*margin), Inches(1.0))
        set_textbox_text(head, truncate_to(headline_text, 80), theme, "title", color_key="text")
        reveal.append(head)
        y_offset += 1.2
    else:
        y_offset = 1.5
        
    panel_w = (13.333 - 2 * margin - gutter) / 2
    panel_h = 7.5 - y_offset - margin
    if not headline_text:
        panel_h = 4.5
        y_offset = (7.5 - panel_h) / 2
        
    # Left Panel (Standard)
    left_x = margin
    left_bg = add_solid_rect(slide, left_x, y_offset, panel_w, panel_h, theme, color_key="panel", line=True)
    reveal.append(left_bg)
    
    l_head_text = get_slot(slots, "left_headline", required=True)
    l_sub_text = get_slot(slots, "left_subhead")
    
    l_content_h = 1.0 if not l_sub_text else 2.2
    l_start_y = y_offset + (panel_h - l_content_h) / 2
    
    l_head = slide.shapes.add_textbox(Inches(left_x + 0.5), Inches(l_start_y), Inches(panel_w - 1.0), Inches(1.0))
    set_textbox_text(l_head, truncate_to(l_head_text, 40), theme, "title", color_key="text")
    reveal.append(l_head)
    
    if l_sub_text:
        l_sub = slide.shapes.add_textbox(Inches(left_x + 0.5), Inches(l_start_y + 1.2), Inches(panel_w - 1.0), Inches(1.0))
        set_textbox_text(l_sub, truncate_to(l_sub_text, 100), theme, "subtitle", color_key="muted")
        reveal.append(l_sub)

    # Right Panel (Accent/Highlighted)
    right_x = margin + panel_w + gutter
    right_bg = add_solid_rect(slide, right_x, y_offset, panel_w, panel_h, theme, color_key="accent", line=False)
    reveal.append(right_bg)
    
    r_head_text = get_slot(slots, "right_headline", required=True)
    r_sub_text = get_slot(slots, "right_subhead")
    
    r_content_h = 1.0 if not r_sub_text else 2.2
    r_start_y = y_offset + (panel_h - r_content_h) / 2
    
    r_head = slide.shapes.add_textbox(Inches(right_x + 0.5), Inches(r_start_y), Inches(panel_w - 1.0), Inches(1.0))
    set_textbox_text(r_head, truncate_to(r_head_text, 40), theme, "title", color_key="bg")
    reveal.append(r_head)
    
    if r_sub_text:
        r_sub = slide.shapes.add_textbox(Inches(right_x + 0.5), Inches(r_start_y + 1.2), Inches(panel_w - 1.0), Inches(1.0))
        set_textbox_text(r_sub, truncate_to(r_sub_text, 100), theme, "subtitle", color_key="bg")
        reveal.append(r_sub)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
