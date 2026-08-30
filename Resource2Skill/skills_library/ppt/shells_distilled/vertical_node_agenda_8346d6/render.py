from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="primary", line=False, shape_type=MSO_SHAPE.RECTANGLE)
    
    reveal = []
    
    # Left side: Title and vertical line
    title_text = get_slot(slots, "title", default="AGENDA")
    
    line_x = 3.5
    line_y = 3.0
    line_h = 1.5
    left_line = add_solid_rect(slide, line_x, line_y, 0.05, line_h, theme, color_key="bg", line=False, shape_type=MSO_SHAPE.RECTANGLE)
    reveal.append(left_line)
    
    title_box = slide.shapes.add_textbox(Inches(line_x + 0.3), Inches(line_y + 0.2), Inches(4.0), Inches(1.0))
    set_textbox_text(title_box, title_text, theme, "title_xl", color_key="bg")
    reveal.append(title_box)
    
    # Right side: Agenda items
    items = get_slot(slots, "items", required=True)
    if not isinstance(items, list):
        items = [{"title": "Item 1"}]
    
    num_items = min(len(items), 5)
    
    right_x = 8.5
    start_y = 1.5
    end_y = 6.0
    
    if num_items > 1:
        spacing = (end_y - start_y) / (num_items - 1)
    else:
        spacing = 0
        start_y = 3.75
        
    # Vertical connecting line
    if num_items > 1:
        conn_line = add_solid_rect(slide, right_x - 0.025, start_y, 0.05, end_y - start_y, theme, color_key="bg", line=False, shape_type=MSO_SHAPE.RECTANGLE)
        reveal.append(conn_line)
        
    circle_d = 0.8
    
    for i in range(num_items):
        item = items[i]
        if not isinstance(item, dict):
            item = {"title": str(item)}
            
        cy = start_y + i * spacing
        
        # Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(right_x - circle_d/2), Inches(cy - circle_d/2), Inches(circle_d), Inches(circle_d))
        circle.fill.solid()
        circle.fill.fore_color.rgb = palette_color(theme, "text")
        circle.line.color.rgb = palette_color(theme, "bg")
        circle.line.width = Inches(0.03)
        
        # Number inside circle
        num_box = slide.shapes.add_textbox(Inches(right_x - circle_d/2), Inches(cy - circle_d/2 + 0.15), Inches(circle_d), Inches(circle_d))
        set_textbox_text(num_box, f"{i+1:02d}", theme, "body_bold", color_key="bg")
        
        # Item title
        title_y = cy - 0.25
        item_box = slide.shapes.add_textbox(Inches(right_x + circle_d/2 + 0.3), Inches(title_y), Inches(3.5), Inches(0.5))
        set_textbox_text(item_box, item.get("title", ""), theme, "body", color_key="bg")
        
        reveal.extend([circle, num_box, item_box])
        
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
