from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin_x = 0.6
    margin_y = 0.3
    max_x = 12.7
    max_y = 7.2
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Right side image strip
    img_w = 3.0
    img_x = max_x - img_w
    gap = 0.1
    # 4 images, 3 gaps
    img_h = (max_y - margin_y - (3 * gap)) / 4

    for i in range(4):
        img_y = margin_y + i * (img_h + gap)
        img_slot = get_slot(slots, f"image_{i+1}")
        
        added_pic = False
        if img_slot:
            try:
                pic = slide.shapes.add_picture(img_slot, Inches(img_x), Inches(img_y), Inches(img_w), Inches(img_h))
                reveal.append(pic)
                added_pic = True
            except Exception:
                pass
        
        if not added_pic:
            # Fallback placeholder
            placeholder = add_solid_rect(
                slide, img_x, img_y, img_w, img_h, theme,
                color_key="panel", line=False
            )
            reveal.append(placeholder)

    # Left side text
    text_w = img_x - margin_x - 0.5 # 0.5 inch gap between text and images
    
    headline = get_slot(slots, "headline", required=True)
    body = get_slot(slots, "body")
    
    # Estimate heights to center vertically
    head_h = 1.2
    body_h = 1.5 if body else 0
    total_text_h = head_h + body_h
    start_y = margin_y + (max_y - margin_y - total_text_h) / 2

    head_box = slide.shapes.add_textbox(Inches(margin_x), Inches(start_y), Inches(text_w), Inches(head_h))
    set_textbox_text(head_box, truncate_to(headline, 60), theme, "title_xl", color_key="text")
    reveal.append(head_box)

    if body:
        body_box = slide.shapes.add_textbox(Inches(margin_x), Inches(start_y + head_h), Inches(text_w), Inches(body_h))
        set_textbox_text(body_box, truncate_to(body, 200), theme, "body", color_key="muted")
        reveal.append(body_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
