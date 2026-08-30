from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Full-bleed background using accent color for high impact
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "accent")
    bg.line.fill.background()

    message = truncate_to(get_slot(slots, "message", default="THANK YOU"), 30)

    # Massive centered textbox
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Use 'bg' color for text to contrast strongly against the 'accent' background
    set_textbox_text(tb, message, theme, "title_xl", color_key="bg")

    # Force center alignment for all paragraphs
    for paragraph in tb.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

    add_theme_entrance(slide, bg, theme, delay_ms=0, duration_ms=500, index=0)
    add_theme_entrance(slide, tb, theme, delay_ms=200, duration_ms=800, index=1)
