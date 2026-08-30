from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Dimensions for the centered block
    total_w = 9.0
    h = 4.5
    panel_w = 3.5
    img_w = 5.5
    
    x = (13.333 - total_w) / 2
    y = (7.5 - h) / 2
    
    # Left Panel (Solid Color Block)
    panel = add_solid_rect(slide, x, y, panel_w, h, theme, color_key="primary", line=False)
    reveal.append(panel)
    
    # Text inside the Left Panel
    pad = 0.5
    headline_text = get_slot(slots, "headline")
    if headline_text:
        head_box = slide.shapes.add_textbox(Inches(x + pad), Inches(y + pad), Inches(panel_w - 2*pad), Inches(1.2))
        # Using 'bg' color for text assuming 'primary' is a dark/bold color
        set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title", color_key="bg")
        reveal.append(head_box)
        
    body_text = get_slot(slots, "body")
    if body_text:
        body_box = slide.shapes.add_textbox(Inches(x + pad), Inches(y + pad + 1.3), Inches(panel_w - 2*pad), Inches(h - 2*pad - 1.3))
        set_textbox_text(body_box, truncate_to(body_text, 150), theme, "body", color_key="bg")
        reveal.append(body_box)
        
    # Right Image
    img_path = get_slot(slots, "hero_image")
    if img_path:
        try:
            pic = slide.shapes.add_picture(img_path, Inches(x + panel_w), Inches(y), width=Inches(img_w), height=Inches(h))
            reveal.append(pic)
        except Exception:
            # Fallback if image fails to load
            fallback = add_solid_rect(slide, x + panel_w, y, img_w, h, theme, color_key="muted", line=False)
            reveal.append(fallback)
            
    # Staggered Entrance Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
