from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "headline",   "kind": "text",  "max_chars": 50,  "style": "title", "required": True},
    {"name": "body",       "kind": "text",  "max_chars": 500, "style": "body",  "required": True},
    {"name": "hero_image", "kind": "image", "aspect": "4:3",                    "required": True},
]

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Layout metrics
    left_w = 3.5
    gap = 0.5
    img_x = margin + left_w + gap
    img_w = 13.333 - img_x - margin
    img_h = 6.0
    img_y = (7.5 - img_h) / 2
    
    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 50)
    head_box = slide.shapes.add_textbox(Inches(margin), Inches(1.5), Inches(left_w), Inches(1.0))
    set_textbox_text(head_box, headline, theme, "title", color_key="text")
    reveal.append(head_box)
    
    # Body
    body = truncate_to(get_slot(slots, "body", required=True), 500)
    body_box = slide.shapes.add_textbox(Inches(margin), Inches(2.6), Inches(left_w), Inches(4.0))
    set_textbox_text(body_box, body, theme, "body", color_key="text")
    reveal.append(body_box)
    
    # Image
    img_path = get_slot(slots, "hero_image")
    if img_path:
        try:
            pic = slide.shapes.add_picture(img_path, Inches(img_x), Inches(img_y), Inches(img_w), Inches(img_h))
            reveal.append(pic)
        except Exception:
            placeholder = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="panel")
            reveal.append(placeholder)
    else:
        placeholder = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="panel")
        reveal.append(placeholder)
        
    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
