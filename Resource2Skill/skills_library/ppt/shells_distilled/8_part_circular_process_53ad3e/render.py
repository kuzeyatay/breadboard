import math
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(1.0))
    set_textbox_text(head, headline, theme, "title", color_key="text")
    reveal.append(head)
    
    # Diagram Center
    cx = 13.333 / 2
    cy = 4.0
    radius = 2.2
    
    # Draw 8 pie slices to form a segmented donut
    for i in range(8):
        start_angle = i * 45 + 2
        end_angle = (i + 1) * 45 - 2
        
        left = cx - radius
        top = cy - radius
        width = radius * 2
        height = radius * 2
        
        pie = slide.shapes.add_shape(MSO_SHAPE.PIE, Inches(left), Inches(top), Inches(width), Inches(height))
        pie.fill.solid()
        pie.fill.fore_color.rgb = palette_color(theme, "primary")
        pie.line.fill.background()
        
        try:
            pie.adjustments[0] = start_angle
            pie.adjustments[1] = end_angle
        except Exception:
            pass # Fallback if adjustments are not supported in this exact way
            
        reveal.append(pie)
        
    # Draw center cutout to complete the donut look
    cutout_radius = 1.0
    cutout = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - cutout_radius), Inches(cy - cutout_radius),
        Inches(cutout_radius * 2), Inches(cutout_radius * 2)
    )
    cutout.fill.solid()
    cutout.fill.fore_color.rgb = palette_color(theme, "bg")
    cutout.line.fill.background()
    reveal.append(cutout)
    
    # Center Label
    center_label = get_slot(slots, "center_label")
    if center_label:
        ct_box = slide.shapes.add_textbox(
            Inches(cx - 0.9), Inches(cy - 0.5), Inches(1.8), Inches(1.0)
        )
        set_textbox_text(ct_box, truncate_to(center_label, 20), theme, "body_bold", color_key="text")
        reveal.append(ct_box)
        
    # Items around the perimeter
    items = get_slot(slots, "items", required=True)
    if not isinstance(items, list):
        items = []
        
    text_radius = 3.0
    for i in range(8):
        if i < len(items):
            item = items[i]
            if not isinstance(item, dict):
                item = {"title": str(item), "body": ""}
        else:
            item = {"title": "", "body": ""}
            
        if not item.get("title") and not item.get("body"):
            continue
            
        angle_deg = i * 45 + 22.5
        angle_rad = math.radians(angle_deg)
        
        tx = cx + text_radius * math.cos(angle_rad) - 1.0
        ty = cy + text_radius * math.sin(angle_rad) - 0.4
        
        t_box = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(2.0), Inches(0.8))
        
        text = item.get("title", "")
        if item.get("body"):
            text += "\n" + item.get("body")
            
        set_textbox_text(t_box, truncate_to(text, 60), theme, "caption", color_key="text")
        reveal.append(t_box)
        
    # Bottom decorative bar (inspired by the image's color palette bar)
    bar = add_solid_rect(slide, 0, 7.3, 13.333, 0.2, theme, color_key="accent")
    reveal.append(bar)
    
    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
