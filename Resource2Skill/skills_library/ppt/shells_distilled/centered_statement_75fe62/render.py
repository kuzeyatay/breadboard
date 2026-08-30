from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    statement_text = get_slot(slots, "statement", required=True)
    statement_text = truncate_to(statement_text, 150)

    # Centered text box
    tb_width = 10.0
    tb_height = 3.0
    x = (13.333 - tb_width) / 2
    y = (7.5 - tb_height) / 2

    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(tb_width), Inches(tb_height))
    set_textbox_text(tb, statement_text, theme, "title", color_key="text")

    add_theme_entrance(slide, tb, theme, delay_ms=0, index=0)
