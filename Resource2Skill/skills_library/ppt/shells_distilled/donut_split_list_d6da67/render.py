from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

SLOTS = [
    {"name": "headline", "kind": "text", "style": "title", "max_chars": 40, "required": True},
    {"name": "bullets", "kind": "bullet_list", "bullet_capacity": 4, "required": True}
]

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # Left Graphic (Donut)
    donut_x, donut_y, donut_size = 1.0, 1.25, 5.0
    donut = slide.shapes.add_shape(
        MSO_SHAPE.DONUT, Inches(donut_x), Inches(donut_y), Inches(donut_size), Inches(donut_size)
    )
    donut.fill.solid()
    donut.fill.fore_color.rgb = palette_color(theme, "accent")
    donut.line.fill.background()
    donut.adjustments[0] = 0.25  # Adjust hole size to make it a thick ring
    reveal.append(donut)

    # Headline Background
    head_x, head_y, head_w, head_h = 5.5, 3.4, 6.0, 0.7
    head_bg = add_solid_rect(
        slide, head_x, head_y, head_w, head_h, theme, color_key="accent"
    )
    reveal.append(head_bg)

    # Headline Text
    headline = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(
        Inches(head_x + 0.2), Inches(head_y), Inches(head_w - 0.4), Inches(head_h)
    )
    set_textbox_text(
        head_box, truncate_to(headline, 40), theme, "title", color_key="bg"
    )
    reveal.append(head_box)

    # Bullets
    bullets = get_slot(slots, "bullets", required=True)
    bul_x, bul_y, bul_w, bul_h = 5.5, 4.3, 6.5, 2.5
    bul_box = slide.shapes.add_textbox(
        Inches(bul_x), Inches(bul_y), Inches(bul_w), Inches(bul_h)
    )
    
    if isinstance(bullets, list):
        # Format as a simple bulleted string
        bullet_text = "\n".join(
            f"• {item}" if isinstance(item, str) else f"• {item.get('title', '')}" 
            for item in bullets[:4]
        )
        set_textbox_text(bul_box, bullet_text, theme, "body", color_key="text")
    reveal.append(bul_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
