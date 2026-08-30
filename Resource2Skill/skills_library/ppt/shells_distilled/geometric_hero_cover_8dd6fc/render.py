from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # 1. Background Image or Fallback
    img_path = get_slot(slots, "hero_image")
    if img_path and isinstance(img_path, str):
        try:
            bg = slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            reveal.append(bg)
        except Exception:
            bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="muted", line=False, shape_type=MSO_SHAPE.RECTANGLE)
            reveal.append(bg)
    else:
        bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="muted", line=False, shape_type=MSO_SHAPE.RECTANGLE)
        reveal.append(bg)

    # 2. Geometric Overlays
    # Left solid block
    left_rect = add_solid_rect(slide, 0, 0, 5.0, 7.5, theme, color_key="primary", line=False, shape_type=MSO_SHAPE.RECTANGLE)
    reveal.append(left_rect)

    # Left diagonal edge (Right Triangle)
    left_tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(5.0), Inches(0), Inches(2.5), Inches(7.5))
    left_tri.fill.solid()
    left_tri.fill.fore_color.rgb = palette_color(theme, "primary")
    left_tri.line.color.rgb = palette_color(theme, "primary")
    reveal.append(left_tri)

    # Right diagonal stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(-1.5), Inches(1.0), Inches(10.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = palette_color(theme, "primary")
    stripe.line.color.rgb = palette_color(theme, "primary")
    stripe.rotation = 18.4
    reveal.append(stripe)

    # Bottom right accent triangle
    white_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(7.8), Inches(4.5), Inches(2.8), Inches(3.0))
    white_tri.fill.solid()
    white_tri.fill.fore_color.rgb = palette_color(theme, "bg")
    white_tri.line.color.rgb = palette_color(theme, "bg")
    reveal.append(white_tri)

    # 3. Text Content
    headline_text = get_slot(slots, "headline", default="Geometric Cover")
    head_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(4.0), Inches(1.5))
    set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title_xl", color_key="bg")
    reveal.append(head_box)

    subtitle_text = get_slot(slots, "subtitle")
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(4.0), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subtitle_text, 120), theme, "body", color_key="bg")
        reveal.append(sub_box)

    # 4. Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
