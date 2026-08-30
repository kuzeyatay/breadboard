from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    message = get_slot(slots, "message")
    if message:
        tb = slide.shapes.add_textbox(Inches(1.66), Inches(3.0), Inches(10.0), Inches(1.5))
        set_textbox_text(tb, truncate_to(message, 100), theme, "title", color_key="text")
        add_theme_entrance(slide, tb, theme, index=0)
