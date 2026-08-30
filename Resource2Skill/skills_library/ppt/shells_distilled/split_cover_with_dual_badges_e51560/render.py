from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []

    # Hero Image (Left side)
    hero_img_path = get_slot(slots, "hero_image")
    if hero_img_path:
        try:
            pic = slide.shapes.add_picture(hero_img_path, Inches(0), Inches(0), width=Inches(7.5), height=Inches(7.5))
            reveal.append(pic)
        except Exception:
            pass

    # Right side layout parameters
    right_x = 7.5
    margin = 0.6
    content_w = 13.333 - right_x - margin * 2
    
    # Badges (Top Right)
    badge_y = 0.8
    badge_size = (content_w - 0.4) / 2  # Two badges with 0.4 spacing
    has_badges = False
    
    badge_1_path = get_slot(slots, "badge_1")
    if badge_1_path:
        try:
            b1 = slide.shapes.add_picture(badge_1_path, Inches(right_x + margin), Inches(badge_y), width=Inches(badge_size), height=Inches(badge_size))
            reveal.append(b1)
            has_badges = True
        except Exception:
            pass
            
    badge_2_path = get_slot(slots, "badge_2")
    if badge_2_path:
        try:
            b2 = slide.shapes.add_picture(badge_2_path, Inches(right_x + margin + badge_size + 0.4), Inches(badge_y), width=Inches(badge_size), height=Inches(badge_size))
            reveal.append(b2)
            has_badges = True
        except Exception:
            pass

    # Text Content (Below badges)
    if has_badges:
        y_text = badge_y + badge_size + 0.8
    else:
        y_text = 2.5
    
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(right_x + margin), Inches(y_text), Inches(content_w), Inches(1.2))
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="text")
        reveal.append(head_box)
        y_text += 1.3

    subhead = get_slot(slots, "subhead")
    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(right_x + margin), Inches(y_text), Inches(content_w), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subhead, 100), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
