from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    metrics = get_slot(slots, "metrics", required=True)
    if not isinstance(metrics, list):
        metrics = [{"title": str(metrics)}]

    reveal = []
    num_items = min(len(metrics), 5)
    
    box_width = Inches(6)
    item_height = Inches(1.5)
    total_height = num_items * item_height
    
    start_y = Inches(7.5 / 2) - (total_height / 2)
    center_x = Inches(13.333 / 2)

    for i in range(num_items):
        item = metrics[i]
        if not isinstance(item, dict):
            item = {"title": str(item)}
        
        y = start_y + (i * item_height)
        
        box = slide.shapes.add_textbox(center_x - (box_width/2), y, box_width, item_height)
        set_textbox_text(box, truncate_to(item.get("title", ""), 15), theme, "metric_xl", color_key="accent")
        
        for paragraph in box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            
        reveal.append(box)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
