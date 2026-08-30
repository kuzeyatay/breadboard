from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Large circle on the left (touching top/bottom, right edge near center)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.8), Inches(0), Inches(7.5), Inches(7.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = palette_color(theme, "panel")
    circle.line.fill.background()
    reveal.append(circle)
    
    # Fetch slots
    kicker = get_slot(slots, "kicker")
    headline = get_slot(slots, "headline", required=True)
    subhead = get_slot(slots, "subhead")
    
    # Calculate vertical centering for text block on the right
    h_kicker = 0.5 if kicker else 0
    h_head = 1.5
    h_sub = 1.0 if subhead else 0
    total_h = h_kicker + h_head + h_sub
    
    y = (7.5 - total_h) / 2
    x = 7.5
    w = 5.0
    
    if kicker:
        kicker_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h_kicker))
        set_textbox_text(kicker_box, truncate_to(kicker, 30), theme, "caption", color_key="accent")
        reveal.append(kicker_box)
        y += h_kicker
        
    head_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h_head))
    set_textbox_text(head_box, truncate_to(headline, 80), theme, "title", color_key="text")
    reveal.append(head_box)
    y += h_head
    
    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h_sub))
        set_textbox_text(sub_box, truncate_to(subhead, 150), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)
        
    # Staggered entrance animation
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
