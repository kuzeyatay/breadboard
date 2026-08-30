from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # Background Image
    bg_img = get_slot(slots, "background_image", required=True)
    if bg_img:
        pic = slide.shapes.add_picture(bg_img, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        reveal.append(pic)

    # Center Box
    box_w, box_h = 4.5, 0.8
    box_x = (13.333 - box_w) / 2
    box_y = (7.5 - box_h) / 2 - 0.2

    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(box_x), Inches(box_y), Inches(box_w), Inches(box_h))
    box.fill.solid()
    box.fill.fore_color.rgb = palette_color(theme, "accent")
    box.line.color.rgb = palette_color(theme, "bg")
    box.line.width = Inches(0.02)
    reveal.append(box)

    # Title
    title_text = get_slot(slots, "title", required=True)
    title_box = slide.shapes.add_textbox(Inches(box_x), Inches(box_y + 0.15), Inches(box_w), Inches(box_h - 0.3))
    set_textbox_text(title_box, truncate_to(title_text, 40), theme, "title", color_key="bg")
    for paragraph in title_box.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    reveal.append(title_box)

    # Subtitle
    subtitle_text = get_slot(slots, "subtitle")
    if subtitle_text:
        sub_w = 6.0
        sub_x = (13.333 - sub_w) / 2
        sub_y = box_y + box_h + 0.1
        sub_box = slide.shapes.add_textbox(Inches(sub_x), Inches(sub_y), Inches(sub_w), Inches(0.4))
        set_textbox_text(sub_box, truncate_to(subtitle_text, 60), theme, "caption", color_key="text")
        for paragraph in sub_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        reveal.append(sub_box)

    # Footer
    footer_text = get_slot(slots, "footer")
    if footer_text:
        foot_w = 5.0
        foot_x = 13.333 - foot_w - 0.5
        foot_y = 7.5 - 0.6
        foot_box = slide.shapes.add_textbox(Inches(foot_x), Inches(foot_y), Inches(foot_w), Inches(0.4))
        set_textbox_text(foot_box, truncate_to(footer_text, 80), theme, "caption", color_key="text")
        for paragraph in foot_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.RIGHT
        reveal.append(foot_box)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
