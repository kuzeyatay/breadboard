from __future__ import annotations
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head = slide.shapes.add_textbox(Inches(margin), Inches(0.4), Inches(13.333 - 2*margin), Inches(0.8))
    set_textbox_text(head, truncate_to(headline_text, 80), theme, "title", color_key="text")
    reveal.append(head)
    
    # Top Section (Image + Cards)
    top_y = 1.4
    top_h = 3.2
    
    # Left Image Placeholder
    img_slot = get_slot(slots, "hero_image")
    img_w = 5.0
    if img_slot:
        pic = add_solid_rect(slide, margin, top_y, img_w, top_h, theme, color_key="muted", line=False, shape_type=MSO_SHAPE.RECTANGLE)
        reveal.append(pic)

    # Right Metric Cards
    cards = get_slot(slots, "metric_cards", required=True)
    if not isinstance(cards, list):
        cards = []
    cards = cards[:4]
    
    cards_x = margin + img_w + 0.4
    cards_w = 13.333 - margin - cards_x
    
    if cards:
        gap = 0.15
        card_h = (top_h - (len(cards) - 1) * gap) / len(cards)
        
        for i, card in enumerate(cards):
            if not isinstance(card, dict):
                card = {"title": str(card), "body": ""}
            
            cy = top_y + i * (card_h + gap)
            
            # Card Background
            bg_rect = add_solid_rect(slide, cards_x, cy, cards_w, card_h, theme, color_key="panel", line=False, shape_type=MSO_SHAPE.RECTANGLE)
            
            # Accent Bar
            bar = add_solid_rect(slide, cards_x, cy, 0.15, card_h, theme, color_key="accent", line=False, shape_type=MSO_SHAPE.RECTANGLE)
            
            # Card Title (Source/Label)
            t_box = slide.shapes.add_textbox(Inches(cards_x + 0.3), Inches(cy + 0.05), Inches(cards_w - 0.4), Inches(0.3))
            set_textbox_text(t_box, truncate_to(card.get("title", ""), 50), theme, "caption", color_key="muted")
            
            # Card Body (Metric/Value)
            b_box = slide.shapes.add_textbox(Inches(cards_x + 0.3), Inches(cy + 0.35), Inches(cards_w - 0.4), Inches(card_h - 0.4))
            set_textbox_text(b_box, truncate_to(card.get("body", ""), 80), theme, "body_bold", color_key="text")
            
            reveal.extend([bg_rect, bar, t_box, b_box])

    # Bottom Bullets
    bottom_bullets = get_slot(slots, "bottom_bullets")
    if bottom_bullets and isinstance(bottom_bullets, list):
        bottom_bullets = bottom_bullets[:3]
        bot_y = top_y + top_h + 0.4
        bot_h = 7.5 - margin - bot_y
        
        cols = len(bottom_bullets)
        col_gap = 0.4
        col_w = (13.333 - 2*margin - (cols - 1)*col_gap) / cols
        
        for i, bullet in enumerate(bottom_bullets):
            if not isinstance(bullet, dict):
                bullet = {"title": str(bullet), "body": ""}
            
            bx = margin + i * (col_w + col_gap)
            
            # Bullet Accent Square
            sq = add_solid_rect(slide, bx, bot_y + 0.1, 0.15, 0.15, theme, color_key="accent", line=False, shape_type=MSO_SHAPE.RECTANGLE)
            
            # Bullet Title
            bt_box = slide.shapes.add_textbox(Inches(bx + 0.25), Inches(bot_y), Inches(col_w - 0.25), Inches(0.4))
            set_textbox_text(bt_box, truncate_to(bullet.get("title", ""), 40), theme, "body_bold", color_key="text")
            
            # Bullet Body
            bb_box = slide.shapes.add_textbox(Inches(bx + 0.25), Inches(bot_y + 0.4), Inches(col_w - 0.25), Inches(bot_h - 0.4))
            set_textbox_text(bb_box, truncate_to(bullet.get("body", ""), 150), theme, "body", color_key="muted")
            
            reveal.extend([sq, bt_box, bb_box])

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
