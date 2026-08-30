from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get('spacing', {}).get('margin', 0.8)
    content_w = 13.333 - 2 * margin
    reveal = []
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, 'bg')
    bg.line.fill.background()
    kicker_text = get_slot(slots, 'kicker')
    if kicker_text:
        kicker = slide.shapes.add_textbox(Inches(margin), Inches(0.5), Inches(content_w), Inches(0.4))
        set_textbox_text(kicker, truncate_to(kicker_text, 50), theme, 'caption', color_key='muted')
        reveal.append(kicker)
    banner_y = 1.2
    banner_h = 2.0
    banner_bg = add_solid_rect(slide, margin, banner_y, content_w, banner_h, theme, color_key='panel')
    reveal.append(banner_bg)
    img_w = content_w * 0.5
    img_rect = add_solid_rect(slide, margin, banner_y, img_w, banner_h, theme, color_key='muted')
    reveal.append(img_rect)
    headline_text = get_slot(slots, 'headline', required=True)
    head_x = margin + img_w + 0.5
    head_w = content_w - img_w - 0.8
    head = slide.shapes.add_textbox(Inches(head_x), Inches(banner_y + 0.4), Inches(head_w), Inches(1.2))
    set_textbox_text(head, truncate_to(headline_text, 60), theme, 'title', color_key='text')
    reveal.append(head)
    panel_y = 3.5
    panel_h = 3.0
    panel = add_solid_rect(slide, margin, panel_y, content_w, panel_h, theme, color_key='panel')
    reveal.append(panel)
    author_w = 2.5
    author_center_x = margin + (author_w / 2)
    author_img_size = 0.8
    author_img_x = author_center_x - (author_img_size / 2)
    author_img_y = panel_y + 0.5
    author_img = add_solid_rect(slide, author_img_x, author_img_y, author_img_size, author_img_size, theme, color_key='muted')
    reveal.append(author_img)
    name_text = get_slot(slots, 'author_name', required=True)
    name_box = slide.shapes.add_textbox(Inches(margin + 0.2), Inches(author_img_y + author_img_size + 0.2), Inches(author_w - 0.4), Inches(0.4))
    set_textbox_text(name_box, truncate_to(name_text, 40), theme, 'body_bold', color_key='text')
    reveal.append(name_box)
    role_text = get_slot(slots, 'author_role')
    if role_text:
        role_box = slide.shapes.add_textbox(Inches(margin + 0.2), Inches(author_img_y + author_img_size + 0.5), Inches(author_w - 0.4), Inches(0.3))
        set_textbox_text(role_box, truncate_to(role_text, 40), theme, 'caption', color_key='muted')
        reveal.append(role_box)
    quote_text = get_slot(slots, 'quote', required=True)
    quote_x = margin + author_w + 0.5
    quote_w = content_w - author_w - 1.0
    quote_y = panel_y + 0.5
    quote_h = panel_h - 1.0
    quote_box = slide.shapes.add_textbox(Inches(quote_x), Inches(quote_y), Inches(quote_w), Inches(quote_h))
    set_textbox_text(quote_box, truncate_to(quote_text, 250), theme, 'subtitle', color_key='text')
    reveal.append(quote_box)
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
