from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Decorative Pill
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.6), Inches(3.2), Inches(2.5), Inches(0.8)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = palette_color(theme, "accent")
    pill.line.fill.background()
    pill.adjustments[0] = 0.5  # Fully rounded ends
    pill.rotation = -15.0
    reveal.append(pill)

    # Headline
    headline = get_slot(slots, "headline", default="MEET OUR TEAM")
    headline = truncate_to(headline, 40)
    words = headline.split()

    # Group words into up to 3 lines for the staggered effect
    lines = []
    if len(words) == 1:
        lines = [words[0]]
    elif len(words) == 2:
        lines = [words[0], words[1]]
    elif len(words) >= 3:
        lines = [words[0], words[1], " ".join(words[2:])]
    else:
        lines = ["MEET", "OUR", "TEAM"]

    start_x = 7.6
    start_y = 2.7
    x_step = 1.0
    y_step = 1.0

    for i, line_text in enumerate(lines):
        x = start_x + (i * x_step)
        y = start_y + (i * y_step)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(5.0), Inches(1.2))
        # Force uppercase to match the bold, graphic mood of the original design
        set_textbox_text(tb, line_text.upper(), theme, "title_xl", color_key="text")
        reveal.append(tb)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
