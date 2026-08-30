from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline_text = truncate_to(get_slot(slots, "headline", required=True), 80)
    head = slide.shapes.add_textbox(Inches(1.66), Inches(1.5), Inches(10.0), Inches(2.0))
    set_textbox_text(head, headline_text, theme, "title_xl", color_key="text")
    for paragraph in head.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    reveal.append(head)

    # Search Bar Group
    bar_y = 4.2
    bar_h = 0.8
    input_w = 5.0
    btn_w = 1.0
    total_w = input_w + btn_w
    start_x = (13.333 - total_w) / 2

    # Input Box
    input_box = add_solid_rect(
        slide, start_x, bar_y, input_w, bar_h, theme,
        color_key="bg", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    input_box.line.color.rgb = palette_color(theme, "text")
    input_box.line.width = Inches(0.03)
    reveal.append(input_box)

    # Search Query Text
    query = get_slot(slots, "search_query")
    if query:
        q_box = slide.shapes.add_textbox(Inches(start_x + 0.3), Inches(bar_y + 0.15), Inches(input_w - 0.6), Inches(0.5))
        set_textbox_text(q_box, truncate_to(query, 40), theme, "body", color_key="text")
        reveal.append(q_box)

    # Search Button
    btn_x = start_x + input_w - 0.2  # Overlap slightly to hide right corners of input box
    btn = add_solid_rect(
        slide, btn_x, bar_y, btn_w + 0.2, bar_h, theme,
        color_key="text", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    reveal.append(btn)

    # Magnifying Glass Icon (Circle + Handle)
    r = 0.12
    cx = btn_x + (btn_w + 0.2) / 2 - 0.05
    cy = bar_y + bar_h / 2 - 0.05
    
    glass = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(r*2), Inches(r*2))
    glass.fill.background()
    glass.line.color.rgb = palette_color(theme, "bg")
    glass.line.width = Inches(0.03)
    reveal.append(glass)

    handle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx + r - 0.02), Inches(cy + r - 0.02), Inches(0.15), Inches(0.05))
    handle.fill.solid()
    handle.fill.fore_color.rgb = palette_color(theme, "bg")
    handle.line.fill.background()
    handle.rotation = 45
    reveal.append(handle)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
