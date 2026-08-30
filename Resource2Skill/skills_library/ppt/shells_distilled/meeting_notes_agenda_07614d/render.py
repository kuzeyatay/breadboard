from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_hairline, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    y = margin
    
    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.8))
        set_textbox_text(head_box, truncate_to(headline, 100), theme, "title", color_key="text")
        reveal.append(head_box)
        y += 0.9
        
    # Meta info
    meta_info = get_slot(slots, "meta_info")
    if meta_info:
        meta_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.8))
        set_textbox_text(meta_box, truncate_to(meta_info, 200), theme, "body", color_key="muted")
        reveal.append(meta_box)
        y += 1.0
        
        add_hairline(slide, margin, y, 13.333 - 2 * margin, theme)
        y += 0.3
        
    # Content Title
    content_title = get_slot(slots, "content_title")
    if content_title:
        ct_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.5))
        set_textbox_text(ct_box, truncate_to(content_title, 100), theme, "subtitle", color_key="text")
        reveal.append(ct_box)
        y += 0.6
        
    # Bullets
    bullets = get_slot(slots, "bullets")
    if bullets and isinstance(bullets, list):
        for item in bullets[:8]:  # limit to 8 to avoid overflow
            if isinstance(item, dict):
                title = item.get("title", "")
                body = item.get("body", "")
            else:
                title = str(item)
                body = ""
                
            if title:
                t_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.35))
                set_textbox_text(t_box, f"• {truncate_to(title, 100)}", theme, "body_bold", color_key="text")
                reveal.append(t_box)
                y += 0.4
                
            if body:
                b_box = slide.shapes.add_textbox(Inches(margin + 0.5), Inches(y), Inches(13.333 - 2 * margin - 0.5), Inches(0.35))
                set_textbox_text(b_box, truncate_to(body, 200), theme, "body", color_key="muted")
                reveal.append(b_box)
                y += 0.4
                
            y += 0.1

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
