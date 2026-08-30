from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to, add_emphasis_pulse
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Dimensions
    center_x = 13.333 / 2
    center_y = 7.5 / 2
    ring_size = 5.5
    inner_size = 4.0

    # Vertical Line (Mimicking the magenta line)
    line_w = 0.1
    line_h = 6.5
    vert_line = add_solid_rect(
        slide, center_x - line_w/2, center_y - line_h/2, line_w, line_h, theme,
        color_key="accent", line=False
    )
    reveal.append(vert_line)

    # Outer Ring (Mimicking the orange ring)
    ring = add_solid_rect(
        slide, center_x - ring_size/2, center_y - ring_size/2, ring_size, ring_size, theme,
        color_key="accent", line=False, shape_type=MSO_SHAPE.DONUT
    )
    try:
        ring.adjustments[0] = 0.25 # Set donut thickness
    except:
        pass
    reveal.append(ring)

    # Inner Circle (Panel for text readability)
    inner_circle = add_solid_rect(
        slide, center_x - inner_size/2, center_y - inner_size/2, inner_size, inner_size, theme,
        color_key="panel", line=False, shape_type=MSO_SHAPE.OVAL
    )
    reveal.append(inner_circle)

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(
            Inches(center_x - inner_size/2 + 0.2), Inches(center_y - 0.8),
            Inches(inner_size - 0.4), Inches(1.0)
        )
        set_textbox_text(head_box, truncate_to(headline, 40), theme, "title_xl", color_key="text")
        reveal.append(head_box)

    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(center_x - inner_size/2 + 0.2), Inches(center_y + 0.2),
            Inches(inner_size - 0.4), Inches(0.8)
        )
        set_textbox_text(sub_box, truncate_to(subtitle, 80), theme, "body", color_key="muted")
        reveal.append(sub_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100 + i*100, index=i)
    
    add_emphasis_pulse(slide, ring, theme, delay_ms=1500)
