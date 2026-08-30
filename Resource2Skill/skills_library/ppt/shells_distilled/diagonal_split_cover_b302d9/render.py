from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")

    reveal = []

    # Diagonal accent shapes on the right
    # Outer diagonal (muted)
    ff_builder1 = slide.shapes.build_freeform(Inches(8.0), Inches(0))
    ff_builder1.add_line_segments([
        (Inches(13.333), Inches(0)),
        (Inches(13.333), Inches(7.5)),
        (Inches(6.0), Inches(7.5)),
        (Inches(8.0), Inches(0))
    ])
    poly1 = ff_builder1.convert_to_shape()
    poly1.fill.solid()
    poly1.fill.fore_color.rgb = palette_color(theme, "muted")
    poly1.line.fill.background()
    reveal.append(poly1)

    # Inner diagonal (panel)
    ff_builder2 = slide.shapes.build_freeform(Inches(9.5), Inches(0))
    ff_builder2.add_line_segments([
        (Inches(13.333), Inches(0)),
        (Inches(13.333), Inches(7.5)),
        (Inches(7.5), Inches(7.5)),
        (Inches(9.5), Inches(0))
    ])
    poly2 = ff_builder2.convert_to_shape()
    poly2.fill.solid()
    poly2.fill.fore_color.rgb = palette_color(theme, "panel")
    poly2.line.fill.background()
    reveal.append(poly2)

    # Text content
    margin_x = Inches(1.5)
    
    headline_text = truncate_to(get_slot(slots, "headline", required=True), 100)
    head_box = slide.shapes.add_textbox(margin_x, Inches(2.5), Inches(6.0), Inches(1.5))
    set_textbox_text(head_box, headline_text, theme, "title_xl", color_key="text")
    reveal.append(head_box)
    
    subhead_text = get_slot(slots, "subhead")
    if subhead_text:
        sub_box = slide.shapes.add_textbox(margin_x, Inches(4.5), Inches(6.0), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subhead_text, 80), theme, "subtitle", color_key="text")
        reveal.append(sub_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
