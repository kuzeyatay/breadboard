from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Large accent rectangle
    rect_w, rect_h = 9.5, 4.5
    rect_x = (13.333 - rect_w) / 2
    rect_y = (7.5 - rect_h) / 2
    
    accent_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rect_x), Inches(rect_y), Inches(rect_w), Inches(rect_h))
    accent_rect.fill.solid()
    accent_rect.fill.fore_color.rgb = palette_color(theme, "accent")
    accent_rect.line.fill.background()
    reveal.append(accent_rect)
    
    # Inner rounded rectangle
    inner_w, inner_h = 4.5, 1.5
    inner_x = (13.333 - inner_w) / 2
    inner_y = (7.5 - inner_h) / 2
    
    inner_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(inner_x), Inches(inner_y), Inches(inner_w), Inches(inner_h))
    inner_rect.fill.solid()
    inner_rect.fill.fore_color.rgb = palette_color(theme, "panel")
    inner_rect.line.fill.background()
    if inner_rect.adjustments:
        inner_rect.adjustments[0] = 0.15
    reveal.append(inner_rect)
    
    # Text
    headline = get_slot(slots, "headline", required=True)
    subtitle = get_slot(slots, "subtitle")
    
    if subtitle:
        head_h = 0.8
        sub_h = 0.5
        total_h = head_h + sub_h
        start_y = inner_y + (inner_h - total_h) / 2
        
        head_box = slide.shapes.add_textbox(Inches(inner_x + 0.2), Inches(start_y), Inches(inner_w - 0.4), Inches(head_h))
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="text")
        head_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(head_box)
        
        sub_box = slide.shapes.add_textbox(Inches(inner_x + 0.2), Inches(start_y + head_h), Inches(inner_w - 0.4), Inches(sub_h))
        set_textbox_text(sub_box, truncate_to(subtitle, 80), theme, "body", color_key="muted")
        sub_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(sub_box)
    else:
        head_h = 1.0
        start_y = inner_y + (inner_h - head_h) / 2
        head_box = slide.shapes.add_textbox(Inches(inner_x + 0.2), Inches(start_y), Inches(inner_w - 0.4), Inches(head_h))
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="text")
        head_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(head_box)
        
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
