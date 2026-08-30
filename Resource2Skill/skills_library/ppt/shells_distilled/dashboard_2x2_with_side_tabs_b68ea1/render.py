from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Layout parameters
    x_left = 0.5
    y_top = 0.5
    left_w = 7.8
    row_h = 3.1
    gutter = 0.4
    x_right = x_left + left_w + gutter
    right_w = 13.333 - x_right - 0.5
    y_bottom = y_top + row_h + gutter

    # --- Top Left Panel ---
    p1 = add_solid_rect(slide, x_left, y_top, left_w, row_h, theme, color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    p1.adjustments[0] = 0.05
    reveal.append(p1)

    # --- Bottom Left Panel ---
    p2 = add_solid_rect(slide, x_left, y_bottom, left_w, row_h, theme, color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    p2.adjustments[0] = 0.05
    reveal.append(p2)

    # --- Top Right Panel ---
    tab1 = add_solid_rect(slide, x_right, y_top, 0.6, row_h, theme, color_key="accent", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tab1.adjustments[0] = 0.1
    p3 = add_solid_rect(slide, x_right + 0.3, y_top, right_w - 0.3, row_h, theme, color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    p3.adjustments[0] = 0.05
    reveal.extend([tab1, p3])

    label_tr = get_slot(slots, "label_top_right")
    if label_tr:
        # Vertical text: rotate 270 degrees around center
        tb1 = slide.shapes.add_textbox(Inches(x_right - 1.2), Inches(y_top + 1.3), Inches(3.0), Inches(0.5))
        tb1.rotation = 270
        set_textbox_text(tb1, truncate_to(label_tr, 20), theme, "body_bold", color_key="bg")
        reveal.append(tb1)

    # --- Bottom Right Panel ---
    tab2 = add_solid_rect(slide, x_right, y_bottom, 0.6, row_h, theme, color_key="accent", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tab2.adjustments[0] = 0.1
    p4 = add_solid_rect(slide, x_right + 0.3, y_bottom, right_w - 0.3, row_h, theme, color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    p4.adjustments[0] = 0.05
    reveal.extend([tab2, p4])

    label_br = get_slot(slots, "label_bottom_right")
    if label_br:
        tb2 = slide.shapes.add_textbox(Inches(x_right - 1.2), Inches(y_bottom + 1.3), Inches(3.0), Inches(0.5))
        tb2.rotation = 270
        set_textbox_text(tb2, truncate_to(label_br, 20), theme, "body_bold", color_key="bg")
        reveal.append(tb2)

    # Metrics in Bottom Right
    metrics = get_slot(slots, "metrics")
    if metrics and isinstance(metrics, list):
        m_w = 1.5
        m_y = y_bottom + 0.8
        for i, item in enumerate(metrics[:2]):
            if not isinstance(item, dict):
                item = {"title": str(item), "body": ""}
            m_x = x_right + 0.8 + i * 1.6

            val_box = slide.shapes.add_textbox(Inches(m_x), Inches(m_y), Inches(m_w), Inches(0.8))
            set_textbox_text(val_box, truncate_to(item.get("title", ""), 10), theme, "metric_xl", color_key="accent")

            lbl_box = slide.shapes.add_textbox(Inches(m_x), Inches(m_y + 0.8), Inches(m_w), Inches(0.4))
            set_textbox_text(lbl_box, truncate_to(item.get("body", ""), 20), theme, "body", color_key="muted")

            reveal.extend([val_box, lbl_box])

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
