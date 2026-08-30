from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    # 1. Background (Accent color)
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="accent", line=False)

    # 2. White Triangle (Base on left, pointing right)
    # Create an upward-pointing isosceles triangle and rotate 90 degrees.
    # Unrotated dimensions: width=7.5, height=13.333.
    # Center must be at slide center (6.6665, 3.75).
    left = 6.6665 - (7.5 / 2)
    top = 3.75 - (13.333 / 2)
    
    triangle = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(left), Inches(top), Inches(7.5), Inches(13.333)
    )
    triangle.rotation = 90
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = palette_color(theme, "bg")
    triangle.line.fill.background()

    reveal = [bg, triangle]

    # 3. Text Content
    headline = get_slot(slots, "headline")
    subhead = get_slot(slots, "subhead")

    has_sub = bool(subhead)
    start_y = 2.65 if has_sub else 3.15

    if headline:
        head_box = slide.shapes.add_textbox(Inches(1.0), Inches(start_y), Inches(7.0), Inches(1.2))
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title_xl", color_key="text")
        reveal.append(head_box)
        start_y += 1.4

    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(start_y), Inches(7.0), Inches(0.8))
        set_textbox_text(sub_box, truncate_to(subhead, 100), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)

    # 4. Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
