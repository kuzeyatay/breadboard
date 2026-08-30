from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "hero_image", "kind": "image", "required": True},
    {"name": "headline",   "kind": "text",  "style": "title_xl", "required": True,  "max_chars": 40},
    {"name": "subhead",    "kind": "text",  "style": "subtitle", "required": False, "max_chars": 60},
]

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Hero Image (Left)
    img_path = get_slot(slots, "hero_image")
    if img_path:
        try:
            pic = slide.shapes.add_picture(img_path, Inches(0.6), Inches(0.75), height=Inches(6.0))
            reveal.append(pic)
        except Exception:
            placeholder = add_solid_rect(slide, 0.6, 0.75, 4.5, 6.0, theme, color_key="muted")
            reveal.append(placeholder)
    else:
        placeholder = add_solid_rect(slide, 0.6, 0.75, 4.5, 6.0, theme, color_key="muted")
        reveal.append(placeholder)

    # Text Area (Right)
    text_x = 6.0
    text_w = 6.0

    # Headline
    headline = get_slot(slots, "headline", required=True)
    if headline:
        head_box = slide.shapes.add_textbox(Inches(text_x), Inches(3.0), Inches(text_w), Inches(1.2))
        set_textbox_text(head_box, truncate_to(headline, 40), theme, "title_xl", color_key="text")
        for paragraph in head_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.RIGHT
        reveal.append(head_box)

    # Subhead
    subhead = get_slot(slots, "subhead")
    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(text_x), Inches(4.3), Inches(text_w), Inches(0.8))
        set_textbox_text(sub_box, truncate_to(subhead, 60), theme, "subtitle", color_key="text")
        for paragraph in sub_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.RIGHT
        reveal.append(sub_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
