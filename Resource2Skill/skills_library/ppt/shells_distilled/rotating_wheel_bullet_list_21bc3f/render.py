"""
shells/rotating_wheel_bullets.py

A dynamic layout featuring a rotating wheel graphic on the left,
a prominent call-to-action banner, a directional arrow for the subhead,
and a bulleted list.
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text,
    truncate_to, add_infinite_rotation
)

AMBIENT_CAPABLE = True

SLOTS = [
    {"name": "call_to_action", "kind": "text", "max_chars": 25, "style": "title", "required": False},
    {"name": "headline", "kind": "text", "max_chars": 30, "style": "title_xl", "required": True},
    {"name": "subhead", "kind": "text", "max_chars": 40, "style": "subtitle", "required": False},
    {"name": "bullets", "kind": "bullet_list", "bullet_capacity": 4, "required": True}
]


def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal_shapes = []

    # Left Rotating Wheel (Ambient Motion)
    # Using a BLOCK_ARC so rotation is visually obvious
    wheel = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, Inches(-2.5), Inches(0.5), Inches(6.5), Inches(6.5))
    wheel.fill.solid()
    wheel.fill.fore_color.rgb = palette_color(theme, "accent")
    wheel.line.fill.background()
    wheel.adjustments[0] = 180.0  # Start angle
    wheel.adjustments[1] = 360.0  # End angle
    wheel.adjustments[2] = 0.2    # Thickness
    add_infinite_rotation(slide, wheel, duration_ms=15000, direction="cw")
    reveal_shapes.append(wheel)

    # Inner decorative wheel
    inner_wheel = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(-1.5), Inches(1.5), Inches(4.5), Inches(4.5))
    inner_wheel.fill.solid()
    inner_wheel.fill.fore_color.rgb = palette_color(theme, "panel")
    inner_wheel.line.fill.background()
    inner_wheel.adjustments[0] = 0.1
    add_infinite_rotation(slide, inner_wheel, duration_ms=20000, direction="ccw")
    reveal_shapes.append(inner_wheel)

    # Top CTA Banner
    cta_text = get_slot(slots, "call_to_action")
    if cta_text:
        cta_banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(0.4), Inches(9.833), Inches(1.2))
        cta_banner.fill.solid()
        cta_banner.fill.fore_color.rgb = palette_color(theme, "accent2", fallback="accent")
        cta_banner.line.fill.background()
        set_textbox_text(cta_banner, truncate_to(cta_text, 25), theme, "title", color_key="bg")
        reveal_shapes.append(cta_banner)

    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.0), Inches(7.0), Inches(0.8))
    set_textbox_text(head_box, truncate_to(headline_text, 30), theme, "title_xl", color_key="accent")
    reveal_shapes.append(head_box)

    # Subhead Arrow Banner
    subhead_text = get_slot(slots, "subhead")
    if subhead_text:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.0), Inches(2.9), Inches(8.5), Inches(1.0))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = palette_color(theme, "accent")
        arrow.line.fill.background()
        arrow.adjustments[0] = 0.5  # Arrow head width
        arrow.adjustments[1] = 0.5  # Arrow neck width
        set_textbox_text(arrow, truncate_to(subhead_text, 40), theme, "subtitle", color_key="bg")
        reveal_shapes.append(arrow)

    # Bullets
    bullets = get_slot(slots, "bullets", required=True)
    if bullets:
        y_cursor = 4.3
        for i, item in enumerate(bullets[:4]):
            if isinstance(item, dict):
                text = item.get("title", "")
            else:
                text = str(item)
            
            bullet_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_cursor), Inches(7.0), Inches(0.6))
            set_textbox_text(bullet_box, f"• {truncate_to(text, 60)}", theme, "title", color_key="text")
            reveal_shapes.append(bullet_box)
            y_cursor += 0.7

    # Entrances
    for i, shape in enumerate(reveal_shapes):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
