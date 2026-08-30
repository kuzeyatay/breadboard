from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    gutter = theme.get("spacing", {}).get("gutter", 0.4)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    y = margin
    total_w = 13.333 - 2 * margin
    
    # Top hairline
    add_hairline(slide, margin, y, total_w, theme)
    y += 0.1
    
    # Kicker
    kicker = get_slot(slots, "kicker")
    if kicker:
        k_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(total_w), Inches(0.3))
        set_textbox_text(k_box, truncate_to(kicker, 80), theme, "caption", color_key="muted")
        reveal.append(k_box)
    y += 0.4
    
    # Bottom hairline
    add_hairline(slide, margin, y, total_w, theme)
    y += 0.2
    
    # Headline
    headline = get_slot(slots, "headline", required=True)
    h_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(total_w), Inches(1.2))
    set_textbox_text(h_box, truncate_to(headline, 120), theme, "title", color_key="text")
    reveal.append(h_box)
    y += 1.3
    
    # Byline
    byline = get_slot(slots, "byline")
    if byline:
        b_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(total_w), Inches(0.3))
        set_textbox_text(b_box, truncate_to(byline, 50), theme, "caption", color_key="accent")
        reveal.append(b_box)
    y += 0.5
    
    # 3 Columns
    col_w = (total_w - 2 * gutter) / 3
    col_h = 7.5 - y - margin
    
    # Col 1
    col1_text = get_slot(slots, "col1_text")
    if col1_text:
        c1_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(col_w), Inches(col_h))
        set_textbox_text(c1_box, truncate_to(col1_text, 800), theme, "body", color_key="text")
        reveal.append(c1_box)
        
    # Col 2
    x2 = margin + col_w + gutter
    img_slot = get_slot(slots, "image")
    cy = y
    if img_slot:
        img_h = col_w * 0.6
        img_rect = add_solid_rect(slide, x2, cy, col_w, img_h, theme, color_key="muted")
        reveal.append(img_rect)
        cy += img_h + 0.1
        
        caption = get_slot(slots, "caption")
        if caption:
            cap_box = slide.shapes.add_textbox(Inches(x2), Inches(cy), Inches(col_w), Inches(0.4))
            set_textbox_text(cap_box, truncate_to(caption, 150), theme, "caption", color_key="muted")
            reveal.append(cap_box)
            cy += 0.5
            
    col2_text = get_slot(slots, "col2_text")
    if col2_text:
        c2_box = slide.shapes.add_textbox(Inches(x2), Inches(cy), Inches(col_w), Inches(7.5 - cy - margin))
        set_textbox_text(c2_box, truncate_to(col2_text, 400), theme, "body", color_key="text")
        reveal.append(c2_box)
        
    # Col 3
    x3 = x2 + col_w + gutter
    col3_text = get_slot(slots, "col3_text")
    if col3_text:
        c3_box = slide.shapes.add_textbox(Inches(x3), Inches(y), Inches(col_w), Inches(col_h))
        set_textbox_text(c3_box, truncate_to(col3_text, 800), theme, "body", color_key="text")
        reveal.append(c3_box)
        
    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
