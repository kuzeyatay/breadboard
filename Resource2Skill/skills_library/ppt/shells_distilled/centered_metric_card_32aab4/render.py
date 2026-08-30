from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(1.66), Inches(1.2), Inches(10.0), Inches(1.0))
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="text")
        head_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(head_box)

    # Metric Card
    card_w, card_h = 8.0, 3.5
    card_x = (13.333 - card_w) / 2
    card_y = 2.8

    card = add_solid_rect(
        slide, card_x, card_y, card_w, card_h, theme,
        color_key="bg", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    card.line.color.rgb = palette_color(theme, "accent")
    card.line.width = Inches(0.04)
    card.adjustments[0] = 0.1
    reveal.append(card)

    # Metric Value & Unit
    metric_val = get_slot(slots, "metric_value", required=True)
    metric_unit = get_slot(slots, "metric_unit")

    if metric_unit:
        val_w = 5.0
        unit_w = 2.5
        val_x = card_x + 0.2
        unit_x = val_x + val_w
        
        val_box = slide.shapes.add_textbox(Inches(val_x), Inches(card_y + 0.6), Inches(val_w), Inches(2.0))
        set_textbox_text(val_box, truncate_to(metric_val, 15), theme, "metric_xl", color_key="accent")
        val_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        unit_box = slide.shapes.add_textbox(Inches(unit_x), Inches(card_y + 1.8), Inches(unit_w), Inches(1.0))
        set_textbox_text(unit_box, truncate_to(metric_unit, 10), theme, "title", color_key="text")
        unit_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        reveal.extend([val_box, unit_box])
    else:
        val_box = slide.shapes.add_textbox(Inches(card_x + 0.5), Inches(card_y + 0.6), Inches(card_w - 1.0), Inches(2.0))
        set_textbox_text(val_box, truncate_to(metric_val, 20), theme, "metric_xl", color_key="accent")
        val_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(val_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
