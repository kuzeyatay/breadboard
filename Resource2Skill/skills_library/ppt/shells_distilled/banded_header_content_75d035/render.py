from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Header Banding Structure
    # Top thin line
    line1 = add_solid_rect(slide, 0, 0.1, 13.333, 0.03, theme, color_key="accent")
    # Main thick banner
    banner = add_solid_rect(slide, 0, 0.18, 13.333, 1.2, theme, color_key="accent")
    # Bottom thin line
    line2 = add_solid_rect(slide, 0, 1.43, 13.333, 0.03, theme, color_key="accent")
    
    reveal.extend([line1, banner, line2])
    
    # Headline (placed inside the banner)
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(0.25), Inches(13.333 - 2 * margin), Inches(1.0))
        # Using 'bg' color for text to contrast against the solid 'accent' banner
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="bg")
        reveal.append(head)
        
    # Body Content (placed in the main white area)
    body_text = get_slot(slots, "body")
    if body_text:
        body = slide.shapes.add_textbox(Inches(margin), Inches(2.0), Inches(13.333 - 2 * margin), Inches(4.5))
        set_textbox_text(body, truncate_to(body_text, 500), theme, "body", color_key="text")
        reveal.append(body)
        
    # Staggered entrance animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
