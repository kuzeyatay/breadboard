def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from _shell_helpers import add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    
    box_w = 5.0
    x = (13.333 - box_w) / 2
    
    headline = get_slot(slots, "headline", required=True)
    subhead = get_slot(slots, "subhead")
    body = get_slot(slots, "body")
    
    # Calculate total height to vertically center the block
    total_h = 1.5
    if subhead:
        total_h += 0.7
    if body:
        total_h += 2.5
    
    y = (7.5 - total_h) / 2
    
    # Headline
    head_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_w), Inches(1.5))
    set_textbox_text(head_box, truncate_to(headline, 80), theme, "title", color_key="text")
    reveal.append(head_box)
    y += 1.5
        
    # Subhead (Accent)
    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_w), Inches(0.5))
        set_textbox_text(sub_box, truncate_to(subhead, 40), theme, "subtitle", color_key="accent")
        reveal.append(sub_box)
        y += 0.7
        
    # Body
    if body:
        body_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_w), Inches(2.5))
        set_textbox_text(body_box, truncate_to(body, 300), theme, "body", color_key="muted")
        reveal.append(body_box)
        
    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)