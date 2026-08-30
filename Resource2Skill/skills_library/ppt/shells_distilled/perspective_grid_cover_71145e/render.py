from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # Top background (Sky)
    top_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(3.75))
    top_bg.fill.solid()
    top_bg.fill.fore_color.rgb = palette_color(theme, "accent")
    top_bg.line.fill.background()
    reveal.append(top_bg)

    # Bottom background (Ground)
    bot_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.75), Inches(13.333), Inches(3.75))
    bot_bg.fill.solid()
    bot_bg.fill.fore_color.rgb = palette_color(theme, "accent2")
    bot_bg.line.fill.background()
    reveal.append(bot_bg)

    # Horizon Line
    horizon = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(3.75), Inches(13.333), Inches(3.75))
    horizon.line.color.rgb = palette_color(theme, "bg")
    horizon.line.width = Inches(0.02)
    reveal.append(horizon)

    # Perspective Lines
    center_x = 13.333 / 2
    center_y = 3.75
    bottom_y = 7.5
    x_points = [0, 3.333, 6.666, 10.0, 13.333]

    for x in x_points:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(center_x), Inches(center_y), Inches(x), Inches(bottom_y))
        line.line.color.rgb = palette_color(theme, "muted")
        line.line.width = Inches(0.01)
        reveal.append(line)

    # Headline
    headline_text = get_slot(slots, "headline")
    if headline_text:
        head = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(1.5))
        set_textbox_text(head, truncate_to(headline_text, 80), theme, "title_xl", color_key="bg")
        reveal.append(head)

    # Subtitle
    subtitle_text = get_slot(slots, "subtitle")
    if subtitle_text:
        sub = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(1.0))
        set_textbox_text(sub, truncate_to(subtitle_text, 120), theme, "subtitle", color_key="bg")
        reveal.append(sub)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
