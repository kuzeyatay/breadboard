from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "headline",   "kind": "text",  "max_chars": 60,  "style": "title_xl", "required": True},
    {"name": "subtitle",   "kind": "text",  "max_chars": 100, "style": "subtitle", "required": False},
    {"name": "hero_image", "kind": "image", "aspect": "16:9",                      "required": False},
]

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(1.66), Inches(2.0), Inches(10.0), Inches(1.0))
    set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title_xl", color_key="text")
    reveal.append(head_box)

    # Thick Accent Divider
    divider = add_solid_rect(
        slide, 4.66, 3.1, 4.0, 0.08, theme,
        color_key="accent", line=False
    )
    reveal.append(divider)

    # Subtitle
    subtitle_text = get_slot(slots, "subtitle")
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(1.66), Inches(3.3), Inches(10.0), Inches(0.8))
        set_textbox_text(sub_box, truncate_to(subtitle_text, 100), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)

    # Hero Image (Optional central graphic)
    hero_image = get_slot(slots, "hero_image")
    if hero_image:
        try:
            pic = slide.shapes.add_picture(hero_image, Inches(4.66), Inches(4.5), width=Inches(4.0))
            reveal.append(pic)
        except Exception:
            # Fallback if image path is invalid in environment
            placeholder = add_solid_rect(slide, 4.66, 4.5, 4.0, 2.0, theme, color_key="panel", line=True)
            reveal.append(placeholder)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
