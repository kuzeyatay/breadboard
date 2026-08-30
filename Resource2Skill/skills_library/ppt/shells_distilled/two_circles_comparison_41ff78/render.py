from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Left Circle (Smaller)
    left_w = 2.5
    left_h = 2.5
    left_x = 3.916 - left_w / 2
    left_y = 3.75 - left_h / 2
    
    left_circle = add_solid_rect(
        slide, left_x, left_y, left_w, left_h, theme,
        color_key="accent", line=False, shape_type=MSO_SHAPE.OVAL
    )
    reveal.append(left_circle)

    left_text = get_slot(slots, "left_text")
    if left_text:
        lb = slide.shapes.add_textbox(Inches(left_x), Inches(left_y + left_h/2 - 0.5), Inches(left_w), Inches(1.0))
        set_textbox_text(lb, truncate_to(left_text, 20), theme, "title", color_key="bg")
        reveal.append(lb)

    # Right Circle (Larger)
    right_w = 3.5
    right_h = 3.5
    right_x = 7.416 - right_w / 2
    right_y = 3.75 - right_h / 2

    right_circle = add_solid_rect(
        slide, right_x, right_y, right_w, right_h, theme,
        color_key="accent", line=False, shape_type=MSO_SHAPE.OVAL
    )
    reveal.append(right_circle)

    right_text = get_slot(slots, "right_text")
    if right_text:
        rb = slide.shapes.add_textbox(Inches(right_x), Inches(right_y + right_h/2 - 0.5), Inches(right_w), Inches(1.0))
        set_textbox_text(rb, truncate_to(right_text, 20), theme, "title", color_key="bg")
        reveal.append(rb)

    # Entrance animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
