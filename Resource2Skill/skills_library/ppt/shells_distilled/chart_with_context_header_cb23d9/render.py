from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "headline", "kind": "text", "style": "title", "required": True},
    {"name": "subtitle", "kind": "text", "style": "subtitle", "required": False},
    {"name": "context", "kind": "text", "style": "caption", "required": False},
    {"name": "chart", "kind": "chart", "required": True},
    {"name": "footer", "kind": "text", "style": "caption", "required": False}
]

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    y = margin

    # Headline
    headline = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.6))
    set_textbox_text(head_box, truncate_to(headline, 100), theme, "title", color_key="text")
    reveal.append(head_box)
    y += 0.7

    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.4))
        set_textbox_text(sub_box, truncate_to(subtitle, 100), theme, "subtitle", color_key="text")
        reveal.append(sub_box)
        y += 0.4

    # Context
    context = get_slot(slots, "context")
    if context:
        ctx_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.3))
        set_textbox_text(ctx_box, truncate_to(context, 100), theme, "caption", color_key="muted")
        reveal.append(ctx_box)
        y += 0.35

    # Divider
    add_hairline(slide, margin, y, 13.333 - 2 * margin, theme)
    y += 0.2

    # Chart Placeholder
    chart_h = 6.8 - y
    chart_w = 13.333 - 2 * margin
    chart_placeholder = add_solid_rect(
        slide, margin, y, chart_w, chart_h, theme,
        color_key="panel", line=True, shape_type=MSO_SHAPE.RECTANGLE
    )
    ph_label = slide.shapes.add_textbox(Inches(margin), Inches(y + chart_h/2 - 0.25), Inches(chart_w), Inches(0.5))
    set_textbox_text(ph_label, "[ Chart Area ]", theme, "body", color_key="muted")
    for paragraph in ph_label.text_frame.paragraphs:
        paragraph.alignment = 2  # Center align
    
    reveal.extend([chart_placeholder, ph_label])

    # Footer
    footer = get_slot(slots, "footer")
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(margin), Inches(7.0), Inches(13.333 - 2 * margin), Inches(0.3))
        set_textbox_text(foot_box, truncate_to(footer, 100), theme, "caption", color_key="muted")
        for paragraph in foot_box.text_frame.paragraphs:
            paragraph.alignment = 2  # Center align
        reveal.append(foot_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
