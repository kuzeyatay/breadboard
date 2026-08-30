def render(slide, slots: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Dimensions based on the visual proportions
    card_w = 4.7
    card_h = 5.9
    gutter = 0.1
    start_x = (13.333 - (2 * card_w + gutter)) / 2
    start_y = (7.5 - card_h) / 2

    # Left Card
    left_x = start_x
    left_card = add_solid_rect(slide, left_x, start_y, card_w, card_h, theme, color_key="primary")
    reveal.append(left_card)

    left_title = truncate_to(get_slot(slots, "left_title", default="Left Option"), 40)
    lt_box = slide.shapes.add_textbox(Inches(left_x + 0.4), Inches(start_y + 0.5), Inches(card_w - 0.8), Inches(0.8))
    set_textbox_text(lt_box, left_title, theme, "title", color_key="bg")
    reveal.append(lt_box)

    left_body = truncate_to(get_slot(slots, "left_body", default="Description for the left option goes here."), 200)
    lb_box = slide.shapes.add_textbox(Inches(left_x + 0.4), Inches(start_y + 1.5), Inches(card_w - 0.8), Inches(card_h - 2.0))
    set_textbox_text(lb_box, left_body, theme, "body", color_key="bg")
    reveal.append(lb_box)

    # Right Card
    right_x = start_x + card_w + gutter
    right_card = add_solid_rect(slide, right_x, start_y, card_w, card_h, theme, color_key="accent")
    reveal.append(right_card)

    right_title = truncate_to(get_slot(slots, "right_title", default="Right Option"), 40)
    rt_box = slide.shapes.add_textbox(Inches(right_x + 0.4), Inches(start_y + 0.5), Inches(card_w - 0.8), Inches(0.8))
    set_textbox_text(rt_box, right_title, theme, "title", color_key="bg")
    reveal.append(rt_box)

    right_body = truncate_to(get_slot(slots, "right_body", default="Description for the right option goes here."), 200)
    rb_box = slide.shapes.add_textbox(Inches(right_x + 0.4), Inches(start_y + 1.5), Inches(card_w - 0.8), Inches(card_h - 2.0))
    set_textbox_text(rb_box, right_body, theme, "body", color_key="bg")
    reveal.append(rb_box)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
