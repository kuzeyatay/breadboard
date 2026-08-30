from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)


def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # Background base (explicitly using palette_color to satisfy strict checks)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    reveal.append(bg)

    # Hero Image
    hero_img = get_slot(slots, "hero_image")
    if hero_img:
        try:
            pic = slide.shapes.add_picture(hero_img, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            reveal.append(pic)
        except Exception:
            pass

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.333), Inches(2.0))
        set_textbox_text(head_box, truncate_to(headline, 30), theme, "title_xl", color_key="text")
        reveal.append(head_box)

    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10.333), Inches(1.5))
        set_textbox_text(sub_box, truncate_to(subtitle, 60), theme, "title", color_key="muted")
        reveal.append(sub_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
