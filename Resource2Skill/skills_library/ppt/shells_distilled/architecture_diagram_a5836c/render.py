from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        set_textbox_text(head, truncate_to(headline, 60), theme, "title", color_key="text")
        reveal.append(head)

    # Footer
    footer = get_slot(slots, "footer")
    if footer:
        foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
        set_textbox_text(foot, truncate_to(footer, 150), theme, "caption", color_key="muted")
        for p in foot.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        reveal.append(foot)

    # Layout metrics
    content_y = 1.3
    content_h = 5.2
    left_x = 0.5
    left_w = 2.0
    right_w = 2.0
    right_x = 13.333 - 0.5 - right_w
    center_x = left_x + left_w + 0.2
    center_w = right_x - center_x - 0.2

    # Left Column
    left_title = get_slot(slots, "left_title")
    left_items = get_slot(slots, "left_items", default=[])
    
    if left_title or left_items:
        if left_title:
            lt_box = add_solid_rect(slide, left_x, content_y, left_w, 0.8, theme, "panel")
            lt_txt = slide.shapes.add_textbox(Inches(left_x), Inches(content_y), Inches(left_w), Inches(0.8))
            set_textbox_text(lt_txt, truncate_to(left_title, 30), theme, "body_bold", color_key="text")
            for p in lt_txt.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
            reveal.extend([lt_box, lt_txt])
            
        if left_items:
            start_y = content_y + 0.9 if left_title else content_y
            avail_h = content_h - 0.9 if left_title else content_h
            item_h = min(0.6, avail_h / max(len(left_items), 1))
            
            for i, item in enumerate(left_items):
                y = start_y + i * (item_h + 0.05)
                text = item.get("title", "") if isinstance(item, dict) else str(item)
                
                ibox = add_solid_rect(slide, left_x, y, left_w, item_h, theme, "bg", line=True)
                itxt = slide.shapes.add_textbox(Inches(left_x), Inches(y), Inches(left_w), Inches(item_h))
                set_textbox_text(itxt, truncate_to(text, 30), theme, "body", color_key="text")
                for p in itxt.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                reveal.extend([ibox, itxt])

    # Right Column
    right_title = get_slot(slots, "right_title")
    right_items = get_slot(slots, "right_items", default=[])
    
    if right_title or right_items:
        if right_title:
            rt_box = add_solid_rect(slide, right_x, content_y, right_w, 0.8, theme, "panel")
            rt_txt = slide.shapes.add_textbox(Inches(right_x), Inches(content_y), Inches(right_w), Inches(0.8))
            set_textbox_text(rt_txt, truncate_to(right_title, 30), theme, "body_bold", color_key="text")
            for p in rt_txt.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
            reveal.extend([rt_box, rt_txt])
            
        if right_items:
            start_y = content_y + 0.9 if right_title else content_y
            avail_h = content_h - 0.9 if right_title else content_h
            item_h = min(0.6, avail_h / max(len(right_items), 1))
            
            for i, item in enumerate(right_items):
                y = start_y + i * (item_h + 0.05)
                text = item.get("title", "") if isinstance(item, dict) else str(item)
                
                ibox = add_solid_rect(slide, right_x, y, right_w, item_h, theme, "bg", line=True)
                itxt = slide.shapes.add_textbox(Inches(right_x), Inches(y), Inches(right_w), Inches(item_h))
                set_textbox_text(itxt, truncate_to(text, 30), theme, "body", color_key="text")
                for p in itxt.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                reveal.extend([ibox, itxt])

    # Center Layers
    center_layers = get_slot(slots, "center_layers", default=[])
    if center_layers:
        layer_spacing = 0.1
        total_spacing = layer_spacing * (len(center_layers) - 1)
        layer_h = (content_h - total_spacing) / max(len(center_layers), 1)
        
        for i, layer in enumerate(center_layers):
            y = content_y + i * (layer_h + layer_spacing)
            
            header_h = 0.5
            body_h = layer_h - header_h
            
            color_key = "primary" if i % 2 == 0 else "secondary"
            
            # Header
            h_box = add_solid_rect(slide, center_x, y, center_w, header_h, theme, color_key)
            h_txt = slide.shapes.add_textbox(Inches(center_x), Inches(y), Inches(center_w), Inches(header_h))
            set_textbox_text(h_txt, truncate_to(layer.get("title", ""), 50), theme, "body_bold", color_key="bg")
            for p in h_txt.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                
            # Body
            b_box = add_solid_rect(slide, center_x, y + header_h, center_w, body_h, theme, "panel")
            b_txt = slide.shapes.add_textbox(Inches(center_x + 0.2), Inches(y + header_h + 0.1), Inches(center_w - 0.4), Inches(body_h - 0.2))
            set_textbox_text(b_txt, truncate_to(layer.get("body", ""), 200), theme, "body", color_key="text")
            
            reveal.extend([h_box, h_txt, b_box, b_txt])

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
