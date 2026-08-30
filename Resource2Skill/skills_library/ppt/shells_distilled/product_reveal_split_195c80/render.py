from __future__ import annotations
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Left side: Hero Image
    img_path = get_slot(slots, "hero_image")
    img_x, img_y, img_w, img_h = 2.5, 1.5, 2.5, 4.5
    if img_path:
        try:
            pic = slide.shapes.add_picture(img_path, Inches(img_x), Inches(img_y), width=Inches(img_w), height=Inches(img_h))
            reveal.append(pic)
        except Exception:
            ph = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="panel")
            reveal.append(ph)
    else:
        ph = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="panel")
        reveal.append(ph)

    # Right side: Text
    text_x = 6.5
    text_w = 5.5

    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(text_x), Inches(2.5), Inches(text_w), Inches(1.0))
        set_textbox_text(head_box, truncate_to(headline, 50), theme, "title_xl", color_key="text")
        reveal.append(head_box)

    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(text_x), Inches(3.8), Inches(text_w), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subtitle, 50), theme, "title", color_key="text")
        reveal.append(sub_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
