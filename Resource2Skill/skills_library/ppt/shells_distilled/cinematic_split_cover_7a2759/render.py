from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # Base Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Image Area (Placeholder for hero_image slot)
    img_rect = add_solid_rect(
        slide, 0, 0.8, 13.333, 5.2, theme, 
        color_key="muted", line=False, shape_type=MSO_SHAPE.RECTANGLE
    )
    reveal.append(img_rect)

    # Top Bar (Panel)
    top_bar = add_solid_rect(
        slide, 0, 0, 13.333, 0.8, theme, 
        color_key="panel", line=False, shape_type=MSO_SHAPE.RECTANGLE
    )
    reveal.append(top_bar)

    kicker = get_slot(slots, "kicker")
    if kicker:
        k_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.4))
        set_textbox_text(k_box, truncate_to(kicker, 50), theme, "caption", color_key="text")
        reveal.append(k_box)

    # Bottom Bar (Panel)
    bottom_bar = add_solid_rect(
        slide, 0, 6.0, 13.333, 1.5, theme, 
        color_key="panel", line=False, shape_type=MSO_SHAPE.RECTANGLE
    )
    reveal.append(bottom_bar)

    headline = get_slot(slots, "headline", required=True)
    h_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.6))
    set_textbox_text(h_box, truncate_to(headline, 80), theme, "title", color_key="text")
    reveal.append(h_box)

    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        s_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4))
        set_textbox_text(s_box, truncate_to(subtitle, 100), theme, "body", color_key="muted")
        reveal.append(s_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
