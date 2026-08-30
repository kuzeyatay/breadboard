from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 1.0)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Grid layout parameters
    col_w = 4.5
    gutter_x = 1.0
    row_h = 1.8
    start_y = 3.0

    # Center the grid horizontally
    total_grid_w = (2 * col_w) + gutter_x
    start_x = (13.333 - total_grid_w) / 2

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 80)
    head = slide.shapes.add_textbox(Inches(start_x), Inches(1.0), Inches(13.333 - start_x - margin), Inches(1.2))
    set_textbox_text(head, headline, theme, "title", color_key="text")
    reveal.append(head)

    # Agenda Items
    items = get_slot(slots, "agenda_items", required=True)
    if not isinstance(items, list):
        items = []

    max_items = 4
    items = items[:max_items]

    for i, item in enumerate(items):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}

        r = i // 2
        c = i % 2

        x = start_x + c * (col_w + gutter_x)
        y = start_y + r * row_h

        # Number Badge (Circle)
        badge_size = 1.0
        badge = add_solid_rect(
            slide, x, y, badge_size, badge_size, theme,
            color_key="primary", line=False, shape_type=MSO_SHAPE.OVAL
        )
        set_textbox_text(badge, f"{i+1:02d}", theme, "title", color_key="bg")
        reveal.append(badge)

        # Item Title
        text_x = x + badge_size + 0.3
        text_w = col_w - badge_size - 0.3

        t_box = slide.shapes.add_textbox(Inches(text_x), Inches(y), Inches(text_w), Inches(0.5))
        set_textbox_text(t_box, truncate_to(item.get("title", ""), 40), theme, "body_bold", color_key="text")
        reveal.append(t_box)

        # Item Body
        b_box = slide.shapes.add_textbox(Inches(text_x), Inches(y + 0.5), Inches(text_w), Inches(1.0))
        set_textbox_text(b_box, truncate_to(item.get("body", ""), 100), theme, "body", color_key="muted")
        reveal.append(b_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
