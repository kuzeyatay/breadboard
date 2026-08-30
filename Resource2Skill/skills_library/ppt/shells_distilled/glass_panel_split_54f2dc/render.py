from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # 1. Background
    bg_img = get_slot(slots, "background_image")
    if bg_img:
        bg = slide.shapes.add_picture(bg_img, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
        reveal.append(bg)
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = palette_color(theme, "bg")
        bg.line.fill.background()
        reveal.append(bg)

    # 2. Central Rounded Panel
    panel_w, panel_h = 10.0, 5.5
    panel_x = (13.333 - panel_w) / 2
    panel_y = (7.5 - panel_h) / 2

    panel = add_solid_rect(
        slide, panel_x, panel_y, panel_w, panel_h, theme,
        color_key="panel", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    if panel.adjustments:
        panel.adjustments[0] = 0.08
    reveal.append(panel)

    # 3. Hero Image (Left side of panel)
    hero_img = get_slot(slots, "hero_image")
    img_w, img_h = 3.5, 4.5
    img_x = panel_x + 0.5
    img_y = panel_y + 0.5

    if hero_img:
        hero = slide.shapes.add_picture(hero_img, Inches(img_x), Inches(img_y), width=Inches(img_w))
        reveal.append(hero)
    else:
        ph = add_solid_rect(
            slide, img_x, img_y, img_w, img_h, theme,
            color_key="muted", line=False, shape_type=MSO_SHAPE.RECTANGLE
        )
        reveal.append(ph)

    # 4. Text Content (Right side of panel)
    text_x = img_x + img_w + 0.5
    text_w = panel_w - img_w - 1.5
    text_y = panel_y + 0.8

    headline_text = get_slot(slots, "headline")
    if headline_text:
        head_box = slide.shapes.add_textbox(Inches(text_x), Inches(text_y), Inches(text_w), Inches(1.2))
        set_textbox_text(head_box, truncate_to(headline_text, 80), theme, "title", color_key="text")
        reveal.append(head_box)
        text_y += 1.3

    body_text = get_slot(slots, "body")
    if body_text:
        body_box = slide.shapes.add_textbox(Inches(text_x), Inches(text_y), Inches(text_w), Inches(2.5))
        set_textbox_text(body_box, truncate_to(body_text, 300), theme, "body", color_key="text")
        reveal.append(body_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
