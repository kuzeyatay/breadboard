from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.5)
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Headline
    headline = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(margin), Inches(0.4), Inches(13.333 - 2 * margin), Inches(0.8))
    set_textbox_text(head_box, truncate_to(headline, 100), theme, "title", color_key="text")
    reveal.append(head_box)

    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(margin), Inches(1.3), Inches(13.333 - 2 * margin), Inches(0.4))
        set_textbox_text(sub_box, truncate_to(subtitle, 120), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)

    # Divider Line
    add_hairline(slide, margin, 1.8, 13.333 - 2 * margin, theme)

    # Content Areas (Side-by-Side)
    content_y = 2.2
    content_h = 4.5
    gutter = 0.5
    half_w = (13.333 - 2 * margin - gutter) / 2

    # Left Chart Placeholder
    left_chart = add_solid_rect(
        slide, margin, content_y, half_w, content_h, theme,
        color_key="panel", line=True
    )
    reveal.append(left_chart)

    # Right Chart Placeholder
    right_chart = add_solid_rect(
        slide, margin + half_w + gutter, content_y, half_w, content_h, theme,
        color_key="panel", line=True
    )
    reveal.append(right_chart)

    # Footer
    footer = get_slot(slots, "footer")
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(margin), Inches(7.0), Inches(13.333 - 2 * margin), Inches(0.3))
        set_textbox_text(foot_box, truncate_to(footer, 80), theme, "caption", color_key="muted")
        reveal.append(foot_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
