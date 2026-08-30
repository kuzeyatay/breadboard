from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Central Circle
    diameter = 6.5
    x = (13.333 - diameter) / 2.0
    y = (7.5 - diameter) / 2.0

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter))
    circle.fill.solid()
    circle.fill.fore_color.rgb = palette_color(theme, "primary")
    circle.line.fill.background()
    reveal.append(circle)

    # Optional Headline inside the circle
    headline = get_slot(slots, "headline")
    if headline:
        box_w = diameter * 0.8
        box_h = 2.0
        box_x = x + (diameter - box_w) / 2.0
        box_y = y + (diameter - box_h) / 2.0
        
        head = slide.shapes.add_textbox(Inches(box_x), Inches(box_y), Inches(box_w), Inches(box_h))
        set_textbox_text(head, truncate_to(headline, 60), theme, "title", color_key="bg")
        reveal.append(head)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=80, duration_ms=500, index=i)
