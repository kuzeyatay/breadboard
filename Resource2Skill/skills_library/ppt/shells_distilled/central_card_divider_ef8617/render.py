from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "headline", "kind": "text", "max_chars": 60, "style": "title", "required": False},
    {"name": "subtitle", "kind": "text", "max_chars": 100, "style": "subtitle", "required": False},
]

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Central Card Dimensions
    card_w = 6.0
    card_h = 6.0
    x = (13.333 - card_w) / 2
    y = (7.5 - card_h) / 2

    # Top Accent Band
    band_h = 0.6
    band = add_solid_rect(
        slide, x, y, card_w, band_h, theme,
        color_key="accent", line=False
    )
    reveal.append(band)

    # Main Card Body
    body_h = card_h - band_h
    body = add_solid_rect(
        slide, x, y + band_h, card_w, body_h, theme,
        color_key="primary", line=False
    )
    reveal.append(body)

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(
            Inches(x + 0.5), Inches(y + band_h + 1.0), Inches(card_w - 1.0), Inches(1.5)
        )
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="bg")
        reveal.append(head_box)

    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(x + 0.5), Inches(y + band_h + 2.5), Inches(card_w - 1.0), Inches(1.0)
        )
        set_textbox_text(sub_box, truncate_to(subtitle, 100), theme, "subtitle", color_key="bg")
        reveal.append(sub_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
