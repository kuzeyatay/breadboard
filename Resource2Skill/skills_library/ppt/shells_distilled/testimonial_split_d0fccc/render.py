from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 1.0)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Kicker
    kicker = get_slot(slots, "kicker")
    if kicker:
        k_box = slide.shapes.add_textbox(Inches(margin), Inches(0.4), Inches(6.0), Inches(0.4))
        set_textbox_text(k_box, truncate_to(kicker, 40), theme, "caption", color_key="muted")
        reveal.append(k_box)
        
    # Hero Image Placeholder
    hero_w = 13.333 - 2 * margin
    hero_h = 2.5
    hero_img = get_slot(slots, "hero_image")
    if hero_img:
        try:
            pic = slide.shapes.add_picture(hero_img, Inches(margin), Inches(1.0), Inches(hero_w), Inches(hero_h))
            reveal.append(pic)
        except Exception:
            ph = add_solid_rect(slide, margin, 1.0, hero_w, hero_h, theme, color_key="accent")
            reveal.append(ph)
    else:
        ph = add_solid_rect(slide, margin, 1.0, hero_w, hero_h, theme, color_key="accent")
        reveal.append(ph)
        
    # Headline (overlaid on hero image)
    headline = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(margin + hero_w * 0.5), Inches(1.5), Inches(hero_w * 0.45), Inches(1.5))
    set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="bg")
    reveal.append(head_box)
    
    # Bottom Panel
    panel_y = 3.8
    panel_h = 3.0
    panel = add_solid_rect(slide, margin, panel_y, hero_w, panel_h, theme, color_key="panel")
    reveal.append(panel)
    
    # Author Image
    author_img = get_slot(slots, "author_image")
    author_x = margin + 0.5
    author_y = panel_y + 0.5
    if author_img:
        try:
            pic = slide.shapes.add_picture(author_img, Inches(author_x), Inches(author_y), Inches(1.0), Inches(1.0))
            reveal.append(pic)
        except Exception:
            ph = add_solid_rect(slide, author_x, author_y, 1.0, 1.0, theme, color_key="muted")
            reveal.append(ph)
    else:
        ph = add_solid_rect(slide, author_x, author_y, 1.0, 1.0, theme, color_key="muted")
        reveal.append(ph)
        
    # Author Name
    author_name = get_slot(slots, "author_name", required=True)
    name_box = slide.shapes.add_textbox(Inches(author_x), Inches(author_y + 1.2), Inches(2.5), Inches(0.4))
    set_textbox_text(name_box, truncate_to(author_name, 30), theme, "body_bold", color_key="text")
    reveal.append(name_box)
    
    # Author Title
    author_title = get_slot(slots, "author_title")
    if author_title:
        title_box = slide.shapes.add_textbox(Inches(author_x), Inches(author_y + 1.6), Inches(2.5), Inches(0.4))
        set_textbox_text(title_box, truncate_to(author_title, 40), theme, "caption", color_key="muted")
        reveal.append(title_box)
        
    # Quote
    quote = get_slot(slots, "quote", required=True)
    quote_x = margin + 3.5
    quote_w = hero_w - 4.0
    quote_box = slide.shapes.add_textbox(Inches(quote_x), Inches(panel_y + 0.5), Inches(quote_w), Inches(2.0))
    set_textbox_text(quote_box, truncate_to(quote, 200), theme, "subtitle", color_key="text")
    reveal.append(quote_box)
    
    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
