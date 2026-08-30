from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    gutter = theme.get("spacing", {}).get("gutter", 0.3)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    y = margin
    reveal = []

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)
        y += 1.2

    # Items
    items = get_slot(slots, "items", required=True)
    if not items:
        items = [{"title": f"Item {i+1}", "body": ""} for i in range(10)]

    max_items = min(len(items), 10)
    rows = (max_items + 1) // 2
    col_w = (13.333 - 2 * margin - gutter) / 2
    row_h = (7.5 - y - margin) / max(rows, 1)

    for i in range(max_items):
        item = items[i]
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}

        col = i % 2
        row = i // 2

        x = margin + col * (col_w + gutter)
        current_y = y + row * row_h

        # Large Number
        num_box = slide.shapes.add_textbox(Inches(x), Inches(current_y), Inches(1.2), Inches(row_h * 0.9))
        set_textbox_text(num_box, str(i + 1), theme, "title_xl", color_key="accent")
        reveal.append(num_box)

        # Text content
        text_x = x + 1.3
        text_w = col_w - 1.3
        t_box = slide.shapes.add_textbox(Inches(text_x), Inches(current_y + 0.1), Inches(text_w), Inches(row_h * 0.8))
        
        content = item.get("title", "")
        if item.get("body"):
            content += "\n" + item.get("body", "")
            
        set_textbox_text(t_box, truncate_to(content, 120), theme, "body", color_key="text")
        reveal.append(t_box)

    # Entrance animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
