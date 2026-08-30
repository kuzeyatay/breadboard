from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, add_emphasis_pulse, add_sequential_reveal,
    get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(4.5), Inches(1.0))
    set_textbox_text(head, truncate_to(headline_text, 40), theme, "title", color_key="text")
    add_theme_entrance(slide, head, theme, delay_ms=0, index=0)
    
    # Hero Image
    hero_img = get_slot(slots, "hero_image")
    hero_pic = None
    img_x, img_y, img_w, img_h = 6.2, 1.5, 6.0, 4.8
    if hero_img:
        try:
            hero_pic = slide.shapes.add_picture(hero_img, Inches(img_x), Inches(img_y), width=Inches(img_w), height=Inches(img_h))
        except Exception:
            hero_pic = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="panel", line=False)
    else:
        hero_pic = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="panel", line=False)
        
    if hero_pic:
        add_theme_entrance(slide, hero_pic, theme, delay_ms=100, index=1)
        add_emphasis_pulse(slide, hero_pic, theme, delay_ms=2000)
    
    # Metrics
    metrics = get_slot(slots, "metrics", required=True)
    if not isinstance(metrics, list):
        metrics = []
        
    y_offset = 2.5
    metric_shapes = []
    
    for i, item in enumerate(metrics[:2]):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
            
        # Title
        t_box = slide.shapes.add_textbox(Inches(1.0), Inches(y_offset), Inches(4.5), Inches(0.5))
        set_textbox_text(t_box, truncate_to(item.get("title", ""), 40), theme, "subtitle", color_key="text")
        metric_shapes.append(t_box)
        
        # Body
        b_box = slide.shapes.add_textbox(Inches(1.0), Inches(y_offset + 0.5), Inches(4.5), Inches(1.0))
        set_textbox_text(b_box, truncate_to(item.get("body", ""), 120), theme, "body", color_key="muted")
        metric_shapes.append(b_box)
        
        # Progress Bar Track
        track = add_solid_rect(slide, 1.75, y_offset + 1.6, 3.0, 0.06, theme, color_key="panel", line=False)
        metric_shapes.append(track)
        
        # Progress Bar Fill (alternating alignment for visual flair)
        if i % 2 == 0:
            fill = add_solid_rect(slide, 3.25, y_offset + 1.6, 1.5, 0.06, theme, color_key="accent", line=False)
        else:
            fill = add_solid_rect(slide, 1.75, y_offset + 1.6, 2.0, 0.06, theme, color_key="accent", line=False)
        metric_shapes.append(fill)
        
        y_offset += 2.3

    if metric_shapes:
        add_sequential_reveal(slide, metric_shapes, theme)
