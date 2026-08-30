from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline_text = get_slot(slots, "headline")
    if headline_text:
        head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(10), Inches(1.0))
        set_textbox_text(head, truncate_to(headline_text, 60), theme, "title", color_key="text")
        reveal.append(head)

    # Accent Bar (nod to the scroll graphic)
    bar_w = 8.0
    bar_h = 0.2
    bar_x = (13.333 - bar_w) / 2
    bar_y = 2.8
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bar_x), Inches(bar_y), Inches(bar_w), Inches(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = palette_color(theme, "accent")
    bar.line.fill.background()
    reveal.append(bar)

    # Quote Text
    quote_text = get_slot(slots, "quote", required=True)
    quote_box = slide.shapes.add_textbox(Inches(margin * 2), Inches(bar_y + 0.6), Inches(13.333 - margin * 4), Inches(2.0))
    set_textbox_text(quote_box, truncate_to(quote_text, 200), theme, "title", color_key="text")
    reveal.append(quote_box)

    # Author
    author_text = get_slot(slots, "author")
    if author_text:
        author_box = slide.shapes.add_textbox(Inches(margin * 2), Inches(bar_y + 2.8), Inches(13.333 - margin * 4), Inches(0.8))
        set_textbox_text(author_box, truncate_to(author_text, 60), theme, "body_bold", color_key="muted")
        reveal.append(author_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
