from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.5)
    gutter = theme.get("spacing", {}).get("gutter", 0.3)
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2*margin), Inches(0.8))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)

    content_y = margin + 1.0
    content_h = 7.5 - content_y - margin

    left_w = (13.333 - 2*margin - gutter) * 0.6
    right_w = (13.333 - 2*margin - gutter) * 0.4

    # Main Chart Panel (Left)
    main_panel = add_solid_rect(
        slide, margin, content_y, left_w, content_h, theme,
        color_key="panel", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    main_panel.adjustments[0] = 0.03
    reveal.append(main_panel)

    # Secondary Chart Panel (Top Right)
    right_x = margin + left_w + gutter
    top_right_h = (content_h - gutter) * 0.5
    sec_panel = add_solid_rect(
        slide, right_x, content_y, right_w, top_right_h, theme,
        color_key="panel", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    sec_panel.adjustments[0] = 0.05
    reveal.append(sec_panel)

    # Metrics Panel (Bottom Right)
    bottom_right_y = content_y + top_right_h + gutter
    bottom_right_h = content_h - top_right_h - gutter
    metrics_panel = add_solid_rect(
        slide, right_x, bottom_right_y, right_w, bottom_right_h, theme,
        color_key="panel", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    metrics_panel.adjustments[0] = 0.05
    reveal.append(metrics_panel)

    # Metrics Grid
    metrics = get_slot(slots, "metrics")
    if metrics and isinstance(metrics, list):
        pad = 0.3
        m_gutter = 0.2
        item_w = (right_w - 2*pad - m_gutter) / 2
        item_h = (bottom_right_h - 2*pad - m_gutter) / 2

        for i, item in enumerate(metrics[:4]):
            if not isinstance(item, dict):
                item = {"title": str(item), "body": ""}

            row = i // 2
            col = i % 2
            ix = right_x + pad + col * (item_w + m_gutter)
            iy = bottom_right_y + pad + row * (item_h + m_gutter)

            # Metric sub-panel
            sub_panel = add_solid_rect(
                slide, ix, iy, item_w, item_h, theme,
                color_key="bg", line=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
            )
            sub_panel.adjustments[0] = 0.1
            reveal.append(sub_panel)

            # Value
            val_box = slide.shapes.add_textbox(Inches(ix + 0.1), Inches(iy + 0.1), Inches(item_w - 0.2), Inches(item_h * 0.5))
            set_textbox_text(val_box, truncate_to(item.get("title", ""), 10), theme, "metric_xl", color_key="accent")
            reveal.append(val_box)

            # Label
            lbl_box = slide.shapes.add_textbox(Inches(ix + 0.1), Inches(iy + item_h * 0.55), Inches(item_w - 0.2), Inches(item_h * 0.35))
            set_textbox_text(lbl_box, truncate_to(item.get("body", ""), 20), theme, "caption", color_key="muted")
            reveal.append(lbl_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
