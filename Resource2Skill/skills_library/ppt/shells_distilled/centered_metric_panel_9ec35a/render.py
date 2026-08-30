from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text
)

SLOTS = [
    {"name": "metric", "kind": "text", "style": "metric_xl", "required": True}
]

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Center Panel Dimensions
    panel_w = 9.0
    panel_h = 4.0
    panel_x = (13.333 - panel_w) / 2
    panel_y = (7.5 - panel_h) / 2

    # Panel
    panel = add_solid_rect(
        slide, panel_x, panel_y, panel_w, panel_h, theme,
        color_key="panel", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
    )
    panel.adjustments[0] = 0.15  # Rounded corners
    
    # Style the panel border to be prominent (accent color)
    panel.line.color.rgb = palette_color(theme, "accent")
    panel.line.width = Inches(0.06)
    reveal.append(panel)

    # Metric Text
    metric_text = get_slot(slots, "metric", required=True)
    
    t_box = slide.shapes.add_textbox(Inches(panel_x), Inches(panel_y), Inches(panel_w), Inches(panel_h))
    set_textbox_text(t_box, metric_text, theme, "metric_xl", color_key="accent")
    
    # Center text vertically and horizontally
    tf = t_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        
    reveal.append(t_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
