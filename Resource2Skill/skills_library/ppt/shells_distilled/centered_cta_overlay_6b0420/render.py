from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # Background Image or Solid Fill
    bg_img = get_slot(slots, "hero_image")
    if bg_img:
        pic = slide.shapes.add_picture(bg_img, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        reveal.append(pic)
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = palette_color(theme, "bg")
        bg.line.fill.background()
        reveal.append(bg)

    # Kicker (Top Center)
    kicker_text = get_slot(slots, "kicker")
    if kicker_text:
        k_box = slide.shapes.add_textbox(Inches(2.66), Inches(1.0), Inches(8.0), Inches(0.5))
        set_textbox_text(k_box, truncate_to(kicker_text, 50), theme, "caption", color_key="text")
        k_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(k_box)

    # Headline (Center)
    headline_text = get_slot(slots, "headline", required=True)
    h_box = slide.shapes.add_textbox(Inches(1.66), Inches(2.8), Inches(10.0), Inches(1.5))
    set_textbox_text(h_box, truncate_to(headline_text, 40), theme, "title_xl", color_key="text")
    h_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.append(h_box)

    # Button (Below Headline)
    btn_text = get_slot(slots, "button_text")
    if btn_text:
        btn_w, btn_h = 4.0, 0.8
        btn_x = (13.333 - btn_w) / 2
        btn_y = 4.8
        
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(btn_x), Inches(btn_y), Inches(btn_w), Inches(btn_h))
        btn.fill.solid()
        btn.fill.fore_color.rgb = palette_color(theme, "panel")
        btn.line.fill.background()
        btn.adjustments[0] = 0.2
        
        btn_box = slide.shapes.add_textbox(Inches(btn_x), Inches(btn_y + 0.15), Inches(btn_w), Inches(0.5))
        set_textbox_text(btn_box, truncate_to(btn_text, 30), theme, "body_bold", color_key="text")
        btn_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        reveal.extend([btn, btn_box])

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
