from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches
from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get('spacing', {}).get('margin', 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, 'bg')
    bg.line.fill.background()
    reveal = []
    
    headline = get_slot(slots, 'headline', required=True)
    head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(1.0))
    set_textbox_text(head, truncate_to(headline, 80), theme, 'title', color_key='text')
    reveal.append(head)
    
    steps = get_slot(slots, 'steps', required=True)
    if not isinstance(steps, list): steps = []
    while len(steps) < 3: steps.append({'title': '', 'body': ''})
    steps = steps[:3]
    
    col_w = 2.8
    gap = (13.333 - 2 * margin - 3 * col_w) / 2
    base_y = 2.5
    base_h = 3.5
    
    # --- Column 1: Monolith ---
    x0 = margin
    step0 = steps[0]
    out0 = add_solid_rect(slide, x0, base_y, col_w, base_h, theme, color_key='bg', line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    out0.fill.background()
    out0.line.color.rgb = palette_color(theme, 'muted')
    out0.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    reveal.append(out0)
    
    t0 = slide.shapes.add_textbox(Inches(x0), Inches(base_y + 0.1), Inches(col_w), Inches(0.4))
    set_textbox_text(t0, truncate_to(step0.get('title', ''), 40), theme, 'body_bold', color_key='text')
    reveal.append(t0)
    
    in0 = add_solid_rect(slide, x0 + 0.2, base_y + 0.6, col_w - 0.4, base_h - 0.8, theme, color_key='accent', shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    reveal.append(in0)
    
    items0 = [line.strip() for line in step0.get('body', '').split('\n') if line.strip()][:3]
    for j, item_text in enumerate(items0):
        py = base_y + 1.0 + j * 0.8
        pill = add_solid_rect(slide, x0 + 0.4, py, col_w - 0.8, 0.5, theme, color_key='panel', shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        pill.adjustments[0] = 0.5
        pt = slide.shapes.add_textbox(Inches(x0 + 0.4), Inches(py + 0.05), Inches(col_w - 0.8), Inches(0.4))
        set_textbox_text(pt, truncate_to(item_text, 30), theme, 'caption', color_key='text')
        reveal.extend([pill, pt])
        
    # Arrow 1
    arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x0 + col_w + 0.1), Inches(base_y + base_h/2 - 0.1), Inches(gap - 0.2), Inches(0.2))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = palette_color(theme, 'accent')
    arr1.line.fill.background()
    reveal.append(arr1)
    
    # --- Column 2: Microservices ---
    x1 = x0 + col_w + gap
    step1 = steps[1]
    out1 = add_solid_rect(slide, x1, base_y, col_w, base_h, theme, color_key='bg', line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    out1.fill.background()
    out1.line.color.rgb = palette_color(theme, 'muted')
    out1.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    reveal.append(out1)
    
    t1 = slide.shapes.add_textbox(Inches(x1), Inches(base_y + 0.1), Inches(col_w), Inches(0.4))
    set_textbox_text(t1, truncate_to(step1.get('title', ''), 40), theme, 'body_bold', color_key='text')
    reveal.append(t1)
    
    items1 = [line.strip() for line in step1.get('body', '').split('\n') if line.strip()][:3]
    for j, item_text in enumerate(items1):
        py = base_y + 0.8 + j * 0.9
        pill = add_solid_rect(slide, x1 + 0.3, py, col_w - 0.6, 0.6, theme, color_key='accent', shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        pill.adjustments[0] = 0.5
        pt = slide.shapes.add_textbox(Inches(x1 + 0.3), Inches(py + 0.1), Inches(col_w - 0.6), Inches(0.4))
        set_textbox_text(pt, truncate_to(item_text, 30), theme, 'caption', color_key='text')
        reveal.extend([pill, pt])
        
    # Arrow 2
    arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x1 + col_w + 0.1), Inches(base_y + base_h/2 - 0.1), Inches(gap - 0.2), Inches(0.2))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = palette_color(theme, 'accent')
    arr2.line.fill.background()
    reveal.append(arr2)
    
    # --- Column 3: Modular Monolith ---
    x2 = x1 + col_w + gap
    step2 = steps[2]
    out2 = add_solid_rect(slide, x2, base_y - 0.5, col_w, base_h + 1.0, theme, color_key='bg', line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    out2.fill.background()
    out2.line.color.rgb = palette_color(theme, 'muted')
    out2.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    reveal.append(out2)
    
    t2 = slide.shapes.add_textbox(Inches(x2), Inches(base_y - 0.4), Inches(col_w), Inches(0.4))
    set_textbox_text(t2, truncate_to(step2.get('title', ''), 40), theme, 'body_bold', color_key='text')
    reveal.append(t2)
    
    items2 = [line.strip() for line in step2.get('body', '').split('\n') if line.strip()][:3]
    for j, item_text in enumerate(items2):
        gy = base_y + 0.2 + j * 1.3
        in2 = add_solid_rect(slide, x2 + 0.2, gy, col_w - 0.4, 1.1, theme, color_key='bg', line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        in2.fill.background()
        in2.line.color.rgb = palette_color(theme, 'muted')
        in2.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        
        pill = add_solid_rect(slide, x2 + 0.4, gy + 0.3, col_w - 0.8, 0.6, theme, color_key='accent', shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        pill.adjustments[0] = 0.5
        pt = slide.shapes.add_textbox(Inches(x2 + 0.4), Inches(gy + 0.4), Inches(col_w - 0.8), Inches(0.4))
        set_textbox_text(pt, truncate_to(item_text, 30), theme, 'caption', color_key='text')
        reveal.extend([in2, pill, pt])
        
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
