from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Top Right Accent Rectangle
    accent = add_solid_rect(
        slide, 11.0, 0.3, 1.5, 2.7, theme,
        color_key="accent", line=False
    )
    reveal.append(accent)
    
    # Left Text Content
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(6.0), Inches(1.5))
        set_textbox_text(head_box, truncate_to(headline, 60), theme, "title_xl", color_key="text")
        reveal.append(head_box)
        
    subhead = get_slot(slots, "subhead")
    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(6.0), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subhead, 100), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)
        
    # Hero Image Placeholder (Bottom Right)
    img_placeholder = add_solid_rect(
        slide, 8.0, 3.0, 3.5, 3.5, theme,
        color_key="panel", line=True, shape_type=MSO_SHAPE.OVAL
    )
    reveal.append(img_placeholder)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
