from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    img_path = get_slot(slots, "hero_image")
    
    if img_path:
        try:
            pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            add_theme_entrance(slide, pic, theme, index=0)
        except Exception:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = palette_color(theme, "accent")
            bg.line.fill.background()
            add_theme_entrance(slide, bg, theme, index=0)
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = palette_color(theme, "bg")
        bg.line.fill.background()
        add_theme_entrance(slide, bg, theme, index=0)

    headline = get_slot(slots, "headline")
    if headline:
        margin = theme.get("spacing", {}).get("margin", 0.6)
        tb = slide.shapes.add_textbox(
            Inches(margin), Inches(3.0), Inches(13.333 - 2 * margin), Inches(1.5)
        )
        set_textbox_text(tb, truncate_to(headline, 80), theme, "title_xl", color_key="bg")
        add_theme_entrance(slide, tb, theme, index=1)
