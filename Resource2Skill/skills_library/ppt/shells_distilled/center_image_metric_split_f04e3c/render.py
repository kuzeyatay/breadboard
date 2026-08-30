from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Metric
    metric_text = get_slot(slots, "metric", required=True)
    metric_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(4.5), Inches(3.5))
    set_textbox_text(metric_box, truncate_to(metric_text, 10), theme, "metric_xl", color_key="accent")
    reveal.append(metric_box)
    
    # Image
    img_slot = get_slot(slots, "hero_image")
    if img_slot:
        pic = slide.shapes.add_picture(img_slot, Inches(4.5), Inches(2.5), width=Inches(3.0))
        reveal.append(pic)
        
    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(6.0), Inches(2.5), Inches(5.5), Inches(2.0))
    set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title", color_key="text")
    reveal.append(head_box)
    
    # Footnote
    footnote_text = get_slot(slots, "footnote")
    if footnote_text:
        foot_box = slide.shapes.add_textbox(Inches(6.0), Inches(4.5), Inches(5.5), Inches(0.5))
        set_textbox_text(foot_box, truncate_to(footnote_text, 50), theme, "caption", color_key="muted")
        reveal.append(foot_box)
        
    # Footer
    footer_text = get_slot(slots, "footer")
    if footer_text:
        footer_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(10.333), Inches(0.5))
        set_textbox_text(footer_box, truncate_to(footer_text, 50), theme, "body_bold", color_key="accent")
        reveal.append(footer_box)
        
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
