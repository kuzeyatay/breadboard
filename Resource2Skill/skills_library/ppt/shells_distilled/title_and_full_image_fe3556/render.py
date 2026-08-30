from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    y = margin

    # Headline
    headline = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.8))
    set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="text")
    reveal.append(head_box)
    y += 0.8

    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(0.5))
        set_textbox_text(sub_box, truncate_to(subtitle, 100), theme, "body", color_key="muted")
        reveal.append(sub_box)
        y += 0.6

    # Diagram / Image
    diagram_path = get_slot(slots, "diagram")
    img_y = y + 0.2
    img_h = 7.5 - img_y - margin
    img_w = 13.333 - 2 * margin
    img_x = margin

    if diagram_path:
        try:
            pic = slide.shapes.add_picture(diagram_path, Inches(img_x), Inches(img_y), width=Inches(img_w), height=Inches(img_h))
            reveal.append(pic)
        except Exception:
            diagram_path = None # Fallback to placeholder

    if not diagram_path:
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(img_x), Inches(img_y), Inches(img_w), Inches(img_h))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = palette_color(theme, "panel")
        placeholder.line.color.rgb = palette_color(theme, "muted")
        reveal.append(placeholder)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
