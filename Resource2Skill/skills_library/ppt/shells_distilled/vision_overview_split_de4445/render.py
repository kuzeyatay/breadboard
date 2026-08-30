from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "headline", "kind": "text", "max_chars": 60, "style": "title", "required": True},
    {"name": "subtitle", "kind": "text", "max_chars": 60, "style": "subtitle", "required": False},
    {"name": "callout_left", "kind": "text", "max_chars": 150, "style": "body", "required": False},
    {"name": "callout_right", "kind": "text", "max_chars": 150, "style": "body", "required": False},
]

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Bottom Bar
    bottom_bar = add_solid_rect(slide, 0, 7.0, 13.333, 0.5, theme, color_key="text")
    reveal.append(bottom_bar)

    bottom_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.333), Inches(6.666), Inches(0.666), Inches(0.666))
    bottom_circle.fill.solid()
    bottom_circle.fill.fore_color.rgb = palette_color(theme, "text")
    bottom_circle.line.fill.background()
    reveal.append(bottom_circle)

    # Title
    headline_text = get_slot(slots, "headline")
    if headline_text:
        head_box = slide.shapes.add_textbox(Inches(3.666), Inches(0.8), Inches(6.0), Inches(0.8))
        set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title", color_key="text")
        reveal.append(head_box)

    # Subtitle
    subtitle_text = get_slot(slots, "subtitle")
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(3.666), Inches(1.6), Inches(6.0), Inches(0.4))
        set_textbox_text(sub_box, truncate_to(subtitle_text, 60), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)

    # Decorative Line
    line = add_solid_rect(slide, 6.166, 2.1, 1.0, 0.04, theme, color_key="text")
    reveal.append(line)

    # Magnifying Glass Lens
    lens = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(2.8), Inches(3.6), Inches(3.6))
    lens.fill.solid()
    lens.fill.fore_color.rgb = palette_color(theme, "bg")
    lens.line.color.rgb = palette_color(theme, "text")
    lens.line.width = Inches(0.08)
    reveal.append(lens)

    # Magnifying Glass Handle
    handle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.85), Inches(5.25), Inches(0.7), Inches(2.5))
    handle.fill.solid()
    handle.fill.fore_color.rgb = palette_color(theme, "text")
    handle.line.fill.background()
    handle.rotation = 315
    reveal.append(handle)

    # Handle Accent
    handle_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.1), Inches(5.35), Inches(0.1), Inches(2.3))
    handle_accent.fill.solid()
    handle_accent.fill.fore_color.rgb = palette_color(theme, "bg")
    handle_accent.line.fill.background()
    handle_accent.rotation = 315
    reveal.append(handle_accent)

    # Left Callout
    callout_left_text = get_slot(slots, "callout_left")
    if callout_left_text:
        left_rect = add_solid_rect(slide, 0.5, 3.5, 3.8, 1.2, theme, color_key="panel")
        
        left_tail = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.3), Inches(4.3), Inches(0.4), Inches(0.4))
        left_tail.fill.solid()
        left_tail.fill.fore_color.rgb = palette_color(theme, "panel")
        left_tail.line.fill.background()
        left_tail.rotation = 90
        
        left_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(3.4), Inches(1.0))
        set_textbox_text(left_box, truncate_to(callout_left_text, 150), theme, "body", color_key="text")
        
        reveal.extend([left_rect, left_tail, left_box])

    # Right Callout
    callout_right_text = get_slot(slots, "callout_right")
    if callout_right_text:
        right_rect = add_solid_rect(slide, 8.8, 3.5, 4.0, 1.2, theme, color_key="text")
        
        right_tail = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(8.4), Inches(4.3), Inches(0.4), Inches(0.4))
        right_tail.fill.solid()
        right_tail.fill.fore_color.rgb = palette_color(theme, "text")
        right_tail.line.fill.background()
        right_tail.rotation = 270
        
        right_box = slide.shapes.add_textbox(Inches(9.0), Inches(3.6), Inches(3.6), Inches(1.0))
        set_textbox_text(right_box, truncate_to(callout_right_text, 150), theme, "body", color_key="bg")
        
        reveal.extend([right_rect, right_tail, right_box])

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
