from __future__ import annotations
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, "bg")
    add_solid_rect(slide, 3.0, 1.2, 10.333, 4.8, theme, "panel")
    reveal = []
    kicker = get_slot(slots, "kicker")
    if kicker:
        k_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.0), Inches(0.5))
        set_textbox_text(k_box, truncate_to(kicker, 50), theme, "subtitle", color_key="text")
        reveal.append(k_box)
    hero_image = get_slot(slots, "hero_image")
    if hero_image:
        try:
            pic = slide.shapes.add_picture(hero_image, Inches(0.5), Inches(1.0), Inches(5.0), Inches(5.0))
            reveal.append(pic)
        except Exception:
            hero_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(0.8), Inches(6.0), Inches(6.0))
            hero_circle.fill.solid()
            hero_circle.fill.fore_color.rgb = palette_color(theme, "accent")
            hero_circle.line.fill.background()
            reveal.append(hero_circle)
    else:
        hero_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(0.8), Inches(6.0), Inches(6.0))
        hero_circle.fill.solid()
        hero_circle.fill.fore_color.rgb = palette_color(theme, "accent")
        hero_circle.line.fill.background()
        reveal.append(hero_circle)
    title_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(4.8), Inches(3.0), Inches(3.0))
    title_circle.fill.solid()
    title_circle.fill.fore_color.rgb = palette_color(theme, "bg")
    title_circle.line.color.rgb = palette_color(theme, "secondary")
    title_circle.line.width = Inches(0.05)
    reveal.append(title_circle)
    main_title = get_slot(slots, "main_title")
    if main_title:
        mt_box = slide.shapes.add_textbox(Inches(3.8), Inches(5.3), Inches(2.4), Inches(2.0))
        set_textbox_text(mt_box, truncate_to(main_title, 30), theme, "title", color_key="text")
        reveal.append(mt_box)
    sec_num = get_slot(slots, "section_number")
    if sec_num:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.8), Inches(0.5), Inches(0.8))
        pill.adjustments[0] = 0.5
        pill.fill.solid()
        pill.fill.fore_color.rgb = palette_color(theme, "bg")
        pill.line.color.rgb = palette_color(theme, "accent")
        pill.line.width = Inches(0.02)
        reveal.append(pill)
        num_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.95), Inches(0.5), Inches(0.5))
        set_textbox_text(num_box, truncate_to(sec_num, 5), theme, "body_bold", color_key="text")
        reveal.append(num_box)
    sec_title = get_slot(slots, "section_title")
    if sec_title:
        st_box = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.8))
        set_textbox_text(st_box, truncate_to(sec_title, 60), theme, "title", color_key="text")
        reveal.append(st_box)
    add_hairline(slide, 5.5, 2.7, 7.833, theme)
    bullets = get_slot(slots, "bullets")
    if bullets:
        b_box = slide.shapes.add_textbox(Inches(7.2), Inches(2.9), Inches(5.5), Inches(3.0))
        bullet_text = ""
        for item in bullets:
            text = item.get("title", str(item)) if isinstance(item, dict) else str(item)
            bullet_text += f"\u2022 {text}\n"
        set_textbox_text(b_box, bullet_text.strip(), theme, "body", color_key="text")
        reveal.append(b_box)
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
