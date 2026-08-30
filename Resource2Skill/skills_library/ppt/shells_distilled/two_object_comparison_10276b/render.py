from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    margin = theme.get("spacing", {}).get("margin", 0.6)

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head_box = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2*margin), Inches(1.0))
        set_textbox_text(head_box, truncate_to(headline, 80), theme, "title", color_key="text")
        head_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(head_box)

    # Left Image
    img_left = get_slot(slots, "image_left")
    if img_left:
        try:
            pic_l = slide.shapes.add_picture(img_left, Inches(2.0), Inches(2.5), height=Inches(3.0))
            reveal.append(pic_l)
        except Exception:
            pass

    # Right Image
    img_right = get_slot(slots, "image_right")
    if img_right:
        try:
            pic_r = slide.shapes.add_picture(img_right, Inches(8.333), Inches(2.5), height=Inches(3.0))
            reveal.append(pic_r)
        except Exception:
            pass

    # Left Label
    label_l = get_slot(slots, "label_left")
    if label_l:
        box_l = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(4.0), Inches(0.5))
        set_textbox_text(box_l, truncate_to(label_l, 40), theme, "body", color_key="text")
        box_l.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(box_l)

    # Right Label
    label_r = get_slot(slots, "label_right")
    if label_r:
        box_r = slide.shapes.add_textbox(Inches(7.833), Inches(5.8), Inches(4.0), Inches(0.5))
        set_textbox_text(box_r, truncate_to(label_r, 40), theme, "body", color_key="text")
        box_r.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        reveal.append(box_r)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
