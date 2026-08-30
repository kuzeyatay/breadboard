from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin_x = 0.6
    margin_y = 0.6
    gutter = theme.get("spacing", {}).get("gutter", 0.3)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    y = margin_y
    
    # Headline
    headline = get_slot(slots, "headline", default="Gallery")
    head_box = slide.shapes.add_textbox(Inches(margin_x), Inches(y), Inches(12.1), Inches(0.8))
    set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="text")
    reveal.append(head_box)
    y += 0.8
    
    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(margin_x), Inches(y), Inches(12.1), Inches(0.5))
        set_textbox_text(sub_box, truncate_to(subtitle, 100), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)
        y += 0.6
    else:
        y += 0.2
        
    # Grid configuration
    total_w = 12.7 - margin_x
    col_w = (total_w - 2 * gutter) / 3
    col_y = [y, y, y]
    
    available_h = 7.2 - y
    net_h = available_h - gutter
    
    # Masonry proportions to mimic the staggered layout
    h1 = net_h * (3.0 / 5.2)
    h4 = net_h * (2.2 / 5.2)
    h2 = net_h * (2.2 / 5.2)
    h5 = net_h * (3.0 / 5.2)
    h3 = net_h * 0.5
    h6 = net_h * 0.5
    
    heights = [h1, h2, h3, h4, h5, h6]
    cols = [0, 1, 2, 0, 1, 2]
    
    for i in range(1, 7):
        img_slot = get_slot(slots, f"image_{i}")
        c = cols[i-1]
        h = heights[i-1]
        x = margin_x + c * (col_w + gutter)
        cy = col_y[c]
        
        drawn = False
        if img_slot:
            try:
                pic = slide.shapes.add_picture(img_slot, Inches(x), Inches(cy), width=Inches(col_w), height=Inches(h))
                reveal.append(pic)
                drawn = True
            except Exception:
                pass
                
        if not drawn and (i <= 3 or img_slot):
            rect = add_solid_rect(slide, x, cy, col_w, h, theme, color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            rect.adjustments[0] = 0.05
            reveal.append(rect)
            drawn = True
            
        if drawn:
            # Mimic the circular '+' button from the original UI
            btn_size = 0.3
            btn_x = x + col_w - btn_size - 0.15
            btn_y = cy + h - btn_size - 0.15
            btn = add_solid_rect(slide, btn_x, btn_y, btn_size, btn_size, theme, color_key="accent", shape_type=MSO_SHAPE.OVAL)
            reveal.append(btn)
            
        col_y[c] += h + gutter

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=500, index=i)
