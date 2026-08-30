from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []
    
    margin_x = 2.0
    content_w = 13.333 - (2 * margin_x)

    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(margin_x), Inches(1.2), Inches(content_w), Inches(1.0))
    set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title_xl", color_key="text")
    reveal.append(head_box)

    # Hero Image
    img_path = get_slot(slots, "hero_image")
    img_y = 2.5
    img_h = 4.2
    
    if img_path:
        try:
            pic = slide.shapes.add_picture(img_path, Inches(margin_x), Inches(img_y), width=Inches(content_w), height=Inches(img_h))
            reveal.append(pic)
        except Exception:
            # Fallback to placeholder if image fails to load
            pic = add_solid_rect(slide, margin_x, img_y, content_w, img_h, theme, color_key="muted")
            reveal.append(pic)
    else:
        # Placeholder
        pic = add_solid_rect(slide, margin_x, img_y, content_w, img_h, theme, color_key="muted")
        reveal.append(pic)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
