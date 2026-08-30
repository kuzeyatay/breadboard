from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    # Main Window Panel
    panel_w = 11.5
    panel_h = 6.5
    panel_x = (13.333 - panel_w) / 2
    panel_y = (7.5 - panel_h) / 2
    
    window = add_solid_rect(
        slide, panel_x, panel_y, panel_w, panel_h, theme, 
        color_key="panel", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, line=True
    )
    window.line.color.rgb = palette_color(theme, "accent")
    window.line.width = Inches(0.03)
    window.adjustments[0] = 0.03
    
    reveal = [window]
    
    # Kicker
    kicker = get_slot(slots, "kicker")
    if kicker:
        k_box = slide.shapes.add_textbox(
            Inches(panel_x + 0.6), Inches(panel_y + 0.4), Inches(panel_w - 4.0), Inches(0.4)
        )
        set_textbox_text(k_box, truncate_to(kicker, 100), theme, "caption", color_key="muted")
        reveal.append(k_box)
        
    # Items
    items = get_slot(slots, "items", required=True)
    if not isinstance(items, list):
        raise ValueError(f"'items' slot must be a list, got {items!r}")
        
    item_y = panel_y + 1.0
    item_w = panel_w - 4.0
    item_spacing = 1.45
    
    for i, item in enumerate(items[:3]):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
            
        # Title
        t_box = slide.shapes.add_textbox(
            Inches(panel_x + 0.6), Inches(item_y), Inches(item_w), Inches(0.4)
        )
        title_text = f"{i+1}. {item.get('title', '')}"
        set_textbox_text(t_box, truncate_to(title_text, 80), theme, "body_bold", color_key="text")
        reveal.append(t_box)
        
        # Body
        b_box = slide.shapes.add_textbox(
            Inches(panel_x + 0.8), Inches(item_y + 0.4), Inches(item_w - 0.2), Inches(1.0)
        )
        set_textbox_text(b_box, truncate_to(item.get("body", ""), 200), theme, "body", color_key="muted")
        reveal.append(b_box)
        
        item_y += item_spacing
        
    # Side Image
    side_image = get_slot(slots, "side_image")
    if side_image:
        img_w = 2.4
        img_h = 3.6
        img_x = panel_x + panel_w - img_w - 0.4
        img_y = panel_y + (panel_h - img_h) / 2
        
        # Image border/glow effect
        img_bg = add_solid_rect(
            slide, img_x - 0.08, img_y - 0.08, img_w + 0.16, img_h + 0.16, theme, 
            color_key="accent", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE
        )
        img_bg.adjustments[0] = 0.1
        reveal.append(img_bg)
        
        try:
            pic = slide.shapes.add_picture(
                side_image, Inches(img_x), Inches(img_y), width=Inches(img_w), height=Inches(img_h)
            )
            reveal.append(pic)
        except Exception:
            pass
            
    # Action Bar (Footer)
    action_text = get_slot(slots, "action_text")
    if action_text:
        action_y = panel_y + panel_h - 0.8
        action_box = add_solid_rect(
            slide, panel_x + 0.6, action_y, panel_w - 1.2, 0.5, theme,
            color_key="bg", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, line=True
        )
        action_box.line.color.rgb = palette_color(theme, "muted")
        action_box.adjustments[0] = 0.3
        reveal.append(action_box)
        
        a_txt = slide.shapes.add_textbox(
            Inches(panel_x + 0.8), Inches(action_y + 0.05), Inches(panel_w - 1.6), Inches(0.4)
        )
        set_textbox_text(a_txt, truncate_to(action_text, 50), theme, "body", color_key="muted")
        reveal.append(a_txt)
        
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
