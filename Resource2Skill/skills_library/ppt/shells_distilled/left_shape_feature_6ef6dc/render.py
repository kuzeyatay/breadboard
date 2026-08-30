from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")
    reveal = [bg]
    
    # Decorative Circle (Left)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.0), Inches(2.75), Inches(2.0), Inches(2.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = palette_color(theme, "accent")
    circle.line.color.rgb = palette_color(theme, "text")
    circle.line.width = Inches(1/72) # 1 pt outline as shown in the UI
    reveal.append(circle)
    
    # Headline
    headline_text = get_slot(slots, "headline", default="Process Step")
    head_box = slide.shapes.add_textbox(Inches(5.0), Inches(2.75), Inches(7.0), Inches(1.0))
    set_textbox_text(head_box, truncate_to(headline_text, 60), theme, "title", color_key="text")
    reveal.append(head_box)
    
    # Body
    body_text = get_slot(slots, "body")
    if body_text:
        body_box = slide.shapes.add_textbox(Inches(5.0), Inches(3.8), Inches(7.0), Inches(2.0))
        set_textbox_text(body_box, truncate_to(body_text, 200), theme, "body", color_key="muted")
        reveal.append(body_box)
        
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
