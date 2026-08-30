from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")
    reveal = []

    # Left Graphic - Horizontal bar
    bar = add_solid_rect(slide, 0, 3.45, 4.0, 0.6, theme, color_key="accent")
    reveal.append(bar)

    accent_label = get_slot(slots, "accent_label")
    if accent_label:
        lbl_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(2.8), Inches(0.5))
        set_textbox_text(lbl_box, truncate_to(accent_label, 20), theme, "body_bold", color_key="bg")
        reveal.append(lbl_box)

    # Small triangle indicator on the bar
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.5), Inches(3.6), Inches(0.2), Inches(0.3))
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = palette_color(theme, "bg")
    tri.line.fill.background()
    reveal.append(tri)

    # Left Graphic - Donut
    donut = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(2.0), Inches(1.75), Inches(4.0), Inches(4.0))
    donut.fill.solid()
    donut.fill.fore_color.rgb = palette_color(theme, "secondary")
    donut.line.fill.background()
    try:
        donut.adjustments[0] = 0.75  # Make the ring thinner
    except:
        pass
    reveal.append(donut)

    # Right Text - Headline
    headline = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(6.5), Inches(3.0), Inches(6.0), Inches(0.8))
    set_textbox_text(head_box, truncate_to(headline, 50), theme, "title", color_key="text")
    reveal.append(head_box)

    # Right Text - Bullets
    bullets = get_slot(slots, "bullets", required=True)
    if bullets:
        bullet_text = ""
        for i, item in enumerate(bullets):
            if isinstance(item, dict):
                text = item.get("title", "")
            else:
                text = str(item)
            bullet_text += f"• {text}"
            if i < len(bullets) - 1:
                bullet_text += "\n"
        
        body_box = slide.shapes.add_textbox(Inches(6.5), Inches(4.0), Inches(6.0), Inches(2.5))
        set_textbox_text(body_box, bullet_text, theme, "body", color_key="muted")
        reveal.append(body_box)

    # Staggered entrance animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
