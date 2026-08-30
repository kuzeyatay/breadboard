from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin_x = 1.5
    content_w = 10.333
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Top rule
    line1 = add_solid_rect(slide, margin_x, 1.0, content_w, 0.02, theme, color_key="text")
    reveal.append(line1)
    
    # Title
    title_text = truncate_to(get_slot(slots, "title", required=True), 60)
    title_box = slide.shapes.add_textbox(Inches(margin_x), Inches(1.15), Inches(content_w), Inches(0.6))
    set_textbox_text(title_box, title_text, theme, "title", color_key="text")
    reveal.append(title_box)
    
    # Bottom rule
    line2 = add_solid_rect(slide, margin_x, 1.85, content_w, 0.02, theme, color_key="text")
    reveal.append(line2)
    
    # Left Column: Metric
    metric_label = truncate_to(get_slot(slots, "metric_label", required=True), 30)
    label_box = slide.shapes.add_textbox(Inches(margin_x), Inches(2.8), Inches(4.0), Inches(0.8))
    set_textbox_text(label_box, metric_label, theme, "title", color_key="text")
    reveal.append(label_box)
    
    metric_value = truncate_to(get_slot(slots, "metric_value", required=True), 10)
    value_box = slide.shapes.add_textbox(Inches(margin_x), Inches(3.8), Inches(4.0), Inches(2.0))
    set_textbox_text(value_box, metric_value, theme, "metric_xl", color_key="text")
    reveal.append(value_box)
    
    # Right Column: Image
    image_path = get_slot(slots, "image", required=True)
    img_x, img_y, img_w, img_h = 5.833, 2.3, 6.0, 4.5
    
    if image_path:
        try:
            pic = slide.shapes.add_picture(image_path, Inches(img_x), Inches(img_y), width=Inches(img_w), height=Inches(img_h))
            reveal.append(pic)
        except Exception:
            rect = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="muted")
            reveal.append(rect)
    else:
        rect = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="muted")
        reveal.append(rect)
        
    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
