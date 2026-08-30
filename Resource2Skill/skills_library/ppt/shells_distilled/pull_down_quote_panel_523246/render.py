from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(1.2))
    set_textbox_text(head, truncate_to(headline_text, 80), theme, "title", color_key="text")
    reveal.append(head)

    # Quote Panel Dimensions
    panel_w = 8.5
    panel_h = 4.0
    panel_x = (13.333 - panel_w) / 2.0
    panel_y = 2.2

    # Top Bar (mimicking the roller mechanism)
    bar_w = panel_w + 0.4
    bar_x = panel_x - 0.2
    bar_h = 0.35
    bar = add_solid_rect(
        slide, bar_x, panel_y, bar_w, bar_h, theme,
        color_key="accent", line=False
    )
    reveal.append(bar)

    # Main Quote Panel (the "blind")
    panel = add_solid_rect(
        slide, panel_x, panel_y + bar_h, panel_w, panel_h - bar_h, theme,
        color_key="panel", line=False
    )
    reveal.append(panel)

    # Quote Text
    quote_text = get_slot(slots, "quote", required=True)
    # Automatically wrap in quotes if not present
    if not quote_text.startswith('"') and not quote_text.startswith('\''):
        quote_text = f'"{quote_text}"'
        
    q_box = slide.shapes.add_textbox(
        Inches(panel_x + 0.6), Inches(panel_y + bar_h + 0.6),
        Inches(panel_w - 1.2), Inches(panel_h - bar_h - 1.2)
    )
    set_textbox_text(q_box, truncate_to(quote_text, 300), theme, "subtitle", color_key="text")
    reveal.append(q_box)

    # Staggered entrance animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=150, duration_ms=600, index=i)
