from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, add_emphasis_pulse, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)

    # Central Window Shape Dimensions
    card_w = 6.0
    card_h = 5.5
    card_x = (13.333 - card_w) / 2
    card_y = 1.5
    border_w = 0.25

    # Outer Frame (Accent)
    frame = add_solid_rect(slide, card_x, card_y, card_w, card_h, theme, color_key="accent", line=False)
    reveal.append(frame)

    # Inner Panels
    panel_w = card_w - 2 * border_w
    panel_h = (card_h - 3 * border_w) / 2

    # Top Panel
    top_panel = add_solid_rect(slide, card_x + border_w, card_y + border_w, panel_w, panel_h, theme, color_key="panel", line=False)
    reveal.append(top_panel)

    # Bottom Panel
    bottom_panel = add_solid_rect(slide, card_x + border_w, card_y + 2 * border_w + panel_h, panel_w, panel_h, theme, color_key="panel", line=False)
    reveal.append(bottom_panel)

    # Top Content
    top_title = get_slot(slots, "top_title")
    if top_title:
        tt_box = slide.shapes.add_textbox(Inches(card_x + border_w + 0.3), Inches(card_y + border_w + 0.3), Inches(panel_w - 0.6), Inches(0.5))
        set_textbox_text(tt_box, truncate_to(top_title, 40), theme, "subtitle", color_key="text")
        reveal.append(tt_box)

    top_body = get_slot(slots, "top_body")
    if top_body:
        tb_box = slide.shapes.add_textbox(Inches(card_x + border_w + 0.3), Inches(card_y + border_w + 0.9), Inches(panel_w - 0.6), Inches(panel_h - 1.2))
        set_textbox_text(tb_box, truncate_to(top_body, 150), theme, "body", color_key="muted")
        reveal.append(tb_box)

    # Bottom Content
    bottom_title = get_slot(slots, "bottom_title")
    if bottom_title:
        bt_box = slide.shapes.add_textbox(Inches(card_x + border_w + 0.3), Inches(card_y + 2 * border_w + panel_h + 0.3), Inches(panel_w - 0.6), Inches(0.5))
        set_textbox_text(bt_box, truncate_to(bottom_title, 40), theme, "subtitle", color_key="text")
        reveal.append(bt_box)

    bottom_body = get_slot(slots, "bottom_body")
    if bottom_body:
        bb_box = slide.shapes.add_textbox(Inches(card_x + border_w + 0.3), Inches(card_y + 2 * border_w + panel_h + 0.9), Inches(panel_w - 0.6), Inches(panel_h - 1.2))
        set_textbox_text(bb_box, truncate_to(bottom_body, 150), theme, "body", color_key="muted")
        reveal.append(bb_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
    
    add_emphasis_pulse(slide, frame, theme, delay_ms=1500)
