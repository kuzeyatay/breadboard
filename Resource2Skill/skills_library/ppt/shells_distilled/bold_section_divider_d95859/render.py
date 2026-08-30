from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "accent")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline_text = truncate_to(get_slot(slots, "headline", required=True), 60)
    head = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(2.0))
    set_textbox_text(head, headline_text, theme, "title_xl", color_key="bg")
    reveal.append(head)

    # Decorative circles at the bottom
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), Inches(6.0), Inches(3.0), Inches(3.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = palette_color(theme, "bg")
    c1.line.fill.background()
    reveal.append(c1)

    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.0), Inches(6.5), Inches(4.0), Inches(4.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = palette_color(theme, "bg")
    c2.line.fill.background()
    reveal.append(c2)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
