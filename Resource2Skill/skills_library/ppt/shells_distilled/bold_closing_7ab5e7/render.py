from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_hairline, add_solid_rect, add_theme_entrance, add_emphasis_pulse,
    get_slot, palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "accent")
    bg.line.fill.background()
    
    headline_text = get_slot(slots, "headline")
    if not headline_text:
        headline_text = "THANKS"
    headline_text = truncate_to(headline_text, 30)
    
    # Centered text box
    tb_w = 10.0
    tb_h = 2.0
    tb_x = (13.333 - tb_w) / 2
    tb_y = (7.5 - tb_h) / 2
    
    head = slide.shapes.add_textbox(Inches(tb_x), Inches(tb_y), Inches(tb_w), Inches(tb_h))
    set_textbox_text(head, headline_text, theme, "title_xl", color_key="bg")
    
    # Center align the text
    for paragraph in head.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        
    add_theme_entrance(slide, head, theme, delay_ms=0, index=0)
    add_emphasis_pulse(slide, head, theme, delay_ms=2000)
