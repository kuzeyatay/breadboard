from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

SLOTS = [
    {"name": "brand", "kind": "text", "style": "subtitle", "max_chars": 40, "required": False},
    {"name": "media", "kind": "image", "aspect": "16:9", "required": True},
    {"name": "caption", "kind": "text", "style": "body", "max_chars": 100, "required": False},
]

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Brand (Top Left)
    brand_text = get_slot(slots, "brand")
    if brand_text:
        brand_box = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(4), Inches(0.5))
        set_textbox_text(brand_box, truncate_to(brand_text, 40), theme, "subtitle", color_key="text")
        reveal.append(brand_box)

    # Media Placeholder (Center)
    # Simulating the image area with a panel since no image helper is provided
    img_w = 11.5
    img_h = 5.5
    img_x = (13.333 - img_w) / 2
    img_y = (7.5 - img_h) / 2 + 0.2

    img_rect = add_solid_rect(slide, img_x, img_y, img_w, img_h, theme, color_key="panel")
    reveal.append(img_rect)

    # Caption (Bottom Overlay)
    caption_text = get_slot(slots, "caption")
    if caption_text:
        cap_w = 8.0
        cap_h = 0.6
        cap_x = img_x + 0.5
        cap_y = img_y + img_h - cap_h - 0.5

        # Caption background (dark box)
        cap_bg = add_solid_rect(slide, cap_x, cap_y, cap_w, cap_h, theme, color_key="text")
        reveal.append(cap_bg)

        # Caption text (light text)
        cap_box = slide.shapes.add_textbox(Inches(cap_x + 0.2), Inches(cap_y + 0.1), Inches(cap_w - 0.4), Inches(cap_h - 0.2))
        set_textbox_text(cap_box, truncate_to(caption_text, 100), theme, "body", color_key="bg")
        reveal.append(cap_box)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
