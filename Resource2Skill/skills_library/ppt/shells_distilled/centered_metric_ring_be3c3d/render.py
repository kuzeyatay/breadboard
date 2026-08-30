from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline_text = get_slot(slots, "headline", required=True)
    head_box = slide.shapes.add_textbox(Inches(1.666), Inches(1.2), Inches(10), Inches(0.8))
    set_textbox_text(head_box, truncate_to(headline_text, 50), theme, "title", color_key="text")
    reveal.append(head_box)

    # Subtitle
    subtitle_text = get_slot(slots, "subtitle")
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(1.666), Inches(2.0), Inches(10), Inches(0.5))
        set_textbox_text(sub_box, truncate_to(subtitle_text, 80), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)

    # Donut Chart (Decorative Ring)
    donut_x, donut_y, donut_size = 4.916, 3.0, 3.5
    donut = add_solid_rect(
        slide, donut_x, donut_y, donut_size, donut_size, theme,
        color_key="accent", line=False, shape_type=MSO_SHAPE.DONUT
    )
    # Make the donut hole larger for a thinner ring aesthetic
    try:
        donut.adjustments[0] = 0.8
    except Exception:
        pass
    reveal.append(donut)

    # Metric Text (Centered inside the ring)
    metric_text = get_slot(slots, "metric", required=True)
    metric_box = slide.shapes.add_textbox(Inches(donut_x), Inches(donut_y + 1.25), Inches(donut_size), Inches(1.0))
    set_textbox_text(metric_box, truncate_to(metric_text, 10), theme, "metric_xl", color_key="text")
    reveal.append(metric_box)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
