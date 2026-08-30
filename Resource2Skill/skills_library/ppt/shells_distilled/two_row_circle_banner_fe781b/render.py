from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to


def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")

    # Top decorative header
    header_main = add_solid_rect(slide, 0, 0, 13.333, 1.2, theme, color_key="primary")
    header_line = add_solid_rect(slide, 0, 1.25, 13.333, 0.05, theme, color_key="primary")

    reveal = [header_main, header_line]

    rows = get_slot(slots, "rows", required=True)
    if not isinstance(rows, list):
        rows = []
    
    # Ensure we only process up to 2 rows to match the layout
    rows = rows[:2]

    start_y = 2.8
    spacing_y = 2.2
    circle_d = 1.8
    banner_w = 6.5
    banner_h = 1.4
    circle_x = 2.5
    banner_x = 3.4

    # Alternate colors for the circles
    colors = ["accent", "primary"]

    for i, item in enumerate(rows):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
        
        y = start_y + i * spacing_y
        color_key = colors[i % len(colors)]

        # Banner (placed first so it's behind the circle)
        banner_y = y + (circle_d - banner_h) / 2
        banner = add_solid_rect(
            slide, banner_x, banner_y, banner_w, banner_h, theme, color_key="panel"
        )

        # Circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(circle_x), Inches(y), Inches(circle_d), Inches(circle_d)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = palette_color(theme, color_key)
        circle.line.color.rgb = palette_color(theme, "muted")
        circle.line.width = Inches(0.08)

        # Circle Text (Title)
        title_box = slide.shapes.add_textbox(
            Inches(circle_x + 0.1), Inches(y + 0.4), Inches(circle_d - 0.2), Inches(1.0)
        )
        set_textbox_text(
            title_box, truncate_to(item.get("title", "Label"), 30), theme, "body_bold", color_key="bg"
        )

        # Banner Text (Body/Bullets)
        body_box = slide.shapes.add_textbox(
            Inches(banner_x + 0.8), Inches(banner_y + 0.2), Inches(banner_w - 1.0), Inches(banner_h - 0.4)
        )
        set_textbox_text(
            body_box, truncate_to(item.get("body", ""), 200), theme, "body", color_key="text"
        )

        reveal.extend([banner, circle, title_box, body_box])

    # Staggered entrance animation
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
