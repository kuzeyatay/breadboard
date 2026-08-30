from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = add_solid_rect(slide, 0, 0, 13.333, 7.5, theme, color_key="bg")
    reveal = []

    # Top Title & Subtitle
    headline = get_slot(slots, "headline", "Our Services & Solutions")
    head_box = slide.shapes.add_textbox(Inches(2.0), Inches(0.6), Inches(9.33), Inches(0.8))
    set_textbox_text(head_box, truncate_to(headline, 60), theme, "title", color_key="text")
    reveal.append(head_box)

    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(2.0), Inches(1.4), Inches(9.33), Inches(0.4))
        set_textbox_text(sub_box, truncate_to(subtitle, 40), theme, "subtitle", color_key="muted")
        reveal.append(sub_box)
        
        line = add_solid_rect(slide, 6.16, 1.9, 1.0, 0.04, theme, color_key="primary")
        reveal.append(line)

    # Left Side Content
    side_title = get_slot(slots, "side_title", "CONTACT\nINFORMATION")
    side_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(3.5), Inches(1.2))
    set_textbox_text(side_title_box, truncate_to(side_title, 40), theme, "title", color_key="text")
    reveal.append(side_title_box)

    side_body = get_slot(slots, "side_body")
    if side_body:
        side_body_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(3.0), Inches(2.0))
        set_textbox_text(side_body_box, truncate_to(side_body, 200), theme, "body", color_key="muted")
        reveal.append(side_body_box)

    # Scale Graphic
    # Base and Bottom Bar
    base = add_solid_rect(slide, 8.0, 6.6, 1.0, 1.0, theme, color_key="primary", shape_type=MSO_SHAPE.OVAL)
    bottom_bar = add_solid_rect(slide, 0, 7.1, 13.333, 0.4, theme, color_key="primary")
    reveal.extend([base, bottom_bar])

    # Center Post
    post = add_solid_rect(slide, 8.4, 2.8, 0.2, 4.3, theme, color_key="primary")
    reveal.append(post)

    # Beam
    beam = add_solid_rect(slide, 5.65, 2.975, 5.7, 0.15, theme, color_key="primary")
    beam.rotation = 350 # -10 degrees
    reveal.append(beam)

    # Pivot
    pivot_outer = add_solid_rect(slide, 8.25, 2.8, 0.5, 0.5, theme, color_key="primary", shape_type=MSO_SHAPE.OVAL)
    pivot_inner = add_solid_rect(slide, 8.35, 2.9, 0.3, 0.3, theme, color_key="bg", shape_type=MSO_SHAPE.OVAL)
    reveal.extend([pivot_outer, pivot_inner])

    # Left Pan & Strings
    tri_left = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.9), Inches(3.54), Inches(1.6), Inches(1.6))
    tri_left.fill.background()
    tri_left.line.color.rgb = palette_color(theme, "primary")
    tri_left.line.width = Inches(0.04)
    
    pan_left = add_solid_rect(slide, 5.15, 4.59, 1.1, 2.2, theme, color_key="primary", shape_type=MSO_SHAPE.FLOWCHART_DELAY)
    pan_left.rotation = 90
    
    left_item = get_slot(slots, "left_item", "CONTACT\nINFORMATION")
    text_left = slide.shapes.add_textbox(Inches(4.7), Inches(5.3), Inches(2.0), Inches(0.8))
    set_textbox_text(text_left, truncate_to(left_item, 30), theme, "caption", color_key="bg")
    
    reveal.extend([tri_left, pan_left, text_left])

    # Right Pan & Strings
    tri_right = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(2.56), Inches(1.6), Inches(1.6))
    tri_right.fill.background()
    tri_right.line.color.rgb = palette_color(theme, "primary")
    tri_right.line.width = Inches(0.04)
    
    pan_right = add_solid_rect(slide, 10.75, 3.61, 1.1, 2.2, theme, color_key="primary", shape_type=MSO_SHAPE.FLOWCHART_DELAY)
    pan_right.rotation = 90
    
    right_item = get_slot(slots, "right_item", "CONTACT\nINFORMATION")
    text_right = slide.shapes.add_textbox(Inches(10.3), Inches(4.3), Inches(2.0), Inches(0.8))
    set_textbox_text(text_right, truncate_to(right_item, 30), theme, "caption", color_key="bg")
    
    reveal.extend([tri_right, pan_right, text_right])

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
