from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    headline_text = get_slot(slots, "headline")
    if headline_text:
        head_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.8))
        set_textbox_text(head_box, truncate_to(headline_text, 80), theme, "title", color_key="text")
        reveal.append(head_box)
        start_y = 1.8
        body_h = 4.5
    else:
        start_y = 1.2
        body_h = 5.1
        
    center_x = 13.333 / 2
    col_w = 4.0
    header_h = 0.6
    
    left_x = center_x - col_w
    right_x = center_x
    
    # Left Column
    left_title_text = get_slot(slots, "left_title", required=True)
    left_body_text = get_slot(slots, "left_body")
    
    left_header_rect = add_solid_rect(slide, left_x, start_y, col_w, header_h, theme, color_key="panel", line=True)
    left_header_box = slide.shapes.add_textbox(Inches(left_x), Inches(start_y + 0.1), Inches(col_w), Inches(header_h))
    set_textbox_text(left_header_box, truncate_to(left_title_text, 40), theme, "body_bold", color_key="text")
    if left_header_box.text_frame.paragraphs:
        left_header_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.extend([left_header_rect, left_header_box])
    
    left_body_rect = add_solid_rect(slide, left_x, start_y + header_h, col_w, body_h, theme, color_key="bg", line=True)
    reveal.append(left_body_rect)
    if left_body_text:
        left_body_box = slide.shapes.add_textbox(Inches(left_x + 0.2), Inches(start_y + header_h + 0.2), Inches(col_w - 0.4), Inches(body_h - 0.4))
        set_textbox_text(left_body_box, truncate_to(left_body_text, 300), theme, "body", color_key="text")
        reveal.append(left_body_box)
        
    # Right Column
    right_title_text = get_slot(slots, "right_title", required=True)
    right_body_text = get_slot(slots, "right_body")
    
    right_header_rect = add_solid_rect(slide, right_x, start_y, col_w, header_h, theme, color_key="panel", line=True)
    right_header_box = slide.shapes.add_textbox(Inches(right_x), Inches(start_y + 0.1), Inches(col_w), Inches(header_h))
    set_textbox_text(right_header_box, truncate_to(right_title_text, 40), theme, "body_bold", color_key="text")
    if right_header_box.text_frame.paragraphs:
        right_header_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    reveal.extend([right_header_rect, right_header_box])
    
    right_body_rect = add_solid_rect(slide, right_x, start_y + header_h, col_w, body_h, theme, color_key="bg", line=True)
    reveal.append(right_body_rect)
    if right_body_text:
        right_body_box = slide.shapes.add_textbox(Inches(right_x + 0.2), Inches(start_y + header_h + 0.2), Inches(col_w - 0.4), Inches(body_h - 0.4))
        set_textbox_text(right_body_box, truncate_to(right_body_text, 300), theme, "body", color_key="text")
        reveal.append(right_body_box)
        
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
