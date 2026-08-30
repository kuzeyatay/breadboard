from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Headline
    headline = get_slot(slots, "headline", required=True)
    head = slide.shapes.add_textbox(Inches(margin), Inches(1.0), Inches(13.333 - 2 * margin), Inches(1.0))
    set_textbox_text(head, headline, theme, "title", color_key="text")
    reveal.append(head)
    
    # Items
    items = get_slot(slots, "items", required=True)
    if not isinstance(items, list):
        items = []
    items = items[:3]
    
    y_start = 2.6
    y_spacing = 1.5
    
    text_x = 2.5
    text_w = 3.5
    bar_x = 6.5
    bar_w = 5.0
    
    for i, item in enumerate(items):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
            
        y = y_start + i * y_spacing
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(text_x), Inches(y), Inches(text_w), Inches(0.4))
        set_textbox_text(t_box, truncate_to(item.get("title", ""), 40), theme, "title", color_key="accent")
        reveal.append(t_box)
        
        # Body
        b_box = slide.shapes.add_textbox(Inches(text_x), Inches(y + 0.4), Inches(text_w), Inches(0.8))
        set_textbox_text(b_box, truncate_to(item.get("body", ""), 120), theme, "body", color_key="muted")
        reveal.append(b_box)
        
        # Bar Track (Gray)
        track = add_solid_rect(slide, bar_x, y + 0.4, bar_w, 0.2, theme, color_key="panel", line=False)
        reveal.append(track)
        
        # Bar Fill (Colored, sitting on top)
        fill_w = bar_w * (0.3 + i * 0.25) # Staggered fill lengths for visual interest
        fill = add_solid_rect(slide, bar_x, y + 0.2, fill_w, 0.2, theme, color_key="accent", line=False)
        reveal.append(fill)
        
        # Nodes
        num_nodes = 4
        for j in range(num_nodes):
            node_x = bar_x + j * (bar_w / (num_nodes - 1))
            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(node_x - 0.06), Inches(y + 0.44), Inches(0.12), Inches(0.12))
            node.fill.solid()
            node.fill.fore_color.rgb = palette_color(theme, "bg")
            node.line.color.rgb = palette_color(theme, "muted")
            node.line.width = 12700 # 1 pt
            reveal.append(node)

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, duration_ms=400, index=i)
