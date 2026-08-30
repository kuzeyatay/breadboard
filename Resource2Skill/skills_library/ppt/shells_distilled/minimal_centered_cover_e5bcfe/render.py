from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    headline_text = get_slot(slots, "headline", required=True)
    subtitle_text = get_slot(slots, "subtitle")

    # Layout parameters
    box_width = 10.0
    box_x = (13.333 - box_width) / 2
    
    # Calculate vertical centering
    total_height = 1.5
    if subtitle_text:
        total_height += 0.8
        
    start_y = (7.5 - total_height) / 2

    reveal = []

    # Headline
    head_box = slide.shapes.add_textbox(Inches(box_x), Inches(start_y), Inches(box_width), Inches(1.5))
    set_textbox_text(head_box, truncate_to(headline_text, 80), theme, "title_xl", color_key="text")
    for paragraph in head_box.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    reveal.append(head_box)

    # Subtitle
    if subtitle_text:
        sub_y = start_y + 1.5
        sub_box = slide.shapes.add_textbox(Inches(box_x), Inches(sub_y), Inches(box_width), Inches(0.8))
        set_textbox_text(sub_box, truncate_to(subtitle_text, 120), theme, "subtitle", color_key="muted")
        for paragraph in sub_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        reveal.append(sub_box)

    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=150 * i, index=i)
