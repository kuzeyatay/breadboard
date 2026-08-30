from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get('spacing', {}).get('margin', 1.0)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, 'bg')
    bg.line.fill.background()
    
    reveal = []
    
    # Decorative diagonal background elements
    diag_main = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(-3.0), Inches(8.0), Inches(14.0))
    diag_main.rotation = 20
    diag_main.fill.solid()
    diag_main.fill.fore_color.rgb = palette_color(theme, 'panel')
    diag_main.line.fill.background()
    reveal.append(diag_main)
    
    diag_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(-3.0), Inches(1.0), Inches(14.0))
    diag_stripe.rotation = 20
    diag_stripe.fill.solid()
    diag_stripe.fill.fore_color.rgb = palette_color(theme, 'accent')
    diag_stripe.line.fill.background()
    reveal.append(diag_stripe)
    
    # Content positioning
    y = 2.5
    
    # Headline
    headline = get_slot(slots, 'headline', required=True)
    head_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(6.0), Inches(2.0))
    set_textbox_text(head_box, truncate_to(headline, 80), theme, 'title_xl', color_key='text')
    reveal.append(head_box)
    
    y += 2.0
    
    # Subhead
    subhead = get_slot(slots, 'subhead')
    if subhead:
        sub_box = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(6.0), Inches(1.0))
        set_textbox_text(sub_box, truncate_to(subhead, 60), theme, 'subtitle', color_key='muted')
        reveal.append(sub_box)
    
    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=600, index=i)
