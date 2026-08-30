from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Graphic (Window Pane)
    x = 3.0
    y = 1.5
    width = 2.5
    height = 4.5
    thickness = 0.3
    
    # Top border
    top = add_solid_rect(slide, x, y, width, thickness, theme, color_key="accent")
    # Bottom border
    bottom = add_solid_rect(slide, x, y + height - thickness, width, thickness, theme, color_key="accent")
    # Left border
    left = add_solid_rect(slide, x, y, thickness, height, theme, color_key="accent")
    # Right border
    right = add_solid_rect(slide, x + width - thickness, y, thickness, height, theme, color_key="accent")
    # Middle bar
    middle = add_solid_rect(slide, x, y + (height - thickness) / 2, width, thickness, theme, color_key="accent")
    
    reveal.extend([top, bottom, left, right, middle])
    
    # Text Content
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(6.5), Inches(2.0), Inches(6.0), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)
        
    body = get_slot(slots, "body")
    if body:
        body_box = slide.shapes.add_textbox(Inches(6.5), Inches(3.2), Inches(6.0), Inches(3.0))
        set_textbox_text(body_box, truncate_to(body, 300), theme, "body", color_key="text")
        reveal.append(body_box)
        
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)
