from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from _shell_helpers import (
    add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to,
)


def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Headline
    headline = truncate_to(get_slot(slots, "headline", required=True), 60)
    head = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(1.2))
    set_textbox_text(head, headline, theme, "title_xl", color_key="text")
    reveal.append(head)

    # Bullets
    bullets = get_slot(slots, "bullets", required=True)
    if not isinstance(bullets, list):
        bullets = []

    y = 3.5
    for i, item in enumerate(bullets[:5]):  # Max 5 bullets to fit vertically
        if not isinstance(item, dict):
            item = {"title": str(item)}

        text = truncate_to(item.get("title", ""), 80)
        if not text:
            continue

        # Simulate bullet point with text character for consistent alignment
        bullet_text = f"●   {text}"
        
        # Positioned to create a centered block of left-aligned text
        b_box = slide.shapes.add_textbox(Inches(3.5), Inches(y), Inches(7.0), Inches(0.8))
        set_textbox_text(b_box, bullet_text, theme, "subtitle", color_key="text")
        reveal.append(b_box)
        
        y += 0.8

    # Staggered entrance
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, duration_ms=500, index=i)
