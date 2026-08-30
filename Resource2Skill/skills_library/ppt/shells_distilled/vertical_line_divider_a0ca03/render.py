from __future__ import annotations
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Vertical line
    line = add_solid_rect(slide, 6.6, 2.0, 0.05, 3.5, theme, color_key="accent")
    reveal.append(line)

    # Headline
    headline = get_slot(slots, "headline", required=True)
    head = slide.shapes.add_textbox(Inches(7.0), Inches(2.5), Inches(5.5), Inches(1.5))
    set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
    reveal.append(head)

    # Subtitle
    subtitle = get_slot(slots, "subtitle")
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(7.0), Inches(4.0), Inches(5.5), Inches(1.5))
        set_textbox_text(sub, truncate_to(subtitle, 120), theme, "body", color_key="muted")
        reveal.append(sub)

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
