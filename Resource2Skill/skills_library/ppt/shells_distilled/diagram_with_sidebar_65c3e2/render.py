from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, add_emphasis_pulse, get_slot,
    palette_color, set_textbox_text, truncate_to
)

SLOTS = [
    {"name": "headline", "kind": "text", "max_chars": 80, "style": "title", "required": True},
    {"name": "subtitle", "kind": "text", "max_chars": 100, "style": "subtitle", "required": False},
    {"name": "details", "kind": "text", "max_chars": 400, "style": "caption", "required": False},
    {"name": "diagram", "kind": "image", "required": True},
    {"name": "callout", "kind": "text", "max_chars": 120, "style": "body_bold", "required": False},
]

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.5)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    reveal = []

    # Left Sidebar Area
    left_x = margin
    left_w = 3.5
    y = margin

    headline_text = get_slot(slots, "headline")
    if headline_text:
        head = slide.shapes.add_textbox(Inches(left_x), Inches(y), Inches(left_w), Inches(1.0))
        set_textbox_text(head, truncate_to(headline_text, 80), theme, "title", color_key="text")
        reveal.append(head)
        y += 1.1

    subtitle_text = get_slot(slots, "subtitle")
    if subtitle_text:
        sub = slide.shapes.add_textbox(Inches(left_x), Inches(y), Inches(left_w), Inches(0.6))
        set_textbox_text(sub, truncate_to(subtitle_text, 100), theme, "subtitle", color_key="muted")
        reveal.append(sub)
        y += 0.7

    details_text = get_slot(slots, "details")
    if details_text:
        det = slide.shapes.add_textbox(Inches(left_x), Inches(y), Inches(left_w), Inches(4.0))
        set_textbox_text(det, truncate_to(details_text, 400), theme, "caption", color_key="text")
        reveal.append(det)

    # Right Diagram Area
    right_x = left_x + left_w + 0.5
    right_w = 13.333 - right_x - margin
    diag_y = margin
    diag_h = 7.5 - 2 * margin

    callout_text = get_slot(slots, "callout")
    callout_box = None
    if callout_text:
        diag_h -= 0.8 # Make room for callout at the bottom
        callout_y = diag_y + diag_h + 0.2
        callout_box = add_solid_rect(slide, right_x, callout_y, right_w, 0.6, theme, color_key="bg", line=True)
        callout_txt = slide.shapes.add_textbox(Inches(right_x + 0.1), Inches(callout_y + 0.1), Inches(right_w - 0.2), Inches(0.4))
        set_textbox_text(callout_txt, truncate_to(callout_text, 120), theme, "body_bold", color_key="accent")
        reveal.extend([callout_box, callout_txt])

    # Diagram Placeholder (since we don't have the actual image bytes in the shell)
    diag_placeholder = add_solid_rect(slide, right_x, diag_y, right_w, diag_h, theme, color_key="panel", line=True)
    diag_label = slide.shapes.add_textbox(Inches(right_x), Inches(diag_y + diag_h/2 - 0.25), Inches(right_w), Inches(0.5))
    set_textbox_text(diag_label, "[Diagram Image Placeholder]", theme, "body", color_key="muted")
    reveal.extend([diag_placeholder, diag_label])

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)

    # Emphasis on callout if present
    if callout_box:
        add_emphasis_pulse(slide, callout_box, theme, delay_ms=1500)
