from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, 'bg')
    bg.line.fill.background()

    # Headline (Centered)
    headline_text = get_slot(slots, 'headline', default='EVOLUTION')
    head_box = slide.shapes.add_textbox(Inches(3.66), Inches(0.6), Inches(6.0), Inches(1.0))
    set_textbox_text(head_box, truncate_to(headline_text, 60), theme, 'title', color_key='text')
    reveal.append(head_box)

    # Divider Line
    line = add_solid_rect(slide, 4.66, 1.6, 4.0, 0.04, theme, color_key='accent')
    reveal.append(line)

    # Left Column: Subtitle
    subtitle_text = get_slot(slots, 'subtitle')
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(4.5), Inches(0.6))
        set_textbox_text(sub_box, truncate_to(subtitle_text, 60), theme, 'subtitle', color_key='text')
        reveal.append(sub_box)

    # Left Column: Body
    body_text = get_slot(slots, 'body')
    if body_text:
        body_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(4.5), Inches(3.5))
        set_textbox_text(body_box, truncate_to(body_text, 600), theme, 'body', color_key='text')
        reveal.append(body_box)

    # Right Column: Mock Chart (if chart slot is present or as a placeholder)
    chart_x = 7.5
    chart_y = 2.5
    chart_w = 4.5
    chart_h = 3.0

    # Gridlines
    for i in range(5):
        gy = chart_y + i * (chart_h / 4)
        gridline = add_solid_rect(slide, chart_x, gy, chart_w, 0.01, theme, color_key='muted')
        reveal.append(gridline)

    # Mock Bars
    bar_data = [(0.8, 1.2), (1.5, 2.0), (2.2, 3.0), (2.8, 4.0)]
    group_w = chart_w / 5
    bar_w = group_w * 0.4

    for i, (h1, h2) in enumerate(bar_data):
        bx = chart_x + 0.2 + i * group_w
        sh1 = (h1 / 4.0) * chart_h
        sh2 = (h2 / 4.0) * chart_h

        bar1 = add_solid_rect(slide, bx, chart_y + chart_h - sh1, bar_w, sh1, theme, color_key='accent2')
        bar2 = add_solid_rect(slide, bx + bar_w, chart_y + chart_h - sh2, bar_w, sh2, theme, color_key='accent')
        reveal.extend([bar1, bar2])

    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, index=i)
