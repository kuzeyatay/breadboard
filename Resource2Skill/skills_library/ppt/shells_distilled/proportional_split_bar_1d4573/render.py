from __future__ import annotations

import re
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

SLOTS = [
    {"name": "headline", "kind": "text", "style": "title", "required": False},
    {"name": "total_label", "kind": "text", "style": "body", "required": False},
    {"name": "segments", "kind": "bullet_list", "required": True}
]

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    
    reveal = []
    
    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2*margin), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)

    segments = get_slot(slots, "segments", required=True)
    if not isinstance(segments, list) or len(segments) == 0:
        segments = [{"title": "100%"}]
        
    # Parse proportions from titles (e.g., "60%" -> 60.0)
    proportions = []
    for item in segments:
        title = item.get("title", "")
        match = re.search(r"(\d+(\.\d+)?)", title)
        if match:
            proportions.append(float(match.group(1)))
        else:
            proportions.append(1.0)
            
    total_prop = sum(proportions)
    if total_prop == 0:
        proportions = [1.0] * len(segments)
        total_prop = sum(proportions)
        
    normalized = [p / total_prop for p in proportions]
    
    # Layout metrics
    bar_w = 10.0
    bar_h = 1.2
    x_start = (13.333 - bar_w) / 2
    y_start = (7.5 - bar_h) / 2
    
    colors = ["primary", "secondary", "accent", "muted", "panel"]
    
    current_x = x_start
    
    for i, (seg, prop) in enumerate(zip(segments, normalized)):
        seg_w = bar_w * prop
        color_key = colors[i % len(colors)]
        
        # Bar segment
        rect = add_solid_rect(
            slide, current_x, y_start, seg_w, bar_h, theme,
            color_key=color_key, line=True
        )
        rect.line.color.rgb = palette_color(theme, "bg")
        rect.line.width = Inches(0.02)
        reveal.append(rect)
        
        # Top label
        title_text = seg.get("title", "")
        if title_text:
            lbl = slide.shapes.add_textbox(
                Inches(current_x), Inches(y_start - 0.7), Inches(seg_w), Inches(0.6)
            )
            # Match label color to segment color for visual grouping
            set_textbox_text(lbl, title_text, theme, "title", color_key=color_key)
            for paragraph in lbl.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
            reveal.append(lbl)
            
        current_x += seg_w
        
    # Bracket
    total_label = get_slot(slots, "total_label")
    
    bracket_y1 = y_start + bar_h + 0.15
    bracket_y2 = bracket_y1 + 0.2
    
    # Left tick
    lt = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x_start), Inches(bracket_y1), Inches(x_start), Inches(bracket_y2))
    lt.line.color.rgb = palette_color(theme, "muted")
    reveal.append(lt)
    
    # Right tick
    rt = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x_start + bar_w), Inches(bracket_y1), Inches(x_start + bar_w), Inches(bracket_y2))
    rt.line.color.rgb = palette_color(theme, "muted")
    reveal.append(rt)
    
    # Bottom line
    bl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x_start), Inches(bracket_y2), Inches(x_start + bar_w), Inches(bracket_y2))
    bl.line.color.rgb = palette_color(theme, "muted")
    reveal.append(bl)
    
    if total_label:
        t_lbl = slide.shapes.add_textbox(
            Inches(x_start), Inches(bracket_y2 + 0.1), Inches(bar_w), Inches(0.5)
        )
        set_textbox_text(t_lbl, total_label, theme, "body", color_key="text")
        for paragraph in t_lbl.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        reveal.append(t_lbl)
        
    # Entrances
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50*i, index=i)
