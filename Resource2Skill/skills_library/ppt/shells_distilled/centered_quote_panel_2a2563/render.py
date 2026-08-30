from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import (
    add_solid_rect, add_theme_entrance, get_slot,
    palette_color, set_textbox_text, truncate_to
)

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    reveal = []

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()

    # Headline
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(margin), Inches(13.333 - 2 * margin), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)

    # Quote Panel
    panel_w = 9.0
    panel_h = 4.5
    panel_x = (13.333 - panel_w) / 2
    panel_y = (7.5 - panel_h) / 2 + 0.4  # Centered, slightly offset down to balance title

    panel = add_solid_rect(
        slide, panel_x, panel_y, panel_w, panel_h, theme,
        color_key="panel", line=False
    )
    reveal.append(panel)

    # Quote Text
    quote = get_slot(slots, "quote", required=True)
    quote_box = slide.shapes.add_textbox(
        Inches(panel_x + 0.8), Inches(panel_y + 0.8),
        Inches(panel_w - 1.6), Inches(panel_h - 2.0)
    )
    
    # Ensure quote marks
    clean_quote = truncate_to(quote, 250)
    formatted_quote = f'"{clean_quote}"' if not clean_quote.startswith('"') else clean_quote
    
    set_textbox_text(quote_box, formatted_quote, theme, "subtitle", color_key="text")
    reveal.append(quote_box)

    # Author Text
    author = get_slot(slots, "author")
    if author:
        author_box = slide.shapes.add_textbox(
            Inches(panel_x + 0.8), Inches(panel_y + panel_h - 1.2),
            Inches(panel_w - 1.6), Inches(0.6)
        )
        set_textbox_text(author_box, f"— {truncate_to(author, 60)}", theme, "body_bold", color_key="accent")
        reveal.append(author_box)

    # Animations
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=100, index=i)
