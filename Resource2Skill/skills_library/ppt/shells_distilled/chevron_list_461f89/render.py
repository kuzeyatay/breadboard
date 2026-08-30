from __future__ import annotations
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from _shell_helpers import add_hairline, add_solid_rect, add_theme_entrance, get_slot, palette_color, set_textbox_text, truncate_to

def render(slide, slots: dict, theme: dict) -> None:
    margin = theme.get("spacing", {}).get("margin", 0.6)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette_color(theme, "bg")
    bg.line.fill.background()
    reveal = []
    y = margin
    headline = get_slot(slots, "headline")
    if headline:
        head = slide.shapes.add_textbox(Inches(margin), Inches(y), Inches(13.333 - 2 * margin), Inches(1.0))
        set_textbox_text(head, truncate_to(headline, 80), theme, "title", color_key="text")
        reveal.append(head)
        y += 1.0
        add_hairline(slide, margin, y, 2.0, theme)
        y += 0.4
    else:
        y += 0.4
    items = get_slot(slots, "items", required=True)
    if not isinstance(items, list):
        items = []
    items = items[:4]
    while len(items) < 4:
        items.append({"title": f"Item {len(items)+1}", "body": "Add description here."})
    available_h = 7.5 - margin - y
    spacing = 0.3
    row_h = (available_h - 3 * spacing) / 4
    row_h = min(row_h, 1.3)
    chevron_w = 1.6
    rect_w = 9.0
    overlap = 0.5
    total_w = chevron_w + rect_w - overlap
    start_x = (13.333 - total_w) / 2
    for i, item in enumerate(items):
        if not isinstance(item, dict):
            item = {"title": str(item), "body": ""}
        row_y = y + i * (row_h + spacing)
        rect_x = start_x + chevron_w - overlap
        rect = add_solid_rect(slide, rect_x, row_y, rect_w, row_h, theme, color_key="panel", line=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        chev = add_solid_rect(slide, start_x, row_y, chevron_w, row_h, theme, color_key="accent", line=False, shape_type=MSO_SHAPE.CHEVRON)
        num_box = slide.shapes.add_textbox(Inches(start_x), Inches(row_y + (row_h - 0.6) / 2), Inches(chevron_w - overlap), Inches(0.6))
        set_textbox_text(num_box, f"0{i+1}", theme, "title", color_key="bg")
        t_box = slide.shapes.add_textbox(Inches(rect_x + overlap + 0.3), Inches(row_y + 0.1), Inches(rect_w - overlap - 0.6), Inches(0.35))
        set_textbox_text(t_box, truncate_to(item.get("title", ""), 60), theme, "body_bold", color_key="text")
        b_box = slide.shapes.add_textbox(Inches(rect_x + overlap + 0.3), Inches(row_y + 0.45), Inches(rect_w - overlap - 0.6), Inches(row_h - 0.55))
        set_textbox_text(b_box, truncate_to(item.get("body", ""), 150), theme, "body", color_key="muted")
        reveal.extend([rect, chev, num_box, t_box, b_box])
    for i, shape in enumerate(reveal):
        add_theme_entrance(slide, shape, theme, delay_ms=50, index=i)