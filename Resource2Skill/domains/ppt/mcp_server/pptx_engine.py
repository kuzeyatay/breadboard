"""
domains/ppt/mcp_server/pptx_engine.py
Python-pptx execution engine — extracted from GUI_pptagent2.

Provides:
  - exec_skill_code(): execute a skill's create_slide() with auto-repair
  - exec_raw_code(): execute arbitrary python-pptx code
  - render_slide_to_png(): PPTX → PNG via LibreOffice
  - render_slide_to_b64(): PPTX → base64 PNG
  - load_skill_library(): load index.json and metadata.json
  - get_skill_code(): extract Python code from a skill's analysis Markdown
"""

from __future__ import annotations

import base64
import inspect
import json
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path


# ---------------------------------------------------------------------------
# Preamble executed before every skill code
# ---------------------------------------------------------------------------

_SKILL_EXEC_PREAMBLE = r"""from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# --- Robust image download with fallback ---
import urllib.request as _orig_urllib_request
import os as _os, tempfile as _tempfile

_FALLBACK_IMAGES = {
    "technology": "https://images.unsplash.com/photo-1518770660439-4636190af475?w=1920&h=1080&fit=crop",
    "business":   "https://images.unsplash.com/photo-1486406146926-c627a92ad1ab?w=1920&h=1080&fit=crop",
    "city":       "https://images.unsplash.com/photo-1477959858617-67f85cf4f1df?w=1920&h=1080&fit=crop",
    "nature":     "https://images.unsplash.com/photo-1441974231531-c6227db76b6e?w=1920&h=1080&fit=crop",
    "team":       "https://images.unsplash.com/photo-1522071820081-009f0129c71c?w=1920&h=1080&fit=crop",
    "abstract":   "https://images.unsplash.com/photo-1557672172-298e090bd0f1?w=1920&h=1080&fit=crop",
    "default":    "https://images.unsplash.com/photo-1518770660439-4636190af475?w=1920&h=1080&fit=crop",
}

_real_urlopen = _orig_urllib_request.urlopen

def _robust_urlopen(url_or_req, *args, **kwargs):
    try:
        return _real_urlopen(url_or_req, *args, timeout=8, **kwargs)
    except Exception as _e:
        _url_str = url_or_req.full_url if hasattr(url_or_req, 'full_url') else str(url_or_req)
        _fb_key = "default"
        for _k in _FALLBACK_IMAGES:
            if _k in _url_str.lower():
                _fb_key = _k
                break
        _fb_url = _FALLBACK_IMAGES[_fb_key]
        if _fb_url not in _url_str:
            try:
                _fb_req = _orig_urllib_request.Request(_fb_url, headers={"User-Agent": "Mozilla/5.0"})
                return _real_urlopen(_fb_req, timeout=8)
            except Exception:
                pass
        import io as _io
        from PIL import Image as _Img, ImageDraw as _IDraw
        import random as _rnd
        _w, _h = 1920, 1080
        _img = _Img.new('RGB', (_w, _h))
        _draw = _IDraw.Draw(_img)
        for _y in range(_h):
            _r = int(15 + 25 * _y / _h)
            _g = int(20 + 35 * _y / _h)
            _b = int(40 + 50 * _y / _h)
            _draw.line([0, _y, _w, _y], fill=(_r, _g, _b))
        for _ in range(8):
            _cx = _rnd.randint(0, _w)
            _cy = _rnd.randint(0, _h)
            _rad = _rnd.randint(40, 200)
            _draw.ellipse([_cx-_rad, _cy-_rad, _cx+_rad, _cy+_rad],
                         outline=(_rnd.randint(40,100), _rnd.randint(80,180), _rnd.randint(150,255)),
                         width=2)
        _buf = _io.BytesIO()
        _img.save(_buf, format='PNG')
        _buf.seek(0)
        return _buf

_orig_urllib_request.urlopen = _robust_urlopen
"""

# ---------------------------------------------------------------------------
# Import rewrite tables
# ---------------------------------------------------------------------------

_IMPORT_REWRITES = {
    "from pptx.oxml import OxmlElement": "from pptx.oxml.xmlchemy import OxmlElement",
    "from pptx.oxml.ns import OxmlElement": "from pptx.oxml.xmlchemy import OxmlElement",
}

_IMPORT_PATCHES = {
    "from pptx.oxml import CT_": "",
    "from pptx.opc.constants import RELATIONSHIP_TYPE": "",
    "from pptx.slide import Slide": "",
    "from pptx.shapes.base import BaseShape": "",
}


# ---------------------------------------------------------------------------
# patch_skill_code
# ---------------------------------------------------------------------------

def patch_skill_code(code: str) -> str:
    """Apply known patches to skill code for python-pptx compatibility."""
    patched = code

    # --- Phase 1: _IMPORT_REWRITES (exact string replacements) ---
    for bad_import, good_import in _IMPORT_REWRITES.items():
        if bad_import in patched:
            patched = patched.replace(bad_import, good_import)

    # --- Phase 2: _IMPORT_PATCHES (comment-out entire import lines) ---
    for bad_import in _IMPORT_PATCHES:
        if bad_import in patched:
            lines = patched.split("\n")
            new_lines = []
            for line in lines:
                if bad_import in line and ("import" in line or "from" in line):
                    new_lines.append(f"# PATCHED: {line.strip()}")
                else:
                    new_lines.append(line)
            patched = "\n".join(new_lines)

    # --- Phase 3: API renames ---
    if "add_freeform_shape" in patched:
        patched = patched.replace("add_freeform_shape", "build_freeform")
    if ".get_or_add_spPr()" in patched:
        patched = patched.replace(".get_or_add_spPr()", ".spPr")

    # table._element → table._tbl
    patched = re.sub(r"\b(table|tbl)\._element\b", r"\1._tbl", patched, flags=re.IGNORECASE)

    # --- Phase 4: MSO shape type fixes ---
    patched = patched.replace("MSO_AUTO_SHAPE_TYPE.LINE", "MSO_SHAPE.RECTANGLE")
    patched = patched.replace("MSO_AUTO_SHAPE_TYPE.TEARDROP", "MSO_SHAPE.OVAL")
    patched = patched.replace("MSO_AUTO_SHAPE_TYPE._Element", "MSO_SHAPE.RECTANGLE")
    patched = patched.replace("MSO_AUTO_SHAPE_TYPE.LINE_INV", "MSO_SHAPE.RECTANGLE")
    patched = patched.replace("MSO_AUTO_SHAPE_TYPE.ROUND_2_DIAG_RECT", "MSO_SHAPE.ROUNDED_RECTANGLE")

    if "MSO_SHAPE.LINE" in patched and "MSO_LINE_DASH_STYLE" not in patched.split("MSO_SHAPE.LINE")[0].split("\n")[-1]:
        patched = re.sub(r"\bMSO_SHAPE\.LINE\b(?!_)", "MSO_SHAPE.RECTANGLE", patched)
    patched = re.sub(r"\bMSO_SHAPE\.ROUND_2_DIAG_RECT\b", "MSO_SHAPE.ROUNDED_RECTANGLE", patched)

    # add_connector enum fix: use STRAIGHT (1)
    patched = re.sub(
        r"add_connector\(\s*MSO_(?:SHAPE|AUTO_SHAPE_TYPE)\.\w+",
        "add_connector(1",
        patched,
    )

    # Read-only color.type assignment
    patched = re.sub(
        r"^(\s*)(.*\.color\.type\s*=.*)",
        r"\1pass  # PATCHED: \2",
        patched, flags=re.MULTILINE,
    )

    # Wrong enum module for MSO_LINE_DASH_STYLE
    patched = patched.replace(
        "from pptx.enum.shapes import MSO_LINE_DASH_STYLE",
        "from pptx.enum.dml import MSO_LINE_DASH_STYLE",
    )
    # Wrong qn import
    patched = patched.replace(
        "from pptx.oxml import qn",
        "from pptx.oxml.ns import qn",
    )

    # --- Phase 5: Comment-out unsupported property accesses ---
    patched = re.sub(r"^(\s*)(.*\.fill\.user_picture\b.*)", r"\1pass  # PATCHED: \2", patched, flags=re.MULTILINE)
    patched = re.sub(r"^(\s*)(.*\.background\.color\b.*)", r"\1pass  # PATCHED: \2", patched, flags=re.MULTILINE)
    patched = re.sub(r"^(\s*)(.*\.shadow\.color\b.*)", r"\1pass  # PATCHED: \2", patched, flags=re.MULTILINE)
    patched = re.sub(r"(\w+)\.line\.solid\(\)", r"\1.line.fill.solid()", patched)

    # --- Phase 6: Dead URL rewrites ---
    patched = re.sub(
        r'https?://source\.unsplash\.com/(?:random/)?(\d+)x(\d+)/?\?(.+?)(?=["\'])',
        r"https://images.unsplash.com/photo-1518770660439-4636190af475?w=\1&h=\2&fit=crop",
        patched,
    )

    # --- Phase 7: f-string XML namespace escaping ---
    lines = patched.split("\n")
    fixed_lines = []
    for line in lines:
        if ('f"' in line or "f'" in line) and ('{http' in line or '{urn' in line):
            line = re.sub(r"\{(https?://[^}]+)\}", r"{{\1}}", line)
            line = re.sub(r"\{(urn:[^}]+)\}", r"{{\1}}", line)
        fixed_lines.append(line)
    patched = "\n".join(fixed_lines)

    # Strip f-string prefix from xmlns strings entirely
    patched = re.sub(
        r'f(["\'])((?:(?!\1).)*xmlns[^"\']*)\\1',
        r"\1\2\1",
        patched,
    )

    # --- Phase 8: MSO_ANCHOR fixes ---
    patched = re.sub(r"MSO_ANCHOR\.MMIDDLE\b", "MSO_ANCHOR.MIDDLE", patched)
    patched = re.sub(r"""['"']MMIDDLE['"]""", "MSO_ANCHOR.MIDDLE", patched)
    patched = re.sub(r"""['"]middle['"]""", "MSO_ANCHOR.MIDDLE", patched)
    patched = re.sub(r"""vertical_anchor\s*=\s*['"]middle['"]""", "vertical_anchor = MSO_ANCHOR.MIDDLE", patched)

    # --- Phase 9: RGBColor fixes ---
    patched = re.sub(r"RGBColor\.from_rgb\(", "RGBColor(", patched)
    patched = re.sub(r"(RGBColor\([^)]+\))\.rgb", r"\1", patched)

    # --- Phase 10: More unsupported property accesses ---
    patched = re.sub(r"^(\s*)(.*\.image_part\b.*)", r"\1pass  # PATCHED: \2", patched, flags=re.MULTILINE)
    patched = re.sub(r"^(\s*)(.*gradient_stops\.add\b.*)", r"\1pass  # PATCHED: \2", patched, flags=re.MULTILINE)
    patched = re.sub(r"^(\s*)(.*\.gradient_stops\[.*\].*)", r"\1pass  # PATCHED: \2", patched, flags=re.MULTILINE)
    patched = re.sub(r"^(\s*)(.*freeform.*\.fill\b.*)", r"\1pass  # PATCHED: \2", patched, flags=re.MULTILINE | re.IGNORECASE)

    # --- Phase 11: xpath namespace kwarg removal ---
    if ".xpath(" in patched and "namespaces=" in patched:
        patched = re.sub(
            r"\.xpath\(([^,]+),\s*namespaces\s*=\s*[^)]+\)",
            r".xpath(\1)",
            patched,
        )

    # --- Phase 12: Remove duplicate RGBColor import (already in preamble) ---
    patched = re.sub(r"from pptx\.dml\.color import RGBColor\s*\n?", "", patched)
    patched = re.sub(
        r"from pptx\.util import (.*)RGBColor(.*)",
        lambda m: f"from pptx.util import {m.group(1)}{m.group(2)}".replace(", ,", ",").strip(", "),
        patched,
    )

    return patched


# ---------------------------------------------------------------------------
# _KNOWN_IMPORT_FIXES  – used by exec_skill_code for NameError auto-repair
# ---------------------------------------------------------------------------

_KNOWN_IMPORT_FIXES = {
    "OxmlElement": "from pptx.oxml.xmlchemy import OxmlElement",
    "qn": "from pptx.oxml.ns import qn",
    "nsmap": "from pptx.oxml.ns import nsmap",
    "parse_xml": "from pptx.oxml import parse_xml",
    "Inches": "from pptx.util import Inches",
    "Pt": "from pptx.util import Pt",
    "Emu": "from pptx.util import Emu",
    "RGBColor": "from pptx.dml.color import RGBColor",
    "MSO_SHAPE": "from pptx.enum.shapes import MSO_SHAPE",
    "MSO_SHAPE_TYPE": "from pptx.enum.shapes import MSO_SHAPE_TYPE",
    "MSO_AUTO_SHAPE_TYPE": "from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE",
    "PP_ALIGN": "from pptx.enum.text import PP_ALIGN",
    "MSO_ANCHOR": "from pptx.enum.text import MSO_ANCHOR",
    "MSO_AUTO_SIZE": "from pptx.enum.text import MSO_AUTO_SIZE",
    "MSO_LINE_DASH_STYLE": "from pptx.enum.dml import MSO_LINE_DASH_STYLE",
    "MSO_THEME_COLOR": "from pptx.enum.dml import MSO_THEME_COLOR",
}


# ---------------------------------------------------------------------------
# extract_code_from_analysis
# ---------------------------------------------------------------------------

def extract_code_from_analysis(analysis: str) -> str | None:
    """Extract the last ```python ... ``` block from a skill's analysis Markdown."""
    blocks = re.findall(r"```python\s*\n(.*?)```", analysis, re.DOTALL)
    if not blocks:
        return None
    return blocks[-1].strip()


# ---------------------------------------------------------------------------
# exec_skill_code
# ---------------------------------------------------------------------------

def exec_skill_code(
    code: str,
    output_pptx_path: str,
    params: dict | None = None,
    max_retries: int = 2,
) -> tuple[bool, int]:
    """Execute a skill's create_slide() with auto-repair on failure.

    Returns: (success, shape_count)
    """
    params = params or {}
    current_code = patch_skill_code(code)

    for attempt in range(max_retries + 1):
        tmp_ns = {"__builtins__": __builtins__}

        saved_cwd = os.getcwd()
        exec_tmpdir = tempfile.mkdtemp(prefix="ppt_skill_exec_")
        try:
            os.chdir(exec_tmpdir)

            # --- exec preamble ---
            try:
                exec(_SKILL_EXEC_PREAMBLE, tmp_ns)
            except Exception:
                # preamble exec failure is non-fatal, continue
                pass

            # --- exec skill code ---
            try:
                exec(current_code, tmp_ns)
            except ImportError as e:
                # Auto-patch bad imports by commenting them out
                bad_module = str(e)
                lines = current_code.split("\n")
                for i, line in enumerate(lines):
                    if "import" not in line:
                        continue
                    if any(w in bad_module for w in bad_module.split("'") if len(w) > 2):
                        lines[i] = f"# PATCHED: {line.strip()}"
                        break
                current_code = "\n".join(lines)
                if attempt < max_retries:
                    continue
                else:
                    return (False, 0)

            # --- locate create_slide (or similar) ---
            create_fn = tmp_ns.get("create_slide")
            if not create_fn:
                # Fallback: look for any create_* callable
                return (False, 0)

            # --- filter params through signature ---
            try:
                sig = inspect.signature(create_fn)
                valid_params = set(sig.parameters.keys()) - {"output_pptx_path"}
                has_kwargs = any(p.kind == inspect.Parameter.VAR_KEYWORD for p in sig.parameters.values())
            except (ValueError, TypeError):
                valid_params = set()
                has_kwargs = True

            filtered_params = params if has_kwargs else {k: v for k, v in params.items() if k in valid_params}

            # --- call create_slide ---
            try:
                create_fn(output_pptx_path, **filtered_params)
            except (ImportError, AttributeError, NameError) as e:
                err_str = str(e)
                if isinstance(e, NameError):
                    m = re.search(r"name '(\w+)' is not defined", err_str)
                    if m:
                        fix = _KNOWN_IMPORT_FIXES.get(m.group(1))
                        if fix:
                            current_code = fix + "\n" + current_code
                elif isinstance(e, AttributeError) and "theme_part" in err_str:
                    lines = current_code.split("\n")
                    for i, line in enumerate(lines):
                        if "theme_part" not in line:
                            continue
                        if "import" not in line:
                            lines[i] = f"pass  # PATCHED: {line.strip()}"
                            break
                    current_code = "\n".join(lines)
                if attempt < max_retries:
                    continue
                return (Path(output_pptx_path).exists(), 0)
            except Exception as e:
                if attempt < max_retries and isinstance(e, (ValueError, TypeError)):
                    for k in list(filtered_params):
                        if "color" in k.lower():
                            del filtered_params[k]
                    continue
                return (Path(output_pptx_path).exists(), 0)

            # --- verify output ---
            if not Path(output_pptx_path).exists():
                return (False, 0)

            try:
                from pptx import Presentation
                prs = Presentation(output_pptx_path)
                if len(prs.slides) > 0:
                    return (True, len(prs.slides[-1].shapes))
            except Exception:
                return (True, 0)
            return (True, 0)

        finally:
            os.chdir(saved_cwd)
            import shutil as _shutil
            _shutil.rmtree(exec_tmpdir, ignore_errors=True)

    return (False, 0)


# ---------------------------------------------------------------------------
# exec_raw_code
# ---------------------------------------------------------------------------

def exec_raw_code(code: str, output_pptx_path: str) -> tuple[bool, str]:
    """Execute raw python-pptx code (non-skill format).

    The code should create a Presentation and save it. We inject OUTPUT_PATH.

    Returns: (success, error_message)
    """
    # Fix slide.add_textbox → slide.shapes.add_textbox etc.
    for method in ("add_textbox", "add_shape", "add_picture", "add_table",
                   "add_chart", "add_group_shape", "add_connector"):
        code = re.sub(rf"\bslide\.{method}\b", f"slide.shapes.{method}", code)

    # Remove duplicate RGBColor import
    code = re.sub(r"from pptx\.dml\.color import RGBColor\s*\n?", "", code)
    # Remove hardcoded OUTPUT_PATH assignment
    code = re.sub(r'^OUTPUT_PATH\s*=\s*["\'].*["\']\s*\n?', "", code, flags=re.MULTILINE)

    # Ensure the code saves the presentation
    if "prs.save" not in code and ".save(" not in code:
        code += "\nprs.save(OUTPUT_PATH)\n"

    full_code = _SKILL_EXEC_PREAMBLE + "\n" + code
    exec_globals = {"OUTPUT_PATH": str(output_pptx_path), "__builtins__": __builtins__}

    saved_cwd = os.getcwd()
    exec_tmpdir = tempfile.mkdtemp(prefix="ppt_raw_exec_")
    try:
        os.chdir(exec_tmpdir)
        exec(full_code, exec_globals)
        if Path(output_pptx_path).exists():
            return (True, "")
        return (False, "Code executed but no PPTX file was created")
    except Exception as e:
        if Path(output_pptx_path).exists():
            return (True, f"Warning: {type(e).__name__}: {e}")
        return (False, f"{type(e).__name__}: {e}")
    finally:
        os.chdir(saved_cwd)
        import shutil as _shutil
        _shutil.rmtree(exec_tmpdir, ignore_errors=True)


# ---------------------------------------------------------------------------
# PPTX → PNG rendering
# ---------------------------------------------------------------------------

def pptx_to_png(pptx_path: str | Path, png_path: str | Path) -> bool:
    """Convert PPTX first slide to PNG via LibreOffice."""
    with tempfile.TemporaryDirectory() as tmpdir:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "png",
             "--outdir", tmpdir, str(pptx_path)],
            capture_output=True, timeout=30,
        )
        if result.returncode == 0:
            pngs = list(Path(tmpdir).glob("*.png"))
            if pngs:
                shutil.copy2(pngs[0], str(png_path))
                return True
    return False


def pptx_to_b64(pptx_path: str | Path) -> str:
    """Convert PPTX first slide to base64 PNG string."""
    with tempfile.TemporaryDirectory() as tmpdir:
        png_tmp = Path(tmpdir) / "slide.png"
        if pptx_to_png(pptx_path, png_tmp):
            return base64.b64encode(png_tmp.read_bytes()).decode()
    return ""


def pptx_to_slide_png(pptx_path: str | Path, slide_index: int, png_path: str | Path) -> bool:
    """Convert a specific slide to PNG via LibreOffice + pdftoppm."""
    with tempfile.TemporaryDirectory() as tmpdir:
        # Step 1: PPTX → PDF
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", tmpdir, str(pptx_path)],
            capture_output=True, timeout=60,
        )
        pdfs = list(Path(tmpdir).glob("*.pdf"))
        if not pdfs:
            return False

        # Step 2: PDF → PNG (specific page)
        first_page = slide_index + 1
        subprocess.run(
            ["pdftoppm", "-png", "-r", "150",
             "-f", str(first_page), "-l", str(first_page),
             str(pdfs[0]), str(Path(tmpdir) / "s")],
            capture_output=True, timeout=60,
        )
        pngs = sorted(Path(tmpdir).glob("s-*.png"))
        if pngs:
            shutil.copy2(pngs[0], str(png_path))
            return True
    return False


def pptx_slide_to_b64(pptx_path: str | Path, slide_index: int) -> str:
    """Convert a specific slide to base64 PNG string."""
    with tempfile.TemporaryDirectory() as tmpdir:
        png_tmp = Path(tmpdir) / "slide.png"
        if pptx_to_slide_png(pptx_path, slide_index, png_tmp):
            return base64.b64encode(png_tmp.read_bytes()).decode()
    return ""


# ---------------------------------------------------------------------------
# Skill library I/O
# ---------------------------------------------------------------------------

def load_skill_index(library_dir: Path) -> list[dict]:
    """Load skills from index.json."""
    index_path = library_dir / "index.json"
    if not index_path.exists():
        return []
    data = json.loads(index_path.read_text("utf-8"))
    return data.get("skills", [])


def load_skill_metadata(library_dir: Path) -> dict:
    """Load metadata.json → {skill_id: metadata_dict}."""
    meta_path = library_dir / "metadata.json"
    if not meta_path.exists():
        return {}
    data = json.loads(meta_path.read_text("utf-8"))
    return data.get("metadata", {})


def get_skill_detail(library_dir: Path, skill_id: str, index: list[dict] | None = None) -> dict | None:
    """Load full skill detail from its JSON file."""
    if not index:
        index = load_skill_index(library_dir)
    for entry in index:
        if entry["skill_id"] == skill_id:
            detail_path = library_dir / entry["detail_path"]
            if detail_path.exists():
                return json.loads(detail_path.read_text("utf-8"))
            return None
    return None
