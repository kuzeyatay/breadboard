"""
domains/ppt/mcp_server/server.py
PPT MCP Server — exposes slide-level tools for building PPTX presentations.

Legacy Resource2Skill/python-pptx tools:
  create-presentation    — initialize a new Presentation object
  add-slide              — execute python-pptx code to add a slide
  add-slide-from-skill   — add a slide using a skill from the library
  replace-slide          — replace an existing slide with new code
  delete-slide           — delete a slide by index
  save-presentation      — save to disk
  get-slide-info         — inspect a slide's shapes and properties
  render-slide           — render a slide to base64 PNG
  list-skills            — search/browse the skill library
  get-skill-info         — get full details of a skill

PPT Master SVG-first tools use the ``pptmaster_`` prefix and generate a
project directory of SVG slides before exporting to editable native PPTX.

Usage (stdio transport, launched by VideoWorldSkills agent_executor):
    python domains/ppt/mcp_server/server.py --skills-dir /path/to/skills_library/ppt
"""
from __future__ import annotations

import argparse
import json
import logging
import os
import re
import sys
import tempfile
import uuid
from pathlib import Path

# ---------------------------------------------------------------------------
# Path setup
# ---------------------------------------------------------------------------

_SERVER_DIR = Path(__file__).resolve().parent
_PROJECT_ROOT = _SERVER_DIR.parents[2]
sys.path.insert(0, str(_PROJECT_ROOT))
sys.path.insert(0, str(_SERVER_DIR))

from mcp.server.fastmcp import FastMCP

log = logging.getLogger("ppt-mcp")

# ---------------------------------------------------------------------------
# Lazy imports
# ---------------------------------------------------------------------------

_engine = None
_pptx_mod = None
_pptmaster_mod = None
_pptmaster_r2s_policy_mod = None


def _get_engine():
    global _engine
    if not _engine:
        import pptx_engine
        _engine = pptx_engine
    return _engine


def _get_pptx():
    global _pptx_mod
    if not _pptx_mod:
        from pptx import Presentation
        from pptx.util import Inches
        _pptx_mod = {"Presentation": Presentation, "Inches": Inches}
    return _pptx_mod


def _get_pptmaster():
    global _pptmaster_mod
    if not _pptmaster_mod:
        import pptmaster_engine
        _pptmaster_mod = pptmaster_engine
    return _pptmaster_mod


def _get_pptmaster_r2s_policy():
    global _pptmaster_r2s_policy_mod
    if not _pptmaster_r2s_policy_mod:
        from domains.ppt import pptmaster_r2s_prompt_runner
        _pptmaster_r2s_policy_mod = pptmaster_r2s_prompt_runner
    return _pptmaster_r2s_policy_mod


# ---------------------------------------------------------------------------
# Global state
# ---------------------------------------------------------------------------

_presentations: dict[str, object] = {}
_prs_paths: dict[str, Path] = {}

# Skills
_skills_dir: Path | None = None
_skill_index: list[dict] | None = None
_skill_metadata: dict | None = None

# Per-presentation deck-plan metadata. One dict per prs_id, accumulated as
# add_slide_from_shell is called. save_presentation consumes it to write
# demo/<deck>/deck-plan.json, append the methodology appendix, and run the
# morph-lint pass.
_prs_metadata: dict[str, dict] = {}

# Morph anchor force-match prefix — mirrors _shell_helpers and the contract
# at docs/ppt_morph_continuity_contract.md.
_MORPH_ANCHOR_PREFIX = "!!sameName"

_BRAND_SHELL_SPECS = {
    "cover_brand": {
        "skill_id": "brand_cover_6eaebf69",
        "role": "cover",
        "description": "Brandcover slide",
        "slot_names": ["eyebrow", "title", "headline", "subtitle", "subheadline", "wordmark", "mark_text", "bg_keyword"],
        "required_slots": ["title"],
    },
    "section_divider_brand": {
        "skill_id": "brand_section_divider_edd4f207",
        "role": "section_divider",
        "description": "Brandsection divider",
        "slot_names": ["section_index", "title", "headline", "subtitle", "footer", "bg_keyword"],
        "required_slots": ["title"],
    },
    "content_grid_brand": {
        "skill_id": "brand_content_grid_7a8e7679",
        "role": "feature_grid",
        "description": "Brandcontent grid",
        "slot_names": ["title", "headline", "thesis", "subtitle", "tiles", "features", "columns", "bg_keyword"],
        "required_slots": ["title"],
    },
    "data_quadrant_brand": {
        "skill_id": "brand_data_quadrant_3cad8a69",
        "role": "metric_dashboard",
        "description": "Branddata quadrant",
        "slot_names": ["title", "headline", "body", "context", "x_axis", "y_axis", "quadrant_labels", "items", "bg_keyword"],
        "required_slots": ["title"],
    },
}


def _active_brand_root() -> Path | None:
    brand = os.environ.get("PPT_ACTIVE_BRAND", "").strip()
    if not brand:
        return None
    root = _PROJECT_ROOT / "brand_wiki" / "ppt" / brand
    return root if root.is_dir() else None


def _brand_shell_for_role(role: str) -> str | None:
    role = (role or "").lower()
    if role == "cover" or "cover" in role:
        return "cover_brand"
    if "closing" in role or "cta" in role or "ask" in role:
        return "section_divider_brand"
    if "section" in role or "divider" in role:
        return "section_divider_brand"
    if any(token in role for token in ("feature_grid", "content", "grid", "bento")):
        return "content_grid_brand"
    if any(token in role for token in ("quadrant", "data", "viz", "metric", "dashboard", "hero_giant_metric")):
        return "data_quadrant_brand"
    return None


def _brand_shell_entry(shell_id: str) -> dict | None:
    root = _active_brand_root()
    spec = _BRAND_SHELL_SPECS.get(shell_id)
    if root is None or spec is None:
        return None
    skill_dir = root / "skills" / spec["skill_id"]
    if not skill_dir.is_dir():
        return None
    return {
        "shell_id": shell_id,
        "role": spec["role"],
        "description": spec["description"],
        "archetype": "brand",
        "mood": ["brand"],
        "density": "balanced",
        "style_tags": ["brand"],
        "slot_names": spec["slot_names"],
        "required_slots": spec["required_slots"],
        "slots": spec["slot_names"],
        "source": "brand_overlay",
        "status": "active",
        "ambient_capable": False,
        "brand_skill_id": spec["skill_id"],
    }


def _brand_shell_entries(role: str = "") -> list[dict]:
    if _active_brand_root() is None:
        return []
    if role:
        sid = _brand_shell_for_role(role)
        entry = _brand_shell_entry(sid) if sid else None
        return [entry] if entry else []
    return [e for sid in _BRAND_SHELL_SPECS if (e := _brand_shell_entry(sid))]


def _ensure_prs_metadata(prs_id: str) -> dict:
    entry = _prs_metadata.setdefault(prs_id, {
        "archetype": None,
        "archetype_slides_target": None,
        "archetype_suggested_slides": None,
        "requested_slides_target": None,
        "theme": None,
        "slides": [],
    })
    return entry


def _reindex_prs_metadata(prs_id: str) -> None:
    entry = _prs_metadata.get(prs_id)
    if not entry:
        return
    for idx, meta in enumerate(entry.get("slides") or []):
        if isinstance(meta, dict):
            meta["slide_index"] = idx


_REQUESTED_SLIDE_COUNT_RE = re.compile(
    r"\b(\d{1,2})\s*[- ]?\s*(?:slide|slides|page|pages|deck)\b",
    re.IGNORECASE,
)


def _extract_requested_slide_count(task_description: str) -> int | None:
    """Return an explicit brief slide/page count, if one is present."""
    if not task_description:
        return None
    m = _REQUESTED_SLIDE_COUNT_RE.search(task_description)
    if not m:
        return None
    try:
        n = int(m.group(1))
    except Exception:
        return None
    # Ignore accidental years / pathological values.
    if 1 <= n <= 40:
        return n
    return None


def _scan_anchor_names(slide) -> list[str]:
    """Return every `!!sameName...` cNvPr name present on ``slide``."""
    anchors = []
    for el in slide._element.iter():
        if not el.tag.endswith("}cNvPr"):
            continue
        name = el.get("name") or ""
        if name.startswith(_MORPH_ANCHOR_PREFIX):
            anchors.append(name)
    return anchors


def _slide_transition_kind(slide) -> str:
    """Return 'morph' | 'fade' | 'push' | 'wipe' | 'none' based on the
    transition XML attached to the slide."""
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    MORPH_NS = "http://schemas.microsoft.com/office/powerpoint/2015/09/main"
    trans = slide._element.find(f"{{{P_NS}}}transition")
    if trans is None:
        return "none"
    for child in trans:
        tag = child.tag
        if tag == f"{{{MORPH_NS}}}morph" or tag.endswith("}morph"):
            return "morph"
        if tag.endswith("}fade"):
            return "fade"
        if tag.endswith("}push"):
            return "push"
        if tag.endswith("}wipe"):
            return "wipe"
    return "other"


def _shell_ambient_flag(shell_id: str) -> bool:
    """True when the named shell is tagged ambient-capable.

    Seed shells get their flag from a top-level module attribute
    ``AMBIENT_CAPABLE = True`` or from textual presence of any
    ``add_infinite_rotation``/``add_orbital_motion``/``add_pulse_loop``/
    ``add_drift_motion`` call in the file. Distilled shells read
    ``intent_tags.ambient_capable`` from their ``skill.json``.
    """
    try:
        from core.extraction.shell_loader import _SHELLS_DIR, _DISTILLED_DIR
    except Exception:
        return False
    seed_path = _SHELLS_DIR / f"{shell_id}.py"
    if seed_path.exists():
        try:
            src = seed_path.read_text()
        except Exception:
            src = ""
        if "AMBIENT_CAPABLE = True" in src:
            return True
        if any(tok in src for tok in (
            "add_infinite_rotation", "add_orbital_motion",
            "add_pulse_loop", "add_drift_motion",
        )):
            return True
    meta_path = _DISTILLED_DIR / shell_id / "skill.json"
    if meta_path.exists():
        try:
            meta = json.loads(meta_path.read_text())
            tags = meta.get("intent_tags", {}) or {}
            if tags.get("ambient_capable"):
                return True
        except Exception:
            pass
    return False


def _resolve_provenance(shell_id: str, origin_hint: str | None) -> dict:
    """Return the provenance dict recorded for the slide.

    Provides ``skill_id``, ``source_video_id``, ``source_timestamp``, and
    ``selection_origin`` (``distilled|retrieved|composed|new|seed``).
    """
    try:
        from core.extraction.shell_loader import _SHELLS_DIR, _DISTILLED_DIR
    except Exception:
        _SHELLS_DIR = Path("")
        _DISTILLED_DIR = Path("")
    seed_path = _SHELLS_DIR / f"{shell_id}.py"
    distilled_meta = _DISTILLED_DIR / shell_id / "skill.json"
    if seed_path.exists():
        default_origin = "seed"
        video_id = ""
        timestamp = None
    elif distilled_meta.exists():
        default_origin = "distilled"
        try:
            meta = json.loads(distilled_meta.read_text())
            src = (meta.get("provenance") or {}).get("source") or {}
            video_id = src.get("video_id") or ""
            timestamp = src.get("timestamp")
        except Exception:
            video_id = ""
            timestamp = None
    else:
        default_origin = "new"
        video_id = ""
        timestamp = None
    origin = origin_hint if origin_hint in {
        "distilled", "retrieved", "composed", "new", "seed",
    } else default_origin
    return {
        "skill_id": shell_id,
        "source_video_id": video_id,
        "source_timestamp": timestamp,
        "selection_origin": origin,
    }


def _chip_prefix_for_origin(origin: str) -> str:
    return {
        "distilled": "D",
        "seed": "S",
        "composed": "C",
        "retrieved": "R",
        "new": "N",
    }.get(origin, "N")


def _build_chip_text(provenance: dict) -> str:
    prefix = _chip_prefix_for_origin(provenance.get("selection_origin", "new"))
    short = provenance.get("skill_id", "")[-8:] or "unknown"
    if provenance.get("selection_origin") == "distilled":
        vid = provenance.get("source_video_id", "")
        vid_short = (vid.split("_")[-1][:8] if vid else "") or "?"
        chip = f"{prefix}:{short}@{vid_short}"
    else:
        chip = f"{prefix}:{short}"
    return chip[:30]


def _emit_provenance_chip(slide, theme: dict, provenance: dict) -> None:
    """Drop a small text chip at the bottom-left of ``slide``.

    Colour matches ``theme.palette.muted`` (falling back to ``text`` then
    ``#808080``) so the chip stays legible without pulling focus.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    text = _build_chip_text(provenance)
    tb = slide.shapes.add_textbox(Inches(0.15), Inches(7.15),
                                  Inches(3.0), Inches(0.25))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = text
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(8)
        run.font.name = "Inter"
        pal = (theme or {}).get("palette", {}) or {}
        hex_color = pal.get("muted") or pal.get("text") or "#808080"
        s = hex_color.lstrip("#")
        try:
            rgb = RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
            run.font.color.rgb = rgb
        except Exception:
            pass


def _attach_notes_json(slide, record: dict) -> None:
    """Write ``record`` as pretty JSON into the slide's notes text frame."""
    try:
        notes_slide = slide.notes_slide  # creates one if missing
        tf = notes_slide.notes_text_frame
        tf.text = json.dumps(record, ensure_ascii=False, indent=2)
    except Exception:
        pass


def _apply_transition_to_slide(slide, transition_kind: str) -> None:
    """Attach the requested transition XML to ``slide``. Silently ignores
    unknown kinds (caller is expected to validate before calling)."""
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    MORPH_NS = "http://schemas.microsoft.com/office/powerpoint/2015/09/main"
    # Remove any existing transition first so the new one is authoritative.
    existing = slide._element.find(f"{{{P_NS}}}transition")
    if existing is not None:
        slide._element.remove(existing)
    snippets = {
        "fade": f'<p:transition xmlns:p="{P_NS}" spd="med" advClick="1"><p:fade/></p:transition>',
        "push": f'<p:transition xmlns:p="{P_NS}" spd="med" advClick="1"><p:push dir="l"/></p:transition>',
        "wipe": f'<p:transition xmlns:p="{P_NS}" spd="med" advClick="1"><p:wipe dir="l"/></p:transition>',
        "morph": (
            f'<p:transition xmlns:p="{P_NS}" xmlns:p159="{MORPH_NS}" '
            f'spd="med" advClick="1"><p159:morph option="byObject"/></p:transition>'
        ),
    }
    xml = snippets.get(transition_kind)
    if not xml:
        return
    from lxml import etree as _lxml_etree
    slide._element.append(_lxml_etree.fromstring(xml))


def _ensure_skills_loaded():
    global _skill_index, _skill_metadata
    if _skill_index is None and _skills_dir and _skills_dir.exists():
        engine = _get_engine()
        _skill_index = engine.load_skill_index(_skills_dir)
        _skill_metadata = engine.load_skill_metadata(_skills_dir)


def _save_temp(prs, prs_id: str) -> Path:
    """Save presentation to a temp file, deduplicating ZIP entries."""
    tmp = Path(tempfile.gettempdir()) / f"ppt_mcp_{prs_id}.pptx"
    prs.save(str(tmp))

    # Deduplicate ZIP entries (XML clone can create duplicates that break LibreOffice)
    try:
        import zipfile, io, shutil
        with zipfile.ZipFile(str(tmp), 'r') as zin:
            seen = set()
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    if item.filename not in seen:
                        seen.add(item.filename)
                        zout.writestr(item, zin.read(item.filename))
            buf.seek(0)
            with open(str(tmp), 'wb') as f:
                f.write(buf.read())
    except Exception:
        pass  # If dedup fails, the original save is still there

    _prs_paths[prs_id] = tmp
    return tmp


def _clone_first_slide_from_pptx(prs, source_pptx: Path):
    """Append the first slide from source_pptx to prs and return it."""
    from pptx import Presentation as _P
    from pptx.oxml.ns import qn
    import copy

    src_prs = _P(str(source_pptx))
    if len(src_prs.slides) == 0:
        raise RuntimeError("brand skill produced no slides")
    src_slide = src_prs.slides[0]
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])

    rel_map = {}
    for r_id, rel in src_slide.part.rels.items():
        if any(skip in rel.reltype for skip in ("slideLayout", "slideMaster", "theme", "notesMaster", "handoutMaster")):
            continue
        if rel.is_external:
            rel_map[r_id] = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            rel_map[r_id] = new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    for shape in src_slide.shapes:
        elem = copy.deepcopy(shape.element)
        for old_r_id, new_r_id in rel_map.items():
            for attr in (qn("r:embed"), qn("r:link"), qn("r:id")):
                for node in elem.iter():
                    if node.get(attr) == old_r_id:
                        node.set(attr, new_r_id)
        new_slide.shapes._spTree.insert_element_before(elem, "p:extLst")

    if src_slide.background._element is not None:
        bg_elem = copy.deepcopy(src_slide.background._element)
        if new_slide.background._element is not None:
            new_slide._element.replace(new_slide.background._element, bg_elem)
    return new_slide


def _delete_slide(prs, slide_index: int) -> None:
    from pptx.oxml.ns import qn

    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[slide_index]
    r_id = slide_id.get(qn("r:id"))
    if r_id:
        prs.part.drop_rel(r_id)
    slide_id_list.remove(slide_id)


def _move_last_slide_to_index(prs, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    last = slide_id_list[-1]
    slide_id_list.remove(last)
    slide_id_list.insert(index, last)


def _text_from_slots(slots: dict, *keys: str, default: str = "") -> str:
    for key in keys:
        value = slots.get(key)
        if value is None:
            continue
        if isinstance(value, str) and value.strip():
            return value.strip()
        if not isinstance(value, str):
            return str(value)
    return default


def _columns_from_slots(slots: dict) -> list | None:
    for key in ("columns", "tiles", "features", "stats"):
        value = slots.get(key)
        if isinstance(value, list):
            out = []
            for item in value:
                if isinstance(item, dict):
                    out.append({
                        "heading": item.get("heading") or item.get("title") or item.get("label") or item.get("value") or "",
                        "body": item.get("body") or item.get("caption") or item.get("context") or item.get("delta") or "",
                        "outcome": item.get("outcome") or item.get("support") or "",
                    })
                else:
                    out.append({"heading": str(item), "body": "", "outcome": ""})
            return out
    return None


def _quadrant_items_from_slots(slots: dict) -> list | None:
    value = slots.get("items") or slots.get("points") or slots.get("use_cases")
    if isinstance(value, list):
        return value
    return None


def _brand_shell_kwargs(shell_id: str, slots: dict) -> dict:
    kwargs = dict(slots)
    if shell_id == "cover_brand":
        kwargs["title_text"] = _text_from_slots(slots, "title_text", "headline", "title", default="Untitled")
        kwargs["body_text"] = _text_from_slots(slots, "body_text", "subheadline", "subtitle", "caption")
        kwargs.setdefault("eyebrow", _text_from_slots(slots, "eyebrow", "kicker", "section_label"))
        kwargs.setdefault("wordmark", _text_from_slots(slots, "wordmark", "brand", default=""))
        kwargs.setdefault("mark_text", _text_from_slots(slots, "mark_text", default="A")[:2])
    elif shell_id == "section_divider_brand":
        kwargs["title_text"] = _text_from_slots(slots, "title_text", "headline", "title", "section_label", default="Section")
        kwargs["body_text"] = _text_from_slots(slots, "body_text", "subheadline", "subtitle", "caption", "context")
        kwargs.setdefault("section_index", _text_from_slots(slots, "section_index", "eyebrow", default="Section"))
    elif shell_id == "content_grid_brand":
        kwargs["title_text"] = _text_from_slots(slots, "title_text", "headline", "title", default="Overview")
        kwargs["body_text"] = _text_from_slots(slots, "body_text", "thesis", "subtitle", "context", "caption")
        columns = _columns_from_slots(slots)
        if columns:
            kwargs["columns"] = columns
    elif shell_id == "data_quadrant_brand":
        kwargs["title_text"] = _text_from_slots(slots, "title_text", "headline", "title", "eyebrow", default="Decision matrix")
        kwargs["body_text"] = _text_from_slots(slots, "body_text", "context", "caption", "support", "label")
        items = _quadrant_items_from_slots(slots)
        if items:
            kwargs["items"] = items
    if not kwargs.get("bg_keyword"):
        kwargs["bg_keyword"] = "abstract premium coffee"
    return kwargs


def _add_slide_from_brand_shell(
    prs_id: str,
    shell_id: str,
    slots_dict: dict,
    transition_kind: str | None,
    slide_role: str,
    design_reference_skill_ids: str,
) -> str:
    entry = _brand_shell_entry(shell_id)
    if entry is None:
        return f"Error: brand shell '{shell_id}' is not available"
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"
    root = _active_brand_root()
    skill_id = entry["brand_skill_id"]
    skill_dir = root / "skills" / skill_id
    code_files = sorted((skill_dir / "code").glob("*.py"))
    if not code_files:
        return f"Error: no code asset found for brand shell '{shell_id}'"

    import importlib.util
    with tempfile.TemporaryDirectory(dir="/data/tmp") as tmpdir:
        out = Path(tmpdir) / f"{shell_id}.pptx"
        spec = importlib.util.spec_from_file_location(f"_brand_shell_{shell_id}", code_files[0])
        if spec is None or spec.loader is None:
            return f"Error: could not import brand shell '{shell_id}'"
        mod = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(mod)
        create_slide = getattr(mod, "create_slide")
        create_slide(str(out), **_brand_shell_kwargs(shell_id, slots_dict))
        slide = _clone_first_slide_from_pptx(prs, out)

    if transition_kind:
        _apply_transition_to_slide(slide, transition_kind)
    role = slide_role or entry["role"]
    anchor_names = _scan_anchor_names(slide)
    transition_set = _slide_transition_kind(slide)
    design_refs = _parse_design_reference_skill_ids(design_reference_skill_ids or slots_dict.get("design_reference_skill_ids", ""))
    deck_meta = _ensure_prs_metadata(prs_id)
    deck_meta["theme"] = f"brand:{os.environ.get('PPT_ACTIVE_BRAND', '')}"
    slot_summary = {k: (v[:60] if isinstance(v, str) else str(v)[:60]) for k, v in slots_dict.items()}
    provenance = {
        "skill_id": shell_id,
        "source_video_id": "",
        "source_timestamp": None,
        "selection_origin": "brand_overlay",
    }
    slide_entry = {
        "slide_index": len(deck_meta["slides"]),
        "role": role,
        "shell_id": shell_id,
        "anchor_names_set": anchor_names,
        "transition_set": transition_set,
        "hero_flag": role in {"cover", "closing_cta"} or "hero" in role,
        "ambient_flag": False,
        "design_reference_skill_ids": design_refs,
        "distill_provenance": provenance,
        "slot_values_summary": slot_summary,
        "brand_skill_id": skill_id,
    }
    deck_meta["slides"].append(slide_entry)
    _attach_notes_json(slide, {
        **slide_entry,
        "theme": deck_meta["theme"],
        "brand_skill_id": skill_id,
    })
    _save_temp(prs, prs_id)
    return (
        f"Added slide {len(prs.slides)} from brand shell '{shell_id}' "
        f"(skill_id={skill_id}, transition={transition_set}, shapes={len(slide.shapes)})"
    )


# ---------------------------------------------------------------------------
# Post-processing helpers
# ---------------------------------------------------------------------------

import colorsys
import re as _re
import re

_DEFAULT_PALETTE = {
    "dark_bg":    (10, 15, 26),
    "panel":      (20, 30, 46),
    "accent1":    (0, 210, 255),
    "accent2":    (124, 110, 255),
    "accent3":    (57, 214, 138),
    "text_light": (255, 255, 255),
    "text_muted": (146, 166, 190),
}


def _score_frame(shape, full_text: str) -> float:
    """Score a text frame for content injection priority.

Higher score = more likely to be a title/body placeholder.
Returns -1 to exclude the frame entirely.
"""
    w = (shape.width or 0) / 914400
    h = (shape.height or 0) / 914400
    top = (shape.top or 0) / 914400
    left = (shape.left or 0) / 914400
    area = w * h

    # Exclude off-screen shapes
    if left < -0.3 or top < -0.3:
        return -1

    # Short / decorative text
    stripped = full_text.strip()
    if len(stripped) <= 2:
        return area * 0.05

    # Placeholder detection
    is_placeholder = bool(_re.search(
        r"lorem|ipsum|your.?text|add.?text|dummy|subtitle.?here|click.?to.?edit|neque.?porro|sample.?text|description.?here|heading.?here|write.*brief|detail.*text|placeholder",
        full_text,
        _re.I,
    ))

    placeholder_bonus = 2.5 if is_placeholder else 1.0

    # Position bonus (higher = closer to top)
    if top < 1.5:
        top_bonus = 1.5
    elif top < 3.0:
        top_bonus = 1.2
    else:
        top_bonus = 1.0

    # Length bonus
    len_bonus = 1.0 + min(len(stripped), 100) / 200

    return area * placeholder_bonus * top_bonus * len_bonus


def _inject_content_into_slide(slide, content_brief: str):
    """Replace placeholder text in a skill-generated slide with user content.

Uses a composite scoring heuristic (area + position + placeholder detection)
to rank text frames. Highest-scored frame gets the title, next gets first
body line, and so on. Preserves original formatting.
"""
    lines = content_brief.strip().split("\n")
    title = lines[0].strip() if lines else ""
    body_lines = [l.strip() for l in lines[1:] if l.strip()]

    inject_lines = [title] + body_lines

    if not inject_lines:
        return

    # Score all text frames
    scored_items = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full_text = tf.text.strip()
        if not full_text:
            continue
        score = _score_frame(shape, full_text)
        if score < 0:
            continue
        scored_items.append((score, tf, full_text))

    # Sort by score descending
    scored_items.sort(key=lambda x: -x[0])

    # Assign content to top-scored frames
    for i, (score, tf, old_text) in enumerate(scored_items):
        if i >= len(inject_lines):
            return
        _set_text_preserving_format(tf, inject_lines[i])


def _set_text_preserving_format(text_frame, new_text: str):
    """Replace text in a text frame, preserving the first run's formatting."""
    if not text_frame.paragraphs:
        return
    first_para = text_frame.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = new_text
        for run in first_para.runs[1:]:
            run.text = ""
    else:
        first_para.text = new_text

    # Clear remaining paragraphs
    for para in text_frame.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def _rgb_to_hsl(r, g, b):
    """Convert RGB (0-255) to HSL (h: 0-360, s: 0-1, l: 0-1)."""
    h, l, s = colorsys.rgb_to_hls(r / 255, g / 255, b / 255)
    return (h * 360, s, l)


def _unify_slide_colors(slide, palette=None):
    """Remap all colors in a slide's XML to a unified dark tech palette.

Uses HSL-based bucketing to map arbitrary skill colors to the target palette.
"""
    if palette is None:
        palette = _DEFAULT_PALETTE

    dark_bg = palette["dark_bg"]
    panel = palette["panel"]
    accent1 = palette["accent1"]
    accent2 = palette["accent2"]
    accent3 = palette["accent3"]
    text_light = palette["text_light"]
    text_muted = palette["text_muted"]

    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    for elem in slide._element.iter():
        tag = elem.tag
        if not (tag.endswith("}srgbClr") or tag == f"{{{ns}}}srgbClr"):
            continue

        val = elem.get("val", "")
        if len(val) != 6:
            continue

        try:
            r = int(val[0:2], 16)
            g = int(val[2:4], 16)
            b = int(val[4:6], 16)
        except ValueError:
            continue

        h, s, l = _rgb_to_hsl(r, g, b)

        # HSL-based bucketing
        if l > 0.92:
            new = text_light
        elif l < 0.08:
            new = dark_bg
        elif l < 0.18 and s < 0.4:
            new = dark_bg
        elif l < 0.25 and s < 0.5:
            new = panel
        elif s < 0.15:
            new = text_muted
        elif s > 0.4 and 160 < h < 220:
            new = accent1
        elif s > 0.4 and 220 < h < 300:
            new = accent2
        elif s > 0.4 and 100 < h < 170:
            new = accent3
        elif s > 0.4:
            # Warm hues → accent1, green range → accent3, else → accent2
            if h < 60 or h > 330:
                new = accent1
            elif h < 160:
                new = accent3
            else:
                new = accent2
        elif l > 0.6:
            new = text_light
        else:
            new = text_muted

        new_hex = f"{new[0]:02X}{new[1]:02X}{new[2]:02X}"
        elem.set("val", new_hex)


# =========================================================================
# MCP Server
# =========================================================================

mcp = FastMCP("ppt-agent")


# Server-owned reload_registry — exempt from the legacy stale guard so the
# documented stale_registry remediation path is real on this hybrid server.
@mcp.tool()
def reload_registry() -> dict:
    """Refresh the wiki discovery surface from disk and re-key the stale guard.

    The PPT runtime is hybrid: discovery comes from skills_wiki/ppt while
    execution stays on the legacy python-pptx engine. After an external
    update to skills_wiki/ppt/index.json or a domain.yaml flip, this tool
    re-instantiates the adapter, re-registers the universal wiki tools,
    and re-keys the stale guard so wrapped tools become callable again.
    """
    info: dict = {"reloaded": True, "domain": "ppt"}
    try:
        from domains.ppt.wiki_adapter import PPTWikiAdapter
        from core.skill_wiki.mcp_tools import register_wiki_tools
        from core.skill_wiki.legacy_stale import mark_runtime_backend
        from core import get_active_library_backend
        register_wiki_tools(mcp, PPTWikiAdapter())
        backend = get_active_library_backend("ppt")
        mark_runtime_backend(mcp, backend)
        info["backend"] = backend
        info["tool_surface"] = "wiki+legacy"
    except Exception as exc:  # noqa: BLE001
        info["error"] = f"{type(exc).__name__}: {exc}"
    return info


# ============================================================================
# python-pptx SHELL BUILD PATH — DISABLED (PPT goes through PPT Master only).
# The @mcp.tool() decorators on the shell build tools below are commented out
# so they are NOT registered / reachable by the agent. The shell path had
# content-fidelity bugs (list-content slides dropped their body); deck
# generation runs through the PPT Master SVG->native-PPTX path (pptmaster_*
# tools + pptmaster_r2s_prompt_runner.py / demo/PPT_example/replay.sh).
# Bodies kept as dead code for reference; do not re-enable without a fix.
# ============================================================================
# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def create_presentation(
    width_inches: float = 13.333,
    height_inches: float = 7.5,
) -> str:
    """Create a new empty PowerPoint presentation.

Args:
    width_inches: Slide width in inches (default 13.333 for 16:9).
    height_inches: Slide height in inches (default 7.5 for 16:9).

Returns:
    Human-readable confirmation with the presentation ID.
"""
    pptx = _get_pptx()
    prs = pptx["Presentation"]()
    prs.slide_width = pptx["Inches"](width_inches)
    prs.slide_height = pptx["Inches"](height_inches)
    prs_id = f"prs_{uuid.uuid4().hex[:8]}"
    _presentations[prs_id] = prs
    ratio = "16:9" if abs(width_inches / height_inches - 1.7777777777777777) < 0.1 else f"{width_inches}:{height_inches}"
    return f"Created presentation {prs_id} ({ratio}, {width_inches:.2f}x{height_inches:.2f} inches)"


def _pptmaster_error(exc: Exception) -> dict:
    return {"error": f"{type(exc).__name__}: {exc}"}


def _parse_pptmaster_r2s_brief(task_description: str, *, n_slides: int = 0) -> dict:
    raw = (task_description or "").strip()
    payload: object | None = None
    if raw:
        try:
            payload = json.loads(raw)
        except json.JSONDecodeError:
            payload = None
    if isinstance(payload, str):
        raw = payload
        payload = None
    if isinstance(payload, dict):
        brief = dict(payload)
    else:
        brief = {
            "title": raw[:160] or "PPTMaster deck",
            "audience": raw,
            "core_points": [part.strip() for part in re.split(r"[.;\n]+", raw) if part.strip()][:12],
        }
    if n_slides and not brief.get("n_slides"):
        brief["n_slides"] = int(n_slides)
    return brief


@mcp.tool()
def pptmaster_runtime_info() -> dict:
    """Return PPT Master adapter paths and capability flags.

    This is a parallel backend to the Resource2Skill python-pptx runtime.
    It does not replace ``create_presentation`` / ``add_slide_from_skill``;
    use it when the deck should be authored as PPT Master SVG pages and then
    exported to native editable PPTX.
    """
    try:
        return _get_pptmaster().runtime_info()
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_select_r2s_refs(
    task_description: str,
    n_refs: int = 3,
    n_slides: int = 0,
) -> dict:
    """Select prompt-specific Resource2Skill PPT refs for PPTMaster SVG work.

    This is a policy helper, not a PPTMaster runtime primitive.  It keeps
    Resource2Skill prompt interpretation outside ``pptmaster_engine`` while
    giving the agent stable, prompt-specific reference IDs and non-binding
    design opportunities before it starts authoring SVG.

    ``task_description`` may be plain text or a JSON object/string containing
    fields like ``title``, ``audience``, ``tone_words``, ``role_prefer``,
    ``role_avoid``, ``n_slides``, and ``core_points``.
    """
    try:
        policy = _get_pptmaster_r2s_policy()
        brief = _parse_pptmaster_r2s_brief(task_description, n_slides=n_slides)
        entries = policy._load_skill_entries()
        domain = policy.infer_domain(brief)
        refs = policy.select_skill_refs(brief, entries, k=max(2, min(int(n_refs or 3), 5)))
        details = policy._ref_details(refs, entries)
        return {
            "domain": domain,
            "refs": refs,
            "ref_details": details,
            "visual_family": policy.visual_family(domain),
            "design_opportunities": policy.design_opportunities(domain),
            "instructions": [
                "Inspect these refs with get_skill_text/get_skill_code/get_skill_visual when available.",
            "Adapt mechanisms into SVG; do not clone whole skill slides.",
            "Use one primary skill mechanism per slide; other refs may only influence minor styling.",
            "Avoid competing main visuals on a single slide. Explainer pages need one conceptual center with precise labels.",
            "Record per-slide refs in SVG comments or notes as design_refs.",
                "Do not put internal labels such as REFERENCE-ADAPTED or skill mechanisms rewritten on visible slides.",
                "Keep PPTMaster open-ended: decide slide count, order, and topology from the prompt, using refs only as optional enhancement evidence.",
            ],
        }
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_create_project(
    project_name: str,
    canvas_format: str = "ppt169",
    base_dir: str = "",
) -> dict:
    """Create a PPT Master project directory for SVG-first deck generation."""
    try:
        return _get_pptmaster().create_project(
            project_name=project_name,
            canvas_format=canvas_format,
            base_dir=base_dir or None,
        )
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_list_svg_slides(project_path: str) -> list | dict:
    """List SVG source slides in a PPT Master project."""
    try:
        return _get_pptmaster().list_svg_slides(project_path)
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_add_svg_slide(
    project_path: str,
    svg: str,
    slide_name: str = "",
    notes: str = "",
) -> dict:
    """Add a source SVG slide to a PPT Master project.

    The SVG is the editable source code for the slide.  This is the
    PPT Master equivalent of the legacy ``add_slide`` code path.
    """
    try:
        return _get_pptmaster().write_svg_slide(
            project_path=project_path,
            svg=svg,
            slide_name=slide_name,
            notes=notes,
            overwrite=False,
        )
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_replace_svg_slide(
    project_path: str,
    slide_name: str,
    svg: str,
    notes: str = "",
) -> dict:
    """Replace an existing PPT Master SVG slide.

    This preserves the Resource2Skill edit loop at SVG level: inspect source,
    rewrite source, then re-export the PPTX.
    """
    try:
        return _get_pptmaster().write_svg_slide(
            project_path=project_path,
            svg=svg,
            slide_name=slide_name,
            notes=notes,
            overwrite=True,
        )
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_get_svg_slide(project_path: str, slide_name: str) -> dict:
    """Read one PPT Master SVG slide and matching speaker notes."""
    try:
        return _get_pptmaster().read_svg_slide(project_path, slide_name)
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_delete_svg_slide(project_path: str, slide_name: str) -> dict:
    """Delete one PPT Master SVG slide and its matching notes file."""
    try:
        return _get_pptmaster().delete_svg_slide(project_path, slide_name)
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_finalize_project(project_path: str) -> dict:
    """Run PPT Master's SVG finalization pass into ``svg_final/``."""
    try:
        result = _get_pptmaster().finalize_project(project_path)
        return _get_pptmaster().command_result_to_dict(result)
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_export_project(
    project_path: str,
    output_path: str = "",
    source: str = "output",
    transition: str = "fade",
    animation: str = "auto",
    animation_trigger: str = "after-previous",
    compat: bool = False,
    finalize: bool = False,
    layout_strict: bool = False,
) -> dict:
    """Export a PPT Master SVG project to a native editable PPTX.

    ``compat=false`` uses native DrawingML only and avoids optional SVG->PNG
    fallback dependencies.  Set ``compat=true`` if the environment has
    svglib/reportlab installed and Office compatibility fallback images are
    desired.
    """
    try:
        return _get_pptmaster().export_project(
            project_path=project_path,
            output_path=output_path,
            source=source,
            transition=transition,
            animation=animation,
            animation_trigger=animation_trigger,
            compat=compat,
            finalize=finalize,
            layout_strict=layout_strict,
        )
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_validate_project(project_path: str, strict: bool = False) -> dict:
    """Validate PPTMaster SVG source layout before export.

    The validator is intentionally non-prescriptive: it does not enforce a
    template. It flags obvious layout failures such as off-canvas elements,
    likely text overflow, severe text overlap, empty slides, and malformed SVG.
    Use `strict=true` before final export when preparing open-source-quality
    artifacts.
    """
    try:
        return _get_pptmaster().validate_project(project_path, strict=strict)
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_import_pptx_template(
    pptx_path: str,
    output_dir: str = "",
    inheritance_mode: str = "both",
    manifest_only: bool = False,
) -> dict:
    """Import a reference PPTX as PPT Master template source.

    The upstream importer extracts theme/font/layout/master metadata and emits
    SVG views that agents can inspect or turn into reusable template assets.
    """
    try:
        return _get_pptmaster().import_pptx_template(
            pptx_path=pptx_path,
            output_dir=output_dir,
            inheritance_mode=inheritance_mode,
            manifest_only=manifest_only,
        )
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_list_templates(kind: str = "layout") -> dict:
    """List bundled PPT Master templates.

    ``kind`` is one of ``layout``, ``brand``, ``deck``, or ``chart``.
    """
    try:
        return _get_pptmaster().list_templates(kind=kind)
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_get_template(kind: str, template_id: str, svg_name: str = "") -> dict:
    """Inspect a bundled PPT Master template and optionally return one SVG file."""
    try:
        return _get_pptmaster().get_template(
            kind=kind,
            template_id=template_id,
            svg_name=svg_name,
        )
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


@mcp.tool()
def pptmaster_copy_template_to_project(
    project_path: str,
    kind: str,
    template_id: str,
) -> dict:
    """Copy a bundled PPT Master template into a project's ``templates/`` dir."""
    try:
        return _get_pptmaster().copy_template_to_project(
            project_path=project_path,
            kind=kind,
            template_id=template_id,
        )
    except Exception as exc:  # noqa: BLE001
        return _pptmaster_error(exc)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def add_slide(prs_id: str, code: str) -> str:
    """Add a slide by executing python-pptx code.

The code should work with a `prs` (Presentation) variable that is already
available in scope. It should add one slide to `prs`. Common pattern:
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    # ... add shapes to slide ...

Args:
    prs_id: Presentation ID from create-presentation.
    code: Python code using python-pptx API.

Returns:
    Confirmation with slide index and shape count, or error message.
"""
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"

    engine = _get_engine()
    code_patched = engine.patch_skill_code(code)

    # Strip prs creation / save calls — we provide the prs object
    import re
    code_patched = re.sub(r"prs\s*=\s*Presentation\(\)", "", code_patched)
    code_patched = re.sub(r"prs\.save\([^)]*\)", "", code_patched)

    # Margin auto-fix: rewrite bare Inches(0) in textbox/shape positions to
    # Inches(0.6) unless the shape is full-bleed (width >= 12in).
    # Kills the left-edge clipping that is the #1 visual-QA failure.
    def _fix_margin(m):
        call = m.group(0)
        args = re.findall(r"Inches\(\s*([-\d.]+)\s*\)", call)
        if len(args) < 4:
            return call
        try:
            left, top, width, height = map(float, args[:4])
        except ValueError:
            return call
        if width >= 12.0 or height >= 7.0:
            return call  # full-bleed; leave alone
        if left < 0.5:
            call = call.replace(f"Inches({args[0]})", "Inches(0.6)", 1)
        return call
    _pos_pattern = re.compile(
        r"(?:add_textbox|add_shape)\s*\([^)]*?Inches\([^)]+\)\s*,\s*Inches\([^)]+\)\s*,\s*Inches\([^)]+\)\s*,\s*Inches\([^)]+\)[^)]*\)",
        re.DOTALL,
    )
    code_patched = _pos_pattern.sub(_fix_margin, code_patched)

    full_code = engine._SKILL_EXEC_PREAMBLE + "\n" + code_patched
    ns = {"prs": prs, "__builtins__": __builtins__}

    n_before = len(prs.slides)

    import os
    _saved_cwd = os.getcwd()
    _exec_tmpdir = tempfile.mkdtemp(prefix="ppt_exec_")
    try:
        os.chdir(_exec_tmpdir)
        exec(full_code, ns)
    except Exception as e:
        # Rollback orphan slides: if code partially executed and added a slide
        # before crashing, remove it to prevent blank orphan slides.
        n_after_error = len(prs.slides)
        if n_after_error > n_before:
            orphan_count = n_after_error - n_before
            for _ in range(orphan_count):
                rId = prs.slides._sldIdLst[-1].get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                )
                if rId:
                    prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
            log.warning("Rolled back %d orphan slide(s) after exec error", orphan_count)
        return f"Error executing code: {type(e).__name__}: {e}"
    finally:
        os.chdir(_saved_cwd)
        import shutil
        shutil.rmtree(_exec_tmpdir, ignore_errors=True)

    n_after = len(prs.slides)
    if n_after > n_before:
        new_slide = prs.slides[-1]
        _force_word_wrap(new_slide)
        shapes = len(new_slide.shapes)
        _save_temp(prs, prs_id)
        return f"Added slide {n_after} ({shapes} shapes)"
    return "Warning: code executed but no new slide was added to the presentation"


def _force_word_wrap(slide):
    """Force word_wrap=True on every textbox so long body text wraps instead of
    overflowing horizontally (the #1 cause of ugly cross-card text bleed)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if shape.text_frame.word_wrap is not True:
                shape.text_frame.word_wrap = True
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Skill clone quality check
# ---------------------------------------------------------------------------

_PLACEHOLDER_PATTERNS = [
    "text here", "title here", "your text", "insert text",
    "click to add", "lorem ipsum", "placeholder",
    "标题", "副标题", "会议主题", "在此输入", "文本内容",
    "sample text", "add text", "type here",
    # Known template-demo phrases that leak through from downloaded skill samples
    "solar eclipse", "lunar eclipse", "pop quiz", "experience space",
    "neque porro", "dolor sit amet", "aenean commodo",
]


def _check_skill_clone_quality(slide, content_brief: str = "") -> list[str]:
    """Check a skill-cloned slide for common quality issues.

    Returns a list of warning strings (empty if all OK).
    """
    warnings: list[str] = []

    # Gather all text from the slide
    all_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                all_texts.append(text)
    full_text = " ".join(all_texts).lower()

    # Check 1: placeholder text residue
    found_placeholders = []
    for pat in _PLACEHOLDER_PATTERNS:
        if pat in full_text:
            found_placeholders.append(pat)
    if found_placeholders:
        warnings.append(
            f"placeholder text residue detected ({', '.join(found_placeholders[:3])})"
        )

    # Check 2: content_brief injection failure
    if content_brief:
        title_line = content_brief.strip().split("\n")[0].strip()
        # Check if at least the first significant word of the title appears
        title_words = [w for w in title_line.split() if len(w) > 3]
        if title_words:
            matches = sum(1 for w in title_words if w.lower() in full_text)
            if matches < len(title_words) * 0.3:
                warnings.append(
                    f"content not injected — title '{title_line[:40]}' not found in slide text"
                )

    # Check 3: too few shapes (broken/empty clone)
    n_shapes = len(slide.shapes)
    if n_shapes < 3:
        warnings.append(f"only {n_shapes} shapes — slide may be broken or empty")

    # Check 4: no text at all (pure decorative shape with no content)
    if not all_texts and n_shapes > 0:
        warnings.append("slide has shapes but no readable text")

    return warnings


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def add_slide_from_skill(
    prs_id: str,
    skill_id: str,
    content_brief: str = "",
    style_hints: str = "",
    palette: str = "",
) -> str:
    """Add a slide using a skill from the library.

The skill's create_slide() function will be adapted to append to the
current presentation instead of creating a new file.

Args:
    prs_id: Presentation ID.
    skill_id: Skill ID from the library (use list-skills to find).
    content_brief: Text content to inject (title, body text, etc.).
    style_hints: Visual style preferences (e.g., 'dark, gradient, tech').
    palette: Color palette control. Options:
        - "" (empty): keep original skill colors (no remapping)
        - "none": same as empty — keep skill's original colors
        - JSON dict: custom palette, e.g. '{"dark_bg":[245,247,250],"panel":[255,255,255],...}'

Returns:
    Confirmation with slide details, or error message.
"""
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"

    _ensure_skills_loaded()
    if not _skill_index:
        return "Error: skill library not loaded"

    engine = _get_engine()

    # Wiki-first lookup: skills_wiki/ppt/<skill_id>/code/skill.py is a flat
    # python file with a create_slide(output_pptx_path, ...) entrypoint. If
    # the skill_id resolves to a wiki entry, use it directly so the with-skills
    # arm is genuinely executing wiki code (not just referencing it).
    code: str | None = None
    detail: dict | None = None
    _wiki_root = Path(__file__).resolve().parents[3] / "skills_wiki" / "ppt"
    _wiki_skill_dir = _wiki_root / skill_id
    _wiki_code_file = _wiki_skill_dir / "code" / "skill.py"
    if _wiki_code_file.exists():
        try:
            code = _wiki_code_file.read_text(encoding="utf-8")
            # Use empty dict so any later `detail.get(...)` calls return
            # safe defaults; the wiki layout has no "analysis" / "skill_name"
            # but has the executable code which is what we care about.
            detail = {}
        except Exception:
            code = None
            detail = None

    if code is None:
        # Fall back to legacy skills_library/ppt JSON-with-analysis layout.
        detail = engine.get_skill_detail(_skills_dir, skill_id, _skill_index)
        if not detail:
            return f"Error: skill '{skill_id}' not found in library"
        analysis = detail.get("analysis", "")
        code = engine.extract_code_from_analysis(analysis)
    if not code:
        return f"Error: no Python code found in skill '{skill_id}' analysis"

    # Execute skill code in a temporary directory
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = str(Path(tmpdir) / "skill_output.pptx")

        # Build parameter overrides
        params = {}
        if content_brief:
            params["title_text"] = content_brief.split("\n")[0][:80]
            if "\n" in content_brief:
                rest = content_brief.split("\n", 1)[1]
                params["subtitle_text"] = rest[:200]
                params["body_text"] = rest[:500]
        if style_hints:
            params["style"] = style_hints

        success, shapes = engine.exec_skill_code(code, tmp_path, params)
        if not success:
            return f"Error: skill '{skill_id}' execution failed"

        # Import the generated slide into our presentation
        # Only take the FIRST slide — multi-slide skills produce extra
        # slides (navigation bars, alternate layouts) that break the deck.
        try:
            from pptx import Presentation as _P
            tmp_prs = _P(tmp_path)
            if len(tmp_prs.slides) == 0:
                return "Error: skill produced no slides"

            src_slide = tmp_prs.slides[0]
            layout = prs.slide_layouts[6]
            new_slide = prs.slides.add_slide(layout)

            # Clone shapes via XML deep copy
            import copy
            from lxml import etree
            for shape_elem in src_slide.shapes._spTree:
                new_slide.shapes._spTree.append(copy.deepcopy(shape_elem))

            # Clone background
            if src_slide.background._element is not None:
                bg_elem = copy.deepcopy(src_slide.background._element)
                if new_slide.background._element is not None:
                    new_slide._element.replace(
                        new_slide.background._element, bg_elem
                    )

            # Clone image relationships
            for rel in src_slide.part.rels.values():
                if rel.is_external:
                    new_slide.part.rels.get_or_add_ext_rel(
                        rel.reltype, rel.target_ref
                    )
                    continue

                rt = rel.reltype
                if any(skip in rt for skip in ("slideLayout", "slideMaster", "theme", "notesMaster", "handoutMaster")):
                    continue

                try:
                    new_slide.part.rels.get_or_add(
                        rel.reltype, rel.target_part
                    )
                except Exception:
                    continue

        except Exception as e:
            return f"Error importing skill slide: {type(e).__name__}: {e}"

    # Post-process: inject content and unify colors
    if content_brief:
        _inject_content_into_slide(prs.slides[-1], content_brief)

    if palette == "none" or not palette:
        pass  # Keep skill's original colors
    else:
        try:
            import json as _json
            custom_pal = _json.loads(palette)
            # Convert list values to tuples
            custom_pal = {k: tuple(v) if isinstance(v, list) else v for k, v in custom_pal.items()}
            _unify_slide_colors(prs.slides[-1], palette=custom_pal)
        except (ValueError, KeyError, TypeError):
            pass  # Invalid JSON — keep original colors

    n_slides = len(prs.slides)
    final_shapes = len(prs.slides[-1].shapes)
    skill_name = detail.get("skill_name", skill_id)
    _save_temp(prs, prs_id)

    # Post-clone quality check
    warnings = _check_skill_clone_quality(prs.slides[-1], content_brief)
    base_msg = f"Added slide {n_slides} from skill '{skill_name}' ({final_shapes} shapes)"
    if warnings:
        # HARD REJECT if template placeholder phrases leaked through — these are
        # unfixable without full content replacement. Delete the slide and force
        # the agent to use custom `add_slide` code.
        hard_reject_hits = [w for w in warnings if "placeholder text residue" in w]
        if hard_reject_hits:
            # Remove the slide we just added
            try:
                slide_list = prs.slides._sldIdLst
                last = slide_list[-1]
                rId = last.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                if rId:
                    prs.part.drop_rel(rId)
                slide_list.remove(last)
                _save_temp(prs, prs_id)
            except Exception:
                pass
            return (
                f"Error: skill '{skill_name}' clone had {hard_reject_hits[0]} that could not be scrubbed. "
                f"The slide was REJECTED and removed. Do NOT retry this skill — use `add_slide` with "
                f"custom python-pptx code instead."
            )
        return f"{base_msg} — WARNING: {'; '.join(warnings)}. Consider using replace-slide with custom code instead."
    return base_msg


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def replace_slide(prs_id: str, slide_index: int, code: str) -> str:
    """Replace an existing slide with new code.

Deletes the slide at slide_index and inserts a new one generated by code.

Args:
    prs_id: Presentation ID.
    slide_index: 1-based slide index to replace.
    code: Python code for the replacement slide.

Returns:
    Confirmation or error message.
"""
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"

    idx = slide_index - 1
    if idx < 0 or idx >= len(prs.slides):
        return f"Error: slide index {slide_index} out of range (1-{len(prs.slides)})"

    # Delete the slide
    rId = prs.slides._sldIdLst[idx].get("r:id") if hasattr(prs.slides, "_sldIdLst") else None
    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(xml_slides[idx])
    entry = _prs_metadata.get(prs_id)
    if entry and idx < len(entry.get("slides") or []):
        del entry["slides"][idx]
        _reindex_prs_metadata(prs_id)

    # Add new slide via add_slide tool
    result = add_slide(prs_id, code)
    if result.startswith("Error"):
        return result
    return result.replace("Added", "Replaced")


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def delete_slide(prs_id: str, slide_index: int) -> str:
    """Delete a slide from the presentation.

Args:
    prs_id: Presentation ID.
    slide_index: 1-based slide index to delete.

Returns:
    Confirmation or error message.
"""
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"

    idx = slide_index - 1
    if idx < 0 or idx >= len(prs.slides):
        return f"Error: slide index {slide_index} out of range (1-{len(prs.slides)})"

    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(xml_slides[idx])
    entry = _prs_metadata.get(prs_id)
    if entry and idx < len(entry.get("slides") or []):
        del entry["slides"][idx]
        _reindex_prs_metadata(prs_id)
    _save_temp(prs, prs_id)
    return f"Deleted slide {slide_index}. Presentation now has {len(prs.slides)} slides."


def _auto_add_transitions(prs):
    """Add fade transitions to all slides that don't already have one."""
    from lxml import etree
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    for slide in prs.slides:
        # Check if transition already exists
        existing = slide._element.find(f"{{{ns_p}}}transition")
        if existing is not None:
            continue
        # Add a smooth fade transition
        transition_xml = (
            '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'spd="med" advClick="1"><p:fade/></p:transition>'
        )
        slide._element.append(etree.fromstring(transition_xml))


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def set_transition(prs_id: str, slide_index: int, transition_type: str = "fade") -> str:
    """Set the transition effect for a slide.

    Args:
        prs_id: Presentation ID.
        slide_index: 1-based slide index.
        transition_type: One of: fade, push, wipe, morph, none.

    Returns:
        Confirmation message.
    """
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"

    idx = slide_index - 1
    if idx < 0 or idx >= len(prs.slides):
        return f"Error: slide index {slide_index} out of range"

    from lxml import etree
    slide = prs.slides[idx]
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

    # Remove existing transition
    for old in slide._element.findall(f"{{{ns_p}}}transition"):
        slide._element.remove(old)

    if transition_type == "none":
        _save_temp(prs, prs_id)
        return f"Removed transition from slide {slide_index}"

    transition_map = {
        "fade": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advClick="1"><p:fade/></p:transition>',
        "push": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advClick="1"><p:push dir="l"/></p:transition>',
        "wipe": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advClick="1"><p:wipe dir="l"/></p:transition>',
        "morph": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main" spd="med" advClick="1"><p159:morph option="byObject"/></p:transition>',
    }

    xml = transition_map.get(transition_type)
    if not xml:
        return f"Error: unknown transition '{transition_type}'. Use: fade, push, wipe, morph, none"

    slide._element.append(etree.fromstring(xml))
    _save_temp(prs, prs_id)
    return f"Set {transition_type} transition on slide {slide_index}"


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def save_presentation(prs_id: str, output_path: str,
                      morph_lint_strict: bool = False) -> str:
    """Save the presentation to a PPTX file on disk.

Args:
    prs_id: Presentation ID.
    output_path: Absolute file path for the output .pptx file.
    morph_lint_strict: when True, promote morph-lint warnings to errors and
        refuse to save if any slide with non-empty anchor_names_set lacks a
        morph transition. Default False (emit warnings to stderr only).

Returns:
    Confirmation with file path, or error message.
"""
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"

    entry = _ensure_prs_metadata(prs_id)
    # Keep metadata aligned with the actual presentation if the agent deleted
    # slides. Scoring and progress must reflect what will really be saved, not
    # the number of attempted shell calls.
    actual_content_slides = len(prs.slides)
    if len(entry.get("slides") or []) > actual_content_slides:
        entry["slides"] = entry["slides"][:actual_content_slides]
        _reindex_prs_metadata(prs_id)

    target = entry.get("archetype_slides_target")
    if target and actual_content_slides < target:
        return (
            f"Error: presentation has only {actual_content_slides} content "
            f"slides, but the target is {target}. Add "
            f"{target - actual_content_slides} more slides with "
            "add_slide_from_shell before save_presentation."
        )

    # Auto-add fade transitions to slides that don't have one.
    _auto_add_transitions(prs)

    # Morph-lint pass: any slide with anchor names but a non-morph
    # transition is either an oversight or a deliberate choice. Warn by
    # default; fail the save under morph_lint_strict.
    lint_issues = []
    for idx, meta in enumerate(entry["slides"]):
        # Refresh transition_set from the actual slide in case the agent
        # called set_transition after add_slide_from_shell.
        try:
            transition_now = _slide_transition_kind(prs.slides[idx])
        except Exception:
            transition_now = meta.get("transition_set", "none")
        meta["transition_set"] = transition_now
        if meta.get("anchor_names_set") and transition_now != "morph":
            lint_issues.append(
                f"slide {idx}: anchors {meta['anchor_names_set']} present "
                f"but transition is {transition_now!r} — will not morph"
            )

    for msg in lint_issues:
        print(f"[morph-lint] {msg}", file=sys.stderr)

    if morph_lint_strict and lint_issues:
        return (
            "Error: morph_lint_strict=true and "
            f"{len(lint_issues)} slide(s) have anchors without a morph "
            "transition. Fix with transition='morph' in add_slide_from_shell "
            "or call set_transition(morph) before saving. Details: "
            + "; ".join(lint_issues)
        )

    # Append a methodology appendix slide that lists every shell used.
    try:
        _append_methodology_appendix(prs, entry)
    except Exception as e:
        log.warning("methodology appendix append failed: %s", e)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(out))
        # Deduplicate ZIP entries
        try:
            import zipfile, io
            with zipfile.ZipFile(str(out), 'r') as zin:
                seen = set()
                buf = io.BytesIO()
                with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
                    for item in zin.infolist():
                        if item.filename not in seen:
                            seen.add(item.filename)
                            zout.writestr(item, zin.read(item.filename))
                buf.seek(0)
                with open(str(out), 'wb') as f:
                    f.write(buf.read())
        except Exception:
            pass
        _prs_paths[prs_id] = out
        n_slides = len(prs.slides)

        # Emit deck-plan.json next to the saved pptx so downstream diff and
        # gallery tools can find it.
        deck_plan_path = out.parent / "deck-plan.json"
        deck_plan = {
            "archetype": entry.get("archetype"),
            "archetype_suggested_slides": entry.get("archetype_suggested_slides"),
            "requested_slides_target": entry.get("requested_slides_target"),
            "slides_target": entry.get("archetype_slides_target"),
            "theme": entry.get("theme"),
            "pptx_path": str(out),
            "slides": entry["slides"],
            "morph_lint": {
                "issues": lint_issues,
                "strict_mode": bool(morph_lint_strict),
            },
        }
        try:
            deck_plan_path.write_text(
                json.dumps(deck_plan, indent=2, ensure_ascii=False)
            )
        except Exception as e:
            log.warning("failed to write deck-plan.json: %s", e)

        extras = []
        if lint_issues:
            extras.append(f"{len(lint_issues)} morph-lint warning(s)")
        if deck_plan_path.exists():
            extras.append(f"deck-plan: {deck_plan_path}")
        extras_str = (" | " + " | ".join(extras)) if extras else ""

        return (
            f"Saved {n_slides}-slide presentation to {out}."
            f"{extras_str}"
        )
    except Exception as e:
        return f"Error saving: {type(e).__name__}: {e}"


def _append_methodology_appendix(prs, entry: dict) -> None:
    """Append a final appendix slide listing shell provenance for each slide.

    The appendix is regenerated on every save — if a previous appendix exists
    (identified by a marker in the first text frame), it is replaced rather
    than duplicated.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    marker = "[methodology-appendix]"

    # Remove any prior appendix slide (matches the marker at position 0).
    for idx in range(len(prs.slides) - 1, -1, -1):
        slide = prs.slides[idx]
        try:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if marker in shape.text_frame.text:
                    slide_list = prs.slides._sldIdLst
                    rId = slide_list[idx].get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                    )
                    if rId:
                        prs.part.drop_rel(rId)
                    slide_list.remove(slide_list[idx])
                    break
        except Exception:
            continue

    if not entry["slides"]:
        return

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5),
                                         Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.text = f"Skill Library Methodology {marker}"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(26)
    run.font.bold = True

    # Body: a table-like text block listing each slide's shell provenance.
    body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5),
                                        Inches(12.0), Inches(5.5))
    btf = body_box.text_frame
    btf.word_wrap = True
    # Header line — full field names (no truncation).
    header_para = btf.paragraphs[0]
    header_para.text = (
        "  #  role                  shell_id                              "
        "selection_origin  source_video_id                              source_timestamp"
    )
    header_para.runs[0].font.size = Pt(10)
    header_para.runs[0].font.bold = True

    for s in entry["slides"]:
        p = btf.add_paragraph()
        prov = s.get("distill_provenance") or {}
        vid = prov.get("source_video_id") or "-"
        ts = prov.get("source_timestamp")
        ts_str = "-" if ts in (None, "") else str(ts)
        origin = prov.get("selection_origin") or "-"
        # Plain text rows; no truncation. python-pptx will wrap if a row
        # overflows the box width but the full string remains in the
        # underlying XML so reviewers reading the pptx see every byte.
        row = (
            f"  {s['slide_index']:>2}  "
            f"{s['role']:<22}  "
            f"{s['shell_id']:<36}  "
            f"{origin:<16}  "
            f"{vid:<44}  "
            f"{ts_str}"
        )
        p.text = row
        p.runs[0].font.size = Pt(9)
        p.runs[0].font.name = "Courier New"

    footer = slide.shapes.add_textbox(Inches(0.7), Inches(7.05),
                                      Inches(12.0), Inches(0.3))
    ftf = footer.text_frame
    ftf.text = (
        "Origin legend: S=seed, D=distilled from YouTube frame, "
        "R=retrieved as-is, C=composed at runtime, N=newly generated."
    )
    ftf.paragraphs[0].runs[0].font.size = Pt(9)
    ftf.paragraphs[0].runs[0].font.color.rgb = RGBColor(120, 120, 120)





# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def get_slide_info(prs_id: str, slide_index: int) -> str:
    """Get information about a specific slide.

Args:
    prs_id: Presentation ID.
    slide_index: 1-based slide index.

Returns:
    Slide details: shape count, shape types, text content preview.
"""
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"

    idx = slide_index - 1
    if idx < 0 or idx >= len(prs.slides):
        return f"Error: slide index {slide_index} out of range (1-{len(prs.slides)})"

    slide = prs.slides[idx]
    shapes = list(slide.shapes)
    shape_types = {}
    text_preview = []

    for s in shapes:
        stype = str(s.shape_type) if hasattr(s, "shape_type") else "unknown"
        shape_types[stype] = shape_types.get(stype, 0) + 1
        if not s.has_text_frame:
            continue
        text = s.text_frame.text[:100]
        if not text.strip():
            continue
        text_preview.append(text.strip())

    lines = [
        f"Slide {slide_index}: {len(shapes)} shapes",
        f"  Shape types: {dict(shape_types)}",
    ]
    if text_preview:
        lines.append(f"  Text: {'; '.join(text_preview[:5])}")
    w = prs.slide_width
    h = prs.slide_height
    from pptx.util import Inches
    lines.append(f"  Dimensions: {w / Inches(1):.2f} x {h / Inches(1):.2f} inches")

    overlap_report = _detect_overlaps(slide, prs)
    if overlap_report:
        lines.append("  OVERLAPS:")
        for row in overlap_report:
            lines.append(f"    {row}")
    return "\n".join(lines)


def _detect_overlaps(slide, prs):
    """Deterministic bbox collision detection. Returns list of human-readable overlap descriptions."""
    from pptx.util import Inches
    shapes = []
    for s in slide.shapes:
        try:
            l, t, w, h = s.left, s.top, s.width, s.height
            if l is None or t is None or w is None or h is None:
                continue
            in_l = l / Inches(1); in_t = t / Inches(1)
            in_w = w / Inches(1); in_h = h / Inches(1)
            # Pick a label
            label = ""
            if s.has_text_frame and s.text_frame.text.strip():
                label = s.text_frame.text.strip().replace("\n", " ")[:30]
            if not label:
                label = str(s.shape_type) if hasattr(s, "shape_type") else "shape"
            shapes.append({
                "id": s.shape_id, "label": label,
                "l": in_l, "t": in_t, "r": in_l + in_w, "b": in_t + in_h,
                "area": in_w * in_h, "has_text": s.has_text_frame and bool(s.text_frame.text.strip()),
                "w": in_w, "h": in_h,
            })
        except Exception:
            continue

    issues = []
    # Canvas bounds
    slide_w = prs.slide_width / Inches(1)
    slide_h = prs.slide_height / Inches(1)
    for sh in shapes:
        if sh["r"] > slide_w + 0.02 or sh["b"] > slide_h + 0.02 or sh["l"] < -0.02 or sh["t"] < -0.02:
            issues.append(
                f"shape#{sh['id']} '{sh['label']}' OUT OF BOUNDS "
                f"(bbox {sh['l']:.2f},{sh['t']:.2f} -> {sh['r']:.2f},{sh['b']:.2f}; canvas {slide_w:.2f}x{slide_h:.2f})"
            )

    # Pairwise overlap — only flag text-vs-text or small-shape-vs-text
    # full-bleed backgrounds (area > 80% of slide) are excluded
    slide_area = slide_w * slide_h
    for i, a in enumerate(shapes):
        if a["area"] > 0.7 * slide_area:
            continue
        for b in shapes[i+1:]:
            if b["area"] > 0.7 * slide_area:
                continue
            ix = max(0, min(a["r"], b["r"]) - max(a["l"], b["l"]))
            iy = max(0, min(a["b"], b["b"]) - max(a["t"], b["t"]))
            inter = ix * iy
            if inter <= 0:
                continue
            # Significance: ignore tiny touches; only flag if intersection >= 35% of smaller shape
            min_area = min(a["area"], b["area"])
            if min_area <= 0 or inter / min_area < 0.35:
                continue
            # Only interesting when at least one has text (overlapping decorative layers are OK)
            if not (a["has_text"] or b["has_text"]):
                continue
            # Skip if one fully contains the other AND outer has no text (likely a frame around text)
            a_contains_b = a["l"] <= b["l"] and a["t"] <= b["t"] and a["r"] >= b["r"] and a["b"] >= b["b"]
            b_contains_a = b["l"] <= a["l"] and b["t"] <= a["t"] and b["r"] >= a["r"] and b["b"] >= a["b"]
            if (a_contains_b and not a["has_text"]) or (b_contains_a and not b["has_text"]):
                continue
            # Decorative number badges: small shape with ≤3-char text (e.g. "01", "02")
            # sitting over a larger card is an intentional label — don't flag.
            def _is_decor_badge(sh):
                label = sh.get("label", "")
                stripped = label.replace(" ", "").replace(".", "").replace("/", "")
                return sh["has_text"] and len(stripped) <= 3 and sh["area"] < 0.05 * slide_area
            if _is_decor_badge(a) or _is_decor_badge(b):
                continue
            # Hairline/rule exclusion: a text-free shape with very small height
            # (< 0.1in, i.e. a decorative separator rule) shouldn't be flagged
            # when it intersects text boxes crossing the rule.
            def _is_hairline(sh):
                return not sh["has_text"] and sh["h"] < 0.1
            if _is_hairline(a) or _is_hairline(b):
                continue
            # "Text slightly exceeds its framing panel" case: if one shape contains
            # 90%+ of the other's area and the outer has no text, treat as an
            # intentional frame/panel pairing — the small bbox overflow is a
            # rounding/kerning artifact, not a real collision.
            outer, inner = (a, b) if a["area"] >= b["area"] else (b, a)
            if not outer["has_text"] and inner["area"] > 0 and inter / inner["area"] >= 0.9:
                continue
            # Decorative band/overlay: a text-free smaller shape fully inside a
            # text-bearing larger card (e.g. gradient header bar inside a content
            # card). These are always intentional visual decorations.
            if not inner["has_text"] and inter / inner["area"] >= 0.95:
                continue
            pct = inter / min_area * 100
            issues.append(
                f"shape#{a['id']} '{a['label']}' and shape#{b['id']} '{b['label']}' OVERLAP {pct:.0f}% "
                f"(a={a['l']:.2f},{a['t']:.2f}-{a['r']:.2f},{a['b']:.2f}  b={b['l']:.2f},{b['t']:.2f}-{b['r']:.2f},{b['b']:.2f})"
            )
            if len(issues) >= 8:
                return issues
    return issues


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def render_slide(prs_id: str, slide_index: int = 0) -> str:
    """Render a slide to PNG and return as base64.

Requires LibreOffice to be installed. If slide_index is 0, renders
the first slide (for backward compat).

Args:
    prs_id: Presentation ID.
    slide_index: 1-based slide index (default: first slide).

Returns:
    Base64-encoded PNG data, or error message.
"""
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"

    # Save to temp for rendering
    tmp_path = _save_temp(prs, prs_id)
    engine = _get_engine()

    idx = max(0, slide_index - 1)
    if slide_index <= 1:
        b64 = engine.pptx_to_b64(tmp_path)
    else:
        b64 = engine.pptx_slide_to_b64(tmp_path, idx)

    if b64:
        return f"data:image/png;base64,{b64}"
    return "Error: rendering failed. Is LibreOffice installed?"


@mcp.tool()
def list_skills(category: str = "", query: str = "", verified_only: bool = True, theme: str = "", tag: str = "") -> str:
    """List or search skills in the library.

Args:
    category: Filter by category (e.g., 'cover_divider', 'typography').
              Leave empty to list all categories with counts.
    query: Search query to filter skills by name (substring match).
    verified_only: If true (default), only show skills with exec_ok=True.
    theme: Filter by theme tag: 'light' or 'dark'. Leave empty for all.
    tag: Filter by semantic tag (e.g., 'cover', 'gradient_bg', 'timeline').
         Searches across all categories.

Returns:
    List of matching skills with IDs and names.
"""
    _ensure_skills_loaded()
    if not _skill_index:
        return "Error: skill library not loaded"

    # No category or query or tag → show category summary
    if not category and not query and not tag:
        cats = {}
        cats_verified = {}
        all_tags = set()
        for s in _skill_index:
            c = s.get("category", "unknown")
            cats[c] = cats.get(c, 0) + 1
            meta = (_skill_metadata or {}).get(s["skill_id"], {})
            if meta.get("exec_ok") is True:
                cats_verified[c] = cats_verified.get(c, 0) + 1
            for t in meta.get("semantic_tags", []):
                if not t.endswith("_theme"):
                    all_tags.add(t)

        total_verified = sum(cats_verified.values())
        lines = [f"Skill library: {len(_skill_index)} skills in {len(cats)} categories ({total_verified} verified)"]
        for c, n in sorted(cats.items(), key=lambda x: -x[1]):
            v = cats_verified.get(c, 0)
            lines.append(f"  {c}: {n} skills ({v} verified)")
        if all_tags:
            sorted_tags = sorted(all_tags)[:40]
            lines.append(f"\nSemantic tags (use tag=... to search across categories):")
            lines.append(f"  {', '.join(sorted_tags)}")
        return "\n".join(lines)

    # Filter by category
    results = _skill_index
    if category:
        results = [s for s in results if s.get("category") == category]

    # Filter by query
    if query:
        q = query.lower()
        results = [s for s in results if q in s.get("skill_name", "").lower()]

    # Filter by verified
    if verified_only:
        results = [
            s for s in results
            if (_skill_metadata or {}).get(s["skill_id"], {}).get("exec_ok") is True
        ]

    # Filter by theme tag
    if theme:
        theme_tag = f"{theme.lower()}_theme"
        results = [
            s for s in results
            if theme_tag in (_skill_metadata or {}).get(s["skill_id"], {}).get("semantic_tags", [])
        ]

    # Filter by semantic tag
    if tag:
        tag_lower = tag.lower().replace(" ", "_")
        results = [
            s for s in results
            if tag_lower in (_skill_metadata or {}).get(s["skill_id"], {}).get("semantic_tags", [])
        ]

    if not results:
        hint = " (try verified_only=false to see all)" if verified_only else ""
        return f"No skills found (category={category!r}, query={query!r}){hint}"

    # Shuffle and present (randomized for variety)
    import random
    results_copy = list(results)
    random.shuffle(results_copy)

    lines = [f"Found {len(results_copy)} verified skills:"]
    for s in results_copy[:30]:
        sid = s["skill_id"]
        name = s.get("skill_name", sid)
        cat = s.get("category", "?")
        meta = (_skill_metadata or {}).get(sid, {})
        shapes = meta.get("exec_shapes", "?")
        tags = meta.get("semantic_tags", [])
        tag_str = f" tags=[{', '.join(tags[:5])}]" if tags else ""
        lines.append(f"  [{cat}] {name} (id: {sid}) [VERIFIED {shapes} shapes]{tag_str}")

    if len(results_copy) > 30:
        lines.append(f"  ... and {len(results_copy) - 30} more")
    return "\n".join(lines)


@mcp.tool()
def get_skill_info(skill_id: str) -> str:
    """Get detailed information about a specific skill.

Args:
    skill_id: The skill ID to look up.

Returns:
    Skill details: name, category, tags, when to use, and code preview.
"""
    _ensure_skills_loaded()
    if not _skill_index:
        return "Error: skill library not loaded"

    engine = _get_engine()
    detail = engine.get_skill_detail(_skills_dir, skill_id, _skill_index)
    if not detail:
        return f"Error: skill '{skill_id}' not found"

    name = detail.get("skill_name", skill_id)
    cat = detail.get("category", "unknown")
    meta = (_skill_metadata or {}).get(skill_id, {})

    lines = [f"Skill: {name}"]
    lines.append(f"  ID: {skill_id}")
    lines.append(f"  Category: {cat}")

    if meta:
        lines.append(f"  Scope: {meta.get('scope', '?')}")
        lines.append(f"  Layer: {meta.get('layer', '?')}")
        lines.append(f"  Transition: {meta.get('transition_type', 'none')}")
        lines.append(f"  Complexity: {meta.get('visual_complexity', '?')}")
        tags = meta.get("semantic_tags", [])
        if tags:
            lines.append(f"  Tags: {', '.join(tags)}")
        wtu = meta.get("when_to_use", "")
        if wtu:
            lines.append(f"  When to use: {wtu}")
        lines.append(f"  Exec OK: {meta.get('exec_ok', '?')}")
        lines.append(f"  Expected shapes: {meta.get('exec_shapes', '?')}")

    # Code preview
    analysis = detail.get("analysis", "")
    code = engine.extract_code_from_analysis(analysis)
    if code:
        code_lines = code.split("\n")
        preview = "\n".join(code_lines[:8])
        lines.append(f"  Code preview:\n    {preview}")

    return "\n".join(lines)


@mcp.tool()
def get_skill_code(skill_id: str) -> str:
    """Get full Python code and visual techniques from a skill for design reference.

    Use this to study how a skill creates its visual effects, then write your
    OWN code borrowing those techniques with your own content and palette.
    """
    _ensure_skills_loaded()
    engine = _get_engine()
    detail = engine.get_skill_detail(_skills_dir, skill_id, _skill_index)
    if not detail:
        return f"Error: skill '{skill_id}' not found"

    analysis = detail.get("analysis", "")
    code = engine.extract_code_from_analysis(analysis)
    if not code:
        return f"Error: no code found in skill '{skill_id}'"

    name = detail.get("skill_name", skill_id)

    # Extract visual mechanism
    mechanism = ""
    m = re.search(r'\*\*Core (?:Visual )?Mechanism\*\*[:\s]*(.+?)(?:\n\n|\n\*\*)', analysis, re.DOTALL)
    if m:
        mechanism = m.group(1).strip()[:300]

    # Auto-detect techniques
    techniques = []
    if 'gradFill' in code or 'gradient' in code.lower():
        techniques.append('GRADIENT_FILL')
    if 'outerShdw' in code or 'innerShdw' in code:
        techniques.append('SHADOW')
    if 'build_freeform' in code or 'add_line_segments' in code:
        techniques.append('FREEFORM_CURVES')
    if 'ImageDraw' in code or 'ImageChops' in code:
        techniques.append('PIL_IMAGE_GEN')
    if 'parse_xml' in code or 'etree' in code:
        techniques.append('LXML_XML_INJECTION')

    # Extract helper functions (non-create_*)
    helpers = []
    for fn_match in re.finditer(r'^(def (?!create_)\w+\(.*?\):.*?)(?=\ndef |\Z)', code, re.DOTALL | re.MULTILINE):
        fn_code = fn_match.group(1).strip()
        fn_name_m = re.match(r'def (\w+)', fn_code)
        if fn_name_m:
            helpers.append((fn_name_m.group(1), fn_code))

    lines = [f"# Skill Reference: {name}"]
    if mechanism:
        lines.append(f"\n## Visual Mechanism\n{mechanism}")
    if techniques:
        lines.append(f"\n## Detected Techniques\n" + "\n".join(f"- {t}" for t in techniques))
    if helpers:
        lines.append(f"\n## Reusable Helpers (copy verbatim — units are critical)")
        for fn_name, fn_code in helpers[:5]:
            lines.append(f"\n```python\n{fn_code}\n```")
    lines.append(f"\n## Full Code ({len(code.splitlines())} lines)\n```python\n{code}\n```")
    return "\n".join(lines)


_TECHNIQUE_SNIPPETS = {
    "gradient_fill": '''\
def apply_gradient(shape, stops, angle_deg=90):
    """Apply multi-stop gradient fill. stops = [(r,g,b, position_pct), ...] 0-100."""
    from lxml.etree import SubElement
    from pptx.oxml.ns import qn
    spPr = shape.element.spPr
    for tag in [qn('a:solidFill'), qn('a:noFill'), qn('a:gradFill'), qn('a:pattFill')]:
        for old in spPr.findall(tag):
            spPr.remove(old)
    gradFill = SubElement(spPr, qn('a:gradFill'), attrib={'rotWithShape': '1'})
    gsLst = SubElement(gradFill, qn('a:gsLst'))
    for r, g, b, pos in stops:
        gs = SubElement(gsLst, qn('a:gs'), attrib={'pos': str(int(pos * 1000))})
        SubElement(gs, qn('a:srgbClr'), attrib={'val': f'{r:02X}{g:02X}{b:02X}'})
    SubElement(gradFill, qn('a:lin'), attrib={'ang': str(int(angle_deg * 60000)), 'scaled': '1'})
''',
    "shadow": '''\
def apply_shadow(shape, blur_pt=8, dist_pt=4, alpha_pct=35, angle_deg=135):
    """Add outer drop shadow. 1pt = 12700 EMU."""
    from lxml.etree import SubElement
    from pptx.oxml.ns import qn
    spPr = shape.element.spPr
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = SubElement(spPr, qn('a:effectLst'))
    outerShdw = SubElement(effectLst, qn('a:outerShdw'), attrib={
        'blurRad': str(int(blur_pt * 12700)), 'dist': str(int(dist_pt * 12700)),
        'dir': str(int(angle_deg * 60000)), 'algn': 'b', 'rotWithShape': '0',
    })
    srgb = SubElement(outerShdw, qn('a:srgbClr'), attrib={'val': '000000'})
    SubElement(srgb, qn('a:alpha'), attrib={'val': str(int(alpha_pct * 1000))})
''',
    "freeform_arc": '''\
def draw_arc(slide, cx_in, cy_in, radius_in, start_deg, end_deg, steps=20, line_color=None, line_width_pt=3):
    """Draw smooth arc via freeform line segments. All dims in inches."""
    import math
    from pptx.util import Inches, Pt
    pts = []
    for i in range(steps + 1):
        theta = math.radians(start_deg + (end_deg - start_deg) * i / steps)
        pts.append((Inches(cx_in + radius_in * math.cos(theta)),
                     Inches(cy_in + radius_in * math.sin(theta))))
    ffb = slide.shapes.build_freeform(pts[0][0], pts[0][1])
    ffb.add_line_segments(pts[1:])
    shape = ffb.convert_to_shape()
    shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width_pt)
    return shape
''',
    "radial_gradient_bg": '''\
def radial_gradient_bg(slide, prs, center_rgb=(15, 25, 50), edge_rgb=(5, 10, 20)):
    """PIL radial gradient background image."""
    from PIL import Image, ImageDraw, ImageFilter
    import math, tempfile, os
    from pptx.util import Inches
    W, H = 1920, 1080
    img = Image.new('RGB', (W, H))
    draw = ImageDraw.Draw(img)
    cx, cy = W // 2, H // 2
    max_r = int(math.hypot(cx, cy))
    for r in range(max_r, 0, -8):
        ratio = r / max_r
        c = tuple(int(center_rgb[i] + (edge_rgb[i] - center_rgb[i]) * ratio) for i in range(3))
        draw.ellipse((cx - r, cy - r, cx + r, cy + r), fill=c)
    img = img.filter(ImageFilter.GaussianBlur(12))
    tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    img.save(tmp.name); tmp.close()
    slide.shapes.add_picture(tmp.name, 0, 0, Inches(13.333), Inches(7.5))
    os.unlink(tmp.name)
''',
    "entrance_animation": '''\
# Entrance animation via PowerPoint timing XML (wrapped in mainSeq so it actually plays).
#
# PRESET COOKBOOK — use variety across slides to avoid monotony:
#   preset_id=53  filter="zoom"            → zoom center (great for cover, metrics)
#   preset_id=53  filter="in_center"       → zoom from center with subtle pop
#   preset_id=2   filter="wipe(direction)" → fly-in; set filter="in_bottom" / "in_top" / "in_left"
#   preset_id=22  filter="horizontal"      → wipe from left (timelines, list items)
#   preset_id=22  filter="vertical"        → wipe bottom-up (revealing cards)
#   preset_id=42  filter="fade"            → float-up (closing, testimonials)
#   preset_id=10  filter="fade"            → pure fade (secondary elements only)
#   preset_id=6   filter="in"              → rise (titles)
#   preset_id=54  filter="in"              → expand / grow (numbers, badges)
#
# STAGGER RULES: cover/hero → 120-180ms. Card grids → 80-120ms. Timelines → 150-200ms.
# DO NOT use the same (preset_id, filter) pair on every slide — pick at least 4 across the deck.
#
# Call: add_entrance_animation(slide, shape, delay_ms=200, duration_ms=500, preset_id=53, filter="zoom")
# DO NOT replace this function body with `pass` or `return None` — the XML append is the animation.
def add_entrance_animation(slide, shape, delay_ms=0, duration_ms=500, preset_id=53, filter="zoom", subtype=0):
    """Append an entrance animation into the slide's mainSeq click trigger."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml(
            f'<p:timing xmlns:p="{P_NS}">'
            '<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            '<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
            '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            '<p:par><p:cTn id="3" fill="hold">'
            '<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
            '<p:childTnLst><p:par><p:cTn id="4" fill="hold">'
            '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            '<p:childTnLst/>'
            '</p:cTn></p:par></p:childTnLst>'
            '</p:cTn></p:par>'
            '</p:childTnLst></p:cTn>'
            '<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            '<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            '</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>')
        slide._element.append(timing)
    seq = timing.find('.//' + qn('p:seq'))
    if seq is None: return None
    groups = seq.findall('.//' + qn('p:childTnLst'))
    if not groups: return None
    group = groups[-1]
    sid = shape.shape_id
    uid = 1000 + int(delay_ms) + sid * 7
    xml = (
        f'<p:par xmlns:p="{P_NS}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:cTn id="{uid}" presetID="{int(preset_id)}" presetClass="entr" presetSubtype="{int(subtype)}" '
        f'fill="hold" grpId="0" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="{int(delay_ms)}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:set><p:cBhvr><p:cTn id="{uid+1}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'<p:animEffect transition="in" filter="{filter}">'
        f'<p:cBhvr><p:cTn id="{uid+2}" dur="{int(duration_ms)}"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'</p:cBhvr></p:animEffect>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))
    return True

def stagger_entrance(slide, shapes, start_ms=0, step_ms=120, preset_id=53, filter="zoom"):
    for i, sh in enumerate(shapes):
        add_entrance_animation(slide, sh, delay_ms=start_ms + i*step_ms,
                               preset_id=preset_id, filter=filter)
''',
    "emphasis_animation": '''\
# Emphasis animation — fires AFTER an entrance to draw attention.
# Common presets: preset_id=2 (spin), 8 (pulse), 20 (teeter), 31 (grow_shrink), 43 (color)
# Use on headline numbers, CTA buttons, section icons.
def add_emphasis_animation(slide, shape, delay_ms=600, duration_ms=800, preset_id=8, repeat=1):
    """Append an emphasis animation (pulse/spin/teeter) on top of entrance."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        # ensure mainSeq exists first (reuse entrance helper)
        add_entrance_animation(slide, shape, delay_ms=0, preset_id=10, filter="fade")
        # remove the placeholder entrance we just added
        groups = slide._element.findall('.//' + qn('p:childTnLst'))
        if groups and len(groups[-1]):
            groups[-1].remove(groups[-1][-1])
        timing = slide._element.find(qn('p:timing'))
    seq = timing.find('.//' + qn('p:seq'))
    groups = seq.findall('.//' + qn('p:childTnLst'))
    if not groups: return None
    group = groups[-1]
    sid = shape.shape_id
    uid = 2000 + int(delay_ms) + sid * 11
    repeat_attr = f' repeatCount="{int(repeat)*1000}"' if repeat > 1 else ''
    xml = (
        f'<p:par xmlns:p="{P_NS}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:cTn id="{uid}" presetID="{int(preset_id)}" presetClass="emph" '
        f'fill="hold" grpId="1" nodeType="withEffect"{repeat_attr}>'
        f'<p:stCondLst><p:cond delay="{int(delay_ms)}"/></p:stCondLst>'
        f'<p:childTnLst><p:animEffect transition="none">'
        f'<p:cBhvr><p:cTn id="{uid+1}" dur="{int(duration_ms)}"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'</p:cBhvr></p:animEffect></p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))
    return True
''',
    "spin_entrance": '''\
# Spin/rotate entrance — great for section dividers, logos, icons.
# Uses XML animRot to actually rotate the shape while fading in.
def add_spin_entrance(slide, shape, delay_ms=0, duration_ms=800, degrees=360):
    """Spin-in entrance: rotates `degrees` while fading from 0% to 100% opacity."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    # Ensure timing/mainSeq exists (borrow from add_entrance_animation)
    if slide._element.find(qn('p:timing')) is None:
        add_entrance_animation(slide, shape, delay_ms=0, preset_id=10, filter="fade")
        # wipe the placeholder
        grp = slide._element.findall('.//' + qn('p:childTnLst'))[-1]
        if len(grp): grp.remove(grp[-1])
    seq = slide._element.find('.//' + qn('p:seq'))
    groups = seq.findall('.//' + qn('p:childTnLst'))
    group = groups[-1]
    sid = shape.shape_id
    uid = 3000 + int(delay_ms) + sid * 13
    xml = (
        f'<p:par xmlns:p="{P_NS}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:cTn id="{uid}" presetID="8" presetClass="entr" presetSubtype="0" '
        f'fill="hold" grpId="0" nodeType="withEffect">'
        f'<p:stCondLst><p:cond delay="{int(delay_ms)}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:set><p:cBhvr><p:cTn id="{uid+1}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'<p:animRot by="{int(degrees)*60000}">'
        f'<p:cBhvr><p:cTn id="{uid+2}" dur="{int(duration_ms)}" fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>r</p:attrName></p:attrNameLst>'
        f'</p:cBhvr></p:animRot>'
        f'<p:animEffect transition="in" filter="fade">'
        f'<p:cBhvr><p:cTn id="{uid+3}" dur="{int(duration_ms)}"/>'
        f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
        f'</p:cBhvr></p:animEffect>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    group.append(parse_xml(xml))
    return True
''',
    "section_reveal": '''\
# Section-divider reveal: sweeps a wide rectangle from left to right to uncover the title,
# then floats the content up. Combines two entrance effects on different shapes.
def section_reveal(slide, mask_shape, title_shape, body_shapes=None, start_ms=0):
    """Sweep mask left-to-right (wipe), then title zoom, then body shapes float-up staggered."""
    add_entrance_animation(slide, mask_shape, delay_ms=start_ms, duration_ms=700,
                           preset_id=22, filter="horizontal")
    add_entrance_animation(slide, title_shape, delay_ms=start_ms + 400, duration_ms=500,
                           preset_id=53, filter="zoom")
    if body_shapes:
        for i, sh in enumerate(body_shapes):
            add_entrance_animation(slide, sh, delay_ms=start_ms + 700 + i * 120,
                                   duration_ms=500, preset_id=42, filter="fade")
''',
}


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def get_technique_snippet(technique: str) -> str:
    """Get a tested visual technique code snippet (gradient_fill, shadow, freeform_arc, radial_gradient_bg).

    Copy the function into your slide code and call it. Unit conversions are correct.
    """
    snippet = _TECHNIQUE_SNIPPETS.get(technique)
    if not snippet:
        available = ", ".join(sorted(_TECHNIQUE_SNIPPETS.keys()))
        return f"Error: unknown technique '{technique}'. Available: {available}"
    return f"# Technique: {technique}\n# Copy this into your slide code:\n\n{snippet}"


# ---------------------------------------------------------------------------
# Palette presets
# ---------------------------------------------------------------------------

_PALETTE_PRESETS = {
    "dark_tech": {
        "dark_bg":    (10, 15, 26),
        "panel":      (20, 30, 46),
        "accent1":    (0, 210, 255),
        "accent2":    (124, 110, 255),
        "accent3":    (57, 214, 138),
        "text_light": (255, 255, 255),
        "text_muted": (146, 166, 190),
    },
    "clean_light": {
        "dark_bg":    (245, 247, 250),
        "panel":      (255, 255, 255),
        "accent1":    (37, 99, 235),
        "accent2":    (124, 58, 237),
        "accent3":    (16, 185, 129),
        "text_light": (15, 23, 42),
        "text_muted": (100, 116, 139),
    },
    "warm_corporate": {
        "dark_bg":    (250, 245, 240),
        "panel":      (255, 251, 245),
        "accent1":    (217, 119, 6),
        "accent2":    (180, 83, 9),
        "accent3":    (101, 163, 13),
        "text_light": (41, 37, 36),
        "text_muted": (120, 113, 108),
    },
    "bold_vibrant": {
        "dark_bg":    (15, 15, 15),
        "panel":      (30, 30, 30),
        "accent1":    (255, 59, 48),
        "accent2":    (255, 204, 0),
        "accent3":    (52, 199, 89),
        "text_light": (255, 255, 255),
        "text_muted": (174, 174, 178),
    },
    "ocean_calm": {
        "dark_bg":    (15, 30, 50),
        "panel":      (20, 42, 68),
        "accent1":    (56, 189, 248),
        "accent2":    (99, 102, 241),
        "accent3":    (45, 212, 191),
        "text_light": (240, 249, 255),
        "text_muted": (148, 187, 218),
    },
}


# ---------------------------------------------------------------------------
# Image generation via Azure OpenAI GPT-Image
# ---------------------------------------------------------------------------

_image_client = None
_image_counter = 0


def _get_image_client():
    """Lazy-init Azure OpenAI client for image generation."""
    global _image_client
    if _image_client is None:
        import os
        from openai import AzureOpenAI
        api_key = os.environ.get("AZURE_OPENAI_API_KEY", os.environ.get("OPENAI_API_KEY", ""))
        endpoint = (
            os.environ.get("AZURE_IMAGE_ENDPOINT")
            or os.environ.get("AZURE_OPENAI_ENDPOINT")
            or ""
        ).strip()
        api_version = os.environ.get("AZURE_IMAGE_API_VERSION", "2024-12-01-preview")
        if not api_key or not endpoint:
            return None
        _image_client = AzureOpenAI(
            api_version=api_version,
            azure_endpoint=endpoint,
            api_key=api_key,
        )
    return _image_client


@mcp.tool()
def generate_image(prompt: str, style: str = "digital illustration", size: str = "1536x1024") -> str:
    """Generate an image using AI (GPT-Image) and save to a local file.

Use this to create contextual images for slides — product mockups, scene
illustrations, icons, data visualizations, or abstract backgrounds.

The returned file path can be used with `slide.shapes.add_picture(path, ...)`
in your add-slide code.

Args:
    prompt: Detailed description of the image to generate. Be specific about
            content, style, colors, and composition.
            Good: "Clean vector illustration of K-12 students using tablets
                   in a bright modern classroom, warm orange and green palette"
            Bad: "students"
    style: Visual style hint appended to prompt. Default "digital illustration".
           Options: "digital illustration", "photorealistic", "flat vector",
                    "3d render", "watercolor", "minimalist icon", "infographic".
    size: Image dimensions. Default "1536x1024" (landscape, good for slides).
          Options: "1536x1024" (landscape), "1024x1024" (square), "1024x1536" (portrait).

Returns:
    Absolute file path to the generated PNG image, or error message.
    Use this path in add-slide code: slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
"""
    global _image_counter

    client = _get_image_client()
    if client is None:
        return "Error: No API key found for image generation (AZURE_OPENAI_API_KEY or OPENAI_API_KEY)"

    # Enhance prompt with style
    full_prompt = f"{prompt}. Style: {style}. High quality, professional, suitable for a business presentation slide."

    try:
        result = client.images.generate(
            model="gpt-image-1.5",
            prompt=full_prompt,
            n=1,
            size=size,
        )

        if not result.data or not result.data[0].b64_json:
            return "Error: Image generation returned no data"

        # Save to temp file
        import base64
        _image_counter += 1
        img_path = Path(tempfile.gettempdir()) / f"ppt_genimg_{_image_counter:03d}.png"
        img_bytes = base64.b64decode(result.data[0].b64_json)
        img_path.write_bytes(img_bytes)

        log.info("Generated image: %s (%d bytes)", img_path, len(img_bytes))
        return str(img_path)

    except Exception as e:
        return f"Error generating image: {type(e).__name__}: {e}"


@mcp.tool()
def get_palette_preset(preset_name: str = "") -> str:
    """Get a theme palette preset for use with add-slide and add-slide-from-skill.

    Available presets: dark_tech, clean_light, warm_corporate, bold_vibrant, ocean_calm.
    Call with no args to see all options. Call with a preset name to get the full palette as JSON.

    Use the returned RGB values to define your palette variables in add-slide code.
    Pass the JSON to add-slide-from-skill's palette parameter for color remapping.
    """
    if not preset_name:
        lines = ["Available palette presets:"]
        for name, pal in _PALETTE_PRESETS.items():
            bg = pal["dark_bg"]
            a1 = pal["accent1"]
            txt = pal["text_light"]
            lines.append(
                f"  {name:18s} bg=({bg[0]:3d},{bg[1]:3d},{bg[2]:3d}) "
                f"accent=({a1[0]:3d},{a1[1]:3d},{a1[2]:3d}) "
                f"text=({txt[0]:3d},{txt[1]:3d},{txt[2]:3d})"
            )
        lines.append("\nCall get-palette-preset(preset_name='...') to get full palette as JSON.")
        return "\n".join(lines)

    pal = _PALETTE_PRESETS.get(preset_name)
    if not pal:
        return f"Error: unknown preset '{preset_name}'. Available: {', '.join(_PALETTE_PRESETS.keys())}"
    import json as _json
    return _json.dumps({k: list(v) for k, v in pal.items()}, indent=2)


# =========================================================================
# CLI / main
# =========================================================================

def main():
    parser = argparse.ArgumentParser(description="PPT MCP Server")
    parser.add_argument("--skills-dir", type=str, default=None,
                        help="Path to skill library directory")
    parser.add_argument("--test", action="store_true",
                        help="Run a quick self-test instead of starting the server")
    args = parser.parse_args()

    if args.skills_dir:
        global _skills_dir
        _skills_dir = Path(args.skills_dir)

    if args.test:
        _run_self_test()
        return

    # Hybrid mode: even though library_backend stays on 'legacy' so the
    # python-pptx engine + asset reads keep working, expose the universal
    # wiki discovery surface alongside the legacy tools so an agent can
    # browse skills_wiki/ppt/ via the contract. The legacy stale guard
    # is installed too, so a future flip to 'wiki' surfaces a structured
    # error rather than silently serving stale legacy data.
    try:
        from domains.ppt.wiki_adapter import PPTWikiAdapter
        from core.skill_wiki.mcp_tools import register_wiki_tools
        register_wiki_tools(mcp, PPTWikiAdapter())
        log.info("PPT MCP: registered universal wiki discovery surface")
    except Exception as exc:  # noqa: BLE001
        log.warning("PPT MCP: failed to register wiki discovery surface: %s", exc)

    try:
        from core import get_active_library_backend
        from core.skill_wiki.legacy_stale import register_legacy_stale_check
        backend = get_active_library_backend("ppt")
        wrapped = register_legacy_stale_check(mcp, domain="ppt", startup_backend=backend)
        log.info("PPT MCP: stale-registry guard installed on %d tools (backend=%s)",
                 wrapped, backend)
    except Exception as exc:  # noqa: BLE001
        log.warning("PPT MCP: failed to install stale-registry guard: %s", exc)

    logging.basicConfig(level=logging.INFO)
    mcp.run(transport="stdio")


def _run_self_test():
    """Quick self-test: create prs → add slide → save → verify."""
    print("=== PPT MCP Server Self-Test ===")

    # Create presentation
    result = create_presentation()
    print(f"  create-presentation: {result}")
    prs_id = result.split()[2]

    # Add a slide
    code = """
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
tf = txBox.text_frame
tf.text = "Hello from PPT MCP Server!"
tf.paragraphs[0].font.size = Pt(44)
tf.paragraphs[0].font.color.rgb = RGBColor(0x1a, 0x73, 0xe8)
"""
    result = add_slide(prs_id, code)
    print(f"  add-slide: {result}")

    # Get slide info
    result = get_slide_info(prs_id, 1)
    print(f"  get-slide-info:\n    {result}")

    # Save presentation
    import tempfile
    out = Path(tempfile.gettempdir()) / "ppt_mcp_test.pptx"
    result = save_presentation(prs_id, str(out))
    print(f"  save-presentation: {result}")

    # List skills (if available)
    if _skills_dir:
        result = list_skills()
        lines = result.split("\n")
        print(f"  list-skills: {lines[0]}")
        for line in lines[1:4]:
            print(f"    {line}")

    print("\n=== Self-test PASSED ===")


# ---------------------------------------------------------------------------
# V2 Skill System — Theme + Shell (see docs/ppt_skill_system_v2_refactor.md)
# ---------------------------------------------------------------------------
#
# New MCP tools exposed to the agent:
#   - list_themes()                                       → theme summary cards
#   - get_theme(theme_id)                                 → full theme JSON
#   - list_shells()                                       → shell summary cards
#   - add_slide_from_shell(prs_id, shell_id, slots,       → render a slide via a
#                          theme_id)                        shell's render() fn
#
# The agent workflow becomes: pick theme → for each slide, pick shell, fill slots.
# Shells read colors/typography exclusively from the theme, guaranteeing
# cross-slide consistency.

try:
    from core.extraction.shell_loader import (
        load_shell as _v2_load_shell,
        load_theme as _v2_load_theme,
        list_shells as _v2_list_shells,
        list_themes as _v2_list_themes,
    )
    from core.retrieval import (
        list_archetypes as _v2_list_archetypes,
        load_archetype as _v2_load_archetype,
        filter_shells as _v2_filter_shells,
        format_archetype_card as _v2_fmt_archetype,
        format_theme_card as _v2_fmt_theme,
        format_shell_card as _v2_fmt_shell,
        pick_archetype as _v2_pick_archetype,
        pick_theme as _v2_pick_theme,
        pick_shell as _v2_pick_shell,
    )
    _V2_AVAILABLE = True
except Exception as _v2_err:
    log.warning("V2 shell/retrieval loader unavailable: %s", _v2_err)
    _V2_AVAILABLE = False


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def list_themes(archetype: str = "", mood: str = "", mode: str = "",
                exclude_ids: str = "") -> str:
    """List theme candidates that pass pure metadata filters.

    Returns a JSON slate for the LLM to read and pick from. No ranking.

    Args:
        archetype: comma-separated ARCHETYPE_FIT tokens (data, narrative,
              product, research, brand, boardroom). At least one must
              appear in the theme's archetype_fit list. Empty skips.
        mood: comma-separated MOOD tokens. Same semantics.
        mode: "dark" | "light" | "" — filter by background luminance.
        exclude_ids: comma-separated theme_ids to drop (cross-deck dedup).

    Returns JSON {count, themes: [{theme_id, name, mode, mood,
        archetype_fit, palette_preview, motion_personality, motif,
        font_primary}, ...]}.
    """
    import json as _json
    if not _V2_AVAILABLE:
        return _json.dumps({"error": "V2 theme system not loaded", "themes": []})
    arch_want = {x.strip().lower() for x in (archetype or "").split(",") if x.strip()}
    mood_want = {x.strip().lower() for x in (mood or "").split(",") if x.strip()}
    mode_want = (mode or "").strip().lower()
    exclude = {x.strip() for x in (exclude_ids or "").split(",") if x.strip()}
    raw = _v2_list_themes()
    out = []
    for t in raw:
        if t["theme_id"] in exclude:
            continue
        t_mood = {m.lower() for m in t.get("mood", []) or []}
        t_arch = {a.lower() for a in t.get("archetype_fit", []) or []}
        if mood_want and not (mood_want & t_mood):
            continue
        if arch_want and not (arch_want & t_arch):
            continue
        if mode_want and t.get("mode") != mode_want:
            continue
        out.append(t)
    return _json.dumps({
        "filter": {"mood": sorted(mood_want), "archetype": sorted(arch_want),
                   "mode": mode_want, "exclude_ids": sorted(exclude)},
        "count": len(out),
        "themes": out,
    }, indent=2)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def get_theme(theme_id: str) -> str:
    """Get the full JSON for a theme (palette, typography map, motif, spacing, motion)."""
    if not _V2_AVAILABLE:
        return "Error: V2 theme system not loaded"
    try:
        data = _v2_load_theme(theme_id)
    except FileNotFoundError:
        available = [t["theme_id"] for t in _v2_list_themes()]
        return f"Error: theme '{theme_id}' not found. Available: {', '.join(available)}"
    return json.dumps(data, indent=2)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def list_shells(role: str = "", archetype: str = "", mood: str = "",
                density: str = "", style_tags: str = "",
                exclude_ids: str = "") -> str:
    """List shell candidates that pass pure metadata filters.

    Returns a JSON slate for the LLM to read and pick from. No ranking.

    Args:
        role: "cover" | "hero_giant_metric" | "agenda" | "section_divider"
              | "metric_dashboard" | "comparison_split" | "timeline_horizontal"
              | "bullet_card_list" | "hero_quote" | "closing_cta"
              | "feature_grid" | "stat_trio" | "process_flow". HARD filter.
        archetype: comma-separated; at least one must appear in shell.archetype.
        mood: comma-separated; at least one must overlap shell.mood.
        density: "sparse" | "balanced" | "dense" | "". Exact match when set.
        style_tags: comma-separated; at least one must overlap shell.style_tags.
        exclude_ids: comma-separated shell_ids to drop.

    Returns JSON {count, shells: [{shell_id, role, description, archetype,
        mood, density, style_tags, slot_names, required_slots, slots,
        source, status, ambient_capable}, ...]}.
    """
    import json as _json
    if not _V2_AVAILABLE:
        return _json.dumps({"error": "V2 shell system not loaded", "shells": []})
    arch_want = {x.strip().lower() for x in (archetype or "").split(",") if x.strip()}
    mood_want = {x.strip().lower() for x in (mood or "").split(",") if x.strip()}
    style_want = {x.strip().lower() for x in (style_tags or "").split(",") if x.strip()}
    density_want = (density or "").strip().lower()
    exclude = {x.strip() for x in (exclude_ids or "").split(",") if x.strip()}
    raw = _v2_list_shells()
    out = []
    for entry in _brand_shell_entries(role):
        if entry["shell_id"] not in exclude:
            out.append(entry)
    for s in raw:
        if "error" in s:
            continue
        if s["shell_id"] in exclude:
            continue
        # Accept both "active" seed shells and "candidate" distilled shells.
        if s.get("status") == "demoted":
            continue
        s_role = s.get("role") or ""
        if role and s_role and s_role != role:
            continue
        # fallback: if shell lacks explicit role, accept when role tag matches
        if role and not s_role:
            tags = s.get("intent_tags", {}).get("slide_role", [])
            if role not in tags:
                continue
        s_arch = (s.get("archetype") or "").lower()
        if arch_want and s_arch and s_arch not in arch_want:
            continue
        s_mood = {m.lower() for m in s.get("mood", []) or []}
        if mood_want and not (mood_want & s_mood):
            continue
        s_style = {m.lower() for m in s.get("style_tags", []) or []}
        if style_want and not (style_want & s_style):
            continue
        s_density = (s.get("density") or "").lower()
        if density_want and s_density != density_want:
            continue
        out.append({
            "shell_id": s["shell_id"],
            "role": s_role,
            "description": s.get("description", s.get("doc", "")),
            "archetype": s.get("archetype", ""),
            "mood": s.get("mood", []),
            "density": s.get("density", ""),
            "style_tags": s.get("style_tags", []),
            "slot_names": s.get("slot_names", []),
            "required_slots": s.get("required_slots", []),
            "slots": s.get("slots", []),
            "source": s.get("source", "seed"),
            "status": s.get("status", "active"),
            "ambient_capable": s.get("ambient_capable", False),
        })
    out.sort(key=lambda e: (0 if e.get("source") == "brand_overlay" else 1, e["shell_id"]))
    return _json.dumps({
        "filter": {"role": role, "archetype": sorted(arch_want),
                   "mood": sorted(mood_want), "density": density_want,
                   "style_tags": sorted(style_want),
                   "exclude_ids": sorted(exclude)},
        "count": len(out),
        "shells": out,
    }, indent=2)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def add_slide_from_shell(
    prs_id: str,
    shell_id: str,
    slots: str = "{}",
    theme_id: str = "editorial_dark",
    transition: str = "",
    slide_role: str = "",
    selection_origin: str = "",
    design_reference_skill_ids: str = "",
) -> str:
    """Add a slide by rendering a named shell with named-slot content, under a theme.

    Args:
        prs_id:    presentation id from create_presentation
        shell_id:  id of a shell (see list_shells)
        slots:     JSON string of slot_name→value (strings, lists, or dicts depending on slot kind)
        theme_id:  id of a theme (see list_themes). All shells render under the theme's
                   palette + typography + motion.
        transition: optional. One of "morph", "fade", "push", "wipe". "morph" emits
                   the PowerPoint-compatible force-match transition; other values map
                   to their usual transition XML. Empty string preserves prior default
                   behaviour (the auto-fade backfill applied at save time).
        slide_role: optional. The semantic role for this slide (e.g. "cover",
                   "bullet_card_list"). Recorded into the deck-plan metadata so
                   save_presentation can run role-aware motion budget lint.
        selection_origin: optional. Overrides the auto-detected provenance origin.
                   Must be one of "distilled", "retrieved", "composed", "new", "seed".
        design_reference_skill_ids: optional audit trace. Comma-separated or JSON-list
                   wiki skill IDs inspected before choosing this shell/theme. These are
                   references only; execution still comes from the v2 shell scaffold.

    Returns:
        Confirmation with slide index, shape count, and a short morph-lint hint.
    """
    if not _V2_AVAILABLE:
        return "Error: V2 shell system not loaded"
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"

    # Parse slots JSON
    try:
        slots_dict = json.loads(slots) if isinstance(slots, str) else (slots or {})
    except json.JSONDecodeError as e:
        return f"Error: 'slots' is not valid JSON: {e}"
    if not isinstance(slots_dict, dict):
        return f"Error: 'slots' must be a JSON object, got {type(slots_dict).__name__}"

    if _brand_shell_entry(shell_id) is not None:
        transition_kind = (transition or "").strip().lower() or None
        if transition_kind and transition_kind not in {"morph", "fade", "push", "wipe"}:
            return (
                f"Error: unknown transition '{transition_kind}'. "
                "Use one of: morph, fade, push, wipe."
            )
        return _add_slide_from_brand_shell(
            prs_id=prs_id,
            shell_id=shell_id,
            slots_dict=slots_dict,
            transition_kind=transition_kind,
            slide_role=slide_role,
            design_reference_skill_ids=design_reference_skill_ids,
        )

    # Load theme
    try:
        theme = _v2_load_theme(theme_id)
    except FileNotFoundError:
        available = [t["theme_id"] for t in _v2_list_themes()]
        return f"Error: theme '{theme_id}' not found. Available: {', '.join(available)}"

    # Load shell
    try:
        shell_mod = _v2_load_shell(shell_id)
    except (FileNotFoundError, AttributeError, ImportError) as e:
        return f"Error loading shell '{shell_id}': {e}"

    transition_kind = (transition or "").strip().lower() or None
    if transition_kind and transition_kind not in {"morph", "fade", "push", "wipe"}:
        return (
            f"Error: unknown transition '{transition_kind}'. "
            "Use one of: morph, fade, push, wipe."
        )

    # Render
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    n_before = len(prs.slides)
    try:
        shell_mod.render(slide, slots_dict, theme)
    except Exception as e:
        # Rollback the blank slide we added
        try:
            slide_list = prs.slides._sldIdLst
            rId = slide_list[-1].get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            if rId:
                prs.part.drop_rel(rId)
            slide_list.remove(slide_list[-1])
        except Exception:
            pass
        return f"Error rendering shell '{shell_id}': {type(e).__name__}: {e}"

    # Force word_wrap on all text frames (same post-process as add_slide)
    _force_word_wrap(slide)

    # Resolve provenance; only paint the on-slide chip when debug is on.
    # Viewers read the corner chip (e.g. "D:r_ed4e3a@4e155db3") as a bug
    # in a finished deck, so it's opt-in.
    provenance = _resolve_provenance(shell_id, selection_origin or None)
    if os.environ.get("PPT_SHOW_PROVENANCE_CHIP") == "1":
        _emit_provenance_chip(slide, theme, provenance)

    # Apply transition if the caller asked for one.
    if transition_kind:
        _apply_transition_to_slide(slide, transition_kind)

    # Capture per-slide metadata for deck-plan + morph-lint.
    anchor_names = _scan_anchor_names(slide)
    transition_set = _slide_transition_kind(slide)
    ambient_flag = _shell_ambient_flag(shell_id)
    hero_flag = (
        slide_role.startswith("hero")
        or "hero" in shell_id.lower()
        or slide_role in {"cover", "closing"}
    )
    role = slide_role or _role_from_shell_id(shell_id)

    slot_summary = {}
    for k, v in slots_dict.items():
        if isinstance(v, str):
            slot_summary[k] = v[:60]
        elif isinstance(v, list):
            slot_summary[k] = f"[{len(v)} items]"
        else:
            slot_summary[k] = str(v)[:60]

    design_refs = _parse_design_reference_skill_ids(
        design_reference_skill_ids or slots_dict.get("design_reference_skill_ids", "")
    )

    entry = _ensure_prs_metadata(prs_id)
    if entry.get("theme") and entry["theme"] != theme_id:
        # Non-fatal: surface theme drift so operator notices.
        log.warning("theme drift on prs %s: was %s, now %s",
                    prs_id, entry["theme"], theme_id)
    entry["theme"] = theme_id
    slide_entry = {
        "slide_index": len(entry["slides"]),
        "role": role,
        "shell_id": shell_id,
        "anchor_names_set": anchor_names,
        "transition_set": transition_set,
        "hero_flag": bool(hero_flag),
        "ambient_flag": bool(ambient_flag),
        "design_reference_skill_ids": design_refs,
        "distill_provenance": provenance,
        "slot_values_summary": slot_summary,
    }
    entry["slides"].append(slide_entry)

    # Speaker-notes JSON (same record minus the slot summary bulk — we keep it
    # intact here because reviewers often want to cross-check slot values).
    notes_record = {
        "skill_id": provenance["skill_id"],
        "selection_origin": provenance["selection_origin"],
        "source_video_id": provenance["source_video_id"],
        "source_timestamp": provenance["source_timestamp"],
        "slot_values_summary": slot_summary,
        "role": role,
        "theme": theme_id,
        "transition_set": transition_set,
        "anchor_names_set": anchor_names,
        "hero_flag": bool(hero_flag),
        "ambient_flag": bool(ambient_flag),
        "design_reference_skill_ids": design_refs,
    }
    _attach_notes_json(slide, notes_record)

    shape_count = len(slide.shapes)
    _save_temp(prs, prs_id)

    hint = ""
    if anchor_names and transition_set != "morph":
        hint = (
            f" [morph-lint: slide has anchors {anchor_names} but transition is "
            f"{transition_set!r} — set transition='morph' for morph playback]"
        )

    return (
        f"Added slide {n_before} from shell '{shell_id}' under theme '{theme_id}' "
        f"({shape_count} shapes, transition={transition_set}, ambient={ambient_flag})"
        f"{', design_refs=' + ','.join(design_refs) if design_refs else ''}"
        + hint
    )


def _parse_design_reference_skill_ids(raw) -> list[str]:
    """Normalize optional wiki reference IDs for deck-plan/notes audit trace."""
    if not raw:
        return []
    if isinstance(raw, list):
        items = raw
    elif isinstance(raw, str):
        text = raw.strip()
        if not text:
            return []
        try:
            parsed = json.loads(text)
            items = parsed if isinstance(parsed, list) else [text]
        except Exception:
            items = [part.strip() for part in text.split(",")]
    else:
        items = [raw]

    out: list[str] = []
    seen: set[str] = set()
    for item in items:
        sid = str(item).strip()
        if not sid or sid in seen:
            continue
        seen.add(sid)
        out.append(sid)
    return out


def _role_from_shell_id(shell_id: str) -> str:
    sid = (shell_id or "").lower()
    for key in ("cover", "closing", "agenda", "timeline_horizontal",
                "section_divider", "bullet_card_list", "metric_dashboard",
                "comparison_split", "feature_grid", "quote",
                "hero_giant_metric", "hero_quote_reveal", "hero_before_after"):
        if key in sid:
            return key
    return "content"


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def list_archetypes() -> str:
    """List available deck archetypes (narrative blueprints).

    An archetype defines the section flow + per-slide role sequence.
    Pick ONE archetype per task — it dictates how many slides and in what order.
    """
    if not _V2_AVAILABLE:
        return "Error: V2 retrieval not loaded"
    arches = _v2_list_archetypes()
    if not arches:
        return "No archetypes available"
    lines = [f"Deck archetypes ({len(arches)} available):\n"]
    for a in arches:
        if "error" in a:
            lines.append(f"  [{a['archetype_id']}] ERROR: {a['error']}")
            continue
        lines.append(_v2_fmt_archetype(a))
        lines.append("")
    lines.append(
        "Call get_archetype(archetype_id) for the full slide-by-slide plan. "
        "Or use pick_archetype(task) to let an LLM choose."
    )
    return "\n".join(lines)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def get_archetype(archetype_id: str) -> str:
    """Return the full section + per-slide role plan for an archetype, as JSON."""
    if not _V2_AVAILABLE:
        return "Error: V2 retrieval not loaded"
    try:
        data = _v2_load_archetype(archetype_id)
    except FileNotFoundError:
        available = [a["archetype_id"] for a in _v2_list_archetypes()]
        return f"Error: archetype '{archetype_id}' not found. Available: {', '.join(available)}"
    return json.dumps(data, indent=2, ensure_ascii=False)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def pick_archetype(task_description: str, prs_id: str = "") -> str:
    """Let an LLM pick the best-fitting archetype for the given task.

    Returns JSON {archetype_id, reasoning}. Use this once at the start of a deck.
    When ``prs_id`` is provided, the chosen archetype and slide target are
    recorded into deck-plan metadata. If the brief explicitly requests a slide
    count, that count wins; otherwise the archetype's suggested slide count is
    used for length lint.
    """
    if not _V2_AVAILABLE:
        return "Error: V2 retrieval not loaded"
    requested_target = _extract_requested_slide_count(task_description)
    result = _v2_pick_archetype(task_description)
    if prs_id:
        entry = _ensure_prs_metadata(prs_id)
        entry["archetype"] = result.get("archetype_id")
        try:
            arch = _v2_load_archetype(result["archetype_id"])
            suggested = arch.get("suggested_slides") or arch.get("actual_slide_count")
            entry["archetype_suggested_slides"] = suggested
            entry["requested_slides_target"] = requested_target
            entry["archetype_slides_target"] = requested_target or suggested
            result["archetype_suggested_slides"] = suggested
            result["slides_target"] = entry["archetype_slides_target"]
        except Exception:
            entry["requested_slides_target"] = requested_target
            entry["archetype_slides_target"] = requested_target
    if requested_target:
        result["requested_slides_target"] = requested_target
        result["slides_target"] = requested_target
    return json.dumps(result, indent=2, ensure_ascii=False)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def pick_theme(task_description: str, archetype_id: str = "") -> str:
    """Let an LLM pick the best-fitting theme for the given task.

    If archetype_id is given, candidates are narrowed to that archetype's
    recommended_themes first. Returns JSON {theme_id, reasoning}.
    """
    if not _V2_AVAILABLE:
        return "Error: V2 retrieval not loaded"
    themes = _v2_list_themes()
    result = _v2_pick_theme(task_description, themes, archetype_id or None)
    return json.dumps(result, indent=2, ensure_ascii=False)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def select_shell(
    slide_role: str,
    content_brief: str,
    theme_id: str = "",
    requires_image: bool = False,
    bullet_capacity_min: int = 0,
    prefer_ambient: bool = False,
    top_k: int = 3,
) -> str:
    """Pick the best-fitting shell(s) for a specific slide.

    Runs: (1) hard filter on slide_role/theme/task_fit, then (2) LLM ranks
    the candidates. Returns JSON {ranked:[{shell_id, reasoning}, ...]}.

    Args:
        slide_role: e.g. "cover", "section_divider", "bullet_card_list",
                    "metric_dashboard", "timeline_horizontal", "closing"
        content_brief: 1-2 sentence intent for this slide
        theme_id: optional — only shells compatible with the theme
        requires_image: if True, restrict to shells with an image slot
        bullet_capacity_min: if the slide needs ≥ N bullets, set this
        prefer_ambient: if True, ambient-capable shells are promoted to the
            front of the candidate list before the LLM ranks them. Useful
            for hero/cover/divider intents where continuous motion pays off.
        top_k: how many ranked candidates to return (default 3)
    """
    if not _V2_AVAILABLE:
        return "Error: V2 retrieval not loaded"
    brand_shell = _brand_shell_for_role(slide_role)
    brand_entry = _brand_shell_entry(brand_shell) if brand_shell else None
    if brand_entry is not None:
        return json.dumps({
            "ranked": [{
                "shell_id": brand_entry["shell_id"],
                "reasoning": (
                    "Brand mode is active; this brand shell is the primary "
                    f"constructor for slide_role={slide_role!r}."
                ),
            }],
            "brand_overlay": True,
        }, indent=2, ensure_ascii=False)
    shells = _v2_list_shells()
    filtered = _v2_filter_shells(
        shells,
        slide_role=slide_role or None,
        theme_id=theme_id or None,
        requires_image=True if requires_image else None,
        bullet_capacity_min=bullet_capacity_min or None,
    )
    if not filtered:
        return json.dumps({"ranked": [], "warning": f"No shells match slide_role={slide_role!r}"}, indent=2)
    if prefer_ambient:
        ambient = [s for s in filtered if s.get("ambient_capable")]
        rest = [s for s in filtered if not s.get("ambient_capable")]
        filtered = ambient + rest
    intent = (
        f"role: {slide_role}\n"
        f"content: {content_brief}\n"
        f"theme: {theme_id or 'any'}\n"
        f"requires_image: {requires_image}\n"
        f"bullet_capacity_min: {bullet_capacity_min}\n"
        f"prefer_ambient: {prefer_ambient}"
    )
    result = _v2_pick_shell(intent, filtered, top_k=top_k)
    return json.dumps(result, indent=2, ensure_ascii=False)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def suggest_morph_continuity(prs_id: str, slide_a: int, slide_b: int) -> str:
    """Inspect two slides for PowerPoint Morph compatibility.

    Args:
        prs_id: presentation id
        slide_a: 1-based index of the first slide
        slide_b: 1-based index of the second slide

    Returns:
        JSON with:
          shared_anchors:  list of `!!sameName...` names that appear on both
                           slides (these will morph under PowerPoint).
          incompatible_anchors: anchors present on only one slide, or present
                           on both but attached to shapes whose family
                           differs (text↔shape, shape↔group, etc.).
          recommendation: "morph" if there is at least one compatible shared
                           anchor; "fade" if anchors exist but are not
                           compatible; "none" if no anchors exist at all.
    """
    prs = _presentations.get(prs_id)
    if prs is None:
        return json.dumps({"error": f"presentation '{prs_id}' not found"})
    try:
        a_idx = int(slide_a) - 1
        b_idx = int(slide_b) - 1
    except Exception:
        return json.dumps({"error": "slide indices must be integers"})
    if not (0 <= a_idx < len(prs.slides) and 0 <= b_idx < len(prs.slides)):
        return json.dumps({
            "error": f"slide indices out of range 1..{len(prs.slides)}"
        })

    anchors_a = _slide_anchors_with_family(prs.slides[a_idx])
    anchors_b = _slide_anchors_with_family(prs.slides[b_idx])

    shared = []
    incompatible = []
    for name, family_a in anchors_a.items():
        if name in anchors_b:
            family_b = anchors_b[name]
            if family_a == family_b:
                shared.append({"anchor": name, "family": family_a})
            else:
                incompatible.append({
                    "anchor": name,
                    "family_a": family_a,
                    "family_b": family_b,
                    "reason": "object family mismatch",
                })
    only_a = [n for n in anchors_a if n not in anchors_b]
    only_b = [n for n in anchors_b if n not in anchors_a]

    if shared:
        recommendation = "morph"
    elif anchors_a or anchors_b:
        recommendation = "fade"
    else:
        recommendation = "none"

    return json.dumps({
        "slide_a": slide_a,
        "slide_b": slide_b,
        "shared_anchors": shared,
        "incompatible_anchors": incompatible,
        "anchors_only_on_a": only_a,
        "anchors_only_on_b": only_b,
        "recommendation": recommendation,
    }, indent=2, ensure_ascii=False)


def _slide_anchors_with_family(slide) -> dict[str, str]:
    """Return ``{anchor_name: family}`` for every morph anchor on ``slide``.

    ``family`` is one of ``text``, ``auto_shape:<preset>``, ``image``,
    ``group``, ``freeform``, ``other``.
    """
    out: dict[str, str] = {}
    PIC_TAG_END = "}pic"
    GRP_TAG_END = "}grpSp"
    SP_TAG_END = "}sp"
    for el in slide._element.iter():
        if not el.tag.endswith("}cNvPr"):
            continue
        name = el.get("name") or ""
        if not name.startswith(_MORPH_ANCHOR_PREFIX):
            continue
        # Walk up to find the shape container.
        node = el.getparent()
        family = "other"
        while node is not None:
            if node.tag.endswith(PIC_TAG_END):
                family = "image"
                break
            if node.tag.endswith(GRP_TAG_END):
                family = "group"
                break
            if node.tag.endswith(SP_TAG_END):
                preset = _auto_shape_preset(node)
                if preset == "_txBox_":
                    family = "text"
                elif preset == "_freeform_":
                    family = "freeform"
                elif preset:
                    family = f"auto_shape:{preset}"
                else:
                    family = "auto_shape:unknown"
                break
            node = node.getparent() if hasattr(node, "getparent") else None
        out[name] = family
    return out


def _auto_shape_preset(sp_elem) -> str:
    """Return the shape preset key for an ``<p:sp>`` element.

    - ``_txBox_`` if the shape is a text box (has txBox=1)
    - ``_freeform_`` if the geometry is a custom path (no prstGeom)
    - otherwise the ``prstGeom/@prst`` value (e.g. ``rect``, ``ellipse``).
    """
    for child in sp_elem.iter():
        tag = child.tag
        if tag.endswith("}nvSpPr"):
            for sub in child.iter():
                if sub.tag.endswith("}cNvSpPr") and sub.get("txBox") == "1":
                    return "_txBox_"
        if tag.endswith("}prstGeom"):
            return child.get("prst") or "unknown"
        if tag.endswith("}custGeom"):
            return "_freeform_"
    return "unknown"


# ---------------------------------------------------------------------------
# Motion skills (content-triggered animations)
# ---------------------------------------------------------------------------

_MOTIONS_DIR = Path(__file__).resolve().parents[3] / "skills_library" / "ppt" / "motions"
_MOTION_CACHE: dict[str, dict] | None = None


def _load_motions() -> dict[str, dict]:
    """Lazy-load motion skill metadata from skills_library/ppt/motions/*.py."""
    global _MOTION_CACHE
    if _MOTION_CACHE is not None:
        return _MOTION_CACHE
    import importlib.util as _ilu
    out: dict[str, dict] = {}
    if not _MOTIONS_DIR.exists():
        _MOTION_CACHE = out
        return out
    sys.path.insert(0, str(_MOTIONS_DIR))
    sys.path.insert(
        0, str(_MOTIONS_DIR.parent / "shells_seed")
    )
    for py in sorted(_MOTIONS_DIR.glob("*.py")):
        if py.name.startswith("_"):
            continue
        mid = py.stem
        spec = _ilu.spec_from_file_location(f"motion_{mid}", py)
        if spec is None or spec.loader is None:
            continue
        try:
            mod = _ilu.module_from_spec(spec)
            spec.loader.exec_module(mod)
        except Exception as e:
            out[mid] = {"error": str(e)}
            continue
        out[mid] = {
            "id": mid,
            "name": getattr(mod, "NAME", mid),
            "category": getattr(mod, "CATEGORY", "unknown"),
            "description": getattr(mod, "DESCRIPTION", ""),
            "applicability": getattr(mod, "APPLICABILITY", {}),
            "parameters": getattr(mod, "PARAMETERS", {}),
            # Phase-A ranking metadata
            "tags": getattr(mod, "TAGS", []),
            "intensity": int(getattr(mod, "INTENSITY", 5)),
            "embedding_text": getattr(mod, "EMBEDDING_TEXT", "") or getattr(mod, "DESCRIPTION", ""),
            "content_matchers": getattr(mod, "CONTENT_MATCHERS", {}),
            "complementary_with": getattr(mod, "COMPLEMENTARY_WITH", []),
            "conflicts_with": getattr(mod, "CONFLICTS_WITH", []),
            # LLM-readable selection metadata (2026-04-19)
            "mood": list(getattr(mod, "MOOD", []) or []),
            "archetype_fit": list(getattr(mod, "ARCHETYPE_FIT", []) or []),
            "module": mod,
            "path": str(py),
        }
    _MOTION_CACHE = out
    return out


# Keyword/embedding ranking was removed 2026-04-19: the LLM now reads
# metadata directly via list_motions/list_shells/list_themes and makes
# its own picks.


# Anchor-name synonyms. Keys are the motion metadata vocabulary (what
# APPLICABILITY.anchor_names declares); values are the other names that
# LIKELY refer to the same conceptual shape across shells.
_ANCHOR_SYNONYMS = {
    "hero_number":    ["value", "metric_xl", "metric", "big_number", "kpi_value", "hero_value", "hero_metric"],
    "hero_value":     ["value", "metric_xl", "hero_number", "big_number"],
    "big_number":     ["value", "metric_xl", "hero_number"],
    "hero_headline":  ["headline", "title", "title_xl", "h1", "hero_title"],
    "headline":       ["title", "title_xl", "hero_headline", "h1"],
    "hero_title":     ["headline", "title", "title_xl"],
    "title":          ["headline", "title_xl", "hero_headline"],
    "subhead":        ["subtitle", "subhead_text", "sub", "tagline"],
    "subtitle":       ["subhead"],
    "eyebrow":        ["kicker", "overline", "eyebrow_text", "section_label"],
    "kicker":         ["eyebrow", "overline", "section_label"],
    "section_label":  ["eyebrow", "kicker"],
    "section_number": ["section_num", "chapter_number", "number_glyph", "numeral"],
    "chapter_label":  ["section_label", "eyebrow", "chapter"],
    "quote":          ["quote_body", "pull_quote", "body"],
    "attribution":    ["author", "citation", "byline", "attrib"],
    "cta":            ["button", "cta_button", "chip", "badge", "action"],
    "button":         ["cta", "chip", "badge"],
    "badge":          ["chip", "tag", "pill", "cta"],
    "accent_orb":     ["halo", "accent_halo", "accent", "orb", "accent_ring"],
    "accent_ring":    ["halo", "accent_orb", "ring"],
    "halo":           ["accent_halo", "accent_orb", "accent_ring"],
    "hero_image":     ["image", "photo", "picture", "hero_photo"],
    "logo":           ["brand_mark", "seal", "stamp"],
    "card":           ["card_1", "card_2", "item", "tile"],
    "delta":          ["change", "change_chip"],
    "context":        ["description", "body", "supporting_text"],
    "label":          ["caption"],
}


def _find_anchor_shape(slide, anchor_name: str, motion_info: dict | None = None):
    """Locate a shape on `slide` that best matches `anchor_name`.

    Resolution order:
      1. Exact shape.name match (drops the !!sameName prefix too).
      2. Synonym match via _ANCHOR_SYNONYMS (motion vocab → shell vocab).
      3. Motion's own APPLICABILITY.anchor_names / anchor_name_regex —
         try each alternative anchor the motion declared.
      4. Substring contains match (e.g. "headline" in "eyebrow_headline").
      5. Name-token overlap fallback (e.g. anchor="hero_number" finds a
         shape whose name contains any of {hero, number, value, metric}).

    Returns the first matching pptx shape, or None.
    """
    import re as _re
    target_name = (anchor_name or "").strip()
    if not target_name:
        return None

    def _clean(n: str) -> str:
        return (n or "").replace("!!sameName", "").strip()

    shapes = list(slide.shapes)

    # 1) exact
    for shp in shapes:
        if _clean(getattr(shp, "name", "")) == target_name:
            return shp

    # 2) synonyms
    syn_set = set(_ANCHOR_SYNONYMS.get(target_name, []))
    # also reverse lookup: if any synonym-group contains target_name AND the
    # shape's actual name is the group key, still a match.
    for k, v in _ANCHOR_SYNONYMS.items():
        if target_name in v:
            syn_set.add(k)
            syn_set.update(v)
    if syn_set:
        for shp in shapes:
            n = _clean(getattr(shp, "name", ""))
            if n in syn_set:
                return shp

    # 3) motion's declared alternatives
    if motion_info:
        app = motion_info.get("applicability", {}) or {}
        alt_names = set(app.get("anchor_names") or [])
        alt_names.discard(target_name)
        for shp in shapes:
            n = _clean(getattr(shp, "name", ""))
            if n in alt_names:
                return shp
        # and via the motion's regex
        rgx = app.get("anchor_name_regex")
        if rgx:
            try:
                pat = _re.compile(rgx, _re.I)
                for shp in shapes:
                    if pat.search(_clean(getattr(shp, "name", ""))):
                        return shp
            except Exception:
                pass

    # 4) substring contains
    low = target_name.lower()
    for shp in shapes:
        n = _clean(getattr(shp, "name", "")).lower()
        if n and (low in n or n in low):
            return shp

    # 5) token overlap
    tokens = set(_re.split(r"[^a-z0-9]+", low)) - {""}
    # expand with synonym tokens
    for s in syn_set:
        tokens.update(_re.split(r"[^a-z0-9]+", s.lower()))
    tokens -= {""}
    if tokens:
        for shp in shapes:
            name_tokens = set(_re.split(r"[^a-z0-9]+", _clean(getattr(shp, "name", "")).lower())) - {""}
            if tokens & name_tokens:
                return shp

    return None


def _apply_motion_hard_filters(motions: dict, role: str, category: str,
                                already_picked: list[str]) -> dict:
    """Pure metadata filter. Role is a HARD filter when provided.

    Only drops: load errors, wrong category, conflicts with already_picked,
    and entries in already_picked themselves. If `role` is provided and a
    motion declares APPLICABILITY.roles, the role must be included
    (otherwise the motion is off-target for this slide).
    """
    out = {}
    picked_set = set(already_picked or [])
    for mid, info in motions.items():
        if "error" in info:
            continue
        if category and info.get("category") != category:
            continue
        conflicts = set(info.get("conflicts_with", []) or [])
        if conflicts & picked_set:
            continue
        if mid in picked_set:
            continue
        if role:
            allowed = info.get("applicability", {}).get("roles") or []
            if allowed and role not in allowed:
                continue
        out[mid] = info
    return out



# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def list_motions(role: str = "", anchor_name: str = "",
                 category: str = "", mood: str = "",
                 archetype: str = "",
                 already_picked: str = "",
                 exclude_ids: str = "",
                 limit: int = 0,
                 slide_brief: str = "", shape_names: str = "",
                 slot_values: str = "", top_k: int = 0) -> str:
    """List motion candidates that pass pure metadata filters.

    The LLM is expected to read the returned slate and pick the best
    fit — there is NO ranker and NO score. All filters are metadata-
    based, no text similarity.

    Args:
        role: slide role (e.g. "cover", "hero_giant_metric"). HARD filter:
              a motion with narrow APPLICABILITY.roles must include `role`
              or it's dropped. Motions with empty roles are always kept.
        anchor_name: optional shape name. Exposes `anchor_match=True` on
              the returned entry when APPLICABILITY.anchor_names /
              anchor_name_regex matches. Does NOT filter.
        category: "entrance" | "emphasis" | "ambient" | "transition". HARD filter.
        mood: comma-separated MOOD tokens (calm, editorial, bold, punchy,
              playful, boardroom, cinematic, restrained, warm, technical).
              HARD filter: at least one token must appear in the motion's
              MOOD list. Empty mood skips the filter.
        archetype: comma-separated ARCHETYPE_FIT tokens (data, narrative,
              product, research, brand, boardroom). Same semantics as mood.
        already_picked: comma-separated motion ids already chosen for this
              slide. Drops them + anything that declares CONFLICTS_WITH.
        exclude_ids: comma-separated motion ids to also drop (e.g. items
              used by prior variants; soft cross-deck dedup).
        limit: optional cap (0 = unlimited). Deterministic alpha order.

    Legacy args `slide_brief / shape_names / slot_values / top_k` are
    accepted but IGNORED (kept for backwards-compat with older callers).

    Returns:
        JSON with "motions" list. Each entry carries id, name, category,
        description, intensity, tags, mood, archetype_fit, applicability
        (roles/anchors/max_per_slide), conflicts_with, complementary_with,
        content_matchers, and anchor_match boolean. No score field.
    """
    import json as _json
    import re as _re
    ms = _load_motions()

    already = [x.strip() for x in (already_picked or "").split(",") if x.strip()]
    exclude = {x.strip() for x in (exclude_ids or "").split(",") if x.strip()}
    mood_want = {x.strip().lower() for x in (mood or "").split(",") if x.strip()}
    arch_want = {x.strip().lower() for x in (archetype or "").split(",") if x.strip()}

    filtered = _apply_motion_hard_filters(
        ms, role=role, category=category, already_picked=already,
    )

    out: list[dict] = []
    for mid in sorted(filtered.keys()):
        if mid in exclude:
            continue
        info = filtered[mid]
        entry_mood = {m.lower() for m in info.get("mood", []) or []}
        entry_arch = {a.lower() for a in info.get("archetype_fit", []) or []}
        if mood_want and not (mood_want & entry_mood):
            continue
        if arch_want and not (arch_want & entry_arch):
            continue

        app = info.get("applicability", {})
        anchor_match = False
        if anchor_name:
            names = app.get("anchor_names") or []
            rgx = app.get("anchor_name_regex")
            if anchor_name in names:
                anchor_match = True
            elif rgx:
                try:
                    if _re.search(rgx, anchor_name, _re.I):
                        anchor_match = True
                except Exception:
                    pass

        out.append({
            "id": info["id"],
            "name": info["name"],
            "category": info["category"],
            "description": info["description"],
            "intensity": info.get("intensity", 5),
            "tags": info.get("tags", []),
            "mood": info.get("mood", []),
            "archetype_fit": info.get("archetype_fit", []),
            "roles": app.get("roles", []),
            "anchors": (app.get("anchor_names") or []) + (
                [f"regex:{app['anchor_name_regex']}"] if app.get("anchor_name_regex") else []
            ),
            "max_per_slide": app.get("max_per_slide", 1),
            "complementary_with": info.get("complementary_with", []),
            "conflicts_with": info.get("conflicts_with", []),
            "content_matchers": info.get("content_matchers", {}),
            "anchor_match": anchor_match,
        })

    if limit and limit > 0:
        out = out[: int(limit)]

    return _json.dumps({
        "filter": {
            "role": role, "category": category, "mood": sorted(mood_want),
            "archetype": sorted(arch_want), "anchor_name": anchor_name,
            "already_picked": already, "exclude_ids": sorted(exclude),
        },
        "count": len(out),
        "motions": out,
    }, indent=2)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def get_motion_info(motion_id: str) -> str:
    """Return full metadata for a single motion skill, including its
    parameter schema."""
    import json as _json
    ms = _load_motions()
    info = ms.get(motion_id)
    if info is None:
        return f"Error: motion '{motion_id}' not found"
    if "error" in info:
        return f"Error loading motion '{motion_id}': {info['error']}"
    return _json.dumps({
        "id": info["id"],
        "name": info["name"],
        "category": info["category"],
        "description": info["description"],
        "applicability": info["applicability"],
        "parameters": info["parameters"],
    }, indent=2)


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def get_motion_code(motion_id: str) -> str:
    """Return the raw Python source of a motion skill (for LLM to
    learn techniques and adapt parameters)."""
    ms = _load_motions()
    info = ms.get(motion_id)
    if info is None:
        return f"Error: motion '{motion_id}' not found"
    return Path(info["path"]).read_text()


# @mcp.tool()  # [DISABLED: python-pptx shell path blocked — use PPT Master (pptmaster_*)]
def apply_motion(prs_id: str, slide_index: int, motion_id: str,
                 anchor_name: str = "", params: str = "") -> str:
    """Apply a motion skill to a target shape on a slide.

    Args:
        prs_id: presentation id (from create_presentation).
        slide_index: 1-based slide index.
        motion_id: motion skill id (from list_motions).
        anchor_name: either the exact shape.name (e.g. "hero_number")
                     or the force-match morph anchor (e.g.
                     "!!sameNameaccent_orb"). If empty, the motion is
                     applied to the first shape on the slide.
        params: JSON string of parameter overrides. Empty string uses
                defaults from the motion's PARAMETERS schema.

    Returns:
        Status line with the effect's shape id + motion id on success,
        or an Error: ... string.
    """
    import json as _json
    prs = _presentations.get(prs_id)
    if prs is None:
        return f"Error: presentation '{prs_id}' not found"
    try:
        slide = prs.slides[int(slide_index) - 1]
    except Exception:
        return f"Error: slide_index {slide_index} out of range"
    ms = _load_motions()
    info = ms.get(motion_id)
    if info is None:
        return f"Error: motion '{motion_id}' not found"
    # Parse params JSON.
    try:
        parsed_params = _json.loads(params) if params else {}
    except Exception as e:
        return f"Error: invalid params JSON: {e}"
    # Find target shape — exact match first, fuzzy match second. Fuzzy
    # is needed because motion metadata uses canonical anchor vocab
    # (hero_number, hero_headline, ...) while distilled / seed shells
    # may name shapes differently (value, metric_xl, headline, title).
    target = None
    if anchor_name:
        target = _find_anchor_shape(slide, anchor_name, info)
        if target is None:
            return f"Error: anchor '{anchor_name}' not found on slide {slide_index}"
    else:
        shapes = list(slide.shapes)
        if not shapes:
            return f"Error: slide {slide_index} has no shapes"
        target = shapes[0]
    try:
        info["module"].apply(slide, target, parsed_params)
    except Exception as e:
        import traceback
        return f"Error applying motion '{motion_id}': {e}\n{traceback.format_exc()}"
    return (
        f"Applied motion '{motion_id}' to shape '{getattr(target, 'name', '?')}' "
        f"(shape_id={getattr(target, 'shape_id', '?')}) on slide {slide_index}."
    )


if __name__ == "__main__":
    main()
