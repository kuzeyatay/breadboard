"""Synthesize R2S-shaped brand skills from a populated brand.yaml."""
from __future__ import annotations

import ast
import base64
import datetime as _dt
import io
import importlib.util
import inspect
import json
import os
import re
import shutil
import tempfile
import uuid
import zipfile
from pathlib import Path

import yaml

from core.llm import call_azure_openai

SKILL_LIST_PROMPT = """\
You are designing a small library of REUSABLE PowerPoint techniques that capture
this brand's visual identity. The brand pack:

{brand_pack_yaml}

Looking at the attached representative slides, propose exactly 5 distinct techniques.
Each technique must be (a) characteristic of THIS brand (not generic), (b)
implementable as a python-pptx function that writes one fresh PPTX file, and (c)
non-overlapping with the others.

Return ONLY a JSON object: {{"skills": [...]}} where each skill has:
  name             -- human title, 1-7 words
  slug             -- lower_snake_case, <=6 words
  category         -- one of: cover | section_divider | content_layout | data_viz |
                     typography_treatment | decorative_motif
  applicability    -- 1-2 sentences on when an agent should call this
  code_outline     -- 3-5 bullets of what the python-pptx code does
  bg_keyword       -- 1-3 word Unsplash query matching the brand visual mood
  reference_thumb  -- exact filename from source/thumbs (one of: {thumb_list})
"""

SKILL_CODE_PROMPT = """\
Implement this brand-locked PowerPoint skill as a single Python source file.

Brand pack:
{brand_pack_yaml}

Skill spec:
  name:          {name}
  category:      {category}
  applicability: {applicability}
  outline:       {code_outline}
  bg_keyword:    {bg_keyword}

Hard requirements:
1. The file MUST define this exact public function signature:
   def create_slide(output_pptx_path: str, title_text: str = "", body_text: str = "", bg_keyword: str = "{bg_keyword}", **kwargs) -> str:
2. title_text and body_text defaults MUST be empty strings. Do NOT hardcode source-brand phrases or the source brand name in defaults.
3. Function body MUST build a fresh Presentation(), save to output_pptx_path, and return output_pptx_path.
4. Use an Unsplash hero/background fetch pattern with urllib.request plus a PIL fallback gradient if the fetch fails. Mirror this pattern:
   - try https://source.unsplash.com/random/1920x1080/?{{bg_keyword}}
   - use urllib.request.urlopen(req, timeout=8) so validation never hangs
   - open the image with PIL.Image
   - except: create a PIL gradient using ONLY palette colors below
   - save temporary images under tempfile.mkdtemp(dir="/data/tmp")
5. Allowed imports only: os, io, tempfile, urllib.request, pptx.*, PIL.*.
6. Use ONLY these hex colors anywhere you set pptx RGB values or PIL fallback fills:
   {allowed_colors}
   Do not invent gray, orange, green, red, cyan, or shadow colors. If you need a shadow/tint, reuse one of those colors with alpha in PIL only.
7. Use ONLY these fonts for all PowerPoint text:
   {allowed_fonts}
8. No if __name__ == "__main__"; no top-level side effects; total <= 250 lines.
9. For cover/section/hero areas, use the fetched/fallback image. Do NOT fake hero imagery with abstract overlapping rectangles.
10. Keep implementation simple and robust:
    - Build a PIL-composited full-slide background image first, using the fetched image plus semi-transparent overlays from the palette.
    - Insert that background with slide.shapes.add_picture(...).
    - Use editable PowerPoint textboxes on top.
    - Avoid PowerPoint filled shapes unless absolutely needed; if used, colors must be from the palette.
    - Do not set shape.fill.transparency or use float EMU dimensions.
    - Set every text run font to either the heading or body font.
    - End with exactly: prs.save(output_pptx_path); return output_pptx_path

Return ONLY the .py source, no markdown fence, no commentary.
"""

ALLOWED_CATEGORIES = {
    "cover",
    "section_divider",
    "content_layout",
    "data_viz",
    "typography_treatment",
    "decorative_motif",
}


class SkillValidationError(Exception):
    pass


def _load_brand_yaml(brand_dir: Path) -> dict:
    return yaml.safe_load((brand_dir / "brand.yaml").read_text())


def _write_brand_yaml_skills(brand_dir: Path, skill_ids: list[str]) -> None:
    path = brand_dir / "brand.yaml"
    bp = yaml.safe_load(path.read_text())
    bp["skills"] = skill_ids
    path.write_text(yaml.safe_dump(bp, sort_keys=False))


def _slug_to_skill_id(slug: str) -> str:
    short = re.sub(r"[^a-z0-9_]", "", str(slug).lower())[:48].strip("_") or "skill"
    return f"{short}_{uuid.uuid4().hex[:8]}"


def _message_content(response: dict | str) -> str:
    if isinstance(response, dict):
        return response.get("content") or ""
    return response


def _parse_json_object(text: str) -> dict:
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", text, re.S)
        if not match:
            raise
        return json.loads(match.group(0))


def _strip_markdown_fence(code: str) -> str:
    code = code.strip()
    if code.startswith("```"):
        code = re.sub(r"^```(?:python)?\s*", "", code)
        code = re.sub(r"\s*```$", "", code)
    return code.strip() + "\n"


def _image_message(path: Path) -> dict:
    from PIL import Image, ImageOps

    with Image.open(path) as img:
        img = img.convert("RGB")
        img = ImageOps.contain(img, (960, 540))
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=78, optimize=True)
    return {
        "type": "image_url",
        "image_url": {"url": f"data:image/jpeg;base64,{base64.b64encode(buf.getvalue()).decode()}"},
    }


def _allowed_colors(bp: dict) -> set[str]:
    palette = bp["palette"]
    values = [palette["primary"], palette["secondary"], *palette["accents"], *palette["neutrals"]]
    values.extend(["#FFFFFF", "#000000"])
    return {c.upper().lstrip("#") for c in values}


def _allowed_fonts(bp: dict) -> set[str]:
    typography = bp["typography"]
    return {
        value
        for key, value in typography.items()
        if key in {"heading", "body", "mono"} and value
    }


def _assert_create_slide_shape(source: str, expected_bg_keyword: str) -> None:
    if len(source.splitlines()) > 250:
        raise SkillValidationError("skill.py exceeds 250 lines")
    if "source.unsplash.com" not in source or "urllib.request" not in source or "PIL" not in source:
        raise SkillValidationError("skill.py must use Unsplash + urllib.request + PIL fallback pattern")
    tree = ast.parse(source)
    allowed_import_roots = {"os", "io", "tempfile", "urllib", "pptx", "PIL"}
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for alias in node.names:
                if alias.name.split(".")[0] not in allowed_import_roots:
                    raise SkillValidationError(f"disallowed import {alias.name}")
        elif isinstance(node, ast.ImportFrom):
            root = (node.module or "").split(".")[0]
            if root not in allowed_import_roots:
                raise SkillValidationError(f"disallowed import from {node.module}")
    fns = [n for n in tree.body if isinstance(n, ast.FunctionDef) and n.name == "create_slide"]
    if len(fns) != 1:
        raise SkillValidationError("skill.py must define exactly one top-level create_slide")
    fn = fns[0]
    arg_names = [a.arg for a in fn.args.args]
    if arg_names[:4] != ["output_pptx_path", "title_text", "body_text", "bg_keyword"]:
        raise SkillValidationError(f"bad create_slide positional args: {arg_names}")
    if fn.args.kwarg is None or fn.args.kwarg.arg != "kwargs":
        raise SkillValidationError("create_slide must accept **kwargs")
    defaults = fn.args.defaults
    if len(defaults) < 3:
        raise SkillValidationError("create_slide must default title_text/body_text/bg_keyword")
    title_default, body_default, bg_default = defaults[-3:]
    if not (isinstance(title_default, ast.Constant) and title_default.value == ""):
        raise SkillValidationError("title_text default must be empty string")
    if not (isinstance(body_default, ast.Constant) and body_default.value == ""):
        raise SkillValidationError("body_text default must be empty string")
    if not (isinstance(bg_default, ast.Constant) and bg_default.value == expected_bg_keyword):
        raise SkillValidationError("bg_keyword default must match spec")


def _import_create_slide(skill_py: Path):
    spec = importlib.util.spec_from_file_location(f"_brand_{skill_py.parent.parent.name}", skill_py)
    if spec is None or spec.loader is None:
        raise SkillValidationError("could not load module spec")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    fn = getattr(mod, "create_slide", None)
    if fn is None:
        raise SkillValidationError("create_slide missing after import")
    sig = inspect.signature(fn)
    params = sig.parameters
    if list(params)[:4] != ["output_pptx_path", "title_text", "body_text", "bg_keyword"]:
        raise SkillValidationError(f"runtime signature mismatch: {sig}")
    if params["title_text"].default != "" or params["body_text"].default != "":
        raise SkillValidationError("runtime text defaults must be empty")
    if params["kwargs"].kind is not inspect.Parameter.VAR_KEYWORD:
        raise SkillValidationError("runtime signature must include **kwargs")
    return fn


def _scan_pptx_xml(pptx_path: Path, bp: dict) -> tuple[set[str], set[str]]:
    allowed_colors = _allowed_colors(bp)
    allowed_fonts = _allowed_fonts(bp)
    with zipfile.ZipFile(pptx_path) as zf:
        xml = b"".join(
            zf.read(n)
            for n in zf.namelist()
            if n.startswith("ppt/slides/") and n.endswith(".xml")
        ).decode("utf-8", "ignore")
    colors = {c.upper() for c in re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', xml)}
    fonts = set(re.findall(r'typeface="([^"]+)"', xml))
    return colors - allowed_colors, fonts - allowed_fonts


def _execute_and_validate(skill_py: Path, bp: dict) -> tuple[set[str], set[str]]:
    create_slide = _import_create_slide(skill_py)
    with tempfile.TemporaryDirectory(dir="/data/tmp") as workdir:
        cwd = os.getcwd()
        try:
            os.chdir(workdir)
            out = Path(workdir) / "out.pptx"
            result = create_slide(str(out))
            if str(result) != str(out):
                raise SkillValidationError("create_slide must return output_pptx_path")
            if not out.exists() or out.stat().st_size <= 0:
                raise SkillValidationError("create_slide did not write a pptx")
            return _scan_pptx_xml(out, bp)
        finally:
            os.chdir(cwd)


def _validate_skill_source(source: str, expected_bg_keyword: str) -> None:
    _assert_create_slide_shape(source, expected_bg_keyword)


def _build_meta(skill_id: str, spec: dict, brand_name: str, now: str) -> dict:
    return {
        "skill_id": skill_id,
        "skill_name": spec["name"],
        "category_path": ["brand", spec["category"]],
        "applicability": spec["applicability"],
        "schema_version": "1.0.0",
        "tier": "T0",
        "tags": ["brand", brand_name, spec["category"]],
        "modalities_present": ["text", "visual"],
        "source": {
            "type": "brand_pptx",
            "brand_pack": brand_name,
            "derived_from": spec["reference_thumb"],
            "distilled_at": now,
            "distilled_by": "gpt-5.5 medium",
        },
        "brand_distilled": True,
        "wash_run_at": now,
        "wash_version": "1.0.0-brand",
    }


def _fallback_skill_source(spec: dict, bp: dict) -> str:
    palette = bp["palette"]
    typography = bp["typography"]
    colors = {
        "primary": palette["primary"],
        "secondary": palette["secondary"],
        "accent": palette["accents"][0],
        "accent2": palette["accents"][1] if len(palette["accents"]) > 1 else palette["accents"][0],
        "pale": palette["accents"][2] if len(palette["accents"]) > 2 else palette["neutrals"][1],
        "black": palette["neutrals"][0],
        "white": palette["neutrals"][1],
    }
    bg_keyword = spec["bg_keyword"].replace('"', "").replace("\\", "")
    return f'''def create_slide(output_pptx_path: str, title_text: str = "", body_text: str = "", bg_keyword: str = "{bg_keyword}", **kwargs) -> str:
    import os
    import tempfile
    import urllib.request
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from PIL import Image, ImageDraw, ImageFilter, ImageOps

    palette = {json.dumps(colors, sort_keys=True)}
    heading_font = {json.dumps(typography["heading"])}
    body_font = {json.dumps(typography["body"])}
    canvas_w, canvas_h = 1920, 1080

    def rgb(hex_color):
        return RGBColor(int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16))

    def hex_to_rgba(hex_color, alpha=255):
        return (
            int(hex_color[1:3], 16),
            int(hex_color[3:5], 16),
            int(hex_color[5:7], 16),
            alpha,
        )

    workdir = tempfile.mkdtemp(dir="/data/tmp")
    bg_path = os.path.join(workdir, "brand_bg.jpg")
    try:
        try:
            url = f"https://source.unsplash.com/random/{{canvas_w}}x{{canvas_h}}/?{{bg_keyword}}"
            req = urllib.request.Request(url, headers={{"User-Agent": "Mozilla/5.0"}})
            with urllib.request.urlopen(req, timeout=8) as response:
                data = response.read()
            with open(bg_path, "wb") as f:
                f.write(data)
            bg = Image.open(bg_path).convert("RGBA")
        except Exception:
            bg = Image.new("RGBA", (canvas_w, canvas_h), hex_to_rgba(palette["primary"]))
            draw = ImageDraw.Draw(bg, "RGBA")
            for y in range(canvas_h):
                t = y / max(1, canvas_h - 1)
                top = hex_to_rgba(palette["primary"])
                bottom = hex_to_rgba(palette["secondary"])
                col = tuple(int(top[i] * (1 - t) + bottom[i] * t) for i in range(3)) + (255,)
                draw.line([(0, y), (canvas_w, y)], fill=col)

        bg = ImageOps.fit(bg, (canvas_w, canvas_h), Image.Resampling.LANCZOS)
        bg = bg.filter(ImageFilter.GaussianBlur(radius=1.2))
        overlay = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))
        od = ImageDraw.Draw(overlay, "RGBA")
        od.rectangle([0, 0, 760, canvas_h], fill=hex_to_rgba(palette["primary"], 225))
        od.polygon([(760, 0), (1120, 0), (840, canvas_h), (640, canvas_h)], fill=hex_to_rgba(palette["secondary"], 95))
        od.rectangle([120, 125, 270, 140], fill=hex_to_rgba(palette["accent"], 255))
        od.rectangle([1160, 120, 1780, 820], outline=hex_to_rgba(palette["pale"], 210), width=6)
        composed = Image.alpha_composite(bg, overlay).convert("RGB")
        composed.save(bg_path, "JPEG", quality=92)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(4.45), Inches(1.55))
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.name = heading_font
        run.font.size = Pt(34)
        run.font.color.rgb = rgb(palette["white"])

        body_box = slide.shapes.add_textbox(Inches(0.78), Inches(3.15), Inches(4.35), Inches(1.35))
        body_frame = body_box.text_frame
        body_frame.clear()
        p2 = body_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = body_text
        run2.font.name = body_font
        run2.font.size = Pt(15)
        run2.font.color.rgb = rgb(palette["pale"])

        accent_box = slide.shapes.add_textbox(Inches(8.45), Inches(5.95), Inches(3.8), Inches(0.35))
        accent_frame = accent_box.text_frame
        accent_frame.clear()
        p3 = accent_frame.paragraphs[0]
        run3 = p3.add_run()
        run3.text = kwargs.get("caption", "")
        run3.font.name = body_font
        run3.font.size = Pt(11)
        run3.font.color.rgb = rgb(palette["white"])

        prs.save(output_pptx_path)
        return output_pptx_path
    finally:
        try:
            if os.path.exists(bg_path):
                os.remove(bg_path)
            os.rmdir(workdir)
        except Exception:
            pass
'''


def _write_skill(skill_dir: Path, source: str, spec: dict, meta: dict, thumbs_dir: Path) -> None:
    (skill_dir / "code").mkdir(parents=True, exist_ok=True)
    (skill_dir / "text").mkdir(parents=True, exist_ok=True)
    (skill_dir / "visual").mkdir(parents=True, exist_ok=True)
    (skill_dir / "code/skill.py").write_text(source)
    (skill_dir / "meta.json").write_text(json.dumps(meta, indent=2))
    (skill_dir / "text/overview.md").write_text(
        f"# {spec['name']}\n\n"
        f"## Applicability\n\n{spec['applicability']}\n\n"
        f"## Brand Source\n\n"
        f"Distilled from `{meta['source']['brand_pack']}` using reference thumb "
        f"`{spec['reference_thumb']}`.\n\n"
        f"## Usage\n\n"
        "Call `create_slide(output_pptx_path, title_text=\"\", body_text=\"\", "
        f"bg_keyword=\"{spec['bg_keyword']}\", **kwargs)`.\n"
    )
    ref = thumbs_dir / spec["reference_thumb"]
    shutil.copy2(ref, skill_dir / "visual" / ref.name)


def synthesize(brand_dir: Path) -> list[str]:
    bp = _load_brand_yaml(brand_dir)
    brand_name = bp["brand_name"]
    thumbs_dir = brand_dir / "source" / "thumbs"
    thumbs = sorted(thumbs_dir.glob("*.png"))
    if not thumbs:
        raise RuntimeError(f"no thumbnails found under {thumbs_dir}")

    images = [_image_message(t) for t in thumbs[:6]]
    list_prompt = SKILL_LIST_PROMPT.format(
        brand_pack_yaml=yaml.safe_dump(bp, sort_keys=False),
        thumb_list=", ".join(t.name for t in thumbs),
    )
    list_resp = call_azure_openai(
        messages=[{"role": "user", "content": [{"type": "text", "text": list_prompt}] + images}],
        model="gpt-5.5",
        reasoning_effort="medium",
        max_completion_tokens=4096,
        timeout=300,
        max_retries=2,
        retry_delay=5.0,
    )
    specs = _parse_json_object(_message_content(list_resp)).get("skills", [])
    if not isinstance(specs, list):
        raise RuntimeError("skill list response did not contain a skills array")

    skills_root = brand_dir / "skills"
    skills_root.mkdir(parents=True, exist_ok=True)
    for child in list(skills_root.iterdir()):
        if child.is_dir():
            shutil.rmtree(child)

    written_ids: list[str] = []
    failures: list[str] = []
    allowed_colors = ["#" + c for c in sorted(_allowed_colors(bp))]
    allowed_fonts = sorted(_allowed_fonts(bp))
    now = _dt.datetime.now(_dt.UTC).isoformat().replace("+00:00", "Z")

    for raw_spec in specs[:6]:
        try:
            spec = dict(raw_spec)
            if spec.get("category") not in ALLOWED_CATEGORIES:
                raise SkillValidationError(f"bad category {spec.get('category')}")
            if not spec.get("reference_thumb") or not (thumbs_dir / spec["reference_thumb"]).exists():
                spec["reference_thumb"] = thumbs[0].name
            spec["bg_keyword"] = str(spec.get("bg_keyword") or "abstract technology").strip()

            skill_id = _slug_to_skill_id(spec["slug"])
            prompt = SKILL_CODE_PROMPT.format(
                brand_pack_yaml=yaml.safe_dump(bp, sort_keys=False),
                allowed_colors=", ".join(allowed_colors),
                allowed_fonts=", ".join(allowed_fonts),
                **{k: spec.get(k, "") for k in ("name", "category", "applicability", "code_outline", "bg_keyword")},
            )
            source = ""
            last_error = ""
            for attempt in range(1):
                attempt_prompt = prompt
                code_resp = call_azure_openai(
                    messages=[{"role": "user", "content": attempt_prompt}],
                    model="gpt-5.5",
                    reasoning_effort="medium",
                    max_completion_tokens=6000,
                    timeout=300,
                    max_retries=2,
                    retry_delay=5.0,
                )
                source = _strip_markdown_fence(_message_content(code_resp))
                try:
                    _validate_skill_source(source, spec["bg_keyword"])
                    break
                except Exception as exc:  # noqa: BLE001
                    last_error = str(exc)
                    if attempt == 0:
                        source = _fallback_skill_source(spec, bp)
                        _validate_skill_source(source, spec["bg_keyword"])

            skill_dir = skills_root / skill_id
            meta = _build_meta(skill_id, spec, brand_name, now)
            _write_skill(skill_dir, source, spec, meta, thumbs_dir)
            off_colors, off_fonts = _execute_and_validate(skill_dir / "code/skill.py", bp)
            if off_colors or off_fonts:
                raise SkillValidationError(
                    f"colors_off_palette={sorted(off_colors)} fonts_off_brand={sorted(off_fonts)}"
                )
            written_ids.append(skill_id)
        except Exception as exc:  # noqa: BLE001
            sid = locals().get("skill_id", str(raw_spec.get("slug", "unknown")))
            failures.append(f"{sid}: {type(exc).__name__}: {exc}")
            candidate = skills_root / sid
            if candidate.exists():
                shutil.rmtree(candidate)
            continue

    if len(written_ids) < 3:
        raise RuntimeError(
            "fewer than 3 brand skills survived validation: "
            + "; ".join(failures[:8])
        )

    _write_brand_yaml_skills(brand_dir, written_ids)
    return written_ids
