"""Static-artifact source connector.

Imports skills from a curated corpus of static files (``.pptx``, ``.xlsx``,
``.html``, ``.blend``, ``.wav``+``.mid``).

Two ingestion modes:

1. **Catalogue mode** (default): each entry in
   ``skills_wiki/<d>/_static_corpus.json`` declares a path and a licence;
   the connector loads that catalogue and yields one candidate per row.

2. **Directory mode** (round-10 / task12): when ``extract_dir`` is
   passed (or ``EXTRACT_EXTENSIONS`` is set on the connector subclass),
   the connector walks the directory tree, applies a per-domain
   extractor, and yields one candidate per file it can summarise. The
   extractor produces ``text_overview`` from the file content; the
   licence comes from a sibling ``LICENSE`` file or
   ``--default-license``.

Unlicensed or blocked rows surface as ``license=None`` /
``"blocked"`` so the cycle quarantines them.
"""
from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from .base import CandidateSkill, SourceConnector


# Per-domain extension allow-list for directory mode.
_DOMAIN_EXTRACT_EXTENSIONS: dict[str, tuple[str, ...]] = {
    "excel": (".xlsx",),
    "web": (".html", ".htm"),
    "ppt": (".pptx",),
    "blender": (".blend",),
    "reaper": (".wav", ".mid", ".rpp"),
}


class StaticArtifactConnector(SourceConnector):
    SOURCE_TYPE = "static_artifact"

    def __init__(
        self,
        *,
        corpus_path: Path | str | None = None,
        extract_dir: Path | str | None = None,
        permissive_licences: tuple[str, ...] = (
            "permissive", "public-domain", "cc0-1.0", "cc-by-4.0", "cc-by-sa-4.0",
            "mit", "apache-2.0",
        ),
        default_license: str | None = None,
    ) -> None:
        self.corpus_path = Path(corpus_path) if corpus_path else None
        self.extract_dir = Path(extract_dir) if extract_dir else None
        self.permissive_licences = {l.lower() for l in permissive_licences}
        self.default_license = default_license

    def acquire(self, *, domain: str, quota: int, **kwargs: Any) -> list[CandidateSkill]:
        if quota <= 0:
            return []
        # Directory mode wins when explicitly requested via kwargs.
        extract_dir = kwargs.get("extract_dir") or self.extract_dir
        if extract_dir:
            return self._acquire_from_dir(
                domain=domain, quota=quota,
                extract_dir=Path(extract_dir),
                default_license=str(kwargs.get("default_license") or self.default_license or ""),
            )
        return self._acquire_from_catalogue(
            domain=domain, quota=quota,
            corpus_path=Path(kwargs.get("corpus_path") or self.corpus_path or ""),
        )

    def _acquire_from_catalogue(
        self, *, domain: str, quota: int, corpus_path: Path,
    ) -> list[CandidateSkill]:
        if not corpus_path.exists():
            return []
        try:
            corpus = json.loads(corpus_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            return []
        if not isinstance(corpus, list):
            return []
        out: list[CandidateSkill] = []
        for row in corpus[:quota]:
            if not isinstance(row, dict):
                continue
            licence = (row.get("license") or "").lower() or None
            if licence and licence not in self.permissive_licences:
                licence_value: str | None = "blocked"
            else:
                licence_value = licence
            skill_id = str(row.get("skill_id") or row.get("name") or row.get("path") or "static_unknown")
            skill_id = skill_id.replace("/", "_").replace(".", "_").lower()
            out.append(CandidateSkill(
                skill_id=skill_id,
                skill_name=str(row.get("name") or skill_id),
                domain=domain,
                source_type=self.SOURCE_TYPE,
                source_url=str(row.get("source_url") or ""),
                source_channel="static",
                license=licence_value,
                text_overview=str(row.get("description") or ""),
                applicability=str(row.get("applicability") or "") or None,
                tags=list(row.get("tags") or []),
                extras={"path": row.get("path"), "format": row.get("format")},
            ))
        return out

    def _acquire_from_dir(
        self, *, domain: str, quota: int, extract_dir: Path, default_license: str,
    ) -> list[CandidateSkill]:
        if not extract_dir.exists() or not extract_dir.is_dir():
            return []
        extensions = _DOMAIN_EXTRACT_EXTENSIONS.get(domain, (".py",))
        # Read sibling LICENSE file if present.
        license_value: str | None = (default_license or "").lower() or None
        for cand in (extract_dir / "LICENSE", extract_dir / "LICENSE.txt"):
            if cand.exists():
                license_value = _classify_license_text(cand.read_text(encoding="utf-8", errors="ignore"))
                break
        out: list[CandidateSkill] = []
        for path in sorted(extract_dir.rglob("*")):
            if not path.is_file():
                continue
            if path.suffix.lower() not in extensions:
                continue
            extractor = _DOMAIN_EXTRACTORS.get(domain, _default_extractor)
            try:
                summary = extractor(path)
            except Exception:  # noqa: BLE001
                continue
            if summary is None:
                continue
            rel = path.relative_to(extract_dir)
            slug = str(rel).replace("/", "_").replace(".", "_").replace("-", "_").lower()
            out.append(CandidateSkill(
                skill_id=f"static_{slug}",
                skill_name=summary.get("name") or path.stem,
                domain=domain,
                source_type=self.SOURCE_TYPE,
                source_url=str(path),
                source_channel="static",
                license=license_value if license_value else None,
                text_overview=summary.get("overview") or "",
                applicability=summary.get("applicability"),
                tags=summary.get("tags") or [],
                extras={"path": str(rel), "format": path.suffix.lstrip(".").lower()},
            ))
            if len(out) >= quota:
                break
        return out


# ---------- per-domain extractors ----------------------------------------


def _default_extractor(path: Path) -> dict[str, Any] | None:
    text = path.read_text(encoding="utf-8", errors="ignore")[:2000]
    if not text.strip():
        return None
    return {
        "name": path.stem,
        "overview": f"# {path.name}\n\n```\n{text}\n```",
        "applicability": f"static {path.suffix.lstrip('.')} artifact: {path.name}",
        "tags": [path.suffix.lstrip(".").lower()],
    }


def _excel_extractor(path: Path) -> dict[str, Any] | None:
    """Summarise a .xlsx by sheet name + dimensions + sample headers."""
    try:
        import openpyxl  # type: ignore[import-untyped]
    except ImportError:
        return _default_extractor(path)
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception:  # noqa: BLE001
        return None
    sheet_lines: list[str] = []
    tags: list[str] = ["xlsx"]
    sheet_names = list(wb.sheetnames)
    for sheet_name in sheet_names[:6]:
        ws = wb[sheet_name]
        max_col = min(ws.max_column or 0, 8)
        max_row = ws.max_row or 0
        headers: list[str] = []
        if max_row > 0 and max_col > 0:
            try:
                first_row = next(ws.iter_rows(min_row=1, max_row=1, max_col=max_col, values_only=True))
                headers = [str(cell) if cell is not None else "" for cell in first_row]
            except StopIteration:
                pass
        sheet_lines.append(
            f"- **{sheet_name}** ({max_row} rows × {max_col} cols) headers: {headers}"
        )
        tags.append(sheet_name.lower())
    overview = (
        f"# {path.name}\n\nExcel workbook with {len(sheet_names)} sheets.\n\n"
        + "\n".join(sheet_lines)
    )
    wb.close()
    return {
        "name": path.stem,
        "overview": overview,
        "applicability": f"reusable xlsx layout: {path.stem}",
        "tags": tags[:6],
    }


def _web_extractor(path: Path) -> dict[str, Any] | None:
    """Summarise an HTML file by <title>, headings, and tag counts."""
    text = path.read_text(encoding="utf-8", errors="ignore")
    if not text.strip():
        return None
    import re as _re
    title_m = _re.search(r"<title[^>]*>([^<]+)</title>", text, _re.IGNORECASE)
    title = (title_m.group(1).strip() if title_m else path.stem)
    h1s = _re.findall(r"<h1[^>]*>([^<]+)</h1>", text, _re.IGNORECASE)[:5]
    tags: list[str] = ["html"]
    for cls_m in _re.finditer(r"class=[\"']([^\"']+)[\"']", text):
        for cls in cls_m.group(1).split():
            if 3 <= len(cls) <= 24 and cls.isalnum() and cls.lower() not in tags:
                tags.append(cls.lower())
                if len(tags) >= 6:
                    break
        if len(tags) >= 6:
            break
    overview = (
        f"# {title}\n\nHTML page from {path.name}\n\n## Headings\n"
        + ("\n".join(f"- {h}" for h in h1s) if h1s else "(no h1 tags)")
        + f"\n\n## Sample markup\n\n```html\n{text[:1500]}\n```"
    )
    return {
        "name": title,
        "overview": overview,
        "applicability": f"reusable HTML/CSS pattern: {title}",
        "tags": tags,
    }


def _ppt_extractor(path: Path) -> dict[str, Any] | None:
    """Summarise a .pptx by slide count + slide titles + sample text."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        return _default_extractor(path)
    try:
        prs = Presentation(str(path))
    except Exception:  # noqa: BLE001
        return None
    titles: list[str] = []
    text_samples: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if len(titles) <= len(text_samples) and len(text) <= 80:
                titles.append(text)
            text_samples.append(text[:120])
            break
    overview = (
        f"# {path.name}\n\nPowerPoint deck with {len(prs.slides)} slides.\n\n"
        + "## Slide titles\n\n"
        + ("\n".join(f"{i+1}. {t}" for i, t in enumerate(titles[:10])) or "(no titles)")
        + f"\n\n## Sample text\n\n{chr(10).join('- ' + s for s in text_samples[:10])}"
    )
    tags = ["pptx", f"slides_{len(prs.slides)}"]
    return {
        "name": path.stem,
        "overview": overview,
        "applicability": f"reusable PPT deck pattern: {path.stem}",
        "tags": tags,
    }


def _blender_extractor(path: Path) -> dict[str, Any] | None:
    """Summarise a .blend file by stat metadata + magic-header check.

    Reads the first 12 bytes to confirm it's a real Blender file
    (BLENDER\0 prefix); does NOT parse the file (would require bpy).
    """
    try:
        head = path.read_bytes()[:12]
    except Exception:  # noqa: BLE001
        return None
    if not head.startswith(b"BLENDER"):
        return None
    size_mb = path.stat().st_size / (1024 * 1024)
    overview = (
        f"# {path.name}\n\nBlender scene file ({size_mb:.2f} MB).\n\n"
        f"Magic header: `{head[:7].decode('ascii')}` "
        f"version: `{head[7:].decode('ascii', errors='replace')}`\n\n"
        f"This is a static-artifact reference; load it via "
        f"`bpy.ops.wm.open_mainfile(filepath=...)` to inspect contents."
    )
    return {
        "name": path.stem,
        "overview": overview,
        "applicability": f"reusable Blender scene reference: {path.stem}",
        "tags": ["blend", "scene"],
    }


def _reaper_extractor(path: Path) -> dict[str, Any] | None:
    """Summarise a .wav / .mid / .rpp file."""
    suffix = path.suffix.lower()
    size_kb = path.stat().st_size / 1024
    if suffix == ".wav":
        try:
            import wave
            with wave.open(str(path), "rb") as wf:
                channels = wf.getnchannels()
                framerate = wf.getframerate()
                frames = wf.getnframes()
                duration_sec = frames / framerate if framerate else 0
            overview = (
                f"# {path.name}\n\nWAV audio: {channels}ch @ {framerate}Hz, "
                f"{duration_sec:.2f}s, {size_kb:.1f} KB."
            )
        except Exception:  # noqa: BLE001
            overview = f"# {path.name}\n\nWAV audio ({size_kb:.1f} KB) — header parse failed."
    elif suffix == ".mid":
        try:
            data = path.read_bytes()[:14]
            track_count = int.from_bytes(data[10:12], "big") if data.startswith(b"MThd") else 0
            overview = (
                f"# {path.name}\n\nMIDI file ({size_kb:.1f} KB), {track_count} tracks."
            )
        except Exception:  # noqa: BLE001
            overview = f"# {path.name}\n\nMIDI file ({size_kb:.1f} KB)."
    else:  # .rpp
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")[:1500]
            overview = f"# {path.name}\n\nReaper project file.\n\n```\n{text}\n```"
        except Exception:  # noqa: BLE001
            overview = f"# {path.name}\n\nReaper project ({size_kb:.1f} KB)."
    return {
        "name": path.stem,
        "overview": overview,
        "applicability": f"reusable {suffix.lstrip('.')} audio artifact: {path.stem}",
        "tags": ["reaper", suffix.lstrip(".")],
    }


_DOMAIN_EXTRACTORS = {
    "excel": _excel_extractor,
    "web": _web_extractor,
    "ppt": _ppt_extractor,
    "blender": _blender_extractor,
    "reaper": _reaper_extractor,
}


def _classify_license_text(text: str) -> str:
    """Classify a LICENSE file body. Returns 'permissive' / 'blocked' / etc."""
    lower = text.lower()
    if "mit license" in lower or " permission is hereby granted" in lower:
        return "mit"
    if "apache license" in lower:
        return "apache-2.0"
    if "public domain" in lower or "cc0" in lower:
        return "public-domain"
    if "creative commons" in lower:
        return "cc-by-4.0"
    if "gpl" in lower or "agpl" in lower or "lgpl" in lower:
        return "blocked"
    return "permissive"


__all__ = ["StaticArtifactConnector"]
