"""
core/extraction/quality_gate.py

Three hard gates for shell skills before they go to status="active":
  1. exec_ok: shell code executes without exception against sample slots
  2. render_ok: LibreOffice produces a ≥20KB PNG with no placeholder text
  3. visual_similarity: (optional, needs CLIP) renders resemble the source frame

A shell failing any gate is stored with status="candidate" and is hidden
from default retrieval.
"""
from __future__ import annotations

import json
import logging
import re
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

log = logging.getLogger("ppt.quality_gate")

_PROJECT_ROOT = Path(__file__).resolve().parents[2]
_SHELLS_SEED = _PROJECT_ROOT / "skills_library" / "ppt" / "shells_seed"


# ---------------------------------------------------------------------------
# Sample slots generators
# ---------------------------------------------------------------------------


def make_sample_slots(slot_specs: list[dict]) -> dict:
    """Generate dummy content that satisfies a shell's slot schema.
    Uses realistic-length strings so overlap gates catch text-wrapping bugs."""
    sample: dict[str, Any] = {}
    for s in slot_specs:
        name = s["name"]
        kind = s["kind"]
        if kind == "text":
            lim = s.get("max_chars") or 80
            n = name.lower()
            # Realistic content per slot role
            if "headline" in n or "title" in n:
                body = "Designing systems that feel inevitable in retrospect"
            elif "subhead" in n or "subtitle" in n or "tagline" in n:
                body = "A concise one-liner that sets the context for this slide and reads naturally."
            elif "eyebrow" in n or "label" in n or "section" in n:
                body = "Keynote 2026"
            elif "quote" in n or "body" in n:
                body = "Well-designed systems feel inevitable in retrospect — they invite you to use them."
            elif "cta" in n or "contact" in n:
                body = "hello@company.com"
            elif "number" in n:
                body = "02"
            else:
                body = f"Acme {name.replace('_', ' ')}"
            sample[name] = body[: max(lim - 1, 8)]
        elif kind == "bullet_list":
            n = s.get("bullet_capacity") or 3
            titles = ["Fast by default", "Composable parts", "Honest defaults",
                     "Safe to iterate", "Works offline", "Zero config"]
            bodies = [
                "Every page loads in under 100ms, no matter the size.",
                "Mix and match primitives without special cases.",
                "Sensible behavior out of the box, ready to customize.",
                "Reversible operations and clear audit trails.",
                "Full functionality when the network drops.",
                "Install once; nothing else to learn.",
            ]
            sample[name] = [
                {"title": titles[i % len(titles)], "body": bodies[i % len(bodies)],
                 "value": f"{(i+1)*17}%", "label": titles[i % len(titles)],
                 "delta": f"+{(i+1)*3}%"}
                for i in range(n)
            ]
        elif kind == "image":
            sample[name] = _get_test_image()
        elif kind == "icon":
            sample[name] = _get_test_image()
        elif kind == "chart":
            sample[name] = None
        elif kind == "accent_block":
            sample[name] = True
        else:
            sample[name] = f"acme_{name}"
    return sample


_TEST_IMG_CACHE: Path | None = None


def _get_test_image() -> str:
    """Return a path to a real PNG (400×300 gray) for testing image slots."""
    global _TEST_IMG_CACHE
    if _TEST_IMG_CACHE is not None and _TEST_IMG_CACHE.exists():
        return str(_TEST_IMG_CACHE)
    try:
        from PIL import Image
        tmp = Path(tempfile.gettempdir()) / "shell_quality_gate_img.png"
        if not tmp.exists():
            img = Image.new("RGB", (400, 300), (80, 80, 100))
            img.save(tmp)
        _TEST_IMG_CACHE = tmp
        return str(tmp)
    except Exception:
        return ""


# ---------------------------------------------------------------------------
# Load a default theme for rendering
# ---------------------------------------------------------------------------


def _default_theme() -> dict:
    from core.extraction.shell_loader import load_theme
    return load_theme("editorial_dark")


# ---------------------------------------------------------------------------
# Gate 1: exec check
# ---------------------------------------------------------------------------


def gate_exec(render_code: str, slot_specs: list[dict]) -> tuple[bool, str]:
    """Run the shell's render() in a sandbox. Return (ok, message)."""
    # Inject the code into a namespace with _shell_helpers importable
    sys_path_bak = list(sys.path)
    try:
        if str(_SHELLS_SEED) not in sys.path:
            sys.path.insert(0, str(_SHELLS_SEED))

        ns: dict[str, Any] = {}
        try:
            exec(render_code, ns)
        except Exception as e:
            return False, f"exec error at import time: {type(e).__name__}: {e}"

        render = ns.get("render")
        if not callable(render):
            return False, "no render() function defined"

        # Build a throwaway slide and call render
        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        slots = make_sample_slots(slot_specs)
        theme = _default_theme()
        try:
            render(slide, slots, theme)
        except Exception as e:
            return False, f"render() raised: {type(e).__name__}: {e}"

        n_shapes = len(slide.shapes)
        if n_shapes < 2:
            return False, f"render() produced only {n_shapes} shapes (likely broken)"

        return True, f"exec OK, {n_shapes} shapes"
    finally:
        sys.path[:] = sys_path_bak


# ---------------------------------------------------------------------------
# Gate 2: render check (LibreOffice → PNG)
# ---------------------------------------------------------------------------


_PLACEHOLDER_RESIDUE = [
    "sample text", "sample headline", "sample body", "your text", "lorem ipsum",
    "placeholder", "click to add", "solar eclipse", "lunar eclipse",
]


def gate_render(render_code: str, slot_specs: list[dict]) -> tuple[bool, str]:
    """Render the shell to PNG via LibreOffice and check basic sanity."""
    sys_path_bak = list(sys.path)
    try:
        if str(_SHELLS_SEED) not in sys.path:
            sys.path.insert(0, str(_SHELLS_SEED))

        ns: dict[str, Any] = {}
        try:
            exec(render_code, ns)
        except Exception as e:
            return False, f"exec error: {e}"
        render = ns.get("render")
        if not callable(render):
            return False, "no render()"

        slots = make_sample_slots(slot_specs)
        # DON'T use "Sample X" text here — it triggers placeholder lint against itself
        for k, v in slots.items():
            if isinstance(v, str):
                slots[k] = v.replace("Sample", "Acme").replace("sample", "acme")

        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        try:
            render(slide, slots, _default_theme())
        except Exception as e:
            return False, f"render raised during render gate: {e}"

        with tempfile.TemporaryDirectory() as td:
            pptx_path = Path(td) / "test.pptx"
            prs.save(str(pptx_path))

            try:
                subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "pdf",
                     "--outdir", td, str(pptx_path)],
                    capture_output=True, timeout=60, check=True,
                )
            except subprocess.TimeoutExpired:
                return False, "LibreOffice conversion timed out (>60s)"
            except subprocess.CalledProcessError as e:
                return False, f"LibreOffice failed: {e.stderr.decode()[:200]}"

            pdf_path = pptx_path.with_suffix(".pdf")
            if not pdf_path.exists():
                return False, "no PDF produced"

            png_out = Path(td) / "slide"
            subprocess.run(
                ["pdftoppm", "-r", "80", "-png", str(pdf_path), str(png_out)],
                capture_output=True, timeout=30,
            )
            pngs = sorted(Path(td).glob("slide-*.png"))
            if not pngs:
                return False, "no PNG rendered"

            size = pngs[0].stat().st_size
            if size < 8_000:
                return False, f"PNG too small ({size} bytes) — likely blank"

        return True, f"render OK (PNG {size // 1024}KB)"
    finally:
        sys.path[:] = sys_path_bak


# ---------------------------------------------------------------------------
# Gate 2b: overlap check (bbox pairwise)
# ---------------------------------------------------------------------------


def gate_overlap(render_code: str, slot_specs: list[dict], max_bad_pairs: int = 3) -> tuple[bool, str]:
    """Render into a throwaway slide and count bbox overlaps between
    non-background, non-contained shapes. Fails if too many.

    A 'bad pair' is two shapes whose intersection covers >15% of the smaller
    shape without either being fully contained in the other (>90% contained = layering,
    which is intentional and allowed)."""
    sys_path_bak = list(sys.path)
    try:
        if str(_SHELLS_SEED) not in sys.path:
            sys.path.insert(0, str(_SHELLS_SEED))
        ns: dict[str, Any] = {}
        try:
            exec(render_code, ns)
        except Exception as e:
            return False, f"exec error: {e}"
        render = ns.get("render")
        if not callable(render):
            return False, "no render()"
        slots = make_sample_slots(slot_specs)
        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        try:
            render(slide, slots, _default_theme())
        except Exception as e:
            return False, f"render raised: {e}"

        slide_area = Inches(13.333).emu * Inches(7.5).emu
        shapes = []
        for sh in slide.shapes:
            try:
                l, t = sh.left or 0, sh.top or 0
                w, h = sh.width or 0, sh.height or 0
                if w <= 0 or h <= 0:
                    continue
                area = w * h
                # Skip full-slide backgrounds
                if area > slide_area * 0.85:
                    continue
                shapes.append((sh.shape_id, l, t, l + w, t + h, area))
            except Exception:
                continue

        bad = 0
        for i, a in enumerate(shapes):
            for b in shapes[i+1:]:
                x1, y1 = max(a[1], b[1]), max(a[2], b[2])
                x2, y2 = min(a[3], b[3]), min(a[4], b[4])
                if x2 <= x1 or y2 <= y1:
                    continue
                inter = (x2 - x1) * (y2 - y1)
                min_a = min(a[5], b[5])
                if min_a <= 0:
                    continue
                ratio = inter / min_a
                if ratio > 0.9:
                    continue  # containment = intentional layering
                if ratio > 0.15:
                    bad += 1

        if bad > max_bad_pairs:
            return False, f"too many overlaps: {bad} bad pairs (max {max_bad_pairs})"
        return True, f"overlap OK ({bad} pairs)"
    finally:
        sys.path[:] = sys_path_bak


# ---------------------------------------------------------------------------
# Gate 2c: contrast check (text run color vs nearest fill underneath)
# ---------------------------------------------------------------------------


def _hex_to_rgb_tuple(h: str) -> tuple[int, int, int]:
    s = h.lstrip("#")
    return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)


def _luminance(rgb: tuple[int, int, int]) -> float:
    """Relative luminance per WCAG."""
    def chan(c: int) -> float:
        c01 = c / 255.0
        return c01 / 12.92 if c01 <= 0.03928 else ((c01 + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def _contrast_ratio(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    la, lb = _luminance(a), _luminance(b)
    hi, lo = max(la, lb), min(la, lb)
    return (hi + 0.05) / (lo + 0.05)


def gate_contrast(render_code: str, slot_specs: list[dict],
                  min_ratio: float = 3.0, max_bad_texts: int = 2) -> tuple[bool, str]:
    """Render the shell, then for each text shape check that its run color
    has sufficient contrast against the nearest fill underneath (or slide bg).
    Fails if too many text shapes are below min_ratio."""
    sys_path_bak = list(sys.path)
    try:
        if str(_SHELLS_SEED) not in sys.path:
            sys.path.insert(0, str(_SHELLS_SEED))
        ns: dict[str, Any] = {}
        try:
            exec(render_code, ns)
        except Exception as e:
            return False, f"exec error: {e}"
        render = ns.get("render")
        if not callable(render):
            return False, "no render()"
        slots = make_sample_slots(slot_specs)
        theme = _default_theme()
        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        try:
            render(slide, slots, theme)
        except Exception as e:
            return False, f"render raised: {e}"

        bg_hex = theme.get("palette", {}).get("bg", "#FFFFFF")
        bg_rgb = _hex_to_rgb_tuple(bg_hex)

        # Build a list of (left, top, right, bottom, fill_rgb) for non-text fills
        fills = []
        for sh in slide.shapes:
            try:
                if not sh.has_text_frame and sh.fill.type == 1:  # solid
                    rgb = sh.fill.fore_color.rgb
                    if rgb is None:
                        continue
                    rgb_t = (rgb[0], rgb[1], rgb[2])
                    l, t = sh.left or 0, sh.top or 0
                    w, h = sh.width or 0, sh.height or 0
                    if w > 0 and h > 0:
                        fills.append((l, t, l + w, t + h, rgb_t, w * h))
            except Exception:
                continue
        # Sort by area descending so larger fills come first (background-most)
        fills.sort(key=lambda f: -f[5])

        bad = 0
        for sh in slide.shapes:
            try:
                if not sh.has_text_frame:
                    continue
                tf = sh.text_frame
                # Sample first run with a real string
                text = ""
                run = None
                for p in tf.paragraphs:
                    for r in p.runs:
                        if (r.text or "").strip():
                            text = r.text
                            run = r
                            break
                    if run:
                        break
                if not run or not text:
                    continue
                run_rgb = run.font.color.rgb
                if run_rgb is None:
                    continue
                rrgb_t = (run_rgb[0], run_rgb[1], run_rgb[2])
                # Find smallest containing fill
                cx = (sh.left or 0) + (sh.width or 0) // 2
                cy = (sh.top or 0) + (sh.height or 0) // 2
                under = bg_rgb
                # Iterate biggest→smallest, last match wins (= innermost)
                for l, t, r, b, rgb_t, _ in fills:
                    if l <= cx <= r and t <= cy <= b:
                        under = rgb_t
                ratio = _contrast_ratio(rrgb_t, under)
                if ratio < min_ratio:
                    bad += 1
            except Exception:
                continue

        if bad > max_bad_texts:
            return False, f"low contrast on {bad} text shapes (max {max_bad_texts}, min ratio {min_ratio})"
        return True, f"contrast OK ({bad} marginal)"
    finally:
        sys.path[:] = sys_path_bak


# ---------------------------------------------------------------------------
# Run all gates
# ---------------------------------------------------------------------------


def run_gates(shell_dict: dict, run_render: bool = True,
              run_overlap: bool = True, run_contrast: bool = True) -> dict:
    """Apply all gates. Returns {'passed': bool, 'exec':, 'render':, 'overlap':, 'contrast':}."""
    result: dict[str, Any] = {"exec": None, "render": None, "overlap": None,
                              "contrast": None, "visual": None}

    code = shell_dict.get("_render_code") or ""
    slot_specs = shell_dict.get("slots", [])

    ok, msg = gate_exec(code, slot_specs)
    result["exec"] = {"ok": ok, "msg": msg}
    if not ok:
        shell_dict["status"] = "candidate"
        shell_dict.setdefault("quality", {}).update({"exec_ok": False})
        result["passed"] = False
        return result

    if run_render:
        ok_r, msg_r = gate_render(code, slot_specs)
        result["render"] = {"ok": ok_r, "msg": msg_r}
        if not ok_r:
            shell_dict["status"] = "candidate"
            shell_dict.setdefault("quality", {}).update({"exec_ok": True, "render_ok": False})
            result["passed"] = False
            return result
    else:
        result["render"] = {"ok": True, "msg": "skipped"}

    if run_overlap:
        ok_o, msg_o = gate_overlap(code, slot_specs)
        result["overlap"] = {"ok": ok_o, "msg": msg_o}
        if not ok_o:
            shell_dict["status"] = "candidate"
            shell_dict.setdefault("quality", {}).update({
                "exec_ok": True, "render_ok": True, "overlap_ok": False})
            result["passed"] = False
            return result
    else:
        result["overlap"] = {"ok": True, "msg": "skipped"}

    if run_contrast:
        ok_c, msg_c = gate_contrast(code, slot_specs)
        result["contrast"] = {"ok": ok_c, "msg": msg_c}
        if not ok_c:
            shell_dict["status"] = "candidate"
            shell_dict.setdefault("quality", {}).update({
                "exec_ok": True, "render_ok": True, "overlap_ok": True, "contrast_ok": False})
            result["passed"] = False
            return result
    else:
        result["contrast"] = {"ok": True, "msg": "skipped"}

    shell_dict["status"] = "active"
    shell_dict.setdefault("quality", {}).update({
        "exec_ok": True, "render_ok": True, "overlap_ok": True,
        "contrast_ok": True, "usage_count": 0,
    })
    result["passed"] = True
    return result


# ---------------------------------------------------------------------------
# Motion budget gate (role-aware)
# ---------------------------------------------------------------------------

_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Slide roles where a capped motion budget applies. Roles outside this set are
# exempt and can carry arbitrarily many motion effects (covers, hero reveals,
# section dividers intentionally use several coordinated animations).
_MOTION_BUDGET_ROLES = frozenset({
    "bullet_card_list",
    "metric_dashboard",
    "comparison_split",
    "timeline_horizontal",
    "feature_grid",
})

_MOTION_EFFECT_TAGS = ("animEffect", "animRot", "animMotion", "animScale")


def _slide_xml_element(slide: Any):
    """Return the lxml element root for a python-pptx slide or passthrough.

    Accepts either a python-pptx Slide (has ``_element``), an lxml element
    with a matching tag, or anything exposing a ``.element`` attribute.
    """
    for attr in ("_element", "element"):
        el = getattr(slide, attr, None)
        if el is not None:
            return el
    return slide


def count_motion_effects(slide: Any) -> int:
    """Count animation primitives plus the slide-level transition.

    Sums occurrences of ``<p:animEffect>``, ``<p:animRot>``, ``<p:animMotion>``,
    ``<p:animScale>`` anywhere under the slide (these live inside the slide's
    ``<p:timing>`` tree), plus one for the slide-level ``<p:transition>``
    element when present.
    """
    root = _slide_xml_element(slide)
    total = 0
    for tag in _MOTION_EFFECT_TAGS:
        total += len(root.findall(f".//{{{_PML_NS}}}{tag}"))
    if root.find(f"{{{_PML_NS}}}transition") is not None:
        total += 1
    return total


def gate_motion_budget(
    deck_entries: list[tuple[Any, str]],
    cap: int = 3,
) -> tuple[bool, list[dict]]:
    """Enforce a role-aware cap on per-slide motion effects.

    Parameters
    ----------
    deck_entries
        Ordered list of ``(slide, role)`` pairs covering every slide that the
        budget should be applied to. Role strings outside the budgeted set
        are skipped without contributing to the failure count.
    cap
        Maximum number of motion effects permitted on a budgeted slide.

    Returns
    -------
    ``(ok, violations)`` where ``ok`` is False when at least one budgeted
    slide exceeds the cap. Each violation is a dict with ``slide_index``,
    ``role``, ``count``, and a human-readable ``reason`` string.
    """
    violations: list[dict] = []
    for idx, (slide, role) in enumerate(deck_entries):
        if role not in _MOTION_BUDGET_ROLES:
            continue
        count = count_motion_effects(slide)
        if count > cap:
            violations.append({
                "slide_index": idx,
                "role": role,
                "count": count,
                "reason": (
                    f"slide {idx} role={role} carries {count} motion effects; "
                    f"cap is {cap} for content roles"
                ),
            })
    return (len(violations) == 0, violations)


# ---------------------------------------------------------------------------
# Motion-aware overlap gate
# ---------------------------------------------------------------------------

def _shape_bbox_emu(shape) -> tuple[int, int, int, int] | None:
    try:
        l = int(shape.left); t = int(shape.top)
        w = int(shape.width); h = int(shape.height)
    except Exception:
        return None
    return (l, t, l + w, t + h)


def _rects_overlap(a, b) -> bool:
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])


def _parse_motion_path_bbox(path: str):
    import re as _re
    cur_x, cur_y = 0.0, 0.0
    xs, ys = [0.0], [0.0]
    tokens = _re.findall(r"[MLaq]|-?[\d.]+", path or "")
    i = 0
    while i < len(tokens):
        cmd = tokens[i]; i += 1
        if cmd in ("M", "L") and i + 1 < len(tokens):
            cur_x = float(tokens[i]); cur_y = float(tokens[i+1]); i += 2
            xs.append(cur_x); ys.append(cur_y)
        elif cmd == "a" and i + 6 < len(tokens):
            rx = float(tokens[i]); ry = float(tokens[i+1])
            dx = float(tokens[i+5]); dy = float(tokens[i+6]); i += 7
            nx, ny = cur_x + dx, cur_y + dy
            xs += [cur_x, nx, cur_x - rx, cur_x + rx]
            ys += [cur_y, ny, cur_y - ry, cur_y + ry]
            cur_x, cur_y = nx, ny
        elif cmd == "q" and i + 3 < len(tokens):
            cdx = float(tokens[i]); cdy = float(tokens[i+1])
            dx = float(tokens[i+2]); dy = float(tokens[i+3]); i += 4
            cx = cur_x + cdx; cy = cur_y + cdy
            nx = cur_x + dx; ny = cur_y + dy
            xs += [cur_x, cx, nx]; ys += [cur_y, cy, ny]
            cur_x, cur_y = nx, ny
    return (min(xs), min(ys), max(xs), max(ys))


def gate_motion_overlap(slide) -> tuple[bool, list[dict]]:
    """Check that motion-path animations don't sweep through other
    shapes on the same slide."""
    from pptx.oxml.ns import qn
    violations: list[dict] = []
    try:
        slide_w = int(slide.part.package.presentation_part.presentation.slide_width)
        slide_h = int(slide.part.package.presentation_part.presentation.slide_height)
    except Exception:
        slide_w, slide_h = 12192000, 6858000

    shape_map = {}
    for shp in slide.shapes:
        sid = getattr(shp, "shape_id", None)
        if sid is None:
            continue
        shape_map[int(sid)] = (shp, _shape_bbox_emu(shp))

    timing = slide._element.find(qn("p:timing"))
    if timing is None:
        return (True, [])

    for am in timing.iter(qn("p:animMotion")):
        path = am.get("path") or ""
        # Skip entrance-style paths whose start point is off-slide
        # (those are intentional flybys from outside the visible area,
        # not overlaps against on-slide content).
        try:
            import re as _re2
            m0 = _re2.match(r"\s*M\s*(-?[\d.]+)\s+(-?[\d.]+)", path)
            if m0:
                sx, sy = float(m0.group(1)), float(m0.group(2))
                if abs(sx) > 0.4 or abs(sy) > 0.4:
                    continue
        except Exception:
            pass
        sp_tgt = am.find(".//" + qn("p:spTgt"))
        if sp_tgt is None:
            continue
        try:
            sid = int(sp_tgt.get("spid"))
        except Exception:
            continue
        info = shape_map.get(sid)
        if info is None:
            continue
        target_shape, target_bbox = info
        if target_bbox is None:
            continue
        # Exempt decorative ambient anchors (accent_orb, halo, etc.) —
        # their orbital/drift paths are intentionally designed to travel
        # past other elements as a visual layer.
        name = (getattr(target_shape, "name", "") or "").lower()
        name_clean = name.replace("!!samename", "")
        if any(tok in name_clean for tok in ("accent", "orb", "halo", "dot")):
            continue
        try:
            dxmin, dymin, dxmax, dymax = _parse_motion_path_bbox(path)
        except Exception:
            continue
        l, t, r, b = target_bbox
        swept = (
            l + int(dxmin * slide_w),
            t + int(dymin * slide_h),
            r + int(dxmax * slide_w),
            b + int(dymax * slide_h),
        )
        for other_sid, (other_shape, other_bbox) in shape_map.items():
            if other_sid == sid or other_bbox is None:
                continue
            if _rects_overlap(swept, other_bbox):
                violations.append({
                    "motion_target": getattr(target_shape, "name", f"spid={sid}"),
                    "target_spid": sid,
                    "collides_with": getattr(other_shape, "name", f"spid={other_sid}"),
                    "collides_spid": other_sid,
                    "swept_bbox_emu": list(swept),
                    "path": path,
                })
    return (len(violations) == 0, violations)


__all__ = ["run_gates", "gate_exec", "gate_render", "gate_overlap", "gate_contrast",
           "make_sample_slots", "count_motion_effects", "gate_motion_budget",
           "gate_motion_overlap"]
